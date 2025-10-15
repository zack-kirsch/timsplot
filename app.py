# =============================================================================
# Library Imports (only necessary for launch)
# =============================================================================
#region

from shiny import App, Inputs, Outputs, Session, reactive, render, ui, module
from shinyswatch import theme
#https://rstudio.github.io/shinythemes/
from shiny.types import ImgData
from faicons import icon_svg
# #https://fontawesome.com/search?o=r&m=free

#endregion

# =============================================================================
# UI
# =============================================================================
#region

app_ui=ui.page_fluid(
    ui.panel_title("timsplot: timsTOF Proteomics Data Visualization (v.2025.09.18)"),
    ui.navset_pill_list(
        ui.nav_panel("File Import",
            ui.row(
                ui.column(7,
                    ui.card(
                        ui.card_header("Upload Search Report"),
                            ui.row(
                                ui.column(6,
                                    ui.popover(
                                        ui.input_action_button("info_btn","Instructions",class_="btn-success",icon=icon_svg("question")),
                                        ui.p("-Note: if using the app in a browser window, zoom out to 90% or farther (some OS scaling settings may cause UI elements to overlap)."),
                                        ui.p("-Select software used for search and upload .tsv, .zip, .parquet, or .csv file (multiple files utilizing the same software can be uploaded)."),
                                        ui.p("-Fill out R.Condition and R.Replicate in the metadata table as needed."),
                                        ui.p("-Select necessary switches under 'Update from Metadata Table' and click the 'Apply Changes' button."),
                                        ui.p("-Note: click on the 'Apply Changes' button after upload even if metadata table was not updated."),
                                        placement="left"
                                    ),
                                    ui.p(),
                                    ui.input_radio_buttons("software_general","Software used for search:",{"spectronaut":"Spectronaut",
                                                                                                            "diann":"DIA-NN",
                                                                                                            "fragpipe":"FragPipe",
                                                                                                            "bps":"Bruker ProteoScape/GlycoScape",
                                                                                                            "spectromine":"Spectromine",
                                                                                                            "peaks":"PEAKS",
                                                                                                            "sage":"Sage",
                                                                                                            }),
                                    ui.hr(),
                                    ui.output_ui("software_ui"),
                                    ui.hr(),
                                    ui.panel_conditional("input.software === 'bps_timsrescore' || input.software === 'bps_timsdiann' || input.software === 'bps_spectronaut' || input.software === 'bps_pulsar'",
                                                         ui.input_radio_buttons("software_bps_report_type","",choices={"qual":"Qualitative","quant":"Quantitative"})
                                    )
                                ),
                                ui.column(6,
                                    ui.input_file("searchreport","Upload search report(s):",accept=[".tsv",".zip",".parquet",".psms.csv"],multiple=True),
                                    ui.input_switch("searchreport_reupload","Reupload of processed search file?"),
                                    ui.output_text("metadata_reminder"),
                                    ui.panel_conditional("input.software_general === 'diann'",
                                                         ui.input_radio_buttons("diann_mbr_switch","Q value filtering option:",choices=["off","Protein.Q.Value","Global.PG.Q.Value"])
                                    )
                                ),
                            )
                    )
                ),
                ui.column(5,
                    ui.card(
                        ui.card_header("Update from Metadata Table"),
                            ui.row(
                                ui.column(3),
                                ui.input_action_button("rerun_metadata","Apply Metadata to Search File",width="300px",class_="btn-primary",icon=icon_svg("rotate")),
                                ui.column(1,ui.output_ui("metadata_applied_ui"))
                            ),
                            ui.row(
                                ui.column(5,
                                    ui.input_switch("remove","Remove selected runs"),
                                ),
                                ui.column(7,
                                    ui.input_switch("cRAP_filter","Filter protein list with cRAP database"),
                                    #ui.tags.a("cRAP Database",href="https://www.thegpm.org/crap/",target="_blank")
                                )
                            ),
                            ui.row(
                                ui.column(1),
                                ui.popover(
                                    ui.input_action_button("searchreport_download_info_btn","Instructions",class_="btn-success",icon=icon_svg("question")),
                                    ui.p("-After hitting 'Apply Changes' button, updated search file can be downloaded for easier reuploads."),
                                    ui.p("-Use reupload switch if the input file has already been processed in timsplot since necessary columns are already renamed/added. Select the same software that was used in the initial upload."),
                                    ui.p("-Condition, replicate, and concentration values from the metadata are automatically applied and conditions are reordered after hitting 'Apply Changes'."),
                                    placement="left"
                                ),
                                ui.column(1),
                                ui.download_button("searchreport_download","Download file for reuploads",width="300px",icon=icon_svg("file-arrow-down")),
                            )
                    ),
                    ui.card(
                        ui.card_header("Enter path for figure downloads (optional)"),
                        ui.row(
                            ui.column(1),
                            ui.popover(
                                ui.input_action_button("imagedownload_btn","Instructions",class_="btn-success",icon=icon_svg("question")),
                                ui.p("-The most recent generation of the figure shown will be the one that is saved to the specified folder. If multiple iterations of the same figure are desired (e.g. average and individual ID counts), it's recommended to manually rename the files once they're generated."),
                                ui.p("-The directory specified cannot be the timsplot folder, since the app reloads when changes are made to that directory."),
                                ui.p("-The filename of the generated image will contain the name of the software, the name of the file uploaded, and a short identifier to the plot (e.g. 'spectronaut_results.tsv_idmetrics.png')"),
                                ui.p("-If multiple files are uploaded, the first file will be used in the generated file name"),
                                placement="left"
                            ),
                                ui.column(1),
                            ),
                        ui.row(
                            ui.input_text_area("download_path","",width="100%",height="65px",placeholder="ex. - C:\\Users\\Data")
                        ),
                    )
                )
            ),
            ui.card(
                ui.card_header("Metadata Tables"),
                ui.row(
                    ui.column(2,
                        ui.popover(
                            ui.input_action_button("metadata_info_btn","Instructions",class_="btn-success",icon=icon_svg("question")),
                            ui.p("-To remove runs, add an 'x' to the 'remove' column."),
                            ui.p("-To reorder conditions, order them numerically in the 'order' column. Runs will be sorted by the HyStar ID at the end of the file name."),
                            ui.p("-When updating the Concentration column, make sure to unify the units that are being input. If inputting concentrations in ng and there are samples in pg, convert values to the same unit (either ng or pg, not both)."),
                            placement="left"
                        )
                    ),
                    ui.column(4,
                        ui.input_file("metadata_upload","(Optional) Upload filled metadata table:",accept=".csv",multiple=False),
                        ui.input_switch("use_uploaded_metadata","Use uploaded metadata table"),
                    ),
                    ui.column(4,
                        ui.download_button("metadata_download","Download metadata table as shown",width="300px",icon=icon_svg("file-arrow-down"))
                    ),
                ),
                ui.row(
                    ui.column(8,
                        ui.output_data_frame("metadata_table")
                    ),
                    ui.column(4,
                        ui.output_data_frame("metadata_condition_table")
                    )
                )
            ),
            icon=icon_svg("folder-open")
        ),
        ui.nav_panel("Settings",
            ui.navset_pill(
                ui.nav_panel("Color Settings",
                    ui.row(
                        ui.column(4,
                            ui.p(),
                            ui.popover(
                                ui.input_action_button("colors_info_btn","Instructions",class_="btn-success",icon=icon_svg("question")),
                                ui.p("-Replicates of the same condition will have the same color."),
                                ui.p("-'Pick for me' options choose colors based on splitting a rainbow spectrum across the number of conditions or by choosing colors from the matplotlib tableau colors."),
                                ui.p("-Custom colors can be input by adding them to the text box below. Make sure to not leave whitespace or extra lines."),
                                ui.p("-Hex codes can be used to specify custom colors."),
                                placement="right"
                            ),
                            ui.input_radio_buttons("coloroptions","Choose coloring option for output plots:",choices={"pickrainbow":"Pick for me (rainbow)","pickmatplot":"Pick for me (matplotlib tableau)","custom":"Custom"},selected="pickmatplot"),
                            ui.input_text_area("customcolors","Input custom color names (or hex codes), one per line:",autoresize=True),
                            ui.row(
                                ui.column(5,
                                    ui.output_table("customcolors_table1")
                                ),
                                ui.column(5,
                                    ui.output_table("conditioncolors"),
                                    ui.output_plot("customcolors_plot")
                                )
                            )
                        ),
                        ui.column(2,
                            ui.p("Matplotlib Tableau Colors:"),
                            ui.output_image("matplotcolors_image"),
                        ),
                        ui.column(5,
                            ui.p("CSS Colors:"),
                            ui.output_image("csscolors_image"),
                        ),
                    ),
                ),
                ui.nav_panel("Control Panel",
                    ui.row(
                        ui.column(4,
                            ui.card(
                                ui.card_header("Plot Parameters"),
                                ui.input_slider("titlefont","Plot title text size",min=10,max=25,value=20,step=1,ticks=True),
                                ui.input_slider("axisfont","Axis title text size",min=10,max=25,value=15,step=1,ticks=True),
                                ui.input_slider("axisfont_labels","Axis label text size",min=10,max=25,value=10,step=1,ticks=True),
                                ui.input_slider("labelfont","Data label text size",min=10,max=25,value=15,step=1,ticks=True),
                                ui.input_slider("legendfont","Legend text size",min=10,max=25,value=10,step=1,ticks=True),
                                ui.input_slider("ypadding","y-axis padding for data labels",min=0,max=1,value=0.3,step=0.05,ticks=True),
                                ui.input_slider("xaxis_label_rotation","x-axis label rotation",min=0,max=360,value=90,step=5,ticks=True),
                                ui.input_radio_buttons("xaxis_label_alignment","x-axis label alignment (if using angled axis labels)",choices=["center","left","right"])
                            )
                        ),
                        ui.column(4,
                            ui.card(
                                ui.card_header("Misc."),
                                ui.input_radio_buttons("peptide_grouping","Grouping key for peptide CVs",choices={"stripped":"Stripped sequence","modified":"Modified sequence"}),
                                # ui.input_radio_buttons("msquantlevel","Quant level to use for plotting (if available)",width="400px",choices={"MS1Quantity":"MS1","MS2Quantity":"MS2"}),
                                # "PG."+input.msquantlevel()
                                # "FG."+input.msquantlevel()
                                ui.p(""),
                                ui.input_switch("dpi_switch","Change DPI to 300 for publication quality",value=False,width="400px"),
                                ui.p("Note: values in the width/height sliders for plots will need to be increased to accommodate the DPI change, default plotting parameters will be too small since Shiny plots based on pixels."),
                            )
                        )
                    )
                ),
                ui.nav_panel("PTM Settings",
                    ui.popover(
                        ui.input_action_button("ptminfo_btn","Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                        ui.p("-The table serves as a dictionary for substitutions that are made in the search file upon import to simplify PTM classification. Different search software will notate each PTM differently and this tool serves as a way of unifying that notation for data visualization."),
                        ui.p("-The PTM Identifier is the value that is given in the search file to indicate a specific PTM. The Replacement value is what the PTM Identifier is replaced within this app."),
                        ui.p("-To add a PTM not shown in the table, input the exact way it is notated in the search file in the first text box. Then to the second text box add what you want the PTM to be referred to as. Then press 'Apply Changes', which will apply the updated PTM naming to an uploaded file."),
                        ui.p("-Recommended format for replacement value: 'modname (loc)'. For example, carbamidomethylation of Cys residues is notated as 'Carbamidomethyl (C)'. Do not include underscores or more than one space."),
                        ui.p("-NOTE: For modifications from FragPipe, a different format is needed because of the format that FragPipe writes modifications. For example, a fixed Cys modification of 57 Da have an identifier of 'C160' and the Replacement would be 'C[Carbamidomethyl (C)]'."),
                        ui.p("-To save user-defined PTMs to the local dictionary file for future launches of the app, hit 'Save Changes' and overwrite the ptmdict file located in the timsplot folder."),
                    ),
                    ui.row(
                        ui.column(4,
                            ui.p("Table of PTM naming substitutions performed on file import:"),
                            ui.output_table("ptmdict_table")
                        ),
                        ui.column(3,
                            ui.input_text("ptm_key_input","Input PTM identifier (from search file/software):"),
                            ui.input_text("ptm_value_input","Input PTM replacement name:"),
                            ui.input_action_button("ptm_apply","Apply changes",class_="btn-primary",icon=icon_svg("rotate"),width="300px"),
                            ui.p("Note: Hitting 'Apply Changes' button will reset the metadata table if files have already been uploaded."),
                            ui.download_button("ptm_save","Save Changes to ptmdict File",class_="btn-success",icon=icon_svg("file-arrow-down"),width="300px"),                                                    
                        )
                    )
                ),
                ui.nav_panel("File Stats",
                    ui.output_table("filestats")
                ),
                ui.nav_panel("Column Check",
                    ui.output_table("column_check")
                ),
                ui.nav_panel("File Preview",
                    ui.output_data_frame("filepreview")
                ),
                ui.nav_panel("Extra Colors",
                    ui.p("Bruker Colors:"),
                    ui.output_image("brukercolors")
                )
            ),
            icon=icon_svg("gear")
        ),
        ui.nav_panel("ID Counts",
            ui.navset_pill(
                ui.nav_panel("Counts per Condition",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("counts_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar graphs of the number of IDs per run with subplots for each category. The dropdown menu can be used to plot each category individually instead of as subplots."),
                                ui.p("Notes:"),
                                ui.p("-Proteins: number of unique values in the ProteinGroups column."),
                                ui.p("-Proteins with >2 Peptides: number of ProteinGroups that have more than 2 unique ModifiedPeptides."),
                                ui.p("-Peptides: number of unique values in the ModifiedPeptide column (agnostic to charge)."),
                                ui.p("-Precursors: number of unique values between the ModifiedPeptide and Charge columns."),
                                placement="right"
                                ),
                            ui.row(
                                ui.input_slider("idmetrics_width","Plot width",min=100,max=7500,step=100,value=1500,ticks=True),
                                ui.input_slider("idmetrics_height","Plot height",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_radio_buttons("idmetrics_individual_average","Choose how to show ID counts:",choices={"individual":"Individual Counts","average":"Average Counts","unique":"Total Unique Counts"},width="250px"),
                                ui.input_selectize("idplotinput","Choose what metric to plot:",choices={"all":"All","proteins":"Proteins","proteins2pepts":"Proteins with >2 Peptides","peptides":"Peptides","precursors":"Precursors"},multiple=False,selected="all",width="250px"),
                                ui.input_radio_buttons("idmetrics_peptides","Pick peptide counts to plot:",choices={"modified":"Modified Sequences","stripped":"Stripped Sequences"}),
                            ),
                            ui.row(
                                ui.output_ui("idmetrics_peplength_switch_ui"),
                                ui.output_ui("idmetrics_peplength_ui")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("idmetricsplot")
                ),
                ui.nav_panel("Venn Diagram",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("venn_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("venn_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("venn_conditionorrun","Plot by condition or individual run?",choices={"condition":"Condition","individual":"Individual Run"}),
                                ui.input_radio_buttons("venn_numcircles","Number of comparisons:",choices={"2":"2","3":"3"},width="200px"),
                                ui.column(2,
                                    ui.download_button("venn_download","Download Venn list",width="300px",icon=icon_svg("file-arrow-down")),
                                    ui.download_button("venn_download_detailed","Download Detailed Venn list",width="300px",icon=icon_svg("file-arrow-down"))
                                )
                            ),
                            ui.row(
                                ui.column(2,
                                    ui.output_ui("venn_run1_ui"),
                                    ui.output_ui("venn_run2_ui"),
                                    ui.output_ui("venn_run3_ui"),
                                ),
                                ui.column(2,
                                    ui.input_text("venn_color1","Color 1:",value="tab:blue"),
                                    ui.input_text("venn_color2","Color 2:",value="tab:orange"),
                                    ui.output_ui("venn_color3_ui")
                                ),
                                ui.column(2,
                                    ui.input_radio_buttons("venn_plotproperty","Metric to compare:",choices={"proteingroups":"Protein Groups","peptides":"Peptides","precursors":"Precursors","peptides_stripped":"Stripped Peptides"}),
                                    ui.output_ui("peptidecore_ui"),
                                ),
                                ui.column(2,
                                    ui.output_ui("venn_ptm_ui"),
                                    ui.output_ui("venn_ptmlist_ui"),
                                ),
                                ui.column(2,
                                    ui.output_ui("venn_specific_length_ui"),
                                    ui.output_ui("venn_peplength_ui"),
                                ),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("venn_plot")
                ),
                ui.nav_panel("CV Plots",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("cvplot_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Violin plots of the CVs per condition for the selected type of ID. Use the switch to remove top 5% outliers. The table displays the mean and median CV values for the shown plot and updates when the switch is toggled."),
                                ui.p("Notes:"),
                                ui.p("-Go to Settings -> Control Panel to change grouping key from stripped sequence to modified sequence (default is stripped sequence)."),
                                ui.p("-Peptide-level CVs are calculated from the top 3 precursors for a given sequence."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("cvplot_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("cvplot_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("proteins_precursors_cvplot","Pick which IDs to plot",choices={"Protein":"Proteins","Precursor":"Precursors","Peptide":"Peptides"},width="200px"),
                                ui.column(3,
                                    ui.input_switch("removetop5percent","Remove top 5%",width="200px"),
                                    ui.input_switch("cvplot_histogram_bins_switch","Show as histogram"),
                                    ui.output_ui("cvplot_histogram_bins_ui")
                                )
                            ),
                        ),
                        width="100%"
                    ),   
                    ui.output_table("cv_table"),
                    ui.output_plot("cvplot")
                ),
                ui.nav_panel("IDs with CV Cutoff",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("cvcutoffplot_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar graphs for the numbers of IDs per condition that fall within CV cutoffs."),
                                ui.p("-The 'Identified' bar is the average number of IDs per condition."),
                                ui.p("Notes:"),
                                ui.p("-Go to Settings -> Control Panel to change grouping key from stripped sequence to modified sequence (default is stripped sequence)."),
                                ui.p("-Peptide-level CVs are calculated from the top 3 precursors for a given sequence."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("countscvcutoff_width","Plot width",min=100,max=7500,step=100,value=900,ticks=True),
                                ui.input_slider("countscvcutoff_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                                ui.input_radio_buttons("proteins_precursors_idcutoffplot","Pick which IDs to plot",choices={"proteins":"Proteins","precursors":"Precursors","peptides":"Peptides"}),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_table("countscvcutoff_table"),
                    ui.output_plot("countscvcutoff")
                ),
                ui.nav_panel("UpSet Plot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("upsetplot_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Plot showing the intersections of IDs found in each run/condition."),
                                ui.p("Notes:"),
                                ui.p("-The UpSet plot can be generated by comparing every condition, a selected condition, or by comparing all of the runs that were uploaded."),
                                ui.p("-The table shows the results of the filtering options according to the plotting options."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("upsetplot_width","Plot width",min=100,max=7500,step=100,value=900,ticks=True),
                                ui.input_slider("upsetplot_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                            ),
                            ui.row(
                                ui.input_selectize("protein_precursor_pick","Pick which IDs to plot:",choices={"Protein":"Protein","Peptide":"Peptide","Precursor":"Precursor"}),
                                ui.input_radio_buttons("upsetfilter","Filtering option:",choices={"nofilter":"No filtering","1run":"IDs in only 1 run/replicate","n-1runs":"IDs in n-1 runs/replicates"}),
                                ui.input_radio_buttons("upset_condition_or_run","Pick how to plot UpSet plot:",choices={"condition":"All conditions","specific_condition":"By specific condition","individual":"All runs"}),
                                ui.output_ui("specific_condition_ui"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_data_frame("upsetplot_counts"),
                    ui.output_plot("upsetplot")
                ),
                ui.nav_panel("UpSet Plot (stats)",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("upsetplotstats_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Ion mobility vs m/z scatterplot/heatmap based on the UpSet plot calculations to visualize IDs that were only found in single runs/conditions."),
                                ui.p("Notes:"),
                                ui.p("-'Specific condition in entire result file' shows the IDs that were only found in the specified condition compared to the IDs in the whole result file."),
                                ui.p("-'From specific condition' shows the IDs that were only found in single runs within the specified condition."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("upsetplotstats_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("upsetplotstats_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                            ),
                            ui.row(
                                ui.input_radio_buttons("upsetplotstats_plottype","Choose how to plot:",choices={"scatter":"Scatterplot","2dhist":"2-D Histogram"}),
                                ui.input_radio_buttons("upsetplotstats_peptide_precursor","ID type to plot:",choices={"Peptide":"Peptide","Precursor":"Precursor"}),
                                ui.input_radio_buttons("upsetplotstats_whattoplot","Choose what to show for single-hit IDs:",choices={"individual":"From entire result file","condition":"Specific condition in entire result file","specific_condition":"From specific condition"},width="350px"),
                                ui.output_ui("upsetplotstats_conditionlist_ui"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("upsetplotstats_singlehitIDplot"),
                ),
                ui.nav_panel("Protein/Peptide Signal Tracker",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("signaltracker_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Scatterplot of protein or peptide signal across runs."),
                                ui.p("Notes:"),
                                ui.p("-Selection of a protein group will show the PG.MS2Quantity across all runs and will show a table of all peptide sequences associated with that protein."),
                                ui.p("-Selecting a peptide sequence will show the FG.MS2Quantity across all runs for all charge and modification states of the selected sequence."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("tracker_width","Plot width",min=100,max=7500,step=100,value=900,ticks=True),
                                ui.input_slider("tracker_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                                ui.column(3,
                                    ui.input_switch("tracker_logscale","Set y-axis scale to log10"),
                                    ui.input_switch("tracker_yaxiszero","Set y-axis minimum to 0"),
                                ),
                                ui.input_radio_buttons("tracker_proteingrouping","Pick protein column to use:",choices=["PG.ProteinNames","PG.ProteinGroups"])
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(6,
                            ui.output_data_frame("protein_df"),
                            ui.output_data_frame("pickedprotein_df")
                        ),
                        ui.column(6,
                            ui.output_plot("tracker_plot")
                        )
                    ),
                ),
                ui.nav_panel("Individual Protein Counts",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                            #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                            #ui.p("Description:"),
                            #ui.p("-"),
                            #ui.p("Notes:"),
                            #ui.p("-"),
                            #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("protein_id_plot_width","Plot width",min=100,max=7500,step=100,value=900,ticks=True),
                                ui.input_slider("protein_id_plot_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                                ui.input_radio_buttons("protein_id_individual_average","Choose how to show ID counts:",choices={"individual":"Individual Counts","average":"Average Counts"},width="250px"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(5,
                            ui.output_data_frame("protein_ID_df"),
                        ),
                        ui.column(6,
                            ui.output_plot("protein_id_plot")
                        )
                    )
                ),
            ),
            icon=icon_svg("chart-simple")
        ),
        ui.nav_panel("Metrics",
            ui.navset_pill(
                ui.nav_panel("Charge State",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("chargestate_info_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar plots of detected precursor charges annotated with percents or counts of each charge in the search file(s)."),
                                ui.p("-Per file, the values shown are the number of unique modified peptide sequences per charge state."),
                                ui.p("-Averages shown are average number of unique modified peptide sequences per charge state per sample condition."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("chargestate_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("chargestate_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.column(3,
                                    ui.input_radio_buttons("chargestate_condition_or_run","Plot by condition or by individual run?",choices={"condition":"Condition","individual":"Individual Run"},selected="condition"),
                                    ui.input_radio_buttons("chargestate_counts_percent","Show numbers as counts or percent:",choices=["Percent","Counts"],selected="Percent"),
                                ),
                                ui.column(3,
                                    ui.input_switch("chargestate_stacked","Show as stacked bar graphs",value=False),
                                    ui.output_ui("chargestate_averages_ui"),
                                    ui.input_switch("chargestate_peplength","Plot for specific peptide length"),
                                    ui.output_ui("chargestate_peplength_slider_ui"),
                                )
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("chargestateplot")
                ),
                ui.nav_panel("Peptide Length",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("peptidelength_info_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Line or bar plots of the number of unique unmodified peptide sequences per peptide length."),
                                ui.p("Notes:"),
                                ui.p("-Optional marker automatically picks the maximum among the data."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("peptidelength_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("peptidelength_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.column(3,
                                    ui.input_radio_buttons("peptidelengths_condition_or_run","Plot by condition or by individual run?",choices={"condition":"Condition","individual":"Individual Run"}),
                                    ui.input_radio_buttons("peplengthinput","Line plot or bar plot?",choices={"lineplot":"line plot","barplot":"bar plot"}),
                                ),
                                ui.output_ui("lengthmark_ui"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("peptidelengthplot")
                ),
                ui.nav_panel("Peptides per Protein",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("pepsperprotein_info_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Line or bar plots of the number of unique modified peptide sequences per protein."),      
                                ui.p("Notes:"),
                                ui.p("-The 'X-axis high bound' slider is meant to simplify the view since the bar plots get condensed to a point of showing visual aberrations at high enough x scales."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("pepsperprotein_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("pepsperprotein_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("pepsperprotein_condition_or_run","Plot by condition or by individual run?",choices={"condition":"Condition","individual":"Individual Run"}),
                                ui.input_radio_buttons("pepsperproteininput","Line plot or bar plot?",choices={"lineplot":"line plot","barplot":"bar plot"}),
                                ui.input_slider("pepsperprotein_xrange","X-axis high bound",min=0,max=200,value=50,step=5,ticks=True)
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("pepsperproteinplot")
                ),
                ui.nav_panel("Dynamic Range",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("dynamicrange_info_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Scatter plot of ranked protein signal."),
                                ui.p("-Per sample condition, mean or median protein signal is ranked and shown in the plot."),
                                ui.p("-Contribution to the overall protein signal is shown in the upper part of the plot, with the number of protein groups that contribute to 25%, 50% and 75% of the protein signal."),
                                ui.p("-The table shows the top N proteins according to their mean or median signal."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("dynamicrange_width","Plot width",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_slider("dynamicrange_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                                ui.output_ui("sampleconditions_ui"),
                                ui.input_selectize("meanmedian","Mean or median",choices={"mean":"mean","median":"median"}),
                                ui.input_numeric("top_n","Input top N proteins to display:",value=25,min=5,step=5),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(4,
                            ui.output_data_frame("dynamicrange_proteinrank"),
                        ),
                        ui.column(6,
                            ui.output_plot("dynamicrangeplot")
                        )
                    )
                ),
                ui.nav_panel("Mass Accuracy",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("massaccuracy_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Violin plot or histogram of ppm mass accuracy for each precursor."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("massaccuracy_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("massaccuracy_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("massaccuracy_violin_hist","Plot as a violin plot or histogram?",choices={"violin":"Violin Plot","histogram":"Histogram"}),
                                ui.output_ui("massaccuracy_bins_ui")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("massaccuracy_plot")
                ),
                ui.nav_panel("Peak Width",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("peakwidth_info_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Violin plot of LC peak width for each precursor."),
                                ui.p("Notes:"),
                                ui.p("-If peak width is not exported in the original search file, peak width is calculated using the Start and End values for the retention time."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("peakwidth_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("peakwidth_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_switch("peakwidth_removetop5percent","Remove top 5%"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(5,
                            ui.output_table("peakwidth_table"),
                        ),
                        ui.column(7,
                            ui.output_plot("peakwidthplot")
                        )
                    ),              
                ),
                ui.nav_panel("Missed Cleavages",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("missedcleavages_info_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar plot of numbers of missed cleavages. Peptide IDs that do not follow the enzyme rules are included."),
                                ui.p("Notes:"),
                                ui.p("-Per file or per condition, stripped peptide sequences are checked against the selected enzyme cleavage rules."),
                                ui.p("-The sum of the bars for each run is the number of unique stripped peptide sequences for the given run."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("missedcleavages_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("missedcleavages_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.input_slider("missedcleavages_barwidth","Bar width",min=0.1,max=1,step=0.05,value=0.2,ticks=True),
                                ui.input_radio_buttons("enzyme_rules","Enzyme cleavage rules",choices={"trypsin":"Trypsin (K/R)"}),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("missedcleavages_plot")
                ),
                ui.nav_panel("Data Completeness",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("datacompleteness_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar/scatter plot showing data completeness across runs."),
                                ui.p("-Similar to UpSet plot, this calculates how many runs each unique protein or stripped/modified peptide is detected in and shows counts and percents according to the number of IDs."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("datacompleteness_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("datacompleteness_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("protein_peptide","Pick what metric to plot:",choices={"proteins":"Proteins","peptides":"Modified Peptides","strippedpeptides":"Stripped Peptides"}),
                                ui.column(3,
                                    ui.input_switch("datacompleteness_sampleconditions_switch","Plot for specific condition?",value=False),
                                    ui.output_ui("datacompleteness_sampleconditions_ui"),
                                )

                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("datacompletenessplot")
                ),
            ),
            icon=icon_svg("chart-line")
        ),
        ui.nav_panel("PTMs",
            ui.navset_pill(
                ui.nav_panel("PTMs found",
                    ui.p("PTMs Found in Search File:"),
                    ui.output_text_verbatim("ptmlist_text"),
                    ui.hr(),
                    ui.output_ui("ptmlist_ui")
                ),
                ui.nav_panel("Counts per Condition",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("ptmcounts_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar graphs of the number of IDs containing the specified PTM per run with subplots for each category. The dropdown menu can be used to plot each category individually instead of as subplots."),
                                ui.p("Notes:"),
                                ui.p("-The radio button can be used to toggle whether the bar graphs shown counts or the % of all IDs that contain the selected PTM."),
                                ui.p("-Proteins: number of unique values in the ProteinGroups column."),
                                ui.p("-Proteins with >2 Peptides: number of ProteinGroups that have more than 2 unique ModifiedPeptides."),
                                ui.p("-Peptides: number of unique values in the ModifiedPeptide column (agnostic to charge)."),
                                ui.p("-Precursors: number of unique values between the ModifiedPeptide and Charge columns."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("ptmidmetrics_width","Plot width",min=100,max=7500,step=100,value=1500,ticks=True),
                                ui.input_slider("ptmidmetrics_height","Plot height",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.column(3,
                                    ui.input_radio_buttons("ptmidmetrics_individual_average","Choose how to show ID counts:",choices={"individual":"Individual Counts","average":"Average Counts"}),
                                    ui.input_radio_buttons("ptm_counts_vs_enrich","Show counts or % of IDs?",choices={"counts":"Counts","percent":"% of IDs (enrichment)"})
                                ),
                                ui.input_selectize("ptmidplotinput","Choose what metric to plot:",choices={"all":"All","proteins":"Proteins","proteins2pepts":"Proteins with >2 Peptides","peptides":"Peptides","precursors":"Precursors"},multiple=False,selected="all"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("ptmidmetricsplot")
                ),
                ui.nav_panel("CV Plots",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("ptmcvplot_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Violin plots of the CVs per condition for the selected type of ID containing the selected PTM. Use the switch to remove top 5% outliers. The table displays the mean and median CV values for the shown plot and updates when the switch is toggled."),
                                ui.p("Notes:"),
                                ui.p("-Peptide-level CVs are calculated from the top 3 precursors for a given sequence."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("ptmcvplot_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("ptmcvplot_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("ptm_proteins_precursors","Pick which IDs to plot",choices={"Protein":"Proteins","Precursor":"Precursors","Peptide":"Peptides"}),
                                ui.column(3,
                                    ui.input_switch("ptm_removetop5percent","Remove top 5%"),
                                    ui.input_switch("ptm_cvplot_histogram_bins_switch","Show as histogram"),
                                    ui.output_ui("ptm_cvplot_histogram_bins_ui")
                                )
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_table("ptm_cvtable"),
                    ui.output_plot("ptm_cvplot")
                ),
                ui.nav_panel("Mass Accuracy",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("ptm_massaccuracy_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Violin plot or histogram of ppm mass accuracy for each precursor containing the selected PTM."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("ptm_massaccuracy_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("ptm_massaccuracy_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("ptm_massaccuracy_violin_hist","Plot as a violin plot or histogram?",choices={"violin":"Violin Plot","histogram":"Histogram"}),
                                ui.output_ui("ptm_massaccuracy_bins_ui")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("ptm_massaccuracy_plot")
                ),
                ui.nav_panel("PTMs per Precursor",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("ptmsperprecursor_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar graphs for the number of PTMs in each ModifiedPeptide value. Agnostic to PTM identity, a more detailed list of the PTM combinations can be exported in the Export Tables section."),
                                ui.p("-For example, the '1' on the x-axis gives the number of all precursors that contain a single PTM, regardless of the PTM identity."),
                                ui.p("-When plotting for a specific PTM, the '0' on the x-axis includes peptides that may be modified, just not with the selected PTM."),
                                ui.p("Notes:"),
                                ui.p("-For search results with many conditions, the bar width slider helps to avoid overlap of bar groups."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("ptmsperprecursor_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("ptmsperprecursor_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.input_slider("barwidth","Bar width:",min=0.1,max=1,step=0.05,value=0.25,ticks=True),
                                ui.column(3,
                                    ui.input_switch("ptmsperprecursor_specific","Plot for specific PTM"),
                                    ui.output_ui("ptmsperprecursor_ptmlist_ui"),
                                )
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("ptmsperprecursor")
                ),
                ui.nav_panel("Unique PTM Sites",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("ptmsites_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar graphs of the number of unique modification sites for each PTM in each run."),
                                ui.p("-The switch can be used to toggle whether PTMs from the entire search file and every protein are shown or just for a selected protein from the table."),
                                ui.p("-The table on the left shows identified proteins from the search file and the number of unique residues that were modified. The same residues modified by different reagents are counted as unique sites."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("ptmsites_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("ptmsites_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.input_slider("ptmsites_barwidth","Bar width:",min=0.1,max=1,step=0.05,value=0.15,ticks=True),
                                ui.input_radio_buttons("ptmsites_individual_average","Choose how to show ID counts:",choices={"individual":"Individual Counts","average":"Average Counts"}),
                                ui.input_switch("ptmsites_pickprotein","Plot for specific protein",width="250px")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_data_frame("proteinptms_table")
                        ),
                        ui.column(7,
                            ui.output_plot("ptmsites_plot")
                        )
                    )
                ),
                ui.nav_panel("Per-Residue PTM Sites",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("multiresidueptms_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Stacked bar graphs of the number of unique residues that were modified by each PTM in each run, specifically for PTMs that target multiple different residues."),
                                ui.p("-The switch can be used to toggle whether PTMs from the entire search file and every protein are shown or just for a selected protein from the table."),
                                ui.p("-The table on the left shows identified proteins from the search file and the number of unique residues that were modified. The same residues modified by different reagents are counted as unique sites."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("multiresidueptms_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("multiresidueptms_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.output_ui("multiresidueptms_ui"),
                                ui.input_switch("multiresidueptms_pickprotein","Plot for specific protein")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_data_frame("proteinptms_table2")
                        ),
                        ui.column(7,
                            ui.output_plot("multiresidueptms_plot")
                        )
                    )
                ),
            ),
            icon=icon_svg("magnifying-glass")
        ),
        ui.nav_panel("Heatmaps",
            ui.navset_pill(
                ui.nav_panel("RT, m/z, IM Heatmaps",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("heatmap_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-2-dimensional histograms of RT, m/z, and mobility using information from IDs in the search file(s)."),
                                ui.p("Notes:"),
                                ui.p("-The color bar corresponds to the number of IDs in a bin, not precursor intensity."),
                                ui.p("-RT vs. Intensity is a rough line plot approximation of the BPC, but does not account for peak widths since it is based on just the IDs in the search file(s)."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("heatmap_width","Plot width",min=100,max=7500,step=100,value=1400,ticks=True),
                                ui.input_slider("heatmap_height","Plot height",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_radio_buttons("conditiontype","Plot by individual replicate or by condition:",choices={"replicate":"By replicate","condition":"By condition"},width="350px"),
                                ui.output_ui("cond_rep_list_heatmap"),
                            ),
                            ui.row(
                                ui.input_slider("heatmap_numbins_x","Number of x bins:",min=10,max=250,value=100,step=10,ticks=True),
                                ui.input_slider("heatmap_numbins_y","Number of y bins:",min=10,max=250,value=100,step=10,ticks=True),
                                ui.input_selectize("heatmap_cmap","Heatmap Color:",choices={"default":"White_Blue_Red","viridis":"Viridis","plasma":"Plasma","inferno":"Inferno","magma":"Magma","cividis":"Cividis"}),
                            )
                        ),
                        width="100%"
                    ),
                    ui.output_plot("replicate_heatmap")
                ),
                ui.nav_panel("IM vs. m/z Precursor Heatmap",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("heatmap2_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-2-dimensional heatmap of m/z vs mobility using information from IDs in the search file(s)."),
                                ui.p("Notes:"),
                                ui.p("-The color bar corresponds to the number of IDs in a bin, not precursor intensity."),
                                ui.p("-DIA windows may be overlaid using built-in formats or by uploading custom DIA windows."),
                                ui.p("-Distinct charges and PTMs may be plotted for to assess their alignment to overlaid DIA windows."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("chargeptmheatmap_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("chargeptmheatmap_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("chargeptm_conditiontype","Plot by individual replicate or by condition:",choices={"replicate":"By replicate","condition":"By condition"},width="350px"),
                                ui.output_ui("chargeptm_cond_rep_ui"),
                            ),
                            ui.row(        
                                ui.input_slider("chargeptm_numbins_x","Number of m/z bins",min=10,max=250,value=100,step=10,ticks=True),
                                ui.input_slider("chargeptm_numbins_y","Number of mobility bins",min=10,max=250,value=100,step=10,ticks=True),
                                ui.input_selectize("chargeptmheatmap_cmap","Heatmap Color:",choices={"default":"White_Blue_Red","viridis":"Viridis","plasma":"Plasma","inferno":"Inferno","magma":"Magma","cividis":"Cividis"}),
                                ui.input_switch("chargeptmheatmap_axishistogram","Show histograms for each axis")
                            ),
                            ui.row(
                                ui.column(3,
                                    ui.download_button("diawindows_template","Download DIA Window Template",width="300px",icon=icon_svg("file-arrow-down")),
                                    ui.input_file("diawindow_upload","Upload DIA windows as a .csv:"),
                                ),
                                ui.input_selectize("windows_choice","Choose DIA windows to overlay:",choices={"None":"None","lubeck":"Lubeck DIA","phospho":"Phospho DIA","bremen":"Bremen DIA","imported":"Imported DIA windows","diagonal":"Imported Diagonal-PASEF windows"},selected="None"),
                                ui.output_ui("chargestates_chargeptmheatmap_ui"),
                                ui.output_ui("ptm_chargeptmheatmap_ui"),
                            )
                        ),
                        width="100%"
                    ),
                    ui.output_plot("chargeptmheatmap")
                ),
                ui.nav_panel("IM vs. m/z Precursor Scatterplot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("heatmap_scatter_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Scatter plot of m/z vs. ion mobility of identified precursors and a table of the number of precursors per charge state."),
                                ui.p("Notes:"),
                                ui.p("-Distinct charges and PTMs may be plotted as separately-colored data series in the scatter plot."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("chargeptmscatter_width","Plot width",min=100,max=7500,step=100,value=800,ticks=True),
                                ui.input_slider("chargeptmscatter_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.output_ui("chargeptmscatter_cond_rep"),
                                ui.output_ui("ptm_chargeptmscatter_ui"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(2,
                            ui.output_table("chargeptmscatter_table")
                        ),
                        ui.column(2,
                            ui.output_ui("chargestates_chargeptmscatter_ui"),
                        ),
                        ui.column(8,
                            ui.output_plot("chargeptmscatter")
                        )
                    ),
                ),
                ui.nav_panel("#IDs vs RT",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("idsvsrt_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Histogram of the ApexRT values for each run. An estimate of the maximum retention time is used to determine the timespan of each histogram bin."),
                                ui.p("Notes:"),
                                ui.p("-Changing the bin size changes the RT range over which IDs are grouped in the histogram."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("idsvsrt_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("idsvsrt_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.output_ui("binslider_ui"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_ui("ids_vs_rt_checkbox")
                        ),
                        ui.column(8,
                            ui.output_plot("ids_vs_rt")
                        )
                    )
                ),
                ui.nav_panel("Histograms",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                            #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("histogram_width","Plot width",min=100,max=7500,step=100,value=800,ticks=True),
                                ui.input_slider("histogram_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.input_slider("histogram_numbins","Number of bins:",min=10,max=250,value=100,step=10,ticks=True),
                                ui.input_selectize("histogram_pick","Pick property to plot:",choices={"ionmobility":"Ion Mobility","precursormz":"Precursor m/z","precursorintensity":"Precursor Intensity","proteinintensity":"Protein Intensity"}),
                            ),
                        ),
                    width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_ui("histogram_cond_rep_list"),
                        ),
                        ui.column(9,
                            ui.output_plot("histogram_plot")
                        )
                    )
                )
            ),
            icon=icon_svg("chart-area")
        ),
        ui.nav_panel("Statistics",
            ui.navset_pill(
                ui.nav_panel("Volcano Plot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("volcano_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Student's t-test assuming equal variances."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("volcano_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("volcano_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True)
                            ),
                            ui.row(
                                ui.column(4,
                                    ui.output_ui("volcano_condition1"),
                                    ui.output_ui("volcano_condition2"),
                                    ui.download_button("volcano_download","Download protein table",width="300px",icon=icon_svg("file-arrow-down")),
                                    ui.input_switch("show_labels","Show protein labels"),
                                    ui.input_numeric("label_fontsize","Label size:",value=4),
                                ),
                                ui.column(4,
                                    ui.input_slider("volcano_pvalue","log10 pvalue cutoff",min=0.5,max=5.0,value=1.0,step=0.1,ticks=True),
                                    ui.input_slider("volcano_foldchange","log2 fold change cutoff (absolute value)",min=0.1,max=2.0,value=0.5,step=0.1,ticks=True),
                                    ui.input_switch("volcano_h_v_lines","Show lines for pvalue and fold change cutoffs")
                                ),
                                ui.column(4,
                                    ui.input_slider("volcano_xplotrange","Plot x Range:",min=-10,max=10,value=[-2,2],step=0.5,ticks=True,drag_range=True),
                                    ui.input_slider("volcano_yplotrange","Plot y Range:",min=-10,max=10,value=[0,2],step=0.5,ticks=True,drag_range=True),
                                    ui.input_switch("volcano_plotrange_switch","Use sliders for axis ranges")
                                ),
                            )
                        ),
                        width="100%"
                    ),
                    ui.output_plot("volcanoplot")
                ),
                ui.nav_panel("Volcano Plot - Feature Plot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("volcanofeature_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Box-and-whisker plot of Protein Group intensity for selected proteins from the table with results based on the setup in the Volcano Plot tab."),
                                ui.p("Notes:"),
                                ui.p("-Control and Test conditions specified in the Volcano Plot tab."),
                                ui.p("-Select multiple rows in the data table by holding Control and clicking on rows."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("volcano_feature_width","Plot width",min=100,max=7500,step=100,value=800,ticks=True),
                                ui.input_slider("volcano_feature_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(6,
                            ui.output_data_frame("feature_table")
                        ),
                        ui.column(6,
                            ui.output_plot("feature_plot")                                              
                        )
                    )
                ),
                ui.nav_panel("Volcano Plot - Up/Down Regulation",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("volcanoregulation_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Bar graph of p-values or fold changes for up- or down-regulated proteins based on setup in the Volcano Plot tab."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("volcano_regulation_width","Plot width",min=100,max=7500,step=100,value=800,ticks=True),
                                ui.input_slider("volcano_regulation_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.input_selectize("regulation_upordown","Show up- or down-regulated proteins?",choices={"up":"Upregulated","down":"Downregulated"}),
                                ui.input_slider("regulation_topN","Pick top N proteins to show:",min=5,max=50,value=30,step=5,ticks=True),
                                ui.input_radio_buttons("regulation_p_fold","Metric to plot:",choices={"pvalue":"P-values","foldchange":"Fold Change"},width="200px")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("volcano_updownregulation_plot")
                ),
                ui.nav_panel("Correlation",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("correlation_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Scatterplots of correlations between the log10-transformed precursor intensities between replicates of the same condition."),
                                ui.p("-The middle diagonal contains histograms of the intensities for the given replicates."),
                                ui.p("-R^2 goodness-of-fit metric for each plot is listed in the axis directly opposite to it."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("correlations_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("correlations_height","Plot height",min=100,max=7500,step=100,value=800,ticks=True),
                                ui.output_ui("correlations_sampleconditions_ui")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("correlations_plot")
                ),
                ui.nav_panel("Dendrogram/Protein Signal",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("dendrogram_width","Plot width",min=100,max=7500,step=100,value=800,ticks=True),
                                ui.input_slider("dendrogram_height","Plot height",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_selectize("dendrogram_cmap","Heatmap Color:",choices={"viridis":"Viridis","plasma":"Plasma","inferno":"Inferno","magma":"Magma","cividis":"Cividis"},selected="magma"),
                                ui.input_numeric("dendrogram_scaling","Subplot height ratio (useful when changing plot height):",value=3),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("dendrogram_heatmap")
                ),
                ui.nav_panel("PCA",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("pca_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("pca_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_switch("pca_3d_switch","Plot 3rd PC on z-axis"),
                                ui.output_ui("pca_3d_ui"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("pca_plot")
                )
            ),
            icon=icon_svg("network-wired")
        ),
        ui.nav_panel("Immunopeptidomics",
            ui.navset_pill(
                ui.nav_panel("Sequence Motifs",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("sequencemotifs_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Weblogo plot for specific peptide length from selected run."),
                                ui.p("Notes:"),
                                ui.p("-'Information' matrix is the standard Weblogo style. Counts will scale the figure to the number of counts of each residue at each position."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("seqmotif_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("seqmotif_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("seqmotif_plottype","Pick kind of matrix to plot:",choices={"information":"Information","counts":"Counts"}),
                                ui.output_ui("seqmotif_run_ui"),
                                ui.input_slider("seqmotif_peplengths","Pick peptide length to plot:",min=7,max=25,value=9,step=1,ticks=True),
                            ),
                        ),
                        width="100%"
                    ),                                      
                    ui.output_plot("seqmotif_plot")
                ),
                ui.nav_panel("Charge States (Bar)",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("chargestatesbar_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Similar charge state plot as in Metrics section, but instead of separating precursors that were detected with multiple charges, charges are grouped together to show IDs that were detected with multiple charges."),
                                ui.p("Notes:"),
                                ui.p("-Use the checklist to choose which charges specifically to plot in the bar graph and use the switch at the top to activate this feature."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("charge_barchart_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("charge_barchart_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("chargestate_bar_condition_or_run","Plot by condition or by individual run?",choices={"condition":"Condition","individual":"Individual Run"})
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(2,
                            ui.input_switch("chargestate_charges_usepickedcharges","Use picked charges"),
                            ui.output_ui("chargestate_charges_ui")
                        ),
                        ui.column(9,
                            ui.output_plot("charge_barchart")
                        )
                    )
                ),
                ui.nav_panel("Charge States (Stacked)",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("chargestatesstacked_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Stacked representation of the charge state plot in Charge States (bar) with a table showing frequencies of different charges or charge combinations."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("charge_stackedbarchart_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("charge_stackedbarchart_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("chargestate_stacked_condition_or_run","Plot by condition or by individual run?",choices={"condition":"Condition","individual":"Individual Run"}),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_data_frame("charge_stacked_table"),
                    ui.output_plot("charge_stacked_barchart")
                ),
                ui.nav_panel("Charge States per Peptide Length",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("chargestatespeplength_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Stacked bar graph of charge states detected per peptide length. Precursors found in multiple runs per condition are only counted once."),
                                ui.p("Notes:"),
                                ui.p("-Use the slider to adjust the peptide length range over which the bar graph will be plotted."),
                                ui.p("-Use the checklist to choose which charges specifically to plot in the bar graph and use the switch at the top to activate this feature."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("chargestate_peplength_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("chargestate_peplength_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("chargestate_peplength_condition_or_run","Plot by condition or by individual run?",choices={"condition":"Condition","individual":"Individual Run"}),
                                ui.output_ui("chargestate_peplength_plotrange_ui"),
                                ui.column(2,
                                    ui.output_ui("chargestate_peplength_download_ui"),
                                    ui.download_button("chargestate_peplength_download","Download Counts for Selected Run",width="300px",icon=icon_svg("file-arrow-down"))
                                )
                            )
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(2,
                            ui.input_switch("usepickedcharges","Use picked charges"),
                            ui.output_ui("chargestate_peplength_charges_ui")
                        ),
                        ui.column(9,
                            ui.output_plot("chargestate_peplength_plot")
                        ),
                    )
                ),
            ),
            icon=icon_svg("vial-virus")
        ),
        ui.nav_panel("Mixed Proteome",
            ui.navset_pill(
                ui.nav_panel("Info",
                    ui.popover(
                        ui.input_action_button("mixedproteomesetup_btn","Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                        ui.p("Description:"),
                        ui.p("-Setup for mixed proteome calculations."),
                        ui.p("Notes:"),
                        ui.p("-Organisms are automatically detected from the PG.ProteinNames column in the search report."),
                        ui.p("-Edit the table with the desired plotting order for the organisms and their ratios in comparative sample conditions."),
                        ui.p("-Add an 'x' to the Remove column to remove a line completely from being used in later plotting functions."),
                    ),
                    ui.output_data_frame("organismtable"),
                    ui.input_radio_buttons("coloroptions_sumint","Use matplotlib tableau colors or blues/grays?",choices={"matplot":"matplotlib tableau","bluegray":"blues/grays"},width="400px"),
                    ui.p("Organisms detected in search file:"),
                    ui.output_text_verbatim("organisms")
                ),
                ui.nav_panel("Counts per Organism",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("countsperorganism_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("countsperorganism_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                                ui.input_selectize("countsplotinput","Choose what metric to plot:",choices={"proteins":"Proteins","peptides":"Peptides","precursors":"Precursors"},multiple=False),
                                ui.input_slider("countsperorganism_barwidth","Bar width",min=0.1,max=1,step=0.05,value=0.25,ticks=True),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("countsperorganism")
                ),
                ui.nav_panel("Summed Intensities",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                                #),
                            ui.row(
                                ui.input_slider("summedintensities_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("summedintensities_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                                ui.input_selectize("summedintensities_pick","Choose what metric to plot:",choices={"protein":"Proteins","precursor":"Precursors"}),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("summedintensities")
                ),
                ui.nav_panel("Quant Ratios",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            ui.popover(
                                ui.input_action_button("mixedproteome_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                ui.p("Description:"),
                                ui.p("-Plots of number of proteins per organism used in the calculations (those that were found in test and reference conditions), quant ratios, and histogram of quant ratios."),
                                ui.p("Notes:"),
                                ui.p("-Experimental ratios will be shown as a dashed line and theoretical ratios (calculated from the Info tab) will be shown as a solid line"),
                                ui.p("-The table shows the calculated theoretical and experimental quant ratios according to the y-axis log scale specified."),
                                ui.p("-Sliders can be used to adjust the y-axis plot range and CV cutoff % and are activated by their respective switches."),
                                placement="right"
                            ),
                            ui.row(
                                ui.input_slider("quantratios_width","Plot width",min=100,max=7500,step=100,value=1200,ticks=True),
                                ui.input_slider("quantratios_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.column(3,
                                    ui.input_slider("plotrange","Plot Range:",min=-10,max=10,value=[-2,2],step=0.5,ticks=True,width="300px",drag_range=True),
                                    ui.input_switch("plotrange_switch","Use slider for y-axis range?")
                                ),
                                ui.column(3,
                                    ui.input_slider("cvcutofflevel","CV Cutoff Level (%):",min=10,max=50,value=20,step=10,ticks=True,width="300px"),
                                    ui.input_switch("cvcutoff_switch","Include CV cutoff?")
                                ),
                            ),
                            ui.row(
                                ui.column(3,
                                    ui.output_ui("referencecondition"),
                                    ui.output_ui("testcondition"),
                                ),
                                ui.column(3,
                                    ui.input_selectize("quantratios_IDpick","Choose what metric to plot:",choices={"protein":"Proteins","precursor":"Precursors"}),
                                    ui.input_radio_buttons("quantratios_mean_median","Plot using mean or median quant?",choices={"mean":"mean","median":"median"}),
                                ),
                                ui.column(3,
                                    ui.input_switch("mixedproteome_showexperimentalratios","Show experimental ratios (dashed line)",value=True),
                                    ui.input_switch("mixedproteome_showtheoreticalratios","Show theoretical ratios (solid line)",value=True)
                                ),
                                ui.column(2,
                                    ui.input_radio_buttons("x_log_scale","x-axis Log Scale:",choices=["log2","log10"]),
                                    ui.input_radio_buttons("y_log_scale","y-axis Log Scale:",choices=["log2","log10"]), 
                                )
                            )
                        ),
                        width="100%"
                    ),
                    ui.output_table("quantratios_table"),
                    ui.output_plot("quantratios")
                )
            ),
            icon=icon_svg("flask")
        ),
        ui.nav_panel("PRM",
            ui.navset_pill(
                ui.nav_panel("PRM List",
                    ui.download_button("prm_template","Download PRM Template",width="300px",icon=icon_svg("file-arrow-down")),
                    ui.input_file("prm_list","Upload Peptide List:")
                ),
                ui.nav_panel("PRM Table",
                    ui.card(
                        ui.row(
                            ui.input_text("isolationwidth_input","m/z isolation width:"),
                            ui.input_text("rtwindow_input","Retention time window (s):"),
                            ui.input_text("imwindow_input","Ion mobility window (1/k0):")
                        )
                    ),
                    ui.download_button("prm_table_download","Download PRM Table",width="300px",icon=icon_svg("file-arrow-down")),
                    ui.output_data_frame("prm_table")
                ),
                ui.nav_panel("PRM Peptides - Individual Tracker",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("prmpeptracker_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("prmpeptracker_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                                ui.output_ui("prmpeptracker_pick"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("prmpeptracker_plot")
                ),
                ui.nav_panel("PRM Peptides - Intensity Across Runs",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("prmpepintensity_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("prmpepintensity_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True)
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("prmpepintensity_plot"),
                ),
            ),
            icon=icon_svg("crosshairs")
        ),
        ui.nav_panel("Dilution Series",
            ui.navset_pill(
                ui.nav_panel("Dilution Ratios",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("dilutionseries_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("dilutionseries_height","Plot height",min=100,max=7500,step=100,value=700,ticks=True),
                                ui.output_ui("normalizingcondition"),
                            ),
                        ),
                        width="100%"
                    ),                                     
                    ui.output_plot("dilutionseries_plot")
                )
            ),
            icon=icon_svg("vials")
        ),
        ui.nav_panel("Glycoproteomics",
            ui.navset_pill(
                ui.nav_panel("Glyco ID Metrics",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("glycoidcounts_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("glycoIDsplot_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("glycoIDsplot_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True)
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("glycoIDsplot")
                ),
                ui.nav_panel("Venn Diagram",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("glyco_venn_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("glyco_venn_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.input_radio_buttons("glyco_venn_conditionorrun","Plot by condition or individual run?",choices={"condition":"Condition","individual":"Individual Run"}),
                                ui.input_radio_buttons("glyco_venn_numcircles","Pick number of runs to compare:",choices={"2":"2","3":"3"}),
                            ),
                            ui.row(
                                ui.column(2,
                                    ui.output_ui("glyco_venn_run1_ui"),
                                    ui.output_ui("glyco_venn_run2_ui"),
                                    ui.output_ui("glyco_venn_run3_ui"),
                                ),
                                ui.column(2,
                                    ui.input_radio_buttons("glyco_venn_plotproperty","Metric to compare:",choices={"glycoproteins":"Glycoproteins","glycopeptides":"Glycopeptides","glycoPSMs":"GlycoPSMs"}),
                                ),
                                ui.column(2,
                                    ui.download_button("glyco_venn_download","Download Venn list",width="300px",icon=icon_svg("file-arrow-down"))
                                )
                            )
                        ),
                        width="100%"
                    ),
                    ui.output_plot("glyco_venn_plot")
                ),
                ui.nav_panel("Glyco ID Tables",
                    ui.row(
                        ui.column(3,
                            ui.download_button("glycoproteins_download","Download Glycoprotein IDs",width="300px",icon=icon_svg("file-arrow-down")),
                            ui.output_data_frame("glycoproteins_df_view")
                        ),
                        ui.column(4,
                            ui.download_button("glycopeptides_download","Download Glycopeptide IDs",width="300px",icon=icon_svg("file-arrow-down")),
                            ui.output_data_frame("glycopeptides_df_view")
                        ),
                        ui.column(4,
                            ui.download_button("glycoPSMs_download","Download GlycoPSM IDs",width="300px",icon=icon_svg("file-arrow-down")),
                            ui.output_data_frame("glycoPSMs_df_view")
                        ),
                    ),
                ),
                ui.nav_panel("Glycan Tracker",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("glycomodIDsplot_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("glycomodIDsplot_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.output_ui("glycomodlist_ui"),
                                ui.input_radio_buttons("counts_vs_enrich","Show counts or % of IDs?",choices={"counts":"Counts","percent":"% of IDs (enrichment)"}),
                                ui.output_ui("high_mannose_ui")
                            ),
                        ),
                    width="100%"
                    ),
                    ui.output_plot("glycomod_IDs")
                ),
                ui.nav_panel("Peptide Tracker",
                    ui.row(
                        ui.column(6,
                            ui.download_button("selected_glyco_download","Download PSMs for Selected Peptide",width="400px",icon=icon_svg("file-arrow-down")),
                            ui.output_ui("glyco_peplist")
                        )
                    ),
                    ui.output_data_frame("selected_glyco_peplist")
                ),
                ui.nav_panel("Precursor Scatterplot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("glycoscatter_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("glycoscatter_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.output_ui("glycoscatter_ui"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("glycoscatter")
                ),
            ),
            icon=icon_svg("cubes-stacked")
        ),
        ui.nav_panel("MOMA",
            ui.navset_pill(
                ui.nav_panel("MOMA Extraction",
                    ui.card(
                        ui.popover(
                            ui.input_action_button("momainfo_btn","Instructions",class_="btn-success",icon=icon_svg("question"),width="300px"),
                            ui.p("Description:"),
                            ui.p("-MOMA (mobility offset mass aligned) precursors from the search file are found based on m/z, retention time, and ion mobility tolerances set by the sliders."),
                            ui.p("-The possible MOMA events shown in the table can be used to generate extracted ion mobiligrams from user-uploaded raw data corresponding to the runs in the search file."),
                            ui.p("Notes:"),
                            ui.p("-The m/z and retention time sliders are used as tolerances in both finding MOMA events and plotting EIMs."),
                            ui.p("-The 'Group' column is just a placeholder to mark MOMA groups."),
                            placement="right"
                        ),
                        ui.row(
                            ui.input_slider("moma_mztolerance","m/z tolerance (m/z):",min=0,max=0.1,value=0.005,step=0.001,ticks=True),
                            ui.input_slider("moma_rttolerance","Retention time tolerance (s):",min=0.1,max=2,value=0.5,step=0.1,ticks=True),
                        ),
                    ),
                    ui.row(
                        #parameters related to search file
                        ui.column(6,
                            ui.card(
                                ui.card_header("From Search Result File(s)"),
                                ui.output_ui("moma_cond_rep_list"),
                                ui.input_slider("moma_imtolerance","Ion mobility tolerance (1/K0):",min=0.01,max=0.1,value=0.05,step=0.005,ticks=True),
                                ui.p("Table of Possible MOMA Events:"),
                                ui.output_data_frame("moma_events"),
                            )
                        ),
                        #parameters related to raw file
                        ui.column(6,
                            ui.card(
                                ui.card_header("From Raw File(s)"),
                                ui.input_radio_buttons("moma_file_or_folder","Load raw data from:",choices={"individual":"Individual Files","directory":"Directory"}),
                                ui.output_ui("moma_rawfile_input_ui"),
                                ui.output_ui("moma_rawfile_buttons_ui"),
                                ui.row(
                                    ui.column(6,
                                        ui.input_action_button("moma_load_rawfile","Load Raw File",class_="btn-primary",width="300px"),
                                    ),
                                    ui.column(6,
                                        ui.output_text_verbatim("file_uploaded"),
                                    )
                                ),
                                ui.row(
                                    ui.column(6,
                                        ui.input_text("moma_mz","EIM m/z:"),
                                    ),
                                    ui.column(6,
                                        ui.input_text("moma_rt","EIM Retention Time:"),
                                    )
                                ),
                                ui.output_plot("moma_eim")
                            )
                        )
                    ),
                    ui.download_button("momatable_download","Download MOMA Table",width="300px",icon=icon_svg("file-arrow-down")),
                    ui.download_button("precursortable_download","Download Precursor Table",width="300px",icon=icon_svg("file-arrow-down"))
                )
            ),
            icon=icon_svg("user-astronaut")
        ),
        ui.nav_panel("De Novo",
            ui.navset_pill(
                ui.nav_panel("IDs Found in Fasta",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("fasta_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("fasta_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True)
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("fasta_plot")
                ),
                ui.nav_panel("Position Confidence",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("confidence_width","Plot width",min=100,max=7500,step=100,value=800,ticks=True),
                                ui.input_slider("confidence_height","Plot height",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.output_ui("confidence_condition_ui"),
                                ui.input_slider("confidence_lengthslider","Pick peptide length to plot for:",min=7,max=20,step=1,value=9,ticks=True)
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("confidence_plot")
                ),
                ui.nav_panel("Secondary File Import",
                    ui.row(
                        ui.column(7,
                            ui.card(
                                ui.card_header("Upload Search Report"),
                                ui.row(
                                    ui.column(6,
                                        ui.popover(
                                            ui.input_action_button("denovo_import_btn","Instructions",class_="btn-success",icon=icon_svg("question")),
                                            ui.p("-Use main File Import tab for the BPS Novor data, upload data for the software to compare it to in this tab"),
                                            ui.p("-Make sure that the condition names and replicate numbers are the same between the two metadata sheets"),
                                            placement="right"
                                        ),
                                        ui.input_radio_buttons("software_secondary_general","Software used for search:",{"spectronaut":"Spectronaut",
                                                                                                                        "diann":"DIA-NN",
                                                                                                                        "fragpipe":"FragPipe",
                                                                                                                        "bps":"Bruker ProteoScape/GlycoScape",
                                                                                                                        "spectromine":"Spectromine",
                                                                                                                        "peaks":"PEAKS",
                                                                                                                        "sage":"Sage"}),
                                        ui.hr(),
                                        ui.output_ui("software_secondary_ui"),
                                        ui.hr(),
                                        ui.panel_conditional("input.software_secondary === 'bps_timsrescore' || input.software_secondary === 'bps_timsdiann' || input.software_secondary === 'bps_spectronaut' || input.software_secondary === 'bps_pulsar'",
                                                            ui.input_radio_buttons("software_secondary_bps_report_type","",choices={"qual":"Qualitative","quant":"Quantitative"})
                                        )
                                    ),
                                    ui.column(6,
                                        ui.input_file("searchreport_secondary","Upload search report:",accept=[".tsv",".zip",".parquet",".csv"],multiple=True),
                                        ui.input_switch("searchreport_secondary_reupload","Reupload of processed search file?"),
                                        ui.output_text("metadata_reminder_secondary")
                                    )
                                )
                            )
                        ),
                        ui.column(5,
                            ui.card(
                                ui.card_header("Update from Metadata Table"),
                                ui.row(
                                    ui.column(6,
                                        ui.input_action_button("rerun_metadata_secondary","Apply Metadata to Search File",width="300px",class_="btn-primary",icon=icon_svg("rotate")),
                                    ),
                                    ui.column(1),
                                    ui.column(1,
                                        ui.output_ui("denovo_metadata_applied_ui")),
                                    ),
                                    ui.row(
                                        ui.column(7,
                                            ui.input_switch("remove_secondary","Remove selected runs"),
                                            ui.input_switch("cRAP_filter_secondary","Filter protein list with cRAP database"),
                                        ),
                                    )
                                )
                        )
                    ),
                    ui.card(
                        ui.card_header("Metadata Tables"),
                        ui.row(
                            ui.column(4,
                                ui.input_file("metadata_upload_secondary","(Optional) Upload filled metadata table:",accept=".csv",multiple=False),
                                ui.input_switch("use_uploaded_metadata_secondary","Use uploaded metadata table"),
                            ),
                            ui.column(4,
                                ui.download_button("metadata_download_secondary","Download metadata table as shown",width="300px",icon=icon_svg("file-arrow-down"))
                            ),
                        ),
                        ui.row(
                            ui.column(8,
                                ui.output_data_frame("metadata_table_secondary")
                            ),
                            ui.column(4,
                                ui.output_data_frame("metadata_condition_table_secondary")
                            )
                        )
                    ),
                ),
                ui.nav_panel("Compare - Peptide Lengths",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("peplength_compare_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("peplength_compare_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.output_ui("compare_len_samplelist"),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("peplength_compare_plot")
                ),
                ui.nav_panel("Compare - Stripped Peptide IDs",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("denovocompare_venn_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("denovocompare_venn_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.output_ui("denovocompare_venn_samplelist"),
                                ui.column(3,
                                    ui.input_switch("denovocompare_specific_length","Compare specific peptide length?",value=False,width="300px"),
                                    ui.output_ui("denovocompare_specific_length_ui"),
                                ),
                            ui.input_switch("denovocompare_peptidecore","Only consider peptide core (cut first and last 2 residues)",value=False,width="300px"),
                            ),
                            ui.row(
                                ui.download_button("denovocompare_venn_download","Download Peptide List",width="300px",icon=icon_svg("file-arrow-down")),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("denovocompare_venn_plot")
                ),
                ui.nav_panel("Compare - Sequence Motifs",
                    ui.row(
                        ui.column(3,
                            ui.input_radio_buttons("seqmotif_compare_plottype","Pick kind of matrix to plot:",choices={"information":"Information","counts":"Counts"}),
                            ui.input_switch("seqmotif_compare_onlyunique","Use only unique BPS Novor IDs",value=False),
                            ui.input_slider("seqmotif_compare_peplengths","Pick peptide length to plot:",min=7,max=25,value=9,step=1,ticks=True),
                            ui.output_ui("seqmotif_compare_run_ui"),
                        ),
                        ui.column(8,
                            ui.output_plot("seqmotif_compare_plot1"),
                            ui.output_plot("seqmotif_compare_plot2")
                        )
                    )
                ),
                ui.nav_panel("Compare - Sequence Motifs (Stats)",
                    ui.row(
                        ui.column(6,
                            ui.input_switch("seqmotif_compare_onlyunique2","Use only unique BPS Novor IDs",value=False),
                            ui.input_slider("seqmotif_compare_peplengths2","Pick peptide length to plot:",min=7,max=25,value=9,step=1,ticks=True),
                            ui.output_ui("seqmotif_compare_run_ui2"),
                            ui.output_plot("seqmotif_pca")
                        ),
                        ui.column(6,
                            ui.input_slider("seqmotif_3d_azimuth","3D Plot - Azimuth Angle",min=0,max=360,value=0,step=5,ticks=True),
                            ui.input_slider("seqmotif_3d_elevation","3D Plot - Elevation Angle",min=0,max=90,value=0,step=5,ticks=True),
                            ui.output_plot("seqmotif_3d")
                        )
                    )
                ),
            ),
            icon=icon_svg("atom")
        ),
        ui.nav_panel("Two-Software Comparison",
            ui.navset_pill(
                ui.nav_panel("File Import",
                    ui.row(
                        ui.column(8,
                            ui.card(
                                ui.card_header("Upload Search Reports"),
                                ui.row(
                                    ui.column(6,
                                        ui.input_file("compare_searchreport1","Upload first search report:",accept=[".tsv",".zip",".parquet",".csv"],multiple=True),
                                        ui.input_switch("searchreport1_reupload","Reupload of processed search file?"),
                                        ui.input_radio_buttons("software1_general","First search software:",{"spectronaut":"Spectronaut",
                                                                                                            "diann":"DIA-NN",
                                                                                                            "fragpipe":"FragPipe",
                                                                                                            "bps":"Bruker ProteoScape/GlycoScape",
                                                                                                            "spectromine":"Spectromine",
                                                                                                            "peaks":"PEAKS",
                                                                                                            "sage":"Sage"
                                                                                                            }),
                                        ui.hr(),
                                        ui.output_ui("software1_ui"),
                                        ui.hr(),
                                        ui.panel_conditional("input.software1 === 'bps_timsrescore' || input.software1 === 'bps_timsdiann' || input.software1 === 'bps_spectronaut' || input.software1 === 'bps_pulsar'",
                                                            ui.input_radio_buttons("software1_bps_report_type","",choices={"qual":"Qualitative","quant":"Quantitative"})
                                        )
                                    ),
                                    ui.column(6,
                                        ui.input_file("compare_searchreport2","Upload second search report:",accept=[".tsv",".zip",".parquet",".csv"],multiple=True),
                                        ui.input_switch("searchreport2_reupload","Reupload of processed search file?"),
                                        ui.input_radio_buttons("software2_general","Second search software:",{"spectronaut":"Spectronaut",
                                                                                                            "diann":"DIA-NN",
                                                                                                            "fragpipe":"FragPipe",
                                                                                                            "bps":"Bruker ProteoScape/GlycoScape",
                                                                                                            "spectromine":"Spectromine",
                                                                                                            "peaks":"PEAKS",
                                                                                                            "sage":"Sage"
                                                                                                            }),
                                        ui.hr(),
                                        ui.output_ui("software2_ui"),
                                        ui.hr(),
                                        ui.panel_conditional("input.software2 === 'bps_timsrescore' || input.software2 === 'bps_timsdiann' || input.software2 === 'bps_spectronaut' || input.software2 === 'bps_pulsar'",
                                                            ui.input_radio_buttons("software2_bps_report_type","",choices={"qual":"Qualitative","quant":"Quantitative"})
                                        )
                                    )
                                )
                            )
                        ),
                        ui.column(4,
                            ui.card(
                                ui.card_header("Update from Metadata Table"),
                                ui.row(
                                    ui.column(9,
                                        ui.input_action_button("compare_rerun_metadata","Apply Metadata to Search File",width="300px",class_="btn-primary",icon=icon_svg("rotate")),
                                    ),
                                    ui.column(1,
                                        ui.output_ui("twosoftware_metadata_applied_ui")
                                    ),
                                    ),
                                ui.row(
                                    ui.column(10,
                                        ui.input_switch("compare_remove","Remove selected runs"),
                                        ui.input_switch("compare_cRAP_filter","Filter protein list with cRAP database"),
                                    )
                                )
                            ),
                            ui.card(
                                ui.p("Color Options (see Color Settings tab for naming)"),
                                ui.output_data_frame("compare_colortable"),
                                ui.input_text_area("twosoftware_download_path","Enter path for figure downloads (optional)",width="100%",height="65px",placeholder="ex. - C:\\Users\\Data")
                            )
                        )
                    ),
                    ui.card(
                        ui.card_header("Metadata Tables"),
                        ui.row(
                            ui.column(4,
                                ui.input_file("compare_metadata_upload","(Optional) Upload filled metadata table:",accept=".csv",multiple=False),
                                ui.input_switch("compare_use_uploaded_metadata","Use uploaded metadata table"),
                            ),
                            ui.column(4,
                                ui.download_button("compare_metadata_download","Download metadata table as shown",width="300px",icon=icon_svg("file-arrow-down"))
                            ),
                        ),
                        ui.row(
                            ui.column(9,
                                ui.output_data_frame("compare_metadata_table"),
                            ),
                            ui.column(3,
                                ui.output_data_frame("compare_metadata_condition_table")
                            )
                        ),
                    ),
                ),
                ui.nav_panel("ID Counts",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("compare_id_counts_width","Plot width",min=100,max=7500,step=100,value=1500,ticks=True),
                                ui.input_slider("compare_id_counts_height","Plot height",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_radio_buttons("compare_idmetrics_individual_average","Choose how to show ID counts:",choices={"individual":"Individual Counts","average":"Average Counts"},width="250px"),
                                ui.input_radio_buttons("compare_idmetrics_peptides","Pick peptide counts to plot:",choices={"modified":"Modified Sequences","stripped":"Stripped Sequences"}),
                                ui.output_ui("compare_idmetrics_ptm_ui")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.output_plot("compare_id_counts")
                ),
                ui.nav_panel("Venn Diagram",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("compare_venn_width","Plot width",min=100,max=7500,step=100,value=1000,ticks=True),
                                ui.input_slider("compare_venn_height","Plot height",min=100,max=7500,step=100,value=500,ticks=True),
                                ui.column(3,
                                    ui.input_radio_buttons("compare_venn_conditionorrun","Plot by condition or individual run?",choices={"condition":"Condition","individual":"Individual Run"}),
                                    ui.output_ui("compare_venn_run_ui")
                                ),
                                ui.column(3,
                                    ui.input_radio_buttons("compare_venn_plotproperty","Metric to compare:",choices={"proteingroups":"Protein Groups","peptides":"Peptides","precursors":"Precursors","peptides_stripped":"Stripped Peptides",}),
                                    ui.download_button("compare_venn_download","Download Venn list",width="300px",icon=icon_svg("file-arrow-down"))
                                )
                            ),
                            ui.row(
                                ui.column(3,
                                    ui.output_ui("compare_venn_ptm_ui"),
                                    ui.output_ui("compare_venn_ptmlist_ui")
                                ),
                                ui.column(3,
                                    ui.output_ui("compare_venn_specific_length_ui"),
                                    ui.output_ui("compare_venn_peplength_ui")
                                )
                            )
                        ),
                        width="100%"
                    ),
                    ui.output_plot("compare_venn_plot")
                ),
            ),
            icon=icon_svg("code-compare")
        ),
        ui.nav_panel("Export Tables",
            ui.navset_pill(
                ui.nav_panel("Export Tables",
                    ui.row(
                        ui.column(3,
                            ui.card(
                                ui.card_header("Protein ID Metrics and CVs"),
                                ui.download_button("proteinidmetrics_download","Protein ID Metrics",width="300px",icon=icon_svg("file-arrow-down"))
                            ),
                            ui.card(
                                ui.card_header("Protein ID Matrix"),
                                ui.download_button("proteinidmatrix_download","Protein ID Matrix",width="300px",icon=icon_svg("file-arrow-down"))
                            )
                        ),
                        ui.column(3,
                            ui.card(
                                ui.card_header("Precursor ID Metrics and CVs"),
                                ui.download_button("precursoridmetrics_download","Precursor ID Metrics",width="300px",icon=icon_svg("file-arrow-down"))
                            )
                        ),
                        ui.column(3,
                            ui.card(
                                ui.card_header("Peptide ID Metrics and CVs"),
                                ui.download_button("peptidelist","Peptide ID Metrics",width="300px",icon=icon_svg("file-arrow-down")),
                            )
                        ),
                    ),
                    ui.row(
                        ui.column(3,
                            ui.card(
                                ui.card_header("List of PTMs per Precursor"),
                                ui.download_button("ptmlist_download","Download Precursor PTMs",width="300px",icon=icon_svg("file-arrow-down"))
                            )
                        ),
                        ui.column(3,
                            ui.card(
                                ui.card_header("Stripped Peptide IDs (Specific Length)"),
                                ui.download_button("strippedpeptidelist","Download Stripped Peptide IDs",width="300px",icon=icon_svg("file-arrow-down")),
                                ui.input_radio_buttons("peptidelist_condition_or_run","Condition or individual run?",choices={"condition":"condition","individual":"individual run"}),
                                ui.output_ui("peptidelist_dropdown"),
                                ui.input_slider("strippedpeptidelength","Pick specific peptide length to export:",min=7,max=25,value=9,step=1,ticks=True),
                            )
                        ),
                        ui.column(3,
                            ui.card(
                                ui.card_header("List of MOMA Precursors"),
                                ui.download_button("moma_download","Download MOMA List",width="300px",icon=icon_svg("file-arrow-down")),
                                ui.output_ui("cond_rep_list"),
                                ui.input_slider("mztolerance","m/z tolerance (m/z):",min=0.0005,max=0.1,value=0.005,step=0.0005,ticks=True),
                                ui.input_slider("imtolerance","Ion mobility tolerance (1/K0)",min=0.01,max=0.1,value=0.05,step=0.005,ticks=True),
                                ui.input_slider("rttolerance","Retention time tolerance (s):",min=0.1,max=2,value=0.5,step=0.1,ticks=True),
                            )
                        )
                    ),
                ),
                ui.nav_panel("Figures to PPT",
                    ui.p(),
                    ui.row(
                        ui.column(2,
                            ui.popover(
                                ui.input_action_button("figurestoppt_btn","Instructions",class_="btn-success",icon=icon_svg("question")),
                                ui.p("-Paste a path in File Import or Two-Software Comparison tab to automatically save generated figures from a search report as .png files. Image files from the inputted path will be split up into keywords for selection."),
                                ui.p("-Image files containing the selected keywords will fitted to a standard widescreen PowerPoint slide size (13.333 in. by 7.5 in.) keeping aspect ratio of how they were shown in timsplot and exported to a .pptx file."),
                                placement="right"
                            )
                        ),
                        ui.download_button("pptx_download","Download figures based on selected keywords",width="400px",icon=icon_svg("file-arrow-down"))
                    ),
                    ui.p(),
                    ui.row(
                        ui.column(6,
                            ui.card(
                                ui.card_header("Keywords"),
                                ui.output_ui("keys_ui")
                            )
                        ),
                        ui.column(4,
                            ui.card(
                                ui.card_header("Selected Images"),
                                ui.output_text("selected_images")
                            )
                        )
                    ),
                )
            ),
            icon=icon_svg("file-export")
        ),
        ui.nav_panel("Raw Data",
            ui.navset_pill(
                ui.nav_panel("Multi-File Import",
                    ui.popover(
                        ui.input_action_button("rawfileimport_btn","Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                        ui.p("-This section is independent of what has been uploaded to the File Import tab."),
                        ui.p("-Paste the directory containing .d files or the full path of an individual .d file to plot data for."),
                        ui.p("-For TIC, BPC, and Accumulation Time plots, the analysis.tdf file is loaded and a line plot is generated based on columns for the time and the corresponding intensity or accumulation time values."),
                        ui.p("-For EIC and EIM plots, the analysis.tdf_bin file is loaded for a specified file. m/z or IM tolerances adjust how the dataframe is sliced to generate the EIC/EIM."),
                        ui.p("-It should be noted that these plots are not smoothed, so peak widths will not be the same as what may be shown in DataAnalysis."),
                        placement="right"
                    ),
                    ui.p(),
                    ui.row(
                        ui.input_radio_buttons("file_or_folder","Load raw data from:",choices={"individual":"Individual Files","directory":"Directory"}),
                        ui.input_switch("rawfile_download_parent_path","Automatically download figures to parent directory (for individual files, figures will be saved to the directory containing the first raw file listed)",width="450px"),
                    ),
                    ui.output_ui("rawfile_input_ui"),
                    ui.output_text_verbatim("uploadedfiles")
                ),
                ui.nav_panel("TIC Plot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("tic_width","Plot width",min=100,max=7500,step=100,value=1500,ticks=True),
                                ui.input_slider("tic_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.input_switch("stacked_tic","Stack TIC Plots")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_ui("rawfile_checkboxes_tic")
                        ),
                        ui.column(8,
                            ui.output_plot("TIC_plot")
                        )
                    )
                ),
                ui.nav_panel("BPC Plot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("bpc_width","Plot width",min=100,max=7500,step=100,value=1500,ticks=True),
                                ui.input_slider("bpc_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.input_switch("stacked_bpc","Stack BPC Plots")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_ui("rawfile_checkboxes_bpc")
                        ),
                        ui.column(8,
                            ui.output_plot("BPC_plot")
                        )
                    )
                ),
                ui.nav_panel("Accumulation Time",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("accutime_width","Plot width",min=100,max=7500,step=100,value=1500,ticks=True),
                                ui.input_slider("accutime_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.input_switch("stacked_accutime","Stack BPC Plots")
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_ui("rawfile_checkboxes_accutime")
                        ),
                        ui.column(8,
                            ui.output_plot("accutime_plot")
                        )
                    )
                ),
                ui.nav_panel("EIC Plot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                                #),
                            ui.row(
                                ui.input_slider("eic_width","Plot width",min=100,max=7500,step=100,value=1500,ticks=True),
                                ui.input_slider("eic_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.column(3,
                                    ui.input_text("eic_mz_input","Input m/z for EIC:"),
                                    ui.input_text("eic_ppm_input","Input mass error (ppm) for EIC:")
                                ),
                                ui.column(3,
                                    ui.input_switch("include_mobility","Include mobility in EIC filtering"),
                                    ui.output_ui("mobility_input")
                                )
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_ui("rawfile_buttons_eic"),
                            ui.input_action_button("eic_load_rawfile","Load Raw File",width="300px",class_="btn-primary")
                        ),
                        ui.column(8,
                            ui.output_plot("eic")
                        )
                    )
                ),
                ui.nav_panel("EIM Plot",
                    ui.accordion(
                        ui.accordion_panel("Plot Options",
                            #ui.popover(
                                #ui.input_action_button("_btn","Plot Instructions",width="300px",class_="btn-success",icon=icon_svg("question")),
                                #placement="right"
                            #),
                            ui.row(
                                ui.input_slider("eim_width","Plot width",min=100,max=7500,step=100,value=1500,ticks=True),
                                ui.input_slider("eim_height","Plot height",min=100,max=7500,step=100,value=600,ticks=True),
                                ui.column(3,
                                    ui.input_text("eim_mz_input","Input m/z for EIM:"),
                                    ui.input_text("eim_ppm_input","Input mass error (ppm) for EIM:")
                                ),
                            ),
                        ),
                        width="100%"
                    ),
                    ui.row(
                        ui.column(3,
                            ui.output_ui("rawfile_buttons_eim"),
                            ui.input_action_button("eim_load_rawfile","Load Raw File",width="300px",class_="btn-primary")
                        ),
                        ui.column(8,
                            ui.output_plot("eim")
                        )
                    )
                ),
            ),
            icon=icon_svg("desktop")
        ),
        widths=(2,9)
    ),
    theme=theme.cerulean()
)
#endregion

# ============================================================================= Library Imports (all others needed for calculations)
#region
from collections import OrderedDict
import colorsys
from datetime import date
import io
import itertools
from itertools import groupby
import math
import matplotlib
import matplotlib.pyplot as plt
from matplotlib.pyplot import cm
import matplotlib.patches as mpatches
import matplotlib.colors as mcolors
from matplotlib.ticker import MaxNLocator,MultipleLocator
import numpy as np
import os
import pandas as pd
import re
import scipy
import scipy.stats as stats
from tkinter import *
from upsetplot import *
from zipfile import ZipFile

matplotlib.use('Agg')
#endregion

# =============================================================================
# Server
# =============================================================================
#region
def server(input: Inputs, output: Outputs, session: Session):

# ============================================================================= UI calls for use around the app
#region
    #render ui call for dropdown calling sample condition names
    @render.ui
    def sampleconditions_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=sampleconditions
        return ui.input_selectize("conditionname","Pick sample condition:",choices=opts)
    #render ui call for dropdown calling replicate number
    @render.ui
    def replicates_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=np.arange(1,max(maxreplicatelist)+1,1)
        return ui.input.selectize("replicate","Replicate number",opts)
    #render ui call for dropdown calling Cond_Rep column
    @render.ui
    def cond_rep_list():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_selectize("cond_rep","Pick run:",choices=opts)

#endregion

# ============================================================================= Metrics Functions
#region
    # ====================================== Callable metrics functions
    #use the variables_dfs function that imports the searchoutput df to calculate ID metrics
    def idmetrics(searchoutput):
        sampleconditions=searchoutput["R.Condition"].drop_duplicates().tolist()
        maxreplicatelist=searchoutput[["R.Condition","R.Replicate"]].drop_duplicates().groupby("R.Condition").size().reset_index()[0].tolist()
        runnames=searchoutput["Cond_Rep"].drop_duplicates().tolist()

        resultdf=pd.DataFrame(searchoutput[["Cond_Rep","R.FileName","R.Condition","R.Replicate"]].drop_duplicates()).reset_index(drop=True)
        averagedf=pd.DataFrame({"R.Condition":sampleconditions,"N.Replicates":maxreplicatelist})
        numproteins=[]
        numproteins2pepts=[]
        numpeptides=[]
        numstrippedpeptides=[]
        numprecursors=[]
        for i in runnames:
            replicatedata=searchoutput[searchoutput["Cond_Rep"]==i]
            if replicatedata.empty:
                continue
            #identified proteins
            numproteins.append(replicatedata["PG.ProteinGroups"].nunique())
            #identified proteins with 2 peptides
            numproteins2pepts.append(len(replicatedata[["PG.ProteinGroups","EG.ModifiedPeptide"]].drop_duplicates().groupby("PG.ProteinGroups").size().reset_index(name="peptides").query("peptides>1")))
            #identified peptides
            numpeptides.append(replicatedata["EG.ModifiedPeptide"].nunique())
            #identified stripped peptides
            numstrippedpeptides.append(replicatedata["PEP.StrippedSequence"].nunique())
            #identified precursors
            numprecursors.append(len(replicatedata[["EG.ModifiedPeptide","FG.Charge"]].drop_duplicates()))
        resultdf["proteins"]=numproteins
        resultdf["proteins2pepts"]=numproteins2pepts
        resultdf["peptides"]=numpeptides
        resultdf["strippedpeptides"]=numstrippedpeptides
        resultdf["precursors"]=numprecursors
        #avg and stdev values for IDs appended to averagedf dataframe, which holds lists of all the calculated values here
        for i,condition in enumerate(sampleconditions):
            df=resultdf[resultdf["R.Condition"]==condition]
            averagedf.at[i,"proteins_avg"]=round(df["proteins"].mean())
            averagedf.at[i,"proteins_stdev"]=df["proteins"].std()
            averagedf.at[i,"proteins2pepts_avg"]=round(df["proteins2pepts"].mean())
            averagedf.at[i,"proteins2pepts_stdev"]=df["proteins2pepts"].std()
            averagedf.at[i,"peptides_avg"]=round(df["peptides"].mean())
            averagedf.at[i,"peptides_stdev"]=df["peptides"].std()
            averagedf.at[i,"strippedpeptides_avg"]=round(df["strippedpeptides"].mean())
            averagedf.at[i,"strippedpeptides_stdev"]=df["strippedpeptides"].std()
            averagedf.at[i,"precursors_avg"]=round(df["precursors"].mean())
            averagedf.at[i,"precursors_stdev"]=df["precursors"].std()
        #unique counts per condition
        uniqueproteins=[]
        uniqueproteins2peptides=[]
        uniquepeptides=[]
        uniquestrippedpeptides=[]
        uniqueprecursors=[]
        for condition in sampleconditions:
            df=searchoutput[searchoutput["R.Condition"]==condition]
            if df.empty:
                continue
            uniqueproteins.append(df["PG.ProteinGroups"].nunique())
            uniqueproteins2peptides.append(len(df[["PG.ProteinGroups","EG.ModifiedPeptide"]].drop_duplicates().groupby("PG.ProteinGroups").size().reset_index(name="peptides").query("peptides>1")))
            uniquepeptides.append(df["EG.ModifiedPeptide"].nunique())
            uniquestrippedpeptides.append(df["PEP.StrippedSequence"].nunique())
            uniqueprecursors.append(len(df[["EG.ModifiedPeptide","FG.Charge"]].drop_duplicates()))
        uniquecountsdf=pd.DataFrame()
        uniquecountsdf["R.Condition"]=sampleconditions
        uniquecountsdf["proteins"]=uniqueproteins
        uniquecountsdf["proteins2pepts"]=uniqueproteins2peptides
        uniquecountsdf["protein_difference"]=uniquecountsdf["proteins"]-uniquecountsdf["proteins2pepts"]
        uniquecountsdf["peptides"]=uniquepeptides
        uniquecountsdf["strippedpeptides"]=uniquestrippedpeptides
        uniquecountsdf["precursors"]=uniqueprecursors
        return resultdf,averagedf,uniquecountsdf
    #function for doing ID calculations for a picked PTM
    def ptmcounts(searchoutput,ptm):
        sampleconditions=searchoutput["R.Condition"].drop_duplicates().tolist()
        maxreplicatelist=searchoutput[["R.Condition","R.Replicate"]].drop_duplicates().groupby("R.Condition").size().reset_index()[0].tolist()
        runnames=searchoutput["Cond_Rep"].drop_duplicates().tolist()

        ptmresultdf=pd.DataFrame(searchoutput[["Cond_Rep","R.FileName","R.Condition","R.Replicate"]].drop_duplicates()).reset_index(drop=True)
        numptmproteins_list=[]
        numptmproteins2pepts_list=[]
        numptmpeptides_list=[]
        numptmprecursors_list=[]
        enrichptmproteins_list=[]
        enrichptmproteins2pepts_list=[]
        enrichptmpeptides_list=[]
        enrichptmprecursors_list=[]
        for i in runnames:
            replicatedata=searchoutput[searchoutput["Cond_Rep"]==i]
            df=searchoutput[(searchoutput["Cond_Rep"]==i)&(searchoutput["EG.ModifiedPeptide"].str.contains(ptm,regex=False))][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge"]]
            if replicatedata.empty:
                continue
            #identified proteins
            numproteins=replicatedata["PG.ProteinGroups"].nunique()
            ptmproteins=df["PG.ProteinGroups"].nunique()
            enrichproteins=round((ptmproteins/numproteins)*100,1)
            #identified proteins with 2 peptides
            numproteins2pepts=len(replicatedata[["PG.ProteinGroups","EG.ModifiedPeptide"]].drop_duplicates().groupby("PG.ProteinGroups").size().reset_index(name="peptides").query("peptides>1"))
            ptmproteins2pepts=len(df[["PG.ProteinGroups","EG.ModifiedPeptide"]].drop_duplicates().groupby("PG.ProteinGroups").size().reset_index(name="peptides").query("peptides>1"))
            enrichproteins2pepts=round((ptmproteins2pepts/numproteins2pepts)*100,1)
            #identified peptides
            numpeptides=replicatedata["EG.ModifiedPeptide"].nunique()
            ptmpeptides=df["EG.ModifiedPeptide"].nunique()
            enrichpeptides=round((ptmpeptides/numpeptides)*100,1)
            #identified precursors
            numprecursors=len(replicatedata[["EG.ModifiedPeptide","FG.Charge"]].drop_duplicates())
            ptmprecursors=len(df[["EG.ModifiedPeptide","FG.Charge"]].drop_duplicates())
            enrichprecursors=round((ptmprecursors/numprecursors)*100,1)
            numptmproteins_list.append(ptmproteins)
            numptmproteins2pepts_list.append(ptmproteins2pepts)
            numptmpeptides_list.append(ptmpeptides)
            numptmprecursors_list.append(ptmprecursors)
            enrichptmproteins_list.append(enrichproteins)
            enrichptmproteins2pepts_list.append(enrichproteins2pepts)
            enrichptmpeptides_list.append(enrichpeptides)
            enrichptmprecursors_list.append(enrichprecursors)

        ptmresultdf["proteins"]=numptmproteins_list
        ptmresultdf["proteins2pepts"]=numptmproteins2pepts_list
        ptmresultdf["peptides"]=numptmpeptides_list
        ptmresultdf["precursors"]=numptmprecursors_list
        ptmresultdf["proteins_enrich%"]=enrichptmproteins_list
        ptmresultdf["proteins2pepts_enrich%"]=enrichptmproteins2pepts_list
        ptmresultdf["peptides_enrich%"]=enrichptmpeptides_list
        ptmresultdf["precursors_enrich%"]=enrichptmprecursors_list

        ptmaveragedf=pd.DataFrame({"R.Condition":sampleconditions,"N.Replicates":maxreplicatelist})
        #avg and stdev values for IDs appended to averagedf dataframe, which holds lists of all the calculated values here
        columnlist=ptmresultdf.columns.values.tolist()
        for i in columnlist:
            if i=="R.FileName" or i=="Cond_Rep" or i=="R.Condition" or i=="R.Replicate":
                continue
            avglist=[]
            stdevlist=[]
            for j in sampleconditions:
                samplecondition=ptmresultdf[ptmresultdf["R.Condition"]==j]
                avglist.append(round(np.average(samplecondition[i].to_numpy())))
                stdevlist.append(np.std(samplecondition[i].to_numpy()))
            ptmaveragedf[i+"_avg"]=avglist
            ptmaveragedf[i+"_stdev"]=stdevlist
        
        return ptmresultdf,ptmaveragedf
    #function for doing CV calculations, call ptm=None if not for PTM calculation
    def cvmetrics(searchoutput,ptm):
        if ptm is None:
            searchoutput=searchoutput
            if input.peptide_grouping()=="stripped":
                grouping_key="PEP.StrippedSequence"
            if input.peptide_grouping()=="modified":
                grouping_key="EG.ModifiedPeptide"
        else:
            searchoutput=searchoutput[searchoutput["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
            grouping_key="EG.ModifiedPeptide"

        sampleconditions=searchoutput["R.Condition"].drop_duplicates().tolist()
        maxreplicatelist=searchoutput[["R.Condition","R.Replicate"]].drop_duplicates().groupby("R.Condition").size().reset_index()[0].tolist()
        runnames=searchoutput["Cond_Rep"].drop_duplicates().tolist()

        cvcalc_df=pd.DataFrame()
        cvcalc_df["R.Condition"]=sampleconditions

        if "PG.MS2Quantity" in searchoutput.columns:
            #protein-level CVs
            proteincvlist=[]
            proteincvlist95=[]
            proteincvdict={}
            cvproteingroup=searchoutput[["R.Condition","R.Replicate","PG.ProteinGroups","PG.MS2Quantity"]].drop_duplicates().reset_index(drop=True)
            for x,condition in enumerate(sampleconditions):
                if maxreplicatelist[x]==1:
                    emptylist=[]
                    proteincvlist.append(emptylist)
                    proteincvlist95.append(emptylist)
                else:
                    df=pd.DataFrame(cvproteingroup[cvproteingroup["R.Condition"]==condition]).drop(columns=["R.Condition","R.Replicate"])
                    mean=df.groupby("PG.ProteinGroups").mean().rename(columns={"PG.MS2Quantity":"Protein Mean"})
                    std=df.groupby("PG.ProteinGroups").std().rename(columns={"PG.MS2Quantity":"Protein Std"})
                    cvproteintable=pd.concat([mean,std],axis=1)
                    cvproteintable.dropna(inplace=True)
                    cvproteintable=cvproteintable.loc[(cvproteintable!=0).any(axis=1)]
                    cvlist=(cvproteintable["Protein Std"]/cvproteintable["Protein Mean"]*100).tolist()
                    proteincvdict[condition]=pd.DataFrame(cvlist,columns=["CV"])
                    proteincvlist.append(cvlist)
                    top95=np.percentile(cvlist,95)
                    cvlist95=[]
                    for i in cvlist:
                        if i <=top95:
                            cvlist95.append(i)
                    proteincvlist95.append(cvlist95)
            cvcalc_df["Protein CVs"]=proteincvlist
            cvcalc_df["Protein 95% CVs"]=proteincvlist95

            #counts above CV cutoffs
            #protein CVs
            proteinscv20=[]
            proteinscv10=[]
            for x,condition in enumerate(sampleconditions):
                if maxreplicatelist[x]==1:
                    emptylist=[]
                    proteinscv20.append(emptylist)
                    proteinscv10.append(emptylist)
                else:
                    proteinscv20.append(proteincvdict[condition][proteincvdict[condition]["CV"]<20].shape[0])
                    proteinscv10.append(proteincvdict[condition][proteincvdict[condition]["CV"]<10].shape[0])

            cvcalc_df["proteinsCV<20"]=proteinscv20
            cvcalc_df["proteinsCV<10"]=proteinscv10
        else:
            cvcalc_df["Protein CVs"]=[0]*len(cvcalc_df)
            cvcalc_df["Protein 95% CVs"]=[0]*len(cvcalc_df)
            cvcalc_df["proteinsCV<20"]=[0]*len(cvcalc_df)
            cvcalc_df["proteinsCV<10"]=[0]*len(cvcalc_df)

        if "FG.MS2Quantity" in searchoutput.columns:
            #precursor-level CVs
            precursorcvlist=[]
            precursorcvlist95=[]
            precursorcvdict={}
            cvprecursorgroup=searchoutput[["R.Condition","R.Replicate","EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]].drop_duplicates().reset_index(drop=True)
            for x,condition in enumerate(sampleconditions):
                if maxreplicatelist[x]==1:
                    emptylist=[]
                    precursorcvlist.append(emptylist)
                    precursorcvlist95.append(emptylist)
                else:
                    df=pd.DataFrame(cvprecursorgroup[cvprecursorgroup["R.Condition"]==condition]).drop(columns=["R.Condition","R.Replicate"])
                    mean=df.groupby(["EG.ModifiedPeptide","FG.Charge"]).mean().rename(columns={"FG.MS2Quantity":"Precursor Mean"})
                    std=df.groupby(["EG.ModifiedPeptide","FG.Charge"]).std().rename(columns={"FG.MS2Quantity":"Precursor Std"})
                    cvprecursortable=pd.concat([mean,std],axis=1)
                    cvprecursortable.dropna(inplace=True)
                    cvprecursortable=cvprecursortable.loc[(cvprecursortable!=0).any(axis=1)]
                    cvlist=(cvprecursortable["Precursor Std"]/cvprecursortable["Precursor Mean"]*100).tolist()
                    precursorcvdict[condition]=pd.DataFrame(cvlist,columns=["CV"])
                    precursorcvlist.append(cvlist)
                    top95=np.percentile(cvlist,95)
                    cvlist95=[]
                    for i in cvlist:
                        if i <=top95:
                            cvlist95.append(i)
                    precursorcvlist95.append(cvlist95)
            cvcalc_df["Precursor CVs"]=precursorcvlist
            cvcalc_df["Precursor 95% CVs"]=precursorcvlist95

            #peptide-level CVs
            if input.peptide_grouping()=="stripped":
                grouping_key="PEP.StrippedSequence"
            if input.peptide_grouping()=="modified":
                grouping_key="EG.ModifiedPeptide"

            peptidecvlist=[]
            peptidecvlist95=[]
            peptidecvdict={}
            for x,condition in enumerate(sampleconditions):
                df=searchoutput[searchoutput["R.Condition"]==condition]
                placeholderdf=pd.DataFrame()
                if maxreplicatelist[x]==1:
                    emptylist=[]
                    peptidecvlist.append(emptylist)
                    peptidecvlist95.append(emptylist)
                else:
                    for run in runnames:
                        df1=df[df["Cond_Rep"]==run][[grouping_key,"FG.MS2Quantity"]].groupby(grouping_key).head(3).groupby(grouping_key).mean()
                        placeholderdf=pd.concat([placeholderdf,df1],axis=0)
                    mean=placeholderdf.sort_values(grouping_key).groupby(grouping_key).mean()
                    std=placeholderdf.sort_values(grouping_key).groupby(grouping_key).std()
                    cv=(std/mean)*100
                    cv=cv["FG.MS2Quantity"].dropna().tolist()
                    peptidecvlist.append(cv)
                    peptidecvdict[condition]=pd.DataFrame(cv,columns=["CV"])

                    top95=np.percentile(cv,95)
                    cvlist95=[]
                    for i in cv:
                        if i <=top95:
                            cvlist95.append(i)
                    peptidecvlist95.append(cvlist95)
        
            cvcalc_df["Peptide CVs"]=peptidecvlist
            cvcalc_df["Peptide 95% CVs"]=peptidecvlist95

            #precursor CVs
            precursorscv20=[]
            precursorscv10=[]
            for x,condition in enumerate(sampleconditions):
                if maxreplicatelist[x]==1:
                    emptylist=[]
                    precursorscv20.append(emptylist)
                    precursorscv10.append(emptylist)
                else:
                    precursorscv20.append(precursorcvdict[condition][precursorcvdict[condition]["CV"]<20].shape[0])
                    precursorscv10.append(precursorcvdict[condition][precursorcvdict[condition]["CV"]<10].shape[0])
            cvcalc_df["precursorsCV<20"]=precursorscv20
            cvcalc_df["precursorsCV<10"]=precursorscv10

            #peptide CVs
            peptidescv20=[]
            peptidescv10=[]
            for x,condition in enumerate(sampleconditions):
                if maxreplicatelist[x]==1:
                    emptylist=[]
                    peptidescv20.append(emptylist)
                    peptidescv10.append(emptylist)
                else:
                    peptidescv20.append(peptidecvdict[condition][peptidecvdict[condition]["CV"]<20].shape[0])
                    peptidescv10.append(peptidecvdict[condition][peptidecvdict[condition]["CV"]<10].shape[0])

            cvcalc_df["peptidesCV<20"]=peptidescv20
            cvcalc_df["peptidesCV<10"]=peptidescv10
        else:
            cvcalc_df["Precursor CVs"]=[0]*len(cvcalc_df)
            cvcalc_df["Precursor 95% CVs"]=[0]*len(cvcalc_df)
            cvcalc_df["Peptide CVs"]=[0]*len(cvcalc_df)
            cvcalc_df["Peptide 95% CVs"]=[0]*len(cvcalc_df)
            cvcalc_df["precursorsCV<20"]=[0]*len(cvcalc_df)
            cvcalc_df["precursorsCV<10"]=[0]*len(cvcalc_df)
            cvcalc_df["peptidesCV<20"]=[0]*len(cvcalc_df)
            cvcalc_df["peptidesCV<10"]=[0]*len(cvcalc_df)


        return cvcalc_df

    # ====================================== Metrics calc functions
    #take searchoutput df and generate variables and dataframes to be used downstream
    @reactive.calc
    def variables():
        searchoutput=metadata_update()

        sampleconditions=searchoutput["R.Condition"].drop_duplicates().tolist()
        maxreplicatelist=searchoutput[["R.Condition","R.Replicate"]].drop_duplicates().groupby("R.Condition").size().reset_index()[0].tolist()
        numconditions=len(searchoutput["R.Condition"].drop_duplicates())
        numsamples=len(searchoutput[["R.Condition","R.Replicate"]].drop_duplicates())
        runnames=searchoutput["Cond_Rep"].drop_duplicates().tolist()

        #matplotlib.rcParams edits for dpi and xtick alignment done here to simplify the setup globally
        if input.dpi_switch()==True:
            matplotlib.rcParams["figure.dpi"]=300
        else:
            matplotlib.rcParams["figure.dpi"]=100
        matplotlib.rcParams["xtick.alignment"]=input.xaxis_label_alignment()

        return searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames   
    #peptide lengths
    @reactive.calc
    def peptidelengths():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

        peptidelengths_condition=pd.DataFrame()
        peptidelengths_run=pd.DataFrame()

        listoflengths=[]
        for condition in sampleconditions:
            placeholder=searchoutput[searchoutput["R.Condition"]==condition]["PEP.StrippedSequence"].drop_duplicates().reset_index(drop=True).tolist()
            lengths=[]
            for pep in placeholder:
                lengths.append(len(pep))
            listoflengths.append(lengths)
        peptidelengths_condition["Sample Names"]=sampleconditions
        peptidelengths_condition["Peptide Lengths"]=listoflengths

        listoflengths=[]
        for run in runnames:
            placeholder=searchoutput[searchoutput["Cond_Rep"]==run]["PEP.StrippedSequence"].drop_duplicates().reset_index(drop=True).tolist()
            lengths=[]
            for pep in placeholder:
                lengths.append(len(pep))
            listoflengths.append(lengths)
        peptidelengths_run["Sample Names"]=runnames
        peptidelengths_run["Peptide Lengths"]=listoflengths

        return peptidelengths_condition,peptidelengths_run
    #peptides per protein
    @reactive.calc
    def pepsperprotein():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

        pepsperprotein_condition=pd.DataFrame()
        pepsperprotein_run=pd.DataFrame()

        pepsperproteinlist=[]
        for condition in sampleconditions:
            df=searchoutput[searchoutput["R.Condition"]==condition][["PG.ProteinGroups","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True)
            pepsperproteinlist.append(df.groupby(["PG.ProteinGroups"]).size().tolist())
        pepsperprotein_condition["Sample Names"]=sampleconditions
        pepsperprotein_condition["Peptides per Protein"]=pepsperproteinlist

        pepsperproteinlist=[]
        for run in runnames:
            df=searchoutput[searchoutput["Cond_Rep"]==run][["PG.ProteinGroups","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True)
            pepsperproteinlist.append(df.groupby(["PG.ProteinGroups"]).size().tolist())
        pepsperprotein_run["Sample Names"]=runnames
        pepsperprotein_run["Peptides per Protein"]=pepsperproteinlist

        return pepsperprotein_condition,pepsperprotein_run
    #peak widths
    @reactive.calc
    def peakwidths():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

        fwhm_df=pd.DataFrame()
        fwhm_df["Cond_Rep"]=runnames
        peakwidths=[]
        peakwidths_95=[]
        for run in runnames:
            peakwidthlist=(searchoutput[searchoutput["Cond_Rep"]==run]["EG.PeakWidth"]*60).tolist()
            peakwidths.append(peakwidthlist)
            top95=np.percentile(peakwidthlist,95)
            top95list=[]
            for i in peakwidthlist:
                if i <= top95:
                    top95list.append(i)
            peakwidths_95.append(top95list)
        fwhm_df["Peak Width"]=peakwidths
        fwhm_df["Peak Width 95%"]=peakwidths_95

        return fwhm_df

#endregion

# =============================# Sidebar Tabs #================================
# ============================================================================= File Import
#region
    # ====================================== General
    #check for OS, changes what separator to use for file paths
    @reactive.calc
    def os_check():
        import platform
        if platform.system()=="Windows":
            os_var="\\"
        return os_var
    #software choices to simplify the interface
    @render.ui
    def software_ui():
        software=input.software_general()
        if software=="spectronaut":
            opts={"spectronaut":"directDIA / library-based search",
                  "ddalibrary":"DDA Library"}
        if software=="diann":
            opts={"diann":"DIA-NN pre 2.0",
                  "diann2.0":"DIA-NN 2.0"}
        if software=="fragpipe":
            opts={"fragpipe":"FragPipe",
                  "fragpipe_glyco":"FragPipe Glyco",
                  "fragpipe_combined_ion":"FragPipe Quant"}
        if software=="bps":
            opts={"bps_timsrescore":"tims-rescore",
                  "bps_timsdiann":"tims-DIANN",
                  "bps_spectronaut":"Spectronaut",
                  "bps_pulsar":"Pulsar",
                  "bps_denovo":"BPS Novor",
                  "bps_sage":"Sage",
                  "glycoscape":"Glycoscape"}
        if software=="spectromine":
            opts={"spectromine":"Spectromine"}
        if software=="peaks":
            opts={"peaks":"PEAKS"}
        if software=="sage":
            opts={"sage":"Sage",
                  "sage_lfq":"Sage LFQ"}
        return ui.input_radio_buttons("software","",choices=opts)
    #dictionary for name_mod values used in comparative plotting
    @reactive.calc
    def software_dict():
        name_mod_dict={"spectronaut":"Spectronaut",
                  "ddalibrary":"DDA Library",
                  "diann":"DIA-NN",
                  "diann2.0":"DIA-NN 2.0",
                  "fragpipe":"FragPipe",
                  "fragpipe_glyco":"FragPipe Glyco",
                  "bps_timsrescore":"BPS tims-rescore",
                  "bps_timsdiann":"BPS tims-DIANN",
                  "bps_spectronaut":"BPS Spectronaut",
                  "bps_pulsar":"BPS Pulsar",
                  "bps_denovo":"BPS Novor",
                  "bps_sage":"BPS Sage",
                  "spectromine":"Spectromine",
                  "peaks":"PEAKS",
                  "sage":"Sage",
                  "sage_lfq":"Sage LFQ"
        }
        return name_mod_dict
    #input files for each bps input option
    def bps_filedict():
        #input.software()+input.software_bps_report_type() as keys
        peptidefiledict={
            "bps_timsrescore_qual":"pgfdr.peptide.parquet",
            "bps_timsrescore_quant":"consolidation.peptide.parquet",
            "bps_timsdiann_qual":"tims-diann.peptide.parquet",
            "bps_timsdiann_quant":"tims-diann.peptide.parquet",
            "bps_pulsar_qual":"pulsar-id.peptide.parquet",
            "bps_pulsar_quant":"pulsar-archive-merge.peptide.parquet",
            "bps_spectronaut_qual":"spectronaut-id.peptide.parquet",
            "bps_spectronaut_quant":"spectronaut-combine.peptide.parquet",
            #qual and quant for novor needed for dict access
            "bps_denovo_qual":"novor-fasta-mapping-results.peptide.parquet",
            "bps_denovo_quant":"novor-fasta-mapping-results.peptide.parquet",
            "glycoscape":"myriad-merger.glycopsm.parquet",
        }
        proteinfiledict={
            "bps_timsrescore_qual":None,
            "bps_timsrescore_quant":"maxlfq.protein.parquet",
            "bps_timsdiann_qual":"tims-diann.protein.parquet",
            "bps_timsdiann_quant":"tims-diann.protein.parquet",
            "bps_pulsar_qual":"pulsar-id.protein.parquet",
            "bps_pulsar_quant":"pulsar-archive-merge.protein.parquet",
            "bps_spectronaut_qual":"spectronaut-id.protein.parquet",
            "bps_spectronaut_quant":"spectronaut-combine.protein.parquet",
            "bps_denovo_qual":"novor-fasta-mapping-results.protein.parquet",
            "bps_denovo_quant":"novor-fasta-mapping-results.protein.parquet",
            "glycoscape":None,
        }
        return peptidefiledict,proteinfiledict

    # ====================================== Callable input file functions
    #get extension of input file (for multiple files, only looks at the first file name)
    def get_inputfile_type(searchreport):
        if searchreport:
            filename=searchreport[0]["name"]
        else:
            return ""
        fileextension=filename[filename.rfind("."):]
        return fileextension
    #handling for uploads of >1 .tsv, .parquet, or .csv files
    def combine_inputfiles(searchreport,software):
        fileextension=get_inputfile_type(searchreport)

        searchoutput=pd.DataFrame()
        if fileextension==".tsv":
            if software=="fragpipe_combined_ion":
                for i in range(len(searchreport)):
                    if "combined_ion" in searchreport[i]["name"]:
                        combined_ion=pd.read_csv(searchreport[i]["datapath"],sep="\t")
                    if "combined_protein" in searchreport[i]["name"]:
                        combined_protein=pd.read_csv(searchreport[i]["datapath"],sep="\t")

                #remove extraneous columns from protein and modifiedpeptide dfs
                for column in combined_protein.columns:
                    if "Spectral Count" in column:
                        combined_protein=combined_protein.drop(columns=column)
                combined_protein=combined_protein.drop(columns=["Protein Length","Organism","Protein Existence","Description","Protein Probability","Top Peptide Probability","Combined Total Peptides","Indistinguishable Proteins"])
                for column in combined_ion.columns:
                    if "Spectral Count" in column or "Match Type" in column or "Localization" in column:
                        combined_ion=combined_ion.drop(columns=column)
                    
                combined_ion=combined_ion.drop(columns=["Prev AA","Next AA","Start","End","Peptide Length","Compensation Voltage","Assigned Modifications","Protein","Protein Description","Mapped Genes","Mapped Proteins"])

                idcolumns=["Peptide Sequence","Modified Sequence","M/Z","Charge","Protein ID","Entry Name","Gene"]
                combined_ion_intensity=pd.concat([combined_ion[idcolumns],combined_ion.loc[:,combined_ion.columns.str.contains("Intensity")]],axis=1)
                combined_ion_RT=pd.concat([combined_ion[idcolumns],combined_ion.loc[:,combined_ion.columns.str.contains("Retention Time")]],axis=1)
                combined_ion_IM=pd.concat([combined_ion[idcolumns],combined_ion.loc[:,combined_ion.columns.str.contains("Ion Mobility")]],axis=1)

                #unpivot modifiedpeptide df
                columns=["Peptide Sequence","Modified Sequence","M/Z","Charge","Protein ID","Entry Name","Gene"]
                filenames=[]
                filenames_original=[]
                for column in combined_ion_intensity.columns:
                    if column in columns:
                        pass
                    else:
                        filenames_original.append(column)
                        filenames.append(column.split(" Intensity")[0])
                peptide_df=pd.melt(combined_ion_intensity,id_vars=columns,var_name="R.FileName",value_name="FG.MS2Quantity")

                filenames=[]
                filenames_original=[]
                for column in combined_ion_RT.columns:
                    if column in columns:
                        pass
                    else:
                        filenames_original.append(column)
                        filenames.append(column.split(" Intensity")[0])
                peptide_df_RT=pd.melt(combined_ion_RT,id_vars=columns,var_name="R.FileName",value_name="EG.ApexRT")
                peptide_df_RT["EG.ApexRT"]=peptide_df_RT["EG.ApexRT"]/60

                filenames=[]
                filenames_original=[]
                for column in combined_ion_IM.columns:
                    if column in columns:
                        pass
                    else:
                        filenames_original.append(column)
                        filenames.append(column.split(" Intensity")[0])
                peptide_df_IM=pd.melt(combined_ion_IM,id_vars=columns,var_name="R.FileName",value_name="EG.IonMobility")

                peptide_df_concat=pd.concat([peptide_df,peptide_df_RT["EG.ApexRT"],peptide_df_IM["EG.IonMobility"]],axis=1).dropna().reset_index(drop=True)

                #remove "Intensity" from filename entries in both dfs then sort the peptide df by protein ID
                peptide_df_concat[["R.FileName","drop"]]=peptide_df_concat["R.FileName"].str.split(" ",expand=True)
                peptide_df_concat=peptide_df_concat.drop(columns=["drop"])
                peptide_df_concat=peptide_df_concat.sort_values(["Protein ID","R.FileName"],ascending=[True,True]).reset_index(drop=True)

                #remove "Intensity" from the protein df column names
                descriptor_columns=["Protein","Entry Name","Gene","Protein Length","Organism","Protein Existence","Description","Protein Probability","Top Peptide Probability","Combined Total Peptides","Indistinguishable Proteins"]
                intensity_columns=sorted(list(set(combined_protein.columns)-set(descriptor_columns)))
                protein_intensity=combined_protein[intensity_columns].set_index("Protein ID")

                for column in protein_intensity.columns:
                    protein_intensity=protein_intensity.rename(columns={column:column.split(" Intensity")[0]})

                #sort peptide_df by filename
                s=peptide_df_concat.sort_values("R.FileName")["R.FileName"].drop_duplicates().tolist()
                sortedfilenames=sorted(s,key=lambda x: int(re.findall(r'\d+',x)[-1]))
                peptide_df_concat["sortedfilenames"]=pd.Categorical(peptide_df_concat["R.FileName"],categories=sortedfilenames,ordered=True)
                peptide_df_concat=peptide_df_concat.sort_values(["sortedfilenames","Protein ID"],ascending=[True,True])
                peptide_df_concat=peptide_df_concat.drop(columns=["sortedfilenames"]).reset_index(drop=True)

                #sort protein_intensity df columns by filename
                s=protein_intensity.columns.tolist()
                sortedfilenames=sorted(s,key=lambda x: int(re.findall(r'\d+',x)[-1]))
                protein_intensity=protein_intensity[sortedfilenames]

                #map protein intensities per file to the peptide_df
                proteinintensitylist=[]
                for i in range(len(peptide_df_concat)):
                    proteinintensitylist.append(protein_intensity[peptide_df_concat["R.FileName"][i]].loc[peptide_df_concat["Protein ID"][i]])
                peptide_df_concat["PG.MS2Quantity"]=proteinintensitylist
                searchoutput=peptide_df_concat
            else:
                for i in range(len(searchreport)):
                    run=pd.read_csv(searchreport[i]["datapath"],sep="\t")
                    searchoutput=pd.concat([searchoutput,run])
            searchoutput=searchoutput.reset_index(drop=True)
        elif fileextension==".parquet":
            for i in range(len(searchreport)):
                run=pd.read_parquet(searchreport[i]["datapath"])
                searchoutput=pd.concat([searchoutput,run])
            searchoutput=searchoutput.reset_index(drop=True)
        elif fileextension==".csv":
            for i in range(len(searchreport)):
                run=pd.read_csv(searchreport[i]["datapath"])
                searchoutput=pd.concat([searchoutput,run])
            searchoutput=searchoutput.reset_index(drop=True)               
        return searchoutput
    #unpacking .zip files (handling for multiple .zips as well)
    def zip_inputfiles(searchreport,software,software_bps_report_type):
        peptidefiledict,proteinfiledict=bps_filedict()
        os_var=os_check()

        peptidefile=peptidefiledict[software+"_"+software_bps_report_type]
        proteinfile=proteinfiledict[software+"_"+software_bps_report_type]

        searchoutput=pd.DataFrame()
        #chdir to timsplot app.py directory
        os.chdir(os.path.dirname(os.path.realpath(__file__)))
        for i in range(len(searchreport)):
            #unzip inputfile
            bpszip=ZipFile(searchreport[i]["datapath"])
            bpszip.extractall()
            #read metadata file in inputfile
            metadata_bps=pd.read_csv("metadata.csv")
            #runlist and names of subfolders
            runlist=metadata_bps["processing_run_uuid"].tolist()
            #chdir to processing-run subfolder containing folders for each run
            cwd=os.getcwd()+os_var+"processing-run"
            os.chdir(cwd)

            if software_bps_report_type=="qual":
                if software=="bps_denovo" or software=="glycoscape":
                    pass
                else:
                    for run in runlist:
                        os.chdir(cwd)
                        os.chdir(cwd+os_var+run)
                        peptideparquet=pd.read_parquet(peptidefile)
                        #bps_pulsar and bps_spectronaut need mapping from protein file
                        if software=="bps_pulsar" or software=="bps_spectronaut" or software=="bps_timsdiann":
                            #pull protein information, otherwise protein file not needed
                            proteinparquet=pd.read_parquet(proteinfile)
                            #map protein names and quant from protein file to peptide file
                            peptideparquet=peptideparquet.sort_values(by=["protein_group_parent_id"]).reset_index(drop=True)
                            proteinparquet=proteinparquet.sort_values(by=["protein_group_id"]).reset_index(drop=True)
                            #pull protein and gene info from protein parquet file and add to peptideparquet 
                            protein_name=[]
                            protein_group=[]
                            protein_quant=[]
                            #per run, loop through the searchoutput protein id and pull info from proteinparquet from each protein id
                            #map to sorted searchoutput which will already by in run and protein id order
                            for i in range(len(peptideparquet["protein_group_parent_id"].tolist())):
                                protein_name.append(proteinparquet["protein_name"].iloc[peptideparquet["protein_group_parent_id"][i]])
                                protein_group.append(proteinparquet["protein_accession"].iloc[peptideparquet["protein_group_parent_id"][i]])
                                protein_quant.append(proteinparquet["pg_quantity"].iloc[peptideparquet["protein_group_parent_id"][i]])
                            peptideparquet["PG.ProteinNames"]=protein_name
                            peptideparquet["PG.ProteinGroups"]=protein_group
                            peptideparquet["PG.ProteinAccessions"]=protein_group
                            peptideparquet["PG.MS2Quantity"]=protein_quant
                        searchoutput=pd.concat([searchoutput,peptideparquet],ignore_index=True)
            elif software_bps_report_type=="quant":
                if software=="bps_denovo" or software=="glycoscape":
                    pass
                else:
                    #BPS compiles all runs into a single subfolder instead of a subfolder per run
                    run=runlist[0]
                    os.chdir(cwd+os_var+run)
                    peptideparquet=pd.read_parquet(peptidefile)
                    #conditional probably not needed since all quant files need the protein file to be transferred over
                    if proteinfile is not None:
                        proteinparquet=pd.read_parquet(proteinfile)
                        peptideparquet=peptideparquet.sort_values(by=["sample_name","protein_group_parent_id"]).reset_index(drop=True)
                        proteinparquet=proteinparquet.sort_values(by=["sample_name","protein_group_id"]).reset_index(drop=True)
                        proteinparquet_names=proteinparquet.sort_values(by=["sample_name","protein_group_id"]).reset_index(drop=True)[["protein_group_id","protein_group_name","protein_accession","protein_name","gene_id"]].set_index("protein_group_id").drop_duplicates()
                        
                        #check if the gene_id column is actual gene identifiers or just an index
                        try:
                            #check if we can convert the gene_id column to int, if it fails, it means we have real gene identifiers
                            proteinparquet_names["gene_id"]=proteinparquet_names["gene_id"].astype(int)
                            proteinparquet_names=proteinparquet_names.drop(columns=["gene_id"])
                        except:
                            pass

                        #tims-rescore is the only software that reports pg_quantity as pg_max_lfq, change name for convenience
                        if "pg_max_lfq" in proteinparquet.columns:
                            proteinparquet=proteinparquet.rename(columns={"pg_max_lfq":"pg_quantity"})

                        #map protein names and quant from protein file to peptide file
                        protein_name=[]
                        protein_group=[]
                        if "gene_id" in proteinparquet_names.columns:
                            gene_id=[]
                        protein_quant=[]
                        #per run, loop through the peptideparquet protein id and pull info from proteinparquet from each protein id
                        #map to sorted peptideparquet which will already be in run and protein id order
                        for run in peptideparquet["sample_name"].drop_duplicates().tolist():
                            peptideparquet_df=peptideparquet[peptideparquet["sample_name"]==run]
                            proteinparquet_df=proteinparquet[proteinparquet["sample_name"]==run]
                            for i in peptideparquet_df["protein_group_parent_id"].tolist():
                                protein_quant.append(proteinparquet_df[proteinparquet_df["protein_group_id"]==i]["pg_quantity"].values[0])
                        for i in range(len(peptideparquet["protein_group_parent_id"].tolist())):
                            protein_name.append(proteinparquet_names["protein_name"].iloc[peptideparquet["protein_group_parent_id"][i]])
                            #protein group and accession should always be the same. tims-rescore uses "|" delimited group names compared to the rest of BPS software
                            protein_group.append(proteinparquet_names["protein_accession"].iloc[peptideparquet["protein_group_parent_id"][i]])
                            if "gene_id" in proteinparquet_names.columns:
                                gene_id.append(proteinparquet_names["gene_id"].iloc[peptideparquet["protein_group_parent_id"][i]])
                        peptideparquet["PG.ProteinNames"]=protein_name
                        peptideparquet["PG.ProteinGroups"]=protein_group
                        peptideparquet["PG.ProteinAccessions"]=protein_group
                        if "gene_id" in proteinparquet_names.columns:
                            peptideparquet["PG.Genes"]=gene_id
                        peptideparquet["PG.MS2Quantity"]=protein_quant
                    searchoutput=peptideparquet
            if software=="bps_denovo":
                denovo_score=65.0
                rank=1
                for run in runlist:
                    os.chdir(cwd)
                    os.chdir(cwd+os_var+run)
                    peptideparquet=pd.read_parquet(peptidefile)
                    proteinparquet=pd.read_parquet(proteinfile)

                    #filter out low denovo scores and precursors that aren't 1-ranked
                    peptideparquet=peptideparquet[(peptideparquet["denovo_score"]>=denovo_score)&(peptideparquet["rank"]==rank)].reset_index(drop=True)

                    #fill NaN in the protein group column with -1
                    peptideparquet["protein_group_parent_id"]=peptideparquet["protein_group_parent_id"].fillna(-1)

                    #pull protein and gene info from protein parquet file and add to peptideparquet 
                    protein_name=[]
                    protein_accession=[]
                    gene_id=[]
                    uncategorized_placeholder="uncat"
                    for i in range(len(peptideparquet["protein_group_parent_id"])):
                        if peptideparquet["protein_group_parent_id"][i]==-1:
                            protein_name.append(uncategorized_placeholder)
                            protein_accession.append(uncategorized_placeholder)
                            gene_id.append(uncategorized_placeholder)
                        else:
                            protein_name.append(proteinparquet["protein_name"].iloc[int(peptideparquet["protein_group_parent_id"][i])])
                            protein_accession.append(proteinparquet["protein_accession"].iloc[int(peptideparquet["protein_group_parent_id"][i])])
                            gene_id.append(proteinparquet["gene_id"].iloc[int(peptideparquet["protein_group_parent_id"][i])])

                    peptideparquet["PG.ProteinGroups"]=protein_name
                    peptideparquet["PG.ProteinAccessions"]=protein_accession
                    peptideparquet["PG.Genes"]=gene_id

                    searchoutput=pd.concat([searchoutput,peptideparquet],ignore_index=True)
            if software=="glycoscape":
                for run in runlist:
                    os.chdir(cwd)
                    os.chdir(cwd+os_var+run)
                    peptideparquet=pd.read_parquet(peptidefile)

                    #convert glycan_composition column entries to format from fragger-glyco
                    convert_dict={"H":"Hex","N":"HexNAc","F":"Fuc","S":"NeuAc","G":"NeuGc"}
                    glycanlist=peptideparquet["glycan_composition"].tolist()
                    glycanlist_updated=[]
                    glycanlist_ambiguous=[]
                    for i,entry in enumerate(glycanlist):
                        if entry is None or entry=="unidentified":
                            glycanlist_updated.append(np.nan)
                            glycanlist_ambiguous.append(np.nan)
                        else:
                            if "," in entry:
                                entry=entry.split(",")[0]
                                glycanlist_ambiguous.append(True)
                            else:
                                glycanlist_ambiguous.append(False)
                            entrylist=list(entry)
                            newglycanentry=[]
                            for ele in entrylist:
                                if ele.isnumeric()==True:
                                    newglycanentry.append("("+ele+")")
                                else:
                                    newglycanentry.append(convert_dict[ele])
                            glycanlist_updated.append("".join(newglycanentry))
                    peptideparquet["Total Glycan Composition"]=glycanlist_updated
                    peptideparquet["Ambiguous Glycan ID"]=glycanlist_ambiguous

                    #unpack protein_list column
                    proteinlist=peptideparquet["protein_list"].str.split(";").tolist()
                    truthlist=[]
                    for entry in proteinlist:
                        if "Reverse" not in entry[0] and "contaminant" not in str(entry):
                            truthlist.append(True)
                        else:
                            truthlist.append(False)
                    peptideparquet=peptideparquet[truthlist]

                    proteingroups=[]
                    proteinnames=[]
                    proteinlist_column=peptideparquet["protein_list"].tolist()
                    for item in proteinlist_column:
                        if item.count(";")==0:
                            templist=item.split("|")
                            proteingroups.append(templist[1])
                            proteinnames.append(templist[2])
                        else:
                            proteingroups_torejoin=[]
                            proteinnames_torejoin=[]
                            for entry in item.split(";"):
                                templist=entry.split("|")
                                proteingroups_torejoin.append(templist[1])
                                proteinnames_torejoin.append(templist[2])
                            proteingroups.append(";".join(proteingroups_torejoin))
                            proteinnames.append(";".join(proteinnames_torejoin))
                    peptideparquet["PG.ProteinGroups"]=proteingroups
                    peptideparquet["PG.ProteinNames"]=proteinnames

                    searchoutput=pd.concat([searchoutput,peptideparquet],ignore_index=True)
        
        #reset index numbering
        searchoutput=searchoutput.reset_index(drop=True)
        #reset cwd back to location of app.py
        os.chdir(os.path.dirname(os.path.realpath(__file__)))
        return searchoutput
    #file import, extensible to De Novo and Two-Software Comparison instead of needing their own versions of this function
    def inputfile(searchreport,software,software_bps_report_type,diann_mbr_switch,reupload):
        fileextension=get_inputfile_type(searchreport)
        ptmdict=ptmdict_apply()
        if searchreport is None:
            return pd.DataFrame()
        if len(searchreport)>1 and fileextension!=".zip":
            searchoutput=combine_inputfiles(searchreport,software)
        else:
            if fileextension==".tsv":
                searchoutput=pd.read_csv(searchreport[0]["datapath"],sep="\t")
            elif fileextension==".parquet":
                searchoutput=pd.read_parquet(searchreport[0]["datapath"])
            elif fileextension==".csv":
                searchoutput=pd.read_csv(searchreport[0]["datapath"])
            elif fileextension==".zip":
                searchoutput=zip_inputfiles(searchreport,software,software_bps_report_type)
        if reupload==True:
            #for data that's been processed, downloaded, and reuploaded to timsplot. Fixes the PTM Protein Locations and PTMs columns since it interprets the lists as a single string
            if "PTM Protein Locations" in searchoutput.columns:
                fixedlist=[]
                for ele in searchoutput["PTM Protein Locations"]:
                    if ele=="[]":
                        fixedlist.append([""])
                    else:
                        fixedlist.append([x.strip() for x in ele.replace("[","").replace("]","").replace("'","").split(",")])
                searchoutput["PTM Protein Locations"]=fixedlist
            if "PTMs" in searchoutput.columns:
                fixedlist=[]
                for ele in searchoutput["PTMs"]:
                    if ele=="[]":
                        fixedlist.append([])
                    else:
                        fixedlist.append([x.strip() for x in ele.replace("[","").replace("]","").replace("'","").split(",")])
                searchoutput["PTMs"]=fixedlist
            #return statement here since the file would have gone through all the following processing before
            return searchoutput
        if software=="spectronaut":
            renamedict={}
            dropcolumns=[]
        if software=="ddalibrary":
            renamedict={
                "ReferenceRun":"R.FileName",
                "PrecursorCharge":"FG.Charge",
                "ModifiedPeptide":"EG.ModifiedPeptide",
                "StrippedPeptide":"PEP.StrippedSequence",
                "IonMobility":"EG.IonMobility",
                "PrecursorMz":"FG.PrecMz",
                "ReferenceRunMS1Response":"FG.MS2Quantity",
                "Protein Name":"PG.ProteinNames"
            }
            dropcolumns=[]
        if software=="diann":
            if diann_mbr_switch=="Protein.Q.Value":
                searchoutput=searchoutput[searchoutput["Protein.Q.Value"]<=0.01]
            elif diann_mbr_switch=="Global.PG.Q.Value":
                searchoutput=searchoutput[searchoutput["Global.PG.Q.Value"]<=0.01]
            elif diann_mbr_switch=="off":
                searchoutput=searchoutput

            searchoutput["EG.PeakWidth"]=searchoutput["RT.Stop"]-searchoutput["RT.Start"]

            searchoutput["Modified.Sequence"]=searchoutput["Modified.Sequence"].str.replace("(","[")
            searchoutput["Modified.Sequence"]=searchoutput["Modified.Sequence"].str.replace(")","]")

            renamedict={
                "Run":"R.FileName",
                "Protein.Group":"PG.ProteinGroups",
                "Protein.Ids":"PG.ProteinAccessions",
                "Protein.Names":"PG.ProteinNames",
                "Genes":"PG.Genes",
                "PG.MaxLFQ":"PG.MS2Quantity",
                "Modified.Sequence":"EG.ModifiedPeptide",
                "Stripped.Sequence":"PEP.StrippedSequence",
                "Precursor.Charge":"FG.Charge",
                #"Protein.Q.Value":"",
                #"Global.PG.Q.Value":"",
                "Precursor.Quantity":"FG.MS2Quantity",
                "RT":"EG.ApexRT",
                "First.Protein.Description":"PG.ProteinDescriptions",
                "IM":"EG.IonMobility"
            }
            dropcolumns=[
                "File.Name","PG.Normalised","PG.MaxLFQ","Genes.Quantity","Genes.Normalised","Genes.MaxLFQ","Genes.MaxLFQ.Unique","Precursor.Id","Q.Value","PEP",
                "Global.Q.Value","PG.Q.Value","GG.Q.Value","Translated.Q.Value","Proteotypic","Precursor.Normalised","Quantity.Quality","RT.Start","RT.Stop",
                "iRT","Predicted.RT","Predicted.iRT","Lib.Q.Value","Lib.PG.Q.Value","Ms1.Profile.Corr","Ms1.Area","Ms1.Normalised","Normalisation.Factor",
                "Evidence","Spectrum.Similarity","Averagine","Mass.Evidence","CScore","Decoy.Evidence","Decoy.CScore","Fragment.Quant.Raw","Fragment.Correlations",
                "MS2.Scan","iIM","Predicted.IM","Predicted.iIM","Rec"
            ]
        if software=="diann2.0":
            if diann_mbr_switch=="Protein.Q.Value":
                searchoutput=searchoutput[searchoutput["Protein.Q.Value"]<=0.01]
            elif diann_mbr_switch=="Global.PG.Q.Value":
                searchoutput=searchoutput[searchoutput["Global.PG.Q.Value"]<=0.01]
            elif diann_mbr_switch=="off":
                searchoutput=searchoutput

            searchoutput["EG.PeakWidth"]=searchoutput["RT.Stop"]-searchoutput["RT.Start"]

            searchoutput["Modified.Sequence"]=searchoutput["Modified.Sequence"].str.replace("(","[")
            searchoutput["Modified.Sequence"]=searchoutput["Modified.Sequence"].str.replace(")","]")
            #check for blank Protein.Group entries
            if len(searchoutput[searchoutput["Protein.Group"]==""]["Protein.Group"].index.tolist())!=0:
                for i in searchoutput[searchoutput["Protein.Group"]==""][["Protein.Ids","Protein.Group"]].index.tolist():
                    searchoutput.iloc[i,searchoutput.columns.get_loc("Protein.Group")]=searchoutput["Protein.Ids"][i]

            renamedict={
                "Run":"R.FileName",
                "Modified.Sequence":"EG.ModifiedPeptide",
                "Stripped.Sequence":"PEP.StrippedSequence",
                "Precursor.Charge":"FG.Charge",
                "Precursor.Mz":"FG.PrecMz",
                "Protein.Group":"PG.ProteinGroups",
                "Protein.Names":"PG.ProteinNames",
                "Genes":"PG.Genes",
                "RT":"EG.ApexRT",
                "IM":"EG.IonMobility",
                "Precursor.Quantity":"FG.MS2Quantity",
                "PG.MaxLFQ":"PG.MS2Quantity"
            }
            dropcolumns=[
                "Run.Index","Channel","Precursor.Id","Precursor.Lib.Index","Decoy","Proteotypic","Protein.Ids","iRT","Predicted.RT","Predicted.iRT",
                "iIM","Predicted.IM","Predicted.iIM","Precursor.Normalised","Ms1.Area","Ms1.Normalised","Ms1.Apex.Area","Ms1.Apex.Mz.Delta",
                "Normalisation.Factor","Quantity.Quality","Empirical.Quality","Normalisation.Noise","Ms1.Profile.Corr","Evidence","Mass.Evidence",
                "Channel.Evidence","Ms1.Total.Signal.Before","Ms1.Total.Signal.After","RT.Start","RT.Stop","FWHM","PG.TopN","Genes.TopN","Genes.MaxLFQ",
                "Genes.MaxLFQ.Unique","PG.MaxLFQ.Quality","Genes.MaxLFQ.Quality","Genes.MaxLFQ.Unique.Quality","Q.Value","PEP","Global.Q.Value","Lib.Q.Value",
                "Peptidoform.Q.Value","Global.Peptidoform.Q.Value","Lib.Peptidoform.Q.Value","PTM.Site.Confidence","Site.Occupancy.Probabilities","Protein.Sites",
                "Lib.PTM.Site.Confidence","Translated.Q.Value","Channel.Q.Value","PG.Q.Value","PG.PEP","GG.Q.Value","Lib.PG.Q.Value"
            ]
        if software=="fragpipe":
            searchoutput["FG.CalibratedMassAccuracy (PPM)"]=(searchoutput["Delta Mass"]/searchoutput["Calculated M/Z"])*10E6

            searchoutput["Retention"]=searchoutput["Retention"]/60

            searchoutput["Modified Peptide"]=searchoutput["Modified Peptide"].str.replace("[","")
            searchoutput["Modified Peptide"]=searchoutput["Modified Peptide"].str.replace("]","")

            peps=searchoutput["Peptide"].tolist()
            modpeps=searchoutput["Modified Peptide"].tolist()
            for i in range(len(peps)):
                if type(modpeps[i])!=str:
                    modpeps[i]=peps[i]
                else:
                    modpeps[i]=modpeps[i]
            searchoutput["Modified Peptide"]=modpeps

            #convert Assigned Modifications to protein sequence-localized modifications
            if "PTM Protein Locations" not in searchoutput.columns:
                #make a column of specific protein sequence locations for PTMs
                modlist=[]
                residuemodlist=[]
                for ele in searchoutput["Assigned Modifications"].tolist():
                    if type(ele)==float:
                        modlist.append([])
                        residuemodlist.append([])
                    else:
                        ptmloclist=[]
                        ptmlocs=[re.split("[( )]",x.strip())[0] for x in ele.split(",")]
                        for string in ptmlocs:
                            if "N-term" in string:
                                ptmloclist.append(1)
                            else:
                                ptmloclist.extend(map(int,re.findall(r"\d+",string)))
                        modlist.append(ptmloclist)
                        
                        residuemodlist_temp=[]
                        ele_split=ele.split(",")
                        for k in range(len(ele_split)):
                            if "N-term" in ele_split[k]:
                                residuemodlist_temp.append(ele_split[k].split("(")[0])
                            else:
                                residuemodlist_temp.append(re.split(r"(\d+)",ele_split[k].strip().split("(")[0])[2])
                        residuemodlist.append(residuemodlist_temp)
                searchoutput["Mod Locations"]=modlist
                searchoutput["Residue Mod Locations"]=residuemodlist

                proteinmodlocations=[]
                for i,mod in enumerate(searchoutput["Mod Locations"].tolist()):
                    if len(mod)==0:
                        proteinmodlocations.append([])
                    elif len(mod)==1:
                        proteinmodlocations.append([searchoutput["Protein Start"][i]+mod[0]-1])
                    else:
                        multimodloc=[]
                        for j in range(len(mod)):
                            multimodloc.append(searchoutput["Protein Start"][i]+mod[j]-1)
                        proteinmodlocations.append(multimodloc)
                searchoutput["Protein Mod Locations"]=proteinmodlocations

                modnotationlist=[]
                for i in range(len(searchoutput)):
                    modnotationlist_temp=[]
                    for j in range(len(searchoutput["Protein Mod Locations"][i])):
                        modnotationlist_temp.append(str(searchoutput["Protein Mod Locations"][i][j])+searchoutput["Residue Mod Locations"][i][j])
                    modnotationlist.append(modnotationlist_temp)
                searchoutput["PTM Protein Locations"]=modnotationlist

                #sorting functions to sort the PTM Protein Locations column correctly
                #https://stackoverflow.com/questions/5967500/how-to-correctly-sort-a-string-with-a-number-inside
                def atoi(text):
                    return int(text) if text.isdigit() else text
                def natural_keys(text):
                    '''
                    alist.sort(key=natural_keys) sorts in human order
                    http://nedbatchelder.com/blog/200712/human_sorting.html
                    (See Toothy's implementation in the comments)
                    '''
                    return [atoi(c) for c in re.split(r'(\d+)',text)]

                #PTMs column is in sequence order, Assigned Modifications column is not
                sortedlist=[]
                for ele in searchoutput["PTM Protein Locations"].tolist():
                    sortedlist.append(sorted(ele,key=natural_keys))
                searchoutput["PTM Protein Locations"]=sortedlist

            renamedict={
                "Spectrum File":"R.FileName",
                "Peptide":"PEP.StrippedSequence",
                "Modified Peptide":"EG.ModifiedPeptide",
                "Charge":"FG.Charge",
                "Retention":"EG.ApexRT",
                "Observed M/Z":"FG.PrecMz",
                "Ion Mobility":"EG.IonMobility",
                "Protein ID":"PG.ProteinGroups",
                "Entry Name":"PG.ProteinNames",
                "Gene":"PG.Genes",
                "Intensity":"FG.MS2Quantity"
            }
            dropcolumns=[
                "Spectrum","Extended Peptide","Prev AA","Next AA","Peptide Length","Observed Mass","Calibrated Observed Mass","Calibrated Observed M/Z","Calculated Peptide Mass","Calculated M/Z",
                "Delta Mass","Expectation","Hyperscore","Nextscore","Number of Enzymatic Termini","Number of Missed Cleavages","Protein Start","Protein End","Assigned Modifications",
                "Observed Modifications","Purity","Is Unique","Protein","Protein Description","Mapped Genes","Mapped Proteins","SpectralSim","RTScore","Probability"
            ]
        if software=="fragpipe_combined_ion":
            renamedict={
                "Peptide Sequence":"PEP.StrippedSequence",
                "Modified Sequence":"EG.ModifiedPeptide",
                "M/Z":"FG.PrecMz",
                "Charge":"FG.Charge",
                "Protein ID":"PG.ProteinGroups",
                "Entry Name":"PG.ProteinNames",
                "Gene":"PG.Genes"
            }
            dropcolumns=[]
        if software=="spectromine":
            searchoutput["PG.ProteinGroups"]=searchoutput["PG.ProteinAccessions"]
            renamedict={
                "PG.Label-Free Quant":"PG.MS2Quantity",
                "P.Label-Free Quant":"FG.MS2Quantity",
                "PEP.QValue":"EG.Qvalue",
                "PSM.ApexIM":"EG.IonMobility",
                "PP.Charge":"FG.Charge",
                "P.MoleculeID":"EG.ModifiedPeptide",
                "PSM.ApexRT":"EG.ApexRT",
                "PSM.CalibratedMS1MZ":"FG.PrecMz",
                "PSM.DeltaMS1MZ(Theor-Cali)":"FG.CalibratedMassAccuracy (PPM)"
            }
            dropcolumns=[]
        if software=="peaks":
            searchoutput_accessionsplit=searchoutput["Accession"].str.split(";",expand=True)[0].str.split("|",expand=True)
            searchoutput["PG.ProteinGroups"]=searchoutput_accessionsplit[0]
            searchoutput["PG.ProteinNames"]=searchoutput_accessionsplit[1]

            searchoutput["Peptide"]=searchoutput["Peptide"].str.replace("(","[")
            searchoutput["Peptide"]=searchoutput["Peptide"].str.replace(")","]")
            searchoutput["Peptide"]=searchoutput["Peptide"].str.replace("+","")
            #fill NaN values in Area with zeroes
            searchoutput["Area"]=searchoutput["Area"].fillna(0)

            renamedict={
                "Peptide":"EG.ModifiedPeptide",
                "ppm":"FG.CalibratedMassAccuracy (PPM)",
                "m/z":"FG.PrecMz",
                "z":"FG.Charge",
                "RT":"EG.ApexRT",
                "Area":"FG.MS2Quantity", #peptide level peak area, listing it as fragment MS2 quant since that's the convention I've gone with across timsplot
                "Source File":"R.FileName"
            }
            dropcolumns=["-10LgP","Mass","Tag Length","CAA (%)","Length","Delta RT","MS2 Correlation","raw m/z","raw RT","Scan","Feature Id","Accession","PTM","AScore","Positional Confidence","tag(>=0.0%)"]
        if software=="sage":
            proton_mass=1.007276466812
            searchoutput["FG.PrecMz"]=(searchoutput["expmass"]+searchoutput["charge"]*proton_mass)/searchoutput["charge"]

            proteinsplit=searchoutput["proteins"].str.split(";",expand=True)[0].str.split("|",expand=True)
            searchoutput["PG.ProteinGroups"]=proteinsplit[1].tolist()
            searchoutput["PG.ProteinNames"]=proteinsplit[2].tolist()

            searchoutput["peptide"]=searchoutput["peptide"].str.replace("+","")

            renamedict={
                "peptide":"EG.ModifiedPeptide",
                "filename":"R.FileName",
                "charge":"FG.Charge",
                "precursor_ppm":"FG.CalibratedMassAccuracy (PPM)",
                "rt":"EG.ApexRT",
                "ion_mobility":"EG.IonMobility",
                "ms2_intensity":"FG.MS2Quantity"
            }
            dropcolumns=[
                "psm_id","num_proteins","scannr","rank","label","peptide_len","missed_cleavages","semi_enzymatic","isotope_error","fragment_ppm",
                "hyperscore","delta_next","delta_best","aligned_rt","predicted_rt","delta_rt_model","predicted_mobility","delta_mobility","matched_peaks",
                "longest_b","longest_y","longest_y_pct","matched_intensity_pct","scored_candidates","poisson","sage_discriminant_score","posterior_error",
                "spectrum_q","peptide_q","protein_q","Rec","proteins","expmass","calcmass"
            ]
        if software=="sage_lfq":
            searchoutput=searchoutput.melt(["peptide","charge","proteins","q_value","score","spectral_angle"],var_name="R.FileName",value_name="FG.MS2Quantity")
            proteinsplit=searchoutput["proteins"].str.split(";",expand=True)[0].str.split("|",expand=True)
            searchoutput["PG.ProteinGroups"]=proteinsplit[1].tolist()
            searchoutput["PG.ProteinNames"]=proteinsplit[2].tolist()

            searchoutput["peptide"]=searchoutput["peptide"].str.replace("+","")
            renamedict={
                "peptide":"EG.ModifiedPeptide",
                "charge":"FG.Charge"
            }
            dropcolumns=[
                "q_value","score","spectral_angle","proteins"
            ]
        if software=="bps_timsrescore":
            if software_bps_report_type=="qual":
                #filter out any reverse or contaminant sequences
                searchoutput=searchoutput[(searchoutput["protein_list"].str.contains("Reverse")==False)&(searchoutput["protein_list"].str.contains("contaminant")==False)&(searchoutput["protein_list"].str.contains("con_")==False)].reset_index(drop=True)

                proteingroups=[]
                proteinnames=[]
                proteinlist_column=searchoutput["protein_list"].tolist()
                for item in proteinlist_column:
                    if item.count(";")==0:
                        templist=item.split("|")
                        proteingroups.append(templist[1])
                        proteinnames.append(templist[2])
                    else:
                        proteingroups_torejoin=[]
                        proteinnames_torejoin=[]
                        for entry in item.split(";"):
                            templist=entry.split("|")
                            proteingroups_torejoin.append(templist[1])
                            proteinnames_torejoin.append(templist[2])
                        proteingroups.append(";".join(proteingroups_torejoin))
                        proteinnames.append(";".join(proteinnames_torejoin))
                searchoutput["PG.ProteinGroups"]=proteingroups
                searchoutput["PG.ProteinNames"]=proteinnames
                renamedict={
                    "sample_name":"R.FileName",
                    "stripped_peptide":"PEP.StrippedSequence",
                    "precursor_mz":"FG.PrecMz",
                    "rt":"EG.ApexRT",
                    "charge":"FG.Charge",
                    "ook0":"EG.IonMobility",
                    "ppm_error":"FG.CalibratedMassAccuracy (PPM)"
                }
                dropcolumns=[
                    "index","processing_run_uuid","ms2_id","candidate_id","protein_group_parent_id","protein_group_name","leading_aa","trailing_aa","mokapot_psm_score",
                    "mokapot_psm_qvalue","mokapot_psm_pep","mokapot_peptide_qvalue","mokapot_peptide_pep","global_peptide_score","x_corr_score","delta_cn_score",
                    "precursor_mh","calc_mh","protein_list","is_contaminant","is_target","number_matched_ions","global_peptide_qvalue"
                ]
            if software_bps_report_type=="quant":
                renamedict={
                    "sample_name":"R.FileName",
                    "stripped_peptide":"PEP.StrippedSequence",
                    "quantity":"FG.MS2Quantity",
                    "precursor_mz":"FG.PrecMz",
                    "charge":"FG.Charge",
                    "rt_apex":"EG.ApexRT",
                    "ook0":"EG.IonMobility"
                }
                dropcolumns=[
                    "index","maxlfq_peptide_index","processing_run_uuid","maxlfq_run_index","candidate_id","protein_group_name","maxlfq_protein_group_index",
                    "global_peptide_qvalue","psm_pep","peptide_pep","is_target","is_contaminant","fwhm_rt_start","fwhm_rt_end"
                ]
        if software=="bps_timsdiann":
            searchoutput["EG.PeakWidth"]=searchoutput["rt_end"]-searchoutput["rt_start"]
            renamedict={
                "sample_name":"R.FileName",
                "stripped_peptide":"PEP.StrippedSequence",
                "precursor_quantity":"FG.MS2Quantity",
                "precursor_mz":"FG.PrecMz",
                "ppm_error":"FG.CalibratedMassAccuracy (PPM)",
                "precursor_charge":"FG.Charge",
                "rt":"EG.ApexRT",
                "measured_ook0":"EG.IonMobility"
            }
            dropcolumns=[
                "index","processing_run_id","precursor_quantity_normalised","precursor_ms1_quantity","precursor_q_value","precursor_quant_quality","calc_mz",
                "irt","predicted_irt","predicted_rt","library_ook0","is_proteotypic","identified_by","precursor_ms1_profile_corr","precursor_fwhm","evidence",
                "c_score","decoy_evidence","decoy_c_score","ms2_scan","fragment_id","fragment_mz","fragment_correlation","fragment_quantity","fragment_quantity_corrected",
                "rt_start","rt_end","Rec"
            ]
        if software=="bps_spectronaut":
            searchoutput["EG.PeakWidth"]=searchoutput["rt_end"]-searchoutput["rt_start"]
            renamedict={
                "sample_name":"R.FileName",
                "stripped_peptide":"PEP.StrippedSequence",
                "precursor_quantity_ms2":"FG.MS2Quantity",
                "precursor_mz":"FG.PrecMz",
                "mass_accuracy":"FG.CalibratedMassAccuracy (PPM)",
                "precursor_charge":"FG.Charge",
                "rt":"EG.ApexRT",
                "measured_ook0":"EG.IonMobility"
            }
            dropcolumns=[
                "index","processing_run_id","protein_group_name","precursor_quantity","precursor_quantity_ms1","precursor_q_value","c_score","c_score_normalised","precursor_mz_calibrated",
                "rt_start","rt_end","irt","predicted_irt","predicted_rt","is_proteotypic","is_imputed","shape_quality","precursor_fwhm","pep_score","Rec"
            ]
        if software=="bps_pulsar":
            renamedict={
                "sample_name":"R.FileName",
                "stripped_peptide":"PEP.StrippedSequence",
                "precursor_quantity":"FG.MS2Quantity",
                "psm_measured_ms1_mz":"FG.PrecMz",
                "precursor_charge":"FG.Charge",
                "psm_rt_apex":"EG.ApexRT",
                "psm_ook0_apex":"EG.IonMobility",
                "psm_ppmerror":"FG.CalibratedMassAccuracy (PPM)"
            }
            dropcolumns=[
                "index","processing_run_id","precursor_id","ms2_frame_id","protein_group_name","ptm_confidence_index","ptm_confidence_values",
                "ptm_confidence_locations","peptide_quantity","psm_q_value","peptide_q_value","psm_n_matched_ms2_ions","peptide_molecular_weight",
                "peptide_n_missed_cleavages","peptide_score","precursor_iso_quantities","precursor_used_for_quant","precursor_theoretical_mz",
                "psm_delta_ms1_mz","psm_ms1_rt","psm_ook0","psm_fwhm","psm_localization_confidence","is_proteotypic","psm_pep_score","psm_score","psm_svalue"
            ]
        if software=="bps_sage":
            searchoutput["peptide"]=searchoutput["peptide"].str.replace("+","")

            renamedict={
                "peptide":"EG.ModifiedPeptide",
                #"proteins":"",
                "filename":"R.FileName",
                "expmass":"FG.PrecMz",
                "charge":"FG.Charge",
                "precursor_ppm":"FG.CalibratedMassAccuracy (PPM)",
                "rt":"EG.ApexRT",
                "ion_mobility":"EG.IonMobility",
                "ms2_intensity":"FG.MS2Quantity",
                "protein_q":"PG.Qvalue",
                "peptide_q":"EG.Qvalue"
            }
            dropcolumns=[
                "psm_id","num_proteins","scannr","rank","label","calcmass","peptide_len","missed_cleavages","semi_enzymatic",
                "isotope_error","fragment_ppm","hyperscore","delta_next","delta_best","aligned_rt","predicted_rt","delta_rt_model",
                "predicted_mobility","delta_mobility","matched_peaks","longest_b","longest_y","longest_y_pct","matched_intensity_pct",
                "scored_candidates","poisson","sage_discriminant_score","posterior_error","spectrum_q", 
                "dl_predicted_rt","dl_rt_diff","dl_predicted_im","dl_im_diff","dl_pearson_all","dl_pearson_b","dl_pearson_y",
                "dl_mse_all","dl_mse_b","dl_mse_y","dl_min_abs_diff_norm","dl_max_abs_diff_norm","dl_abs_diff_q1_norm",
                "dl_abs_diff_q2_norm","dl_abs_diff_q3_norm","dl_mean_abs_diff_norm","dl_std_abs_diff_norm","dl_ionb_min_abs_diff_norm",
                "dl_ionb_max_abs_diff_norm","dl_ionb_abs_diff_q1_norm","dl_ionb_abs_diff_q2_norm","dl_ionb_abs_diff_q3_norm",
                "dl_ionb_mean_abs_diff_norm","dl_ionb_std_abs_diff_norm","dl_iony_min_abs_diff_norm","dl_iony_max_abs_diff_norm",
                "dl_iony_abs_diff_q1_norm","dl_iony_abs_diff_q2_norm","dl_iony_abs_diff_q3_norm","dl_iony_mean_abs_diff_norm",
                "dl_iony_std_abs_diff_norm","dl_dot_product_all","dl_dot_product_b","dl_dot_product_y","dl_cos_similarity_all",
                "dl_cos_similarity_b","dl_cos_similarity_y","dl_pearson_all_unlog","dl_pearson_b_unlog","dl_pearson_y_unlog",
                "dl_spearman_all","dl_spearman_b","dl_spearman_y","dl_mse_all_unlog","dl_mse_b_unlog","dl_mse_y_unlog",
                "dl_min_abs_diff_ion_type","dl_max_abs_diff_ion_type","dl_min_abs_diff","dl_max_abs_diff","dl_abs_diff_q1",
                "dl_abs_diff_q2","dl_abs_diff_q3","dl_mean_abs_diff","dl_std_abs_diff","dl_ionb_min_abs_diff","dl_ionb_max_abs_diff",
                "dl_ionb_abs_diff_q1","dl_ionb_abs_diff_q2","dl_ionb_abs_diff_q3","dl_ionb_mean_abs_diff","dl_ionb_std_abs_diff",
                "dl_iony_min_abs_diff","dl_iony_max_abs_diff","dl_iony_abs_diff_q1","dl_iony_abs_diff_q2","dl_iony_abs_diff_q3",
                "dl_iony_mean_abs_diff","dl_iony_std_abs_diff","dl_dot_product_all_unlog","dl_dot_product_b_unlog",
                "dl_dot_product_y_unlog","dl_cos_similarity_all_unlog","dl_cos_similarity_b_unlog","dl_cos_similarity_y_unlog"
            ]
        if software=="bps_denovo":
            renamedict={
                "sample_name":"R.FileName",
                "stripped_peptide":"PEP.StrippedSequence",
                "precursor_mz":"FG.PrecMz",
                "charge":"FG.Charge",
                "ppm_error":"FG.CalibratedMassAccuracy (PPM)",
                "rt":"EG.ApexRT",
                "ook0":"EG.IonMobility",
                "precursor_intensity":"FG.MS2Quantity"
            }
            dropcolumns=[
                "index","processing_run_id","ms2_id","rank","leading_aa","trailing_aa","precursor_mh","calc_mh","denovo_tag_length",
                "found_in_dbsearch","denovo_matches_db","protein_group_parent_loc","protein_group_parent_list_loc","protein_group_parent_list_id"
            ]
        if software=="fragpipe_glyco":
            #if the Spectrum File column is just a single value, get the file names from the Spectrum column
            if searchoutput["Spectrum"][0].split(".")[0] not in searchoutput["Spectrum File"][0]:
                fragger_filelist=searchoutput["Spectrum"].str.split(".",expand=True).drop(columns=[1,2,3]).drop_duplicates().reset_index(drop=True)
                fragger_filelist.rename(columns={0:"R.FileName"},inplace=True)

                filenamelist=[]
                for run in fragger_filelist["R.FileName"]:
                    fileindex=fragger_filelist[fragger_filelist["R.FileName"]==run].index.values[0]
                    filenamelist.append([fragger_filelist["R.FileName"][fileindex]]*len(searchoutput[searchoutput["Spectrum"].str.contains(run)]))

                searchoutput.insert(0,"R.FileName",list(itertools.chain(*filenamelist)))
                searchoutput.drop(columns=["Spectrum File"],inplace=True)
            else:
                searchoutput.rename(columns={"Spectrum File":"R.FileName"},inplace=True)
            
            #check if intensity column is just zeroes
            if len(searchoutput["Intensity"].drop_duplicates())==1:
                searchoutput.drop(columns=["Intensity"],inplace=True)
            else:
                searchoutput.rename(columns={"Intensity":"FG.MS2Quantity"},inplace=True)

            searchoutput["Retention"]=searchoutput["Retention"]/60
            searchoutput["Is Unique"]=searchoutput["Is Unique"].astype(str)

            #remove the "% glycan m/z" from the Total Glycan Composition
            searchoutput["Total Glycan Composition"]=searchoutput["Total Glycan Composition"].str.split("%",expand=True)[0]

            renamedict={
                "Peptide":"PEP.StrippedSequence",
                "Modified Peptide":"EG.ModifiedPeptide",
                "Charge":"FG.Charge",
                "Retention":"EG.ApexRT",
                "Observed M/Z":"FG.PrecMz",
                "Ion Mobility":"EG.IonMobility",
                "Protein ID":"PG.ProteinGroups",
                "Entry Name":"PG.ProteinNames",
                "Gene":"PG.Genes"
            }
            dropcolumns=[
                "Spectrum","Extended Peptide","Prev AA","Next AA","Peptide Length","Observed Mass","Calibrated Observed Mass",
                "Calibrated Observed M/Z","Calculated Peptide Mass","Calculated M/Z","Delta Mass","Expectation","Hyperscore","Nextscore",
                "Probability","Number of Enzymatic Termini","Number of Missed Cleavages","Protein Start","Protein End",
                "MSFragger Localization","Number Best Positions","Shifted Only Position Scores","Shifted Only Position Ions",
                "Score Best Position","Ions Best Position","Score Second Best Position","Ions Second Best Position","Score All Unshifted",
                "Ions All Unshifted","Score Shifted Best Position","Ions Shifted Best Position","Score Shifted All Positions",
                "Ions Shifted All Positions","Purity","Protein","Mapped Genes","Mapped Proteins"
            ]
        if software=="glycoscape":
            renamedict={
                "sample_name":"R.FileName",
                "stripped_peptide":"PEP.StrippedSequence",
                "observed_precursor_mz":"FG.PrecMz",
                "precursor_charge":"FG.Charge",
                "rt":"EG.ApexRT",
                "ook0":"EG.IonMobility",
                "peptide_ppm_error":"FG.CalibratedMassAccuracy (PPM)"
            }
            dropcolumns=[
                "index","processing_run_uuid","ms2_id","peptide_candidate_id","glycan_candidate_id","protein_group_parent_id","protein_group_name",
                "leading_aa","trailing_aa","hexnac_modification","glycosylation_motif","is_contaminant","is_target","peptide_mh","peptide_calc_mh",
                "peptide_isotope_offset","x_corr_score","delta_cn_score","number_matched_ions","mokapot_psm_score","mokapot_psm_pep","mokapot_peptide_qvalue",
                "mokapot_peptide_pep","global_peptide_score","Y1_mz","Y1_charge","experimental_glycan_mr","glycan_isotope_offset","glycan_composition_mass",
                "filtered_glycan_rank","ambiguous_glycan_composition","glycan_rank","building_blocks_coverage","fucose_evidence","Y5Y1_evidence",
                "has_core","sia_smaller_hn","composition_oxonium_count","composition_oxonium_intensity","spectrum_oxonium_count","oxonium_relative_intensity_sum",
                "fucose_shadow_count","fucose_shadow_intensity_sum","bb_names","bb_oxonium_count","bb_oxonium_intensity","oxonium_ions_names","oxonium_ions_mzs",
                "oxonium_ions_intensity","Y_ions_names","oxonium_relative_intensity_sum","Y_ions_mass_offset","Y_ions_intensity","extra_ions_names",
                "extra_ions_mass_offset","extra_ions_intensity","Rec"
            ]
        #rename columns based on dictionary assignments per software
        if renamedict!={}:
            searchoutput=searchoutput.rename(columns=renamedict)
        #drop listed columns
        if dropcolumns!=[]:
            searchoutput=searchoutput.drop(columns=dropcolumns,errors="ignore")
        #filter out NaNs from PG.MS2Quantity and FG.MS2Quantity columns
        if "PG.MS2Quantity" in searchoutput.columns:
            searchoutput=searchoutput[searchoutput["PG.MS2Quantity"].isnull()==False].reset_index(drop=True)
        if "FG.MS2Quantity" in searchoutput.columns:
            searchoutput=searchoutput[searchoutput["FG.MS2Quantity"].isnull()==False].reset_index(drop=True)
        #build and add EG.ModifiedPeptide column (for BPS inputs with ptm_locations and ptms columns)
        if "EG.ModifiedPeptide" not in searchoutput.columns:
            searchoutput["ptm_locations"]=searchoutput["ptm_locations"].astype(str)
            searchoutput["ptms"]=searchoutput["ptms"].astype(str)
            searchoutput["ptm_locations"]=searchoutput["ptm_locations"].str.replace("[]","-1").str.replace("[","").str.replace("]","")
            modifiedpeptides=[]
            for i,entry in enumerate(searchoutput["ptm_locations"]):
                if entry=="-1":
                    modifiedpeptides.append(searchoutput["PEP.StrippedSequence"][i])
                else:
                    str_to_list=list(searchoutput["PEP.StrippedSequence"][i])
                    if len(searchoutput["ptm_locations"][i])==1:
                        mod_loc=int(searchoutput["ptm_locations"][i])+1
                        mod_add=searchoutput["ptms"][i]
                        str_to_list.insert(mod_loc,mod_add)
                        modifiedpeptides.append("".join(str_to_list))
                    #if theres >1 ptm, we need to reformat some strings so we can insert them in the sequence
                    else:
                        ptmlocs=searchoutput["ptm_locations"][i].strip().split(" ")
                        ptms=searchoutput["ptms"][i].replace("[","").replace("]","").split(" ")
                        for ele in ptmlocs:
                            if ele=="":
                                ptmlocs.remove(ele)
                        ptms_for_loop=[]
                        for ele in ptms:
                            ele=ele.strip()
                            ptms_for_loop.append("["+ele+"]")
                        for j,loc in enumerate(ptmlocs):
                            mod_loc=int(loc)+j+1
                            mod_add=ptms_for_loop[j]
                            str_to_list.insert(mod_loc,mod_add)
                        modifiedpeptides.append("".join(str_to_list))
            searchoutput["EG.ModifiedPeptide"]=modifiedpeptides
        #build and add PEP.StrippedSequence column (for bps_sage, sage_lfq, sage, and peaks)
        if "PEP.StrippedSequence" not in searchoutput.columns:
            searchoutput["PEP.StrippedSequence"]=[re.sub(r"\(.*?\)|\[.*?\]","",pep) for pep in searchoutput["EG.ModifiedPeptide"].tolist()]
        #filter out decoy and contaminant sequences
        if "PG.ProteinGroups" in searchoutput.columns:
            searchoutput=searchoutput[(searchoutput["PG.ProteinGroups"].str.contains("Reverse")==False)&(searchoutput["PG.ProteinGroups"].str.contains("contaminant")==False)&(searchoutput["PG.ProteinGroups"].str.contains("con_")==False)].reset_index(drop=True)
        #replace None values in protein name column, remove leading semicolon from protein name entries
        if "PG.ProteinNames" in searchoutput.columns:
            if len(searchoutput[searchoutput["PG.ProteinNames"].isnull()])>0:
                df=searchoutput[searchoutput["PG.ProteinNames"].isnull()]
                for index in df.index.tolist():
                    if searchoutput.at[index,"PG.ProteinNames"]==None or math.isnan(searchoutput.at[index,"PG.ProteinNames"])==True:
                        searchoutput.at[index,"PG.ProteinNames"]=searchoutput.at[index,"PG.ProteinGroups"]
            if len(searchoutput[searchoutput["PG.ProteinNames"].str.startswith(";")])>0:
                df=searchoutput[searchoutput["PG.ProteinNames"].str.startswith(";")]
                for index in df.index.tolist():
                    searchoutput.at[index,"PG.ProteinNames"]=searchoutput.at[index,"PG.ProteinNames"].strip(";")
        #replace PTM identifiers using ptm dictionary
        searchoutput["EG.ModifiedPeptide"]=searchoutput["EG.ModifiedPeptide"].replace(ptmdict,regex=True)
        #because of the way ptmdict works, un-converted PTM names from FragPipe will not be detectable in find_ptms(). Re-add [ ] for un-converted PTM names
        if software=="fragpipe":
            for i,pep in enumerate(searchoutput["EG.ModifiedPeptide"].tolist()):
                #look for numeric characters (unsubstituted ptm values)
                ptms=re.findall(r'\d+',pep)
                #if there are no ptms, then no replacement should be needed
                if len(ptms)==0:
                    pass
                    #searchoutput.iloc[i,searchoutput.columns.get_loc("EG.ModifiedPeptide")]=pep
                elif len(ptms)==1:
                    searchoutput.iloc[i,searchoutput.columns.get_loc("EG.ModifiedPeptide")]=pep.replace(re.search(ptms[0],pep)[0],"["+re.search(ptms[0],pep)[0]+"]")
                else:
                    modlist=list(set(re.findall(r'\d+',pep)))
                    for mod in modlist:
                        pep=pep.replace(re.search(mod,pep)[0],"["+re.search(mod,pep)[0]+"]")
                    searchoutput.iloc[i,searchoutput.columns.get_loc("EG.ModifiedPeptide")]=pep
        #some software add a peptide length column, adding universally here
        searchoutput["Peptide Length"]=[len(pep) for pep in searchoutput["PEP.StrippedSequence"]]
        #add Condition and Replicate columns for filling in from metadata
        if "R.Condition" not in searchoutput.columns:
            searchoutput.insert(1,"R.Condition",["Not Defined"]*len(searchoutput))
        #updated searchoutput sorting, sorts by the last number in the file name, which should be the run ID from HyStar
        s=searchoutput.sort_values("R.FileName")["R.FileName"].drop_duplicates().tolist()
        sortedfilenames=sorted(s,key=lambda x: int(re.findall(r'\d+',x)[-1]))
        searchoutput["sortedfilenames"]=pd.Categorical(searchoutput["R.FileName"],categories=sortedfilenames,ordered=True)
        searchoutput=searchoutput.sort_values(["R.Condition","sortedfilenames"])
        searchoutput=searchoutput.drop(columns=["sortedfilenames"]).reset_index(drop=True)
        if "R.Replicate" not in searchoutput.columns:
            replicateslist=[]
            for i,file in enumerate(searchoutput["R.FileName"].drop_duplicates().tolist()):
                replicateslist.append([i+1]*len(searchoutput[searchoutput["R.FileName"]==file]))
            replicateslist=list(itertools.chain(*replicateslist))
            searchoutput.insert(2,"R.Replicate",replicateslist)
        return searchoutput

    # ====================================== Callable metadata functions
    #general function to hold the metadata_reminder text to extend to De Novo instead of using a new function
    def metadata_reminder_text(software):
        if software=="spectronaut":
            text="Use the timsplot_spectronaut_report format from GitHub when exporting search results to a .tsv file."
        if software=="diann":
            text="Use the report.tsv file as the file input."
        if software=="diann2.0":
            text="Use the report.parquet file as the file input."
        if software=="ddalibrary":
            text="DDA libraries have limited functionality, can only plot ID metrics."
        if software=="fragpipe":
            text="Use the psm.tsv file as the file input."
        if software=="fragpipe_glyco":
            text="Use the psm.tsv file as the file input. Use the Glycoproteomics tab for processing."
        if software=="fragpipe_combined_ion":
            text="Upload combined_ion.tsv and combined_protein.tsv files from FragPipe output to be joined."
        if software=="bps_timsrescore" or software=="bps_timsdiann" or software=="bps_pulsar" or software=="bps_spectronaut":
            text="Use the .zip file from the artifacts download. File upload can take several minutes because of how protein information is mapped to the peptide file. Make sure to periodically check the timsplot folder and remove metadata.csv and the processing-run folder to free up space."
        if software=="bps_denovo":
            text="Use the .zip file from the artifacts download. Denovo score=65.0 and rank=1 are used to filter IDs on upload. Precursors identified by denovo with no protein grouping are all grouped under 'uncat' when counting unique proteins."
        if software=="bps_sage":
            text="Use the results.tsv file as the file input. As of latest version, most plotting functions are unavailable because protein group names are not present in the Sage output."
        if software=="glycoscape":
            text="Use the .zip file from the artifacts download. Use the Glycoproteomics tab for processing."
        if software=="spectromine":
            text="Use the timsplot_spectromine_report format from GitHub when exporting search results to a .tsv file. NOTE: Spectromine outputs quant as protein group and precursor label-free quant which is MS1 level, not MS2. For simplicity and consistency across the app, these are renamed to PG.MS2Quantity and FG.MS2Quantity."
        if software=="peaks":
            text="Export results from the Peptides menu and upload the db.psms.csv file. NOTE: Peaks outputs a peptide-level peak area that will be renamed to FG.MS2Quantity for consistency across the app."
        if software=="sage":
            text="Use the results.sage.tsv file as the file input."
        if software=="sage_lfq":
            text="Use the lfq.tsv file as the file input."
        return text
    #generate metadata df
    def metadata_gen(use_uploaded_metadata,metadata_upload,searchoutput):
        if use_uploaded_metadata==True:
            if metadata_upload is None:
                metadata=pd.DataFrame(columns=["R.FileName","R.Condition","R.Replicate","remove"])
            else:
                metadata=pd.read_csv(metadata_upload[0]["datapath"],sep=",")
                metadata["remove"]=metadata["remove"].fillna("")
        else:
            #fileupload outputs an empty df if input.searchreport() is None
            if searchoutput.empty==True:
                metadata=pd.DataFrame(columns=["R.FileName","R.Condition","R.Replicate","remove"])
            else:
                metadata=searchoutput[["R.FileName","R.Condition","R.Replicate"]].drop_duplicates().reset_index(drop=True)
                metadata["remove"]=[""]*len(metadata["R.FileName"])
        return metadata
    #generate metadata_condition df
    def metadata_condition_gen(metadata,metadata_calc,use_uploaded_metadata,remove):
        #metadata is the table view, metadata_calc is related to the uploaded metadata sheet
        if len(metadata)==0:
            metadata_condition=pd.DataFrame(columns=["R.Condition","order","Concentration"])
        else:
            metadata_condition=pd.DataFrame(columns=["R.Condition","order","Concentration"])
            #need uploaded_metadata as an arg because of how remove and concentration columns are handled
            if use_uploaded_metadata==True:
                if remove==True:
                    metadata_calc=metadata_calc[metadata_calc["remove"]!="x"]
                    metadata_condition["order"]=np.arange(1,len(metadata_calc["R.Condition"].drop_duplicates())+1,1)
                else:
                    metadata_condition["order"]=metadata_calc[["order"]].drop_duplicates().reset_index(drop=True)
                metadata_condition["R.Condition"]=metadata_calc[["R.Condition"]].drop_duplicates().reset_index(drop=True)
                metadata_condition["Concentration"]=metadata_calc[["Concentration"]].drop_duplicates().reset_index(drop=True)
            else:
                if remove==True:
                    metadata=metadata[metadata["remove"]!="x"]
                metadata_condition["R.Condition"]=metadata[["R.Condition"]].drop_duplicates().reset_index(drop=True)
                metadata_condition["order"]=np.arange(1,len(metadata_condition["R.Condition"])+1,1)
                metadata_condition["Concentration"]=[""]*len(metadata_condition["R.Condition"])
        return metadata_condition
    #generate df for downloading metadata
    def metadata_download_gen(metadata,metadata_condition):
        metadata_download=pd.DataFrame()
        metadata_download["R.FileName"]=metadata["R.FileName"]
        metadata_download["R.Condition"]=metadata["R.Condition"]
        metadata_download["R.Replicate"]=metadata["R.Replicate"]
        metadata_download["remove"]=metadata["remove"]

        #check if conditions have been removed from the main metadata df
        orderlist=[]
        concentrationlist=[]
        if len(metadata["R.Condition"].drop_duplicates().tolist())!=len(metadata_condition["R.Condition"].tolist()):
            for condition in metadata["R.Condition"].drop_duplicates().tolist():
                if condition not in metadata_condition["R.Condition"].tolist():
                    concentrationlist.append([""]*len(metadata[metadata["R.Condition"]==condition]["R.Replicate"]))
                    orderlist.append([0]*len(metadata[metadata["R.Condition"]==condition]["R.Replicate"]))
                if condition in metadata_condition["R.Condition"].tolist():
                    orderlist.append([metadata_condition[metadata_condition["R.Condition"]==condition]["order"].values[0]]*len(metadata[metadata["R.Condition"]==condition]["R.Replicate"]))
                    concentrationlist.append([metadata_condition[metadata_condition["R.Condition"]==condition]["Concentration"].values[0]]*len(metadata[metadata["R.Condition"]==condition]["R.Replicate"]))
        else:
            for run in metadata_condition["R.Condition"]:
                fileindex=metadata_condition[metadata_condition["R.Condition"]==run].index.values[0]
                orderlist.append([metadata_condition["order"][fileindex]]*len(metadata[metadata["R.Condition"]==run]))
                concentrationlist.append([metadata_condition["Concentration"][fileindex]]*len(metadata[metadata["R.Condition"]==run]))
        metadata_download["order"]=list(itertools.chain(*orderlist))
        metadata_download["Concentration"]=list(itertools.chain(*concentrationlist))
        return metadata_download
    #update uploaded search file based on metadata and switches
    def metadata_update_gen(searchoutput,metadata,metadata_condition,remove,cRAP_filter):
        #make the original metadata table based on the search report to compare to the metadata_table.data_view() to check if it has been edited
        #include what the remove column currently looks like in the metadata table
        metadata_original=pd.DataFrame(searchoutput[["R.FileName","R.Condition","R.Replicate"]]).drop_duplicates().reset_index(drop=True)
        metadata_original["remove"]=metadata["remove"]

        #remove checked runs from search file
        if remove==True:
            metadata=metadata[metadata["remove"]!="x"].reset_index(drop=True)
            metadata_original=metadata_original[metadata_original["remove"]!="x"].reset_index(drop=True)
            searchoutput=searchoutput.set_index("R.FileName").loc[metadata["R.FileName"].tolist()].reset_index()

        #check if the original metadata sheet and metadata sheet's current view are different
        #update condition/replicate names/values from metadata. if not filled out, autofill with generic values
        #if yes, rename/renumber conditions and replicates
        if metadata_original.compare(metadata).empty==False:
            RConditionlist=[]
            RReplicatelist=[]
            for i,run in enumerate(searchoutput["R.FileName"].drop_duplicates().tolist()):
                fileindex=metadata[metadata["R.FileName"]==run].index.values[0]
                if metadata["R.Condition"][fileindex]=="":
                    RConditionlist.append(["Not Defined"]*len(searchoutput.set_index("R.FileName").loc[run]))
                else:
                    RConditionlist.append([metadata["R.Condition"][fileindex]]*len(searchoutput.set_index("R.FileName").loc[run]))
                if metadata["R.Replicate"][fileindex]=="":
                    RReplicatelist.append([i+1]*len(searchoutput.set_index("R.FileName").loc[run]))
                else:
                    RReplicatelist.append([metadata["R.Replicate"][fileindex]]*len(searchoutput.set_index("R.FileName").loc[run]))
            searchoutput["R.Condition"]=list(itertools.chain(*RConditionlist))
            searchoutput["R.Replicate"]=list(itertools.chain(*RReplicatelist))
            searchoutput["R.Replicate"]=searchoutput["R.Replicate"].astype(int)

        #filter protein names with cRAP database
        if cRAP_filter==True:
            cRAP=['ADH1_YEAST','ALBU_BOVIN','ALBU_HUMAN','ALDOA_RABIT','AMYS_HUMAN','ANT3_HUMAN','ANXA5_HUMAN','B2MG_HUMAN','BGAL_ECOLI','BID_HUMAN','CAH1_HUMAN','CAH2_BOVIN','CAH2_HUMAN',
                  'CAS1_BOVIN','CAS2_BOVIN','CASB_BOVIN','CASK_BOVIN','CATA_HUMAN','CATD_HUMAN','CATG_HUMAN','CO5_HUMAN','CRP_HUMAN','CTRA_BOVIN','CTRB_BOVIN','CYB5_HUMAN','CYC_HORSE',
                  'CYC_HUMAN','DHE3_BOVIN','EGF_HUMAN','FABPH_HUMAN','GAG_SCVLA','GELS_HUMAN','GFP_AEQVI','GSTA1_HUMAN','GSTP1_HUMAN','HBA_HUMAN','HBB_HUMAN','IGF2_HUMAN','IL8_HUMAN',
                  'K1C10_HUMAN','K1C15_SHEEP','K1C9_HUMAN','K1H1_HUMAN','K1H2_HUMAN','K1H4_HUMAN','K1H5_HUMAN','K1H6_HUMAN','K1H7_HUMAN','K1H8_HUMAN','K1HA_HUMAN','K1HB_HUMAN','K1M1_SHEEP',
                  'K1M2_SHEEP','K22E_HUMAN','K2C1_HUMAN','K2M1_SHEEP','K2M2_SHEEP','K2M3_SHEEP','KCRM_HUMAN','KRA3_SHEEP','KRA33_SHEEP','KRA34_SHEEP','KRA3A_SHEEP','KRA61_SHEEP','KRB2A_SHEEP',
                  'KRB2B_SHEEP','KRB2C_SHEEP','KRB2D_SHEEP','KRHB1_HUMAN','KRHB2_HUMAN','KRHB3_HUMAN','KRHB4_HUMAN','KRHB5_HUMAN','KRHB6_HUMAN','KRUC_SHEEP','LALBA_BOVIN','LALBA_HUMAN',
                  'LEP_HUMAN','LYSC_CHICK','LYSC_HUMAN','LYSC_LYSEN','MYG_HORSE','MYG_HUMAN','NEDD8_HUMAN','NQO1_HUMAN','NQO2_HUMAN','OVAL_CHICK','PDGFB_HUMAN','PEPA_BOVIN','PEPA_PIG',
                  'PEPB_PIG','PEPC_PIG','PLMP_GRIFR','PPIA_HUMAN','PRDX1_HUMAN','RASH_HUMAN','REF_HEVBR','RETBP_HUMAN','RS27A_HUMAN','SODC_HUMAN','SRPP_HEVBR','SSPA_STAAU','SUMO1_HUMAN',
                  'SYH_HUMAN','TAU_HUMAN','THIO_HUMAN','TNFA_HUMAN','TRFE_HUMAN','TRFL_HUMAN','TRY1_BOVIN','TRY2_BOVIN','TRYP_PIG','UB2E1_HUMAN','UBE2C_HUMAN','UBE2I_HUMAN']
            decoylist=[]
            for protein in searchoutput["PG.ProteinNames"].drop_duplicates().tolist():
                if protein in cRAP:
                    decoylist.append(protein)
            searchoutput=searchoutput.set_index("PG.ProteinNames").drop(decoylist).reset_index()

        #handling for if there are still blanks in metadata_condition R.Condition column
        if "" in metadata_condition["R.Condition"].tolist():
            metadata_condition.loc[metadata_condition.index[metadata_condition["R.Condition"]==""].values[0],"R.Condition"]="Not Defined"

        #add concentration values
        if metadata_condition["Concentration"].drop_duplicates().tolist()!=[""]:
            concentrationlist=[]
            for run in searchoutput["R.Condition"].drop_duplicates().tolist():
                fileindex=metadata_condition[metadata_condition["R.Condition"]==run].index.values[0]
                concentrationlist.append([metadata_condition["Concentration"][fileindex]]*len(searchoutput.set_index("R.Condition").loc[run]))
            if "Concentration" in searchoutput.columns:
                searchoutput["Concentration"]=list(itertools.chain(*concentrationlist))
                searchoutput["Concentration"]=searchoutput["Concentration"].astype(float)
            else:
                searchoutput.insert(3,"Concentration",list(itertools.chain(*concentrationlist)))
                searchoutput["Concentration"]=searchoutput["Concentration"].astype(float)

        #generate Cond_Rep column
        searchoutput["R.Condition"]=searchoutput["R.Condition"].apply(str)
        if "Cond_Rep" not in searchoutput.columns:
            searchoutput.insert(0,"Cond_Rep",searchoutput["R.Condition"]+"_"+searchoutput["R.Replicate"].apply(str))
        elif "Cond_Rep" in searchoutput.columns:
            searchoutput["Cond_Rep"]=searchoutput["R.Condition"]+"_"+searchoutput["R.Replicate"].apply(str)
        
        #reorder
        metadata_condition["order"]=metadata_condition["order"].astype(int)
        if len(metadata_condition["order"])==1:
            pass
        else:
            sortedmetadata_bycondition=metadata_condition.sort_values(by="order").reset_index(drop=True)
            searchoutput=searchoutput.set_index("R.Condition").loc[sortedmetadata_bycondition["R.Condition"].tolist()].reset_index()
        return searchoutput

    # ====================================== UI and calc functions (this section is all that's needed to transfer to De Novo and Two-Software Comparison)
    #store uploaded file as a reactive value so we don't need to call inputfile function each time we need to access the uploaded file
    @reactive.calc
    def fileupload():
        searchoutput=inputfile(input.searchreport(),input.software(),input.software_bps_report_type(),input.diann_mbr_switch(),input.searchreport_reupload())
        return searchoutput
    #give a reminder for what to do with search reports from different software
    @render.text
    def metadata_reminder():
        text=metadata_reminder_text(input.software())
        return text
    #store generated metadata df as a reactive value
    @reactive.calc
    def metadata_calc():
        metadata=metadata_gen(input.use_uploaded_metadata(),input.metadata_upload(),fileupload())
        return metadata
    #render metadata table
    @render.data_frame
    def metadata_table():
        metadata=metadata_calc()
        metadata=metadata[["R.FileName","R.Condition","R.Replicate","remove"]]
        if len(metadata)==0:
            #return non-editable metadata if there's nothing there
            return render.DataGrid(metadata,width="100%")
        else:
            return render.DataGrid(metadata,editable=True,width="100%")
    #store metadata_condition df as a reactive value
    @reactive.calc
    def metadata_condition_calc():
        metadata_condition=metadata_condition_gen(metadata_table.data_view(),metadata_calc(),input.use_uploaded_metadata(),input.remove())
        return metadata_condition
    #render metadata_condition table
    @render.data_frame
    def metadata_condition_table():
        metadata_condition=metadata_condition_calc()
        if len(metadata_condition)==0:
            return render.DataGrid(metadata_condition,width="100%")
        else:
            return render.DataGrid(metadata_condition,editable=True,width="100%")
    #download metadata table as shown
    @render.download(filename="metadata_"+str(date.today())+".csv")
    def metadata_download():
        metadata_download=metadata_download_gen(metadata_table.data_view(),metadata_condition_table.data_view())
        with io.BytesIO() as buf:
            metadata_download.to_csv(buf,index=False)
            yield buf.getvalue()
    #update the searchoutput df to match how we edited the metadata sheet
    @reactive.calc
    @reactive.event(input.rerun_metadata,ignore_none=False)
    def metadata_update():
        searchoutput=metadata_update_gen(fileupload(),metadata_table.data_view(),metadata_condition_table.data_view(),input.remove(),input.cRAP_filter())
        return searchoutput
    #checkmark popup to show that the metadata has been applied to the search report
    @render.text
    @reactive.event(input.rerun_metadata)
    def metadata_applied_ui():
        return '\u2705'

    # ====================================== Downloads
    @render.download(filename=lambda: f"{input.searchreport()[0]['name']}"+".tsv")
    def searchreport_download():
        searchoutput=metadata_update()
        with io.BytesIO() as buf:
            searchoutput.to_csv(buf,index=False,sep="\t")
            yield buf.getvalue()
    #function for downloading figures to a folder outside of the timsplot directory
    def imagedownload(plotname):
        if input.download_path()!="":
            os_var=os_check()
            filename=input.software()+"_"+input.searchreport()[0]["name"]+"_"+plotname
            plt.savefig(input.download_path()+os_var+filename+".png")
        else:
            pass

#endregion

# ============================================================================= Settings
#region
    # ====================================== Color Settings
    #render the color grids as images
    @render.image
    def matplotcolors_image():
        os.chdir(os.path.dirname(os.path.realpath(__file__)))
        cwd=os.getcwd()
        os_var=os_check()
        foldername="images"
        joined=[cwd,foldername]
        matplotcolors_imagefile=os_var.join([cwd,"images","matplotlib_tabcolors.png"])        
        img: ImgData={"src":matplotcolors_imagefile}
        return img
    #CSS color options
    @render.image
    def csscolors_image():
        os.chdir(os.path.dirname(os.path.realpath(__file__)))
        cwd=os.getcwd()
        os_var=os_check()
        foldername="images"
        joined=[cwd,foldername]
        csscolors_imagefile=os_var.join([cwd,"images","css_colors.png"])        
        img: ImgData={"src":csscolors_imagefile}
        return img

    #color options for plotting 
    @reactive.calc
    def colorpicker():
        import seaborn as sns

        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        if input.coloroptions()=="pickrainbow":
            if numconditions==2:
                rainbowlist=cm.gist_rainbow(np.linspace(0,1,6))
                plotcolors=[]
                for i in range(len(rainbowlist)):
                    plotcolors.append(sns.desaturate(rainbowlist[i],0.75))
                indices=[0,3]
                plotcolors=[plotcolors[x] for x in indices]
            elif numconditions==1:
                plotcolors=sns.desaturate(cm.gist_rainbow(np.random.random_sample()),0.75)
            else:
                samplelinspace=np.linspace(0,1,numconditions)
                rainbowlist=cm.gist_rainbow(samplelinspace)
                plotcolors=[]
                for i in range(len(rainbowlist)):
                    plotcolors.append(sns.desaturate(rainbowlist[i],0.75))
        elif input.coloroptions()=="pickmatplot":
            matplottabcolors=list(mcolors.TABLEAU_COLORS)
            plotcolors=[]
            if numconditions > len(matplottabcolors):
                dif=numconditions-len(matplottabcolors)
                plotcolors=matplottabcolors
                for i in range(dif):
                    plotcolors.append(matplottabcolors[i])
            elif numconditions==1:
                plotcolors=matplottabcolors[np.random.randint(low=0,high=len(list(mcolors.TABLEAU_COLORS)))]
            else:
                for i in range(numconditions):
                    plotcolors.append(matplottabcolors[i])
        elif input.coloroptions()=="custom":
            if numconditions==1:
                plotcolors=input.customcolors()
            else:
                plotcolors=list(input.customcolors().split("\n"))
        return plotcolors
    #loop to extend list for replicates, used in the place of calling colorpicker() to have replicates of the same condition plotted with the same color
    @reactive.calc
    def replicatecolors():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        plotcolors=colorpicker()
        if numconditions==1 and len(maxreplicatelist)==1:
            return plotcolors
        else:
            replicateplotcolors=[]
            for i in range(numconditions):
                x=maxreplicatelist[i]
                for ele in range(x):
                    replicateplotcolors.append(plotcolors[i])
            return replicateplotcolors
    #show a table of the sample conditions 
    @render.table()
    def customcolors_table1():
        try:
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            conditioncolordf1=pd.DataFrame({"Condition":sampleconditions})
            return conditioncolordf1
        except:
            conditioncolordf1=pd.DataFrame({"Condition":[]})
            return conditioncolordf1

    #add an empty table header to match the conditions table
    @render.table
    def conditioncolors():
        conditioncolors_table=pd.DataFrame({"Color per run":[]})
        return conditioncolors_table
    
    #plot rectangles in line with the sample conditions they'll be associated with in downstream plotting
    @reactive.effect
    def _():
        #added a variable to help make sure that the rectangles that are shown for the color choices line up roughly with the table of sample conditions
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        height=40*numconditions
        @render.plot(width=75,height=height)
        def customcolors_plot():
            try:
                searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
                plotcolors=colorpicker()
                
                fig,ax=plt.subplots(nrows=len(sampleconditions))
                fig.set_tight_layout(True)
                for i in range(len(sampleconditions)):
                    if len(sampleconditions)==1:
                        rect=matplotlib.patches.Rectangle(xy=(0,0),width=0.5,height=0.5,facecolor=plotcolors)
                        ax.add_patch(rect)
                        ax.axis("off")
                        ax.set_ylim(bottom=0,top=0.5)
                        ax.set_xlim(left=0,right=0.5)
                    else:
                        rect=matplotlib.patches.Rectangle(xy=(0,0),width=0.5,height=0.5,facecolor=plotcolors[i])
                        ax[i].add_patch(rect)
                        ax[i].axis("off")
                        ax[i].set_ylim(bottom=0,top=0.5)
                        ax[i].set_xlim(left=0,right=0.5)
            except:
                pass

    # ====================================== PTM Settings
    #import ptmdict file locally (include in download from GitHub)
    @reactive.calc
    def ptmdict_setup():
        os.chdir(os.path.dirname(os.path.realpath(__file__)))
        ptmdict_download=pd.read_csv("ptmdict.csv",header=None).set_index(0)
        ptmdict=ptmdict_download.to_dict(orient="dict")[1]
        return ptmdict
    #show ptmdict as a table to show current substitution list 
    @render.table
    def ptmdict_table():
        try:
            ptmdict=ptmdict_apply()
        except:
            ptmdict=ptmdict_setup()
        ptmdict=pd.DataFrame.from_dict(ptmdict,orient="index").reset_index().rename(columns={"index":"PTM Identifier",0:"Replacement"})
        return ptmdict
    #apply any user-defined addition to the ptmdict, works even if you don't go to the PTM settings tab before plotting anything
    #to overwrite current values in the table or adding new values, will have to go back to file import and hit apply changes
    @reactive.calc
    @reactive.event(input.ptm_apply,ignore_none=False)
    def ptmdict_apply():
        ptmdict=ptmdict_setup()
        key=input.ptm_key_input()
        value=input.ptm_value_input()
        if key=="" and value=="":
            ptmdict=ptmdict
        else:
            ptmdict[key]=value
        return ptmdict
    #save any user-defined mods, overwrite the main file to do so
    @render.download(filename="ptmdict.csv")
    def ptm_save():
        ptmdict=ptmdict_apply()
        ptmdict=pd.DataFrame.from_dict(ptmdict,orient="index").reset_index()

        with io.BytesIO() as buf:
            ptmdict.to_csv(buf,index=False,header=False)
            yield buf.getvalue()

    # ====================================== File Stats
    #stats about the input file
    @render.table
    def filestats():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

        variablenames=["# samples","# conditions","Conditions","Replicates per Condition"]
        variablecalls=[numsamples,numconditions,sampleconditions,maxreplicatelist]

        filestatsdf=pd.DataFrame({"Property":variablenames,"Values":variablecalls})
        
        return filestatsdf

    # ====================================== Column Check
    #column check
    @render.table
    def column_check():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        columnnames=searchoutput.columns.tolist()
        expectedcolumns=["R.FileName","R.Condition","R.Replicate","Concentration",
                         "PG.ProteinGroups","PG.ProteinAccessions","PG.ProteinNames","PG.Genes",
                         "PEP.StrippedSequence","EG.ModifiedPeptide","FG.Charge","FG.PrecMz",
                         "PG.MS2Quantity","FG.MS2Quantity","EG.IonMobility","EG.ApexRT","EG.PeakWidth"]
        columnnames=searchoutput.columns.tolist()
        in_report=[]
        for i in expectedcolumns:
            if i in columnnames:
                in_report.append("TRUE")
            else:
                in_report.append("FALSE")
        columncheck_df=pd.DataFrame({"Expected Column":expectedcolumns,"in_report":in_report})
        return columncheck_df

    # ====================================== File Preview
    #preview of searchoutput table
    @render.data_frame
    def filepreview():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        return render.DataGrid(searchoutput,editable=False,width="100%")
 
    # ====================================== Extra Colors
    @render.image
    def brukercolors():
        os.chdir(os.path.dirname(os.path.realpath(__file__)))
        cwd=os.getcwd()
        os_var=os_check()
        foldername="images"
        joined=[cwd,foldername]
        brukercolors_imagefile=os_var.join([cwd,"images","bruker_colors.png"])        
        img: ImgData={"src":brukercolors_imagefile}
        return img
#endregion  

# ============================================================================= ID Counts
#region
    #store idmetrics call in a reactive.calc
    @reactive.calc
    def idmetrics_calc():
        resultdf,averagedf,uniquecountsdf=idmetrics(metadata_update())
        return resultdf,averagedf,uniquecountsdf
    @reactive.calc
    def cvcalc():
        cvcalc_df=cvmetrics(metadata_update(),None)
        return cvcalc_df
    # ====================================== Counts per Condition
    @render.ui
    def idmetrics_peplength_switch_ui():
        if input.idplotinput()=="peptides":
            return ui.input_switch("idmetrics_peplength_switch","Show only specific peptide length?",width="315px")
    @render.ui
    def idmetrics_peplength_ui():
        if input.idmetrics_peplength_switch()==True:
            return ui.input_slider("idmetrics_peplength_slider","Peptide length to plot for:",min=5,max=20,value=9,step=1,ticks=True,width="275px")
    @reactive.effect
    def _():
        @render.plot(width=input.idmetrics_width(),height=input.idmetrics_height())
        def idmetricsplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            resultdf,averagedf,uniquecountsdf=idmetrics_calc()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.idplotinput()=="all":
                fig,ax=plt.subplots(nrows=2,ncols=2,sharex=True,sharey="row",figsize=(input.idmetrics_width()*(1/plt.rcParams['figure.dpi']),input.idmetrics_height()*(1/plt.rcParams['figure.dpi'])))
                ax1=ax[0,0]
                ax2=ax[0,1]
                ax3=ax[1,0]
                ax4=ax[1,1]
                if input.idmetrics_peptides()=="modified":
                    pepcolumn="peptides"
                    peptitle="Peptides"
                elif input.idmetrics_peptides()=="stripped":
                    pepcolumn="strippedpeptides"
                    peptitle="Peptide Sequences"

                if input.idmetrics_individual_average()=="individual":
                    idmetricscolor=replicatecolors()
                    ax1.bar(resultdf["Cond_Rep"],resultdf["proteins"],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax2.bar(resultdf["Cond_Rep"],resultdf["proteins2pepts"],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax3.bar(resultdf["Cond_Rep"],resultdf[pepcolumn],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax4.bar(resultdf["Cond_Rep"],resultdf["precursors"],width=0.8,color=idmetricscolor,edgecolor="k")

                    ax1.bar_label(ax1.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax2.bar_label(ax2.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax3.bar_label(ax3.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax4.bar_label(ax4.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)

                    ax1.set_ylim(top=max(resultdf["proteins"].tolist())+y_padding*max(resultdf["proteins"].tolist()))
                    #ax2.set_ylim(top=max(resultdf["proteins2pepts"].tolist())+y_padding*max(resultdf["proteins2pepts"].tolist()))
                    #ax3.set_ylim(top=max(resultdf[pepcolumn].tolist())+(y_padding+0.1)*max(resultdf[pepcolumn].tolist()))
                    ax4.set_ylim(top=max(resultdf["precursors"].tolist())+(y_padding+0.1)*max(resultdf["precursors"].tolist()))
                elif input.idmetrics_individual_average()=="average":
                    idmetricscolor=colorpicker()
                    bars1=ax1.bar(averagedf["R.Condition"],averagedf["proteins_avg"],yerr=averagedf["proteins_stdev"],edgecolor="k",width=0.8,capsize=10,color=idmetricscolor)
                    bars2=ax2.bar(averagedf["R.Condition"],averagedf["proteins2pepts_avg"],yerr=averagedf["proteins2pepts_stdev"],edgecolor="k",width=0.8,capsize=10,color=idmetricscolor)
                    bars3=ax3.bar(averagedf["R.Condition"],averagedf[pepcolumn+"_avg"],yerr=averagedf[pepcolumn+"_stdev"],edgecolor="k",width=0.8,capsize=10,color=idmetricscolor)
                    bars4=ax4.bar(averagedf["R.Condition"],averagedf["precursors_avg"],yerr=averagedf["precursors_stdev"],edgecolor="k",width=0.8,capsize=10,color=idmetricscolor)

                    ax1.bar_label(bars1,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                    ax2.bar_label(bars2,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                    ax3.bar_label(bars3,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                    ax4.bar_label(bars4,label_type="edge",rotation=90,padding=10,fontsize=labelfont)

                    ax1.set_ylim(top=max(averagedf["proteins_avg"].tolist())+y_padding*max(averagedf["proteins_avg"].tolist()))
                    #ax2.set_ylim(top=max(averagedf["proteins2pepts_avg"].tolist())+y_padding*max(averagedf["proteins2pepts_avg"].tolist()))
                    #ax3.set_ylim(top=max(averagedf[pepcolumn+"_avg"].tolist())+y_padding*max(averagedf[pepcolumn+"_avg"].tolist()))
                    ax4.set_ylim(top=max(averagedf["precursors_avg"].tolist())+y_padding*max(averagedf["precursors_avg"].tolist()))
                elif input.idmetrics_individual_average()=="unique":
                    idmetricscolor=colorpicker()
                    idmetricscolor_light=[]
                    if numconditions==1:
                        c=colorsys.rgb_to_hls(*mcolors.to_rgb(idmetricscolor))
                        idmetricscolor_light.append(colorsys.hls_to_rgb(c[0],1-0.3*(1-c[1]),c[2]))
                    else:
                        for color in idmetricscolor:
                            c=colorsys.rgb_to_hls(*mcolors.to_rgb(color))
                            idmetricscolor_light.append(colorsys.hls_to_rgb(c[0],1-0.3*(1-c[1]),c[2]))
                    ax1.bar(uniquecountsdf["R.Condition"],uniquecountsdf["proteins"],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax2.bar(uniquecountsdf["R.Condition"],uniquecountsdf["proteins2pepts"],width=0.8,color=idmetricscolor,edgecolor="k",label=">2 unique peptides")
                    ax2.bar(uniquecountsdf["R.Condition"],uniquecountsdf["protein_difference"],bottom=uniquecountsdf["proteins2pepts"],width=0.8,color=idmetricscolor_light,edgecolor="k",label="Single unique peptide")
                    ax3.bar(uniquecountsdf["R.Condition"],uniquecountsdf[pepcolumn],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax4.bar(uniquecountsdf["R.Condition"],uniquecountsdf["precursors"],width=0.8,color=idmetricscolor,edgecolor="k")

                    ax1.bar_label(ax1.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax2.bar_label(ax2.containers[0],label_type="center",rotation=0,fontsize=labelfont)
                    ax2.bar_label(ax2.containers[1],label_type="center",rotation=0,fontsize=labelfont)
                    ax3.bar_label(ax3.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax4.bar_label(ax4.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)

                    ax1.set_ylim(top=max(uniquecountsdf["proteins"].tolist())+y_padding*max(resultdf["proteins"].tolist()))
                    #ax2.set_ylim(top=max(uniquecountsdf["proteins"].tolist())+y_padding*max(resultdf["proteins2pepts"].tolist()))
                    #ax3.set_ylim(top=max(uniquecountsdf[pepcolumn].tolist())+(y_padding+0.1)*max(resultdf[pepcolumn].tolist()))
                    ax4.set_ylim(top=max(uniquecountsdf["precursors"].tolist())+(y_padding+0.1)*max(resultdf["precursors"].tolist()))

                    ax2.legend(prop={'size':legendfont})
                    legend=ax2.get_legend()
                    legend.legend_handles[0].set_facecolor("grey")
                    legend.legend_handles[0].set_edgecolor("grey")
                    legend.legend_handles[1].set_facecolor("gainsboro")
                    legend.legend_handles[1].set_edgecolor("gainsboro")

                ax1.set_title("Protein Groups",fontsize=titlefont)
                ax2.set_title("Protein Groups with >2 Peptides",fontsize=titlefont)
                ax3.set_title(peptitle,fontsize=titlefont)
                ax3.set_xlabel("Condition",fontsize=axisfont)
                ax3.tick_params(axis="x",rotation=x_label_rotation)
                ax4.set_title("Precursors",fontsize=titlefont)
                ax4.set_xlabel("Condition",fontsize=axisfont)
                ax4.tick_params(axis="x",rotation=x_label_rotation)

                fig.text(0,0.6,"Counts",ha="left",va="center",rotation="vertical",fontsize=axisfont)

                ax1.set_axisbelow(True)
                ax1.grid(linestyle="--")
                ax1.tick_params(axis="both",labelsize=axisfont_labels)
                ax2.set_axisbelow(True)
                ax2.grid(linestyle="--")
                ax2.tick_params(axis="both",labelsize=axisfont_labels)
                ax3.set_axisbelow(True)
                ax3.grid(linestyle="--")
                ax3.tick_params(axis="both",labelsize=axisfont_labels)
                ax4.set_axisbelow(True)
                ax4.grid(linestyle="--")
                ax4.tick_params(axis="both",labelsize=axisfont_labels)
            else:
                plotinput=input.idplotinput()
                if plotinput=="proteins":
                    titleprop="Protein Groups"
                if plotinput=="proteins2pepts":
                    titleprop="Protein Groups with >2 Peptides"
                if plotinput=="peptides":
                    if input.idmetrics_peptides()=="modified":
                        plotinput="peptides"
                        titleprop="Peptides"
                    elif input.idmetrics_peptides()=="stripped":
                        plotinput="strippedpeptides"
                        titleprop="Peptide Sequences"
                if plotinput=="precursors":
                    titleprop="Precursors"

                fig,ax=plt.subplots(figsize=(input.idmetrics_width()*(1/plt.rcParams['figure.dpi']),input.idmetrics_height()*(1/plt.rcParams['figure.dpi'])))

                if (plotinput=="peptides" or plotinput=="strippedpeptides") and input.idmetrics_peplength_switch()==True:
                    peplen=input.idmetrics_peplength_slider()
                    numpeptides=[]
                    numstrippedpeptides=[]
                    peptideresultdf=pd.DataFrame(searchoutput[["Cond_Rep","R.FileName","R.Condition","R.Replicate"]].drop_duplicates()).reset_index(drop=True)
                    for run in peptideresultdf["Cond_Rep"].tolist():
                        numpeptides.append(searchoutput[(searchoutput["Cond_Rep"]==run)&(searchoutput["Peptide Length"]==peplen)]["EG.ModifiedPeptide"].nunique())
                        numstrippedpeptides.append(searchoutput[(searchoutput["Cond_Rep"]==run)&(searchoutput["Peptide Length"]==peplen)]["PEP.StrippedSequence"].nunique())

                    peptideresultdf["peptides"]=numpeptides
                    peptideresultdf["strippedpeptides"]=numstrippedpeptides

                    peptideaveragedf=pd.DataFrame({"R.Condition":sampleconditions,"N.Replicates":maxreplicatelist})
                    columnlist=peptideresultdf.columns.values.tolist()
                    for i in columnlist:
                        if i=="R.FileName" or i=="Cond_Rep" or i=="R.Condition" or i=="R.Replicate":
                            continue
                        avglist=[]
                        stdevlist=[]
                        for j in sampleconditions:
                            samplecondition=peptideresultdf[peptideresultdf["R.Condition"]==j]
                            avglist.append(round(np.average(samplecondition[i].to_numpy())))
                            stdevlist.append(np.std(samplecondition[i].to_numpy()))
                        peptideaveragedf[i+"_avg"]=avglist
                        peptideaveragedf[i+"_stdev"]=stdevlist
                    if input.idmetrics_individual_average()=="individual":
                        idmetricscolor=replicatecolors()
                        peptideresultdf.plot.bar(ax=ax,x="Cond_Rep",y=plotinput,legend=False,width=0.8,color=idmetricscolor,edgecolor="k",fontsize=axisfont)
                        ax.bar_label(ax.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                        ax.set_ylim(top=max(peptideresultdf[plotinput].tolist())+y_padding*max(peptideresultdf[plotinput].tolist()))
                        titlemod=str(input.idmetrics_peplength_slider())+"mer "
                    elif input.idmetrics_individual_average()=="average":
                        idmetricscolor=colorpicker()
                        bars=ax.bar(peptideaveragedf["R.Condition"],peptideaveragedf[plotinput+"_avg"],yerr=peptideaveragedf[plotinput+"_stdev"],edgecolor="k",width=0.8,capsize=10,color=idmetricscolor)
                        ax.bar_label(bars,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                        ax.set_ylim(top=max(peptideaveragedf[plotinput+"_avg"].tolist())+y_padding*max(peptideaveragedf[plotinput+"_avg"].tolist()))
                        titlemod="Average # "+str(input.idmetrics_peplength_slider())+"mer "
                else:
                    if input.idmetrics_individual_average()=="individual":
                        idmetricscolor=replicatecolors()
                        resultdf.plot.bar(ax=ax,x="Cond_Rep",y=plotinput,legend=False,width=0.8,color=idmetricscolor,edgecolor="k",fontsize=axisfont)
                        ax.bar_label(ax.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                        ax.set_ylim(top=max(resultdf[plotinput].tolist())+y_padding*max(resultdf[plotinput].tolist()))
                        titlemod=""
                    elif input.idmetrics_individual_average()=="average":
                        idmetricscolor=colorpicker()
                        bars=ax.bar(averagedf["R.Condition"],averagedf[plotinput+"_avg"],yerr=averagedf[plotinput+"_stdev"],edgecolor="k",width=0.8,capsize=10,color=idmetricscolor)
                        ax.bar_label(bars,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                        ax.set_ylim(top=max(averagedf[plotinput+"_avg"].tolist())+y_padding*max(averagedf[plotinput+"_avg"].tolist()))
                        titlemod="Average #"
                    elif input.idmetrics_individual_average()=="unique":
                        idmetricscolor=colorpicker()
                        idmetricscolor_light=[]
                        for color in idmetricscolor:
                            c=colorsys.rgb_to_hls(*mcolors.to_rgb(color))
                            idmetricscolor_light.append(colorsys.hls_to_rgb(c[0],1-0.3*(1-c[1]),c[2]))
                        if plotinput=="proteins2pepts":
                            ax.bar(uniquecountsdf["R.Condition"],uniquecountsdf["proteins2pepts"],edgecolor="k",width=0.8,color=idmetricscolor,label=">2 unique peptides")
                            ax.bar(uniquecountsdf["R.Condition"],uniquecountsdf["protein_difference"],bottom=uniquecountsdf["proteins2pepts"],edgecolor="k",width=0.8,color=idmetricscolor_light,label="Single unique peptide")
                            ax.bar_label(ax.containers[0],label_type="center",rotation=0,fontsize=labelfont)
                            ax.bar_label(ax.containers[1],label_type="center",rotation=0,fontsize=labelfont)
                            ax.legend(loc="center left",bbox_to_anchor=(1,0.5),prop={'size':legendfont})
                            legend=ax.get_legend()
                            legend.legend_handles[0].set_facecolor("grey")
                            legend.legend_handles[0].set_edgecolor("grey")
                            legend.legend_handles[1].set_facecolor("gainsboro")
                            legend.legend_handles[1].set_edgecolor("gainsboro")
                        else:
                            ax.bar(uniquecountsdf["R.Condition"],uniquecountsdf[plotinput],edgecolor="k",width=0.8,capsize=10,color=idmetricscolor)
                            ax.bar_label(ax.containers[0],label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                            ax.set_ylim(top=max(uniquecountsdf[plotinput].tolist())+y_padding*max(uniquecountsdf[plotinput].tolist()))
                        titlemod="Total Unique "

                ax.set_ylabel("Counts",fontsize=axisfont)
                ax.set_xlabel("Condition",fontsize=axisfont)
                ax.set_title(titlemod+titleprop,fontsize=titlefont)
                ax.tick_params(axis="x",rotation=x_label_rotation)
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
            fig.set_tight_layout(True)
            imagedownload("idcounts")

    # ====================================== Venn Diagram
    @render.ui
    def venn_run1_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        if input.venn_conditionorrun()=="condition":
            opts=sampleconditions
            return ui.input_selectize("venn_run1_list","Condition 1:",choices=opts)
        if input.venn_conditionorrun()=="individual":
            opts=runnames
            return ui.input_selectize("venn_run1_list","Run 1:",choices=opts)   
    @render.ui
    def venn_run2_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        if input.venn_conditionorrun()=="condition":
            opts=sampleconditions
            return ui.input_selectize("venn_run2_list","Condition 2:",choices=opts)
        if input.venn_conditionorrun()=="individual":
            opts=runnames
            return ui.input_selectize("venn_run2_list","Run 2:",choices=opts)   
    @render.ui
    def venn_run3_ui():
        if input.venn_numcircles()=="3":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            if input.venn_conditionorrun()=="condition":
                opts=sampleconditions
                return ui.input_selectize("venn_run3_list","Condition 3:",choices=opts)
            if input.venn_conditionorrun()=="individual":
                opts=runnames
                return ui.input_selectize("venn_run3_list","Run 3:",choices=opts)
    @render.ui
    def venn_color3_ui():
        if input.venn_numcircles()=="3":
            return ui.input_text("venn_color3","Color 3:",value="tab:green")
    @render.ui
    def venn_specific_length_ui():
        if input.venn_plotproperty()=="peptides" or input.venn_plotproperty()=="precursors" or input.venn_plotproperty()=="peptides_stripped":
            return ui.input_switch("venn_specific_length","Compare specific peptide length?",value=False,width="300px")
    @render.ui
    def venn_peplength_ui():
        if input.venn_specific_length()==True:
            return ui.input_slider("venn_peplength_pick","Pick specific peptide length to compare:",min=3,max=30,value=9,step=1,ticks=True)
    @render.ui
    def peptidecore_ui():
        if input.venn_plotproperty()=="peptides_stripped":
            return ui.input_switch("peptidecore","Only consider peptide core (cut first and last 2 residues, can't download detailed list)",value=False)
    @render.ui
    def venn_ptm_ui():
        if input.venn_plotproperty()=="peptides" or input.venn_plotproperty()=="precursors" or input.venn_plotproperty()=="peptides_stripped":
            return ui.input_switch("venn_ptm","Compare only for specific PTM?",value=False)
    @render.ui
    def venn_ptmlist_ui():
        if input.venn_plotproperty()=="peptides" or input.venn_plotproperty()=="precursors" or input.venn_plotproperty()=="peptides_stripped":
            if input.venn_ptm()==True:
                listofptms=find_ptms()

                return ui.input_selectize("venn_foundptms","Pick PTM to plot data for:",choices=listofptms,selected=listofptms[0])  
    #plot Venn Diagram
    @reactive.effect
    def _():
        @render.plot(width=input.venn_width(),height=input.venn_height())
        def venn_plot():
            from matplotlib_venn import venn2,venn2_circles,venn3,venn3_circles

            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            if input.venn_conditionorrun()=="condition":
                A=searchoutput[searchoutput["R.Condition"]==input.venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                B=searchoutput[searchoutput["R.Condition"]==input.venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                if input.venn_numcircles()=="3":
                    C=searchoutput[searchoutput["R.Condition"]==input.venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            if input.venn_conditionorrun()=="individual":
                A=searchoutput[searchoutput["Cond_Rep"]==input.venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                B=searchoutput[searchoutput["Cond_Rep"]==input.venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                if input.venn_numcircles()=="3":
                    C=searchoutput[searchoutput["Cond_Rep"]==input.venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            
            #add extra columns to df for peptide+charge and peptide lengths
            A["pep_charge"]=A["EG.ModifiedPeptide"]+A["FG.Charge"].astype(str)
            B["pep_charge"]=B["EG.ModifiedPeptide"]+B["FG.Charge"].astype(str)
            if input.venn_numcircles()=="3":
                C["pep_charge"]=C["EG.ModifiedPeptide"]+C["FG.Charge"].astype(str)
            
            titlemodlist=[]

            if input.venn_plotproperty()=="proteingroups":
                a=set(A["PG.ProteinGroups"])
                b=set(B["PG.ProteinGroups"])
                if input.venn_numcircles()=="3":
                    c=set(C["PG.ProteinGroups"])
                titlemod="Protein Groups"
            if input.venn_plotproperty()=="peptides":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    titlemodlist.append(ptm)
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                    titlemodlist.append(str(input.venn_peplength_pick())+"mers")
                if input.venn_ptm()==False and input.venn_specific_length()==False:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                a=set(A["EG.ModifiedPeptide"])
                b=set(B["EG.ModifiedPeptide"])
                if input.venn_numcircles()=="3":
                    c=set(C["EG.ModifiedPeptide"])
                titlemod="Peptides"
            if input.venn_plotproperty()=="precursors":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    titlemodlist.append(ptm)
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                    titlemodlist.append(str(input.venn_peplength_pick())+"mers")
                else:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                a=set(A["pep_charge"])
                b=set(B["pep_charge"])
                if input.venn_numcircles()=="3":
                    c=set(C["pep_charge"])
                titlemod="Precursors"
            if input.venn_plotproperty()=="peptides_stripped":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    titlemodlist.append(ptm)
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                    titlemodlist.append(str(input.venn_peplength_pick())+"mers")
                else:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                if input.peptidecore()==True:
                    A_coreAA=[]
                    B_coreAA=[]
                    for pep in A["PEP.StrippedSequence"].tolist():
                        A_coreAA.append(pep[2:-2])
                    for pep in B["PEP.StrippedSequence"].tolist():
                        B_coreAA.append(pep[2:-2])
                    a=set(A_coreAA)
                    b=set(B_coreAA)
                    if input.venn_numcircles()=="3":
                        C_coreAA=[]
                        for pep in C["PEP.StrippedSequence"].tolist():
                            C_coreAA.append(pep[2:-2])
                        c=set(C_coreAA)
                    titlemod="Stripped Peptides (no terminal AAs)"
                else:
                    a=set(A["PEP.StrippedSequence"])
                    b=set(B["PEP.StrippedSequence"])
                    if input.venn_numcircles()=="3":
                        c=set(C["PEP.StrippedSequence"])
                    titlemod="Stripped Peptides"
            if titlemodlist==[]:
                titlemodlist=""
            else:
                titlemodlist=" ("+", ".join(titlemodlist)+")"
            fig,ax=plt.subplots(figsize=(input.venn_width()*(1/plt.rcParams['figure.dpi']),input.venn_height()*(1/plt.rcParams['figure.dpi'])))
            if input.venn_numcircles()=="2":
                Ab=len(a-b)
                aB=len(b-a)
                AB=len(a&b)
                venn2(subsets=(Ab,aB,AB),set_labels=(input.venn_run1_list(),input.venn_run2_list()),set_colors=(input.venn_color1(),input.venn_color2()),ax=ax)
                venn2_circles(subsets=(Ab,aB,AB),linestyle="dashed",linewidth=0.5)
                plt.title("Venn Diagram for "+titlemod+titlemodlist)
            if input.venn_numcircles()=="3":
                Abc=len(a-b-c)
                aBc=len(b-a-c)
                ABc=len((a&b)-c)
                abC=len(c-a-b)
                AbC=len((a&c)-b)
                aBC=len((b&c)-a)
                ABC=len(a&b&c)
                venn3(subsets=(Abc,aBc,ABc,abC,AbC,aBC,ABC),set_labels=(input.venn_run1_list(),input.venn_run2_list(),input.venn_run3_list()),set_colors=(input.venn_color1(),input.venn_color2(),input.venn_color3()),ax=ax)
                venn3_circles(subsets=(Abc,aBc,ABc,abC,AbC,aBC,ABC),linestyle="dashed",linewidth=0.5)
                plt.title("Venn Diagram for "+titlemod+titlemodlist)
            fig.set_tight_layout(True)
            imagedownload("venndiagram_"+input.venn_plotproperty())
    #download table of Venn Diagram intersections
    @reactive.effect
    def _():
        if input.venn_numcircles()=="2":
            filename=lambda: f"VennList_{input.venn_run1_list()}_vs_{input.venn_run2_list()}_{input.venn_plotproperty()}.csv"
        if input.venn_numcircles()=="3":
            filename=lambda: f"VennList_A-{input.venn_run1_list()}_vs_B-{input.venn_run2_list()}_vs_C-{input.venn_run3_list()}_{input.venn_plotproperty()}.csv"
        @render.download(filename=filename)
        def venn_download():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            if input.venn_conditionorrun()=="condition":
                A=searchoutput[searchoutput["R.Condition"]==input.venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                B=searchoutput[searchoutput["R.Condition"]==input.venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                if input.venn_numcircles()=="3":
                    C=searchoutput[searchoutput["R.Condition"]==input.venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            if input.venn_conditionorrun()=="individual":
                A=searchoutput[searchoutput["Cond_Rep"]==input.venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                B=searchoutput[searchoutput["Cond_Rep"]==input.venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                if input.venn_numcircles()=="3":
                    C=searchoutput[searchoutput["Cond_Rep"]==input.venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            
            #add extra columns to df for peptide+charge and peptide lengths
            A["pep_charge"]=A["EG.ModifiedPeptide"]+A["FG.Charge"].astype(str)
            B["pep_charge"]=B["EG.ModifiedPeptide"]+B["FG.Charge"].astype(str)
            if input.venn_numcircles()=="3":
                C["pep_charge"]=C["EG.ModifiedPeptide"]+C["FG.Charge"].astype(str)

            if input.venn_plotproperty()=="proteingroups":
                a=set(A["PG.ProteinGroups"])
                b=set(B["PG.ProteinGroups"])
                if input.venn_numcircles()=="3":
                    c=set(C["PG.ProteinGroups"])
            if input.venn_plotproperty()=="peptides":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                if input.venn_ptm()==False and input.venn_specific_length()==False:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                a=set(A["EG.ModifiedPeptide"])
                b=set(B["EG.ModifiedPeptide"])
                if input.venn_numcircles()=="3":
                    c=set(C["EG.ModifiedPeptide"])
            if input.venn_plotproperty()=="precursors":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                else:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                a=set(A["pep_charge"])
                b=set(B["pep_charge"])
                if input.venn_numcircles()=="3":
                    c=set(C["pep_charge"])
            if input.venn_plotproperty()=="peptides_stripped":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                else:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                if input.peptidecore()==True:
                    A_coreAA=[]
                    B_coreAA=[]
                    for pep in A["PEP.StrippedSequence"].tolist():
                        A_coreAA.append(pep[2:-2])
                    for pep in B["PEP.StrippedSequence"].tolist():
                        B_coreAA.append(pep[2:-2])
                    a=set(A_coreAA)
                    b=set(B_coreAA)
                    if input.venn_numcircles()=="3":
                        C_coreAA=[]
                        for pep in C["PEP.StrippedSequence"].tolist():
                            C_coreAA.append(pep[2:-2])
                        c=set(C_coreAA)
                else:
                    a=set(A["PEP.StrippedSequence"])
                    b=set(B["PEP.StrippedSequence"])
                    if input.venn_numcircles()=="3":
                        c=set(C["PEP.StrippedSequence"])

            df=pd.DataFrame()
            if input.venn_numcircles()=="2":
                Ab=list(a-b)
                aB=list(b-a)
                AB=list(a&b)
                df=pd.concat([df,pd.Series(Ab,name=input.venn_run1_list())],axis=1)
                df=pd.concat([df,pd.Series(aB,name=input.venn_run2_list())],axis=1)
                df=pd.concat([df,pd.Series(AB,name="Both")],axis=1)
            if input.venn_numcircles()=="3":
                Abc=list(a-b-c)
                aBc=list(b-a-c)
                ABc=list((a&b)-c)
                abC=list(c-a-b)
                AbC=list((a&c)-b)
                aBC=list((b&c)-a)
                ABC=list(a&b&c)
                df=pd.concat([df,pd.Series(Abc,name="A only")],axis=1)
                df=pd.concat([df,pd.Series(aBc,name="B only")],axis=1)
                df=pd.concat([df,pd.Series(ABc,name="A and B, not C")],axis=1)
                df=pd.concat([df,pd.Series(abC,name="C only")],axis=1)
                df=pd.concat([df,pd.Series(AbC,name="A and C, not B")],axis=1)
                df=pd.concat([df,pd.Series(aBC,name="B and C, not A")],axis=1)
                df=pd.concat([df,pd.Series(ABC,name="ABC")],axis=1)
            with io.BytesIO() as buf:
                df.to_csv(buf,index=False)
                yield buf.getvalue()            
    #download detailed table of Venn Diagram intersections
    @reactive.effect
    def _():
        if input.venn_numcircles()=="2":
            filename=lambda: f"VennList_Detailed_{input.venn_run1_list()}_vs_{input.venn_run2_list()}_{input.venn_plotproperty()}.csv"
        if input.venn_numcircles()=="3":
            filename=lambda: f"VennList_Detailed_A-{input.venn_run1_list()}_vs_B-{input.venn_run2_list()}_vs_C-{input.venn_run3_list()}_{input.venn_plotproperty()}.csv"
        @render.download(filename=filename)
        def venn_download_detailed():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            if input.venn_conditionorrun()=="condition":
                A=searchoutput[searchoutput["R.Condition"]==input.venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                B=searchoutput[searchoutput["R.Condition"]==input.venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                if input.venn_numcircles()=="3":
                    C=searchoutput[searchoutput["R.Condition"]==input.venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            if input.venn_conditionorrun()=="individual":
                A=searchoutput[searchoutput["Cond_Rep"]==input.venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                B=searchoutput[searchoutput["Cond_Rep"]==input.venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                if input.venn_numcircles()=="3":
                    C=searchoutput[searchoutput["Cond_Rep"]==input.venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            
            #add extra columns to df for peptide+charge and peptide lengths
            A["pep_charge"]=A["EG.ModifiedPeptide"]+A["FG.Charge"].astype(str)
            B["pep_charge"]=B["EG.ModifiedPeptide"]+B["FG.Charge"].astype(str)
            if input.venn_numcircles()=="3":
                C["pep_charge"]=C["EG.ModifiedPeptide"]+C["FG.Charge"].astype(str)

            if input.venn_plotproperty()=="proteingroups":
                a=set(A["PG.ProteinGroups"])
                b=set(B["PG.ProteinGroups"])
                if input.venn_numcircles()=="3":
                    c=set(C["PG.ProteinGroups"])
            if input.venn_plotproperty()=="peptides":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                if input.venn_ptm()==False and input.venn_specific_length()==False:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                a=set(A["EG.ModifiedPeptide"])
                b=set(B["EG.ModifiedPeptide"])
                if input.venn_numcircles()=="3":
                    c=set(C["EG.ModifiedPeptide"])
            if input.venn_plotproperty()=="precursors":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                else:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                a=set(A["pep_charge"])
                b=set(B["pep_charge"])
                if input.venn_numcircles()=="3":
                    c=set(C["pep_charge"])
            if input.venn_plotproperty()=="peptides_stripped":
                if input.venn_ptm()==True:
                    ptm=input.venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    if input.venn_numcircles()=="3":
                        C=C[C["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                if input.venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.venn_peplength_pick())]
                    if input.venn_numcircles()=="3":
                        C=C[C["Peptide Length"]==int(input.venn_peplength_pick())]
                else:
                    A=A
                    B=B
                    if input.venn_numcircles()=="3":
                        C=C
                if input.peptidecore()==True:
                    A_coreAA=[]
                    B_coreAA=[]
                    for pep in A["PEP.StrippedSequence"].tolist():
                        A_coreAA.append(pep[2:-2])
                    for pep in B["PEP.StrippedSequence"].tolist():
                        B_coreAA.append(pep[2:-2])
                    a=set(A_coreAA)
                    b=set(B_coreAA)
                    if input.venn_numcircles()=="3":
                        C_coreAA=[]
                        for pep in C["PEP.StrippedSequence"].tolist():
                            C_coreAA.append(pep[2:-2])
                        c=set(C_coreAA)
                else:
                    a=set(A["PEP.StrippedSequence"])
                    b=set(B["PEP.StrippedSequence"])
                    if input.venn_numcircles()=="3":
                        c=set(C["PEP.StrippedSequence"])

            df=pd.DataFrame()
            if input.venn_numcircles()=="2":
                Ab=list(a-b)
                aB=list(b-a)
                AB=list(a&b)
                df=pd.concat([df,pd.Series(Ab,name=input.venn_run1_list())],axis=1)
                df=pd.concat([df,pd.Series(aB,name=input.venn_run2_list())],axis=1)
                df=pd.concat([df,pd.Series(AB,name="Both")],axis=1)
            if input.venn_numcircles()=="3":
                Abc=list(a-b-c)
                aBc=list(b-a-c)
                ABc=list((a&b)-c)
                abC=list(c-a-b)
                AbC=list((a&c)-b)
                aBC=list((b&c)-a)
                ABC=list(a&b&c)
                df=pd.concat([df,pd.Series(Abc,name="A only")],axis=1)
                df=pd.concat([df,pd.Series(aBc,name="B only")],axis=1)
                df=pd.concat([df,pd.Series(ABc,name="A and B, not C")],axis=1)
                df=pd.concat([df,pd.Series(abC,name="C only")],axis=1)
                df=pd.concat([df,pd.Series(AbC,name="A and C, not B")],axis=1)
                df=pd.concat([df,pd.Series(aBC,name="B and C, not A")],axis=1)
                df=pd.concat([df,pd.Series(ABC,name="ABC")],axis=1)

            if input.venn_plotproperty()=="proteingroups":
                locindex="PG.ProteinGroups"
                pullcolumns=["Cond_Rep","PG.ProteinGroups","PG.ProteinNames","PG.Genes","PG.MS2Quantity"]
            if input.venn_plotproperty()=="peptides":
                locindex="EG.ModifiedPeptide"
                pullcolumns=["Cond_Rep","EG.ModifiedPeptide","PEP.StrippedSequence","FG.Charge","FG.MS2Quantity","PG.ProteinGroups","PG.ProteinNames","PG.Genes","PG.MS2Quantity"]
            if input.venn_plotproperty()=="peptides_stripped":
                locindex="PEP.StrippedSequence"
                pullcolumns=["Cond_Rep","PEP.StrippedSequence","PG.ProteinGroups","PG.ProteinNames","PG.Genes","PG.MS2Quantity"]
            if input.venn_plotproperty()=="precursors":
                locindex=searchoutput["EG.ModifiedPeptide"]+searchoutput["FG.Charge"].astype(str)
                pullcolumns=["Cond_Rep","EG.ModifiedPeptide","PEP.StrippedSequence","FG.Charge","FG.MS2Quantity","PG.ProteinGroups","PG.ProteinNames","PG.Genes","PG.MS2Quantity"]
            if input.venn_conditionorrun()=="individual":
                querycolumn="Cond_Rep"
            if input.venn_conditionorrun()=="condition":
                querycolumn="R.Condition"

            searchoutput_loc=searchoutput.set_index(locindex)

            if input.venn_numcircles()=="2":
                IDlist=df["Both"].dropna().tolist()
                venndf_all=searchoutput_loc[(searchoutput_loc[querycolumn]==df.columns[0])|(searchoutput_loc[querycolumn]==df.columns[1])].loc[IDlist].reset_index()[pullcolumns].drop_duplicates().reset_index(drop=True)

                venndf_individualID=pd.DataFrame()
                columnlist=df.columns[:-1]
                for i in columnlist:
                    venndf=searchoutput_loc[searchoutput_loc[querycolumn]==i].loc[df[i].dropna().tolist()].reset_index()[pullcolumns].drop_duplicates().reset_index(drop=True)
                    venndf_individualID=pd.concat([venndf_individualID,venndf])
                venndf_detailed=pd.concat([venndf_all,venndf_individualID]).reset_index(drop=True)
            if input.venn_numcircles()=="3":
                A=input.venn_run1_list()
                B=input.venn_run2_list()
                C=input.venn_run3_list()

                columndict={"A only":searchoutput_loc[querycolumn]==A,
                            "B only":searchoutput_loc[querycolumn]==B,
                            "C only":searchoutput_loc[querycolumn]==C,
                            "A and B, not C":(searchoutput_loc[querycolumn]==A)|(searchoutput_loc[querycolumn]==B),
                            "A and C, not B":(searchoutput_loc[querycolumn]==A)|(searchoutput_loc[querycolumn]==C),
                            "B and C, not A":(searchoutput_loc[querycolumn]==B)|(searchoutput_loc[querycolumn]==C)}

                
                IDlist=df["ABC"].dropna().tolist()
                venndf_all=searchoutput_loc[(searchoutput_loc[querycolumn]==A)|(searchoutput_loc[querycolumn]==B)|(searchoutput_loc[querycolumn]==C)].loc[IDlist].reset_index()[pullcolumns].drop_duplicates().reset_index(drop=True)

                venndf_individualID=pd.DataFrame()
                columnlist=df.columns[:-1]
                for i in columnlist:
                    venndf=searchoutput_loc[columndict[i]].loc[df[i].dropna()].reset_index()[pullcolumns].drop_duplicates().reset_index(drop=True)
                    venndf_individualID=pd.concat([venndf_individualID,venndf])
                venndf_detailed=pd.concat([venndf_all,venndf_individualID]).reset_index(drop=True)

            with io.BytesIO() as buf:
                venndf_detailed.to_csv(buf,index=False)
                yield buf.getvalue()            

    # ====================================== CV Plots
    @render.ui
    def cvplot_histogram_bins_ui():
        if input.cvplot_histogram_bins_switch()==True:
            return ui.input_slider("cvplot_histogram_bins_slider","Number of bins:",min=10,max=250,value=100,step=10,ticks=True)
    #plot cv violin plots
    @reactive.effect
    def _():
        @render.plot(width=input.cvplot_width(),height=input.cvplot_height())
        def cvplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            cvcalc_df=cvcalc()

            colors=colorpicker()
            cvplotinput=input.proteins_precursors_cvplot()
            cutoff95=input.removetop5percent()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.cvplot_histogram_bins_switch()==True:
                numbins=input.cvplot_histogram_bins_slider()
                fig,ax=plt.subplots(ncols=len(cvcalc_df),figsize=(input.cvplot_width()*(1/plt.rcParams['figure.dpi']),input.cvplot_height()*(1/plt.rcParams['figure.dpi'])))
                colors=colorpicker()
                if numconditions==1:
                    if cutoff95==True:
                        ax.hist(cvcalc_df[cvplotinput+" 95% CVs"],bins=numbins,color=colors)
                    elif cutoff95==False:
                        ax.hist(cvcalc_df[cvplotinput+" CVs"],bins=numbins,color=colors)
                    ax.set_xlabel(cvplotinput+" % CV",fontsize=axisfont)
                    ax.axvline(x=20,color="black",linestyle="--")
                    ax.set_ylabel("Counts",fontsize=axisfont)
                    ax.set_title(cvcalc_df["R.Condition"].tolist()[0],fontsize=titlefont)
                    ax.grid(linestyle="--")
                    ax.set_axisbelow(True)
                else:
                    for i in range(len(cvcalc_df)):
                        if cutoff95==True:
                            ax[i].hist(cvcalc_df[cvplotinput+" 95% CVs"][i],bins=numbins,color=colors[i])
                        elif cutoff95==False:
                            ax[i].hist(cvcalc_df[cvplotinput+" CVs"][i],bins=numbins,color=colors[i])
                        ax[i].set_xlabel(cvplotinput+" % CV",fontsize=axisfont)
                        ax[i].set_title(cvcalc_df["R.Condition"].tolist()[i],fontsize=titlefont)
                        ax[i].axvline(x=20,color="black",linestyle="--")
                        ax[i].grid(linestyle="--")
                        ax[i].set_axisbelow(True)
                    ax[0].set_ylabel("Counts",fontsize=axisfont)
            else:
                fig,ax=plt.subplots(figsize=(input.cvplot_width()*(1/plt.rcParams['figure.dpi']),input.cvplot_height()*(1/plt.rcParams['figure.dpi'])))
                x=np.arange(len(cvcalc_df["R.Condition"]))
                lineprops=dict(linestyle="--",color="black")
                flierprops=dict(markersize=3)
                if cutoff95==True:
                    bplot=ax.boxplot(cvcalc_df[cvplotinput+" 95% CVs"],medianprops=lineprops,flierprops=flierprops)
                    plot=ax.violinplot(cvcalc_df[cvplotinput+" 95% CVs"],showextrema=False)
                    ax.set_title(cvplotinput+" CVs, 95% Cutoff",fontsize=titlefont)
                elif cutoff95==False:
                    bplot=ax.boxplot(cvcalc_df[cvplotinput+" CVs"],medianprops=lineprops,flierprops=flierprops)
                    plot=ax.violinplot(cvcalc_df[cvplotinput+" CVs"],showextrema=False)
                    ax.set_title(cvplotinput+" CVs",fontsize=titlefont)
                ax.set_xticks(x+1,labels=cvcalc_df["R.Condition"],rotation=x_label_rotation)
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                ax.set_ylabel("CV%",fontsize=axisfont)
                ax.set_xlabel("Condition",fontsize=axisfont)
                ax.grid(linestyle="--")
                ax.set_axisbelow(True)
                ax.axhline(y=20,color="black",linestyle="--")

                if numconditions==1:
                    for z in plot["bodies"]:
                        z.set_facecolor(colors)
                        z.set_edgecolor("black")
                        z.set_alpha(0.7)
                else:
                    for z,color in zip(plot["bodies"],colors):
                        z.set_facecolor(color)
                        z.set_edgecolor("black")
                        z.set_alpha(0.7)
            fig.set_tight_layout(True)
            imagedownload("cvplot")
    #show a table of mean/median CV values per condition 
    @render.table
    def cv_table():
        cvcalc_df=cvcalc()

        cvplotinput=input.proteins_precursors_cvplot()
        cutoff95=input.removetop5percent()

        cvtable_protein=pd.DataFrame()
        cvtable_precursor=pd.DataFrame()
        cvtable_peptide=pd.DataFrame()

        if "Protein CVs" in cvcalc_df.columns:
            protein_meanlist=[]
            protein_medianlist=[]
            protein_meanlist95=[]
            protein_medianlist95=[]
            for run in cvcalc_df["R.Condition"]:
                protein_meanlist.append(np.mean(cvcalc_df[cvcalc_df["R.Condition"]==run]["Protein CVs"].tolist()))
                protein_meanlist95.append(np.mean(cvcalc_df[cvcalc_df["R.Condition"]==run]["Protein 95% CVs"].tolist()))
                protein_medianlist.append(np.median(cvcalc_df[cvcalc_df["R.Condition"]==run]["Protein CVs"].tolist()))
                protein_medianlist95.append(np.median(cvcalc_df[cvcalc_df["R.Condition"]==run]["Protein 95% CVs"].tolist()))
            cvtable_protein["R.Condition"]=cvcalc_df["R.Condition"]
            cvtable_protein["Protein Mean CVs"]=protein_meanlist
            cvtable_protein["Protein Mean CVs 95%"]=protein_meanlist95
            cvtable_protein["Protein Median CVs"]=protein_medianlist
            cvtable_protein["Protein Median CVs 95%"]=protein_medianlist95
        else:
            pass

        if "Precursor CVs" in cvcalc_df.columns:
            precursor_meanlist=[]
            precursor_medianlist=[]
            precursor_meanlist95=[]
            precursor_medianlist95=[]
            for run in cvcalc_df["R.Condition"]:
                precursor_meanlist.append(np.mean(cvcalc_df[cvcalc_df["R.Condition"]==run]["Precursor CVs"].tolist()))
                precursor_meanlist95.append(np.mean(cvcalc_df[cvcalc_df["R.Condition"]==run]["Precursor 95% CVs"].tolist()))
                precursor_medianlist.append(np.median(cvcalc_df[cvcalc_df["R.Condition"]==run]["Precursor CVs"].tolist()))
                precursor_medianlist95.append(np.median(cvcalc_df[cvcalc_df["R.Condition"]==run]["Precursor 95% CVs"].tolist()))
            cvtable_precursor["R.Condition"]=cvcalc_df["R.Condition"]
            cvtable_precursor["Precursor Mean CVs"]=precursor_meanlist
            cvtable_precursor["Precursor Mean CVs 95%"]=precursor_meanlist95
            cvtable_precursor["Precursor Median CVs"]=precursor_medianlist
            cvtable_precursor["Precursor Median CVs 95%"]=precursor_medianlist95         

            peptide_meanlist=[]
            peptide_medianlist=[]
            peptide_meanlist95=[]
            peptide_medianlist95=[]
            for run in cvcalc_df["R.Condition"]:
                peptide_meanlist.append(np.mean(cvcalc_df[cvcalc_df["R.Condition"]==run]["Peptide CVs"].tolist()))
                peptide_meanlist95.append(np.mean(cvcalc_df[cvcalc_df["R.Condition"]==run]["Peptide 95% CVs"].tolist()))
                peptide_medianlist.append(np.median(cvcalc_df[cvcalc_df["R.Condition"]==run]["Peptide CVs"].tolist()))
                peptide_medianlist95.append(np.median(cvcalc_df[cvcalc_df["R.Condition"]==run]["Peptide 95% CVs"].tolist()))
            cvtable_peptide["R.Condition"]=cvcalc_df["R.Condition"]
            cvtable_peptide["Peptide Mean CVs"]=peptide_meanlist
            cvtable_peptide["Peptide Mean CVs 95%"]=peptide_meanlist95
            cvtable_peptide["Peptide Median CVs"]=peptide_medianlist
            cvtable_peptide["Peptide Median CVs 95%"]=peptide_medianlist95
        else:
            pass

        if cvplotinput=="Protein":
            if cutoff95==True:
                return cvtable_protein[["R.Condition","Protein Mean CVs 95%","Protein Median CVs 95%"]]
            if cutoff95==False:
                return cvtable_protein[["R.Condition","Protein Mean CVs","Protein Median CVs"]]
        if cvplotinput=="Precursor":
            if cutoff95==True:
                return cvtable_precursor[["R.Condition","Precursor Mean CVs 95%","Precursor Median CVs 95%"]]
            if cutoff95==False:
                return cvtable_precursor[["R.Condition","Precursor Mean CVs","Precursor Median CVs"]]
        if cvplotinput=="Peptide":
            if cutoff95==True:
                return cvtable_peptide[["R.Condition","Peptide Mean CVs 95%","Peptide Median CVs 95%"]]
            if cutoff95==False:
                return cvtable_peptide[["R.Condition","Peptide Mean CVs","Peptide Median CVs"]]

    # ====================================== IDs with CV Cutoff
    #show percentages of IDs below cutoffs as a table
    @render.table
    def countscvcutoff_table():
        resultdf,averagedf,uniquecountsdf=idmetrics_calc()
        cvcalc_df=cvcalc()
        cvinput=input.proteins_precursors_idcutoffplot()

        conditions=cvcalc_df["R.Condition"].tolist()
        sub20percent=[]
        sub10percent=[]
        for i in range(len(conditions)):
            totalID=averagedf[cvinput+"_avg"][i]
            sub20=cvcalc_df[cvinput+"CV<20"][i]
            sub10=cvcalc_df[cvinput+"CV<10"][i]
            sub20percent.append(round((sub20/totalID)*100,2))
            sub10percent.append(round((sub10/totalID)*100,2))
        return pd.DataFrame({"R.Condition":conditions,"% "+cvinput+" CV<20%":sub20percent,"% "+cvinput+" CV<10%":sub10percent})
    #plot counts with CV cutoffs
    @reactive.effect
    def _():
        @render.plot(width=input.countscvcutoff_width(),height=input.countscvcutoff_height())
        def countscvcutoff():
            resultdf,averagedf,uniquecountsdf=idmetrics_calc()
            cvcalc_df=cvcalc()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            x=np.arange(len(cvcalc_df["R.Condition"]))
            width=0.25

            cvinput=input.proteins_precursors_idcutoffplot()

            fig,ax=plt.subplots(figsize=(input.countscvcutoff_width()*(1/plt.rcParams['figure.dpi']),input.countscvcutoff_width()*(1/plt.rcParams['figure.dpi'])))

            ax.bar(x,averagedf[cvinput+"_avg"],width=width,label="Identified",edgecolor="k",color="#054169")
            ax.bar_label(ax.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)

            ax.bar(x+width,cvcalc_df[cvinput+"CV<20"],width=width,label="CV<20%",edgecolor="k",color="#0071BC")
            ax.bar_label(ax.containers[1],label_type="edge",rotation=90,padding=5,fontsize=labelfont)

            ax.bar(x+(2*width),cvcalc_df[cvinput+"CV<10"],width=width,label="CV<10%",edgecolor="k",color="#737373")
            ax.bar_label(ax.containers[2],label_type="edge",rotation=90,padding=5,fontsize=labelfont)

            ax.set_ylim(top=max(averagedf[cvinput+"_avg"])+y_padding*max(averagedf[cvinput+"_avg"]))
            ax.legend(loc='center left',bbox_to_anchor=(1, 0.5),prop={'size':legendfont})
            ax.set_xticks(x+width,cvcalc_df["R.Condition"],rotation=x_label_rotation)
            ax.set_xlabel("Condition",fontsize=axisfont)
            ax.set_ylabel("Counts",fontsize=axisfont)
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            if cvinput=="proteins":
                ax.set_title("Protein Counts with CV Cutoffs",fontsize=titlefont)
            if cvinput=="precursors":
                ax.set_title("Precursor Counts with CV Cutoffs",fontsize=titlefont)
            if cvinput=="peptides":
                ax.set_title("Peptide Counts with CV Cutoffs",fontsize=titlefont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            fig.set_tight_layout(True)
            imagedownload("idcounts_cvcutoff")
    
    # ====================================== UpSet Plot
    @render.ui
    def specific_condition_ui():
        if input.upset_condition_or_run()=="specific_condition":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=sampleconditions
            return ui.input_selectize("specific_condition_pick","Pick sample condition:",choices=opts)
    @render.data_frame
    def upsetplot_counts():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        if input.protein_precursor_pick()=="Protein":
            proteindict=dict()
            if input.upset_condition_or_run()=="condition":
                if numconditions==1:
                    for run in runnames:
                        proteinlist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","PG.ProteinGroups"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        proteinlist.rename(columns={"PG.ProteinGroups":run},inplace=True)
                        proteindict[run]=proteinlist[run].tolist()
                else:
                    for condition in sampleconditions:
                        proteinlist=searchoutput[searchoutput["R.Condition"]==condition][["R.Condition","PG.ProteinGroups"]].drop_duplicates().reset_index(drop=True).drop(columns=["R.Condition"])
                        proteinlist.rename(columns={"PG.ProteinGroups":condition},inplace=True)
                        proteindict[condition]=proteinlist[condition].tolist()
                proteins=from_contents(proteindict)
            if input.upset_condition_or_run()=="individual":
                for run in runnames:
                    proteinlist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","PG.ProteinGroups"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                    proteinlist.rename(columns={"PG.ProteinGroups":run},inplace=True)
                    proteindict[run]=proteinlist[run].tolist()
                proteins=from_contents(proteindict)
            if input.upset_condition_or_run()=="specific_condition":
                df=searchoutput[searchoutput["R.Condition"]==input.specific_condition_pick()]
                for run in df["Cond_Rep"].drop_duplicates():
                    proteinlist=df[df["Cond_Rep"]==run][["Cond_Rep","PG.ProteinGroups"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                    proteinlist.rename(columns={"PG.ProteinGroups":run},inplace=True)
                    proteindict[run]=proteinlist[run].tolist()
                proteins=from_contents(proteindict)
            plottingdf=proteins
            titlemod="Protein"

        elif input.protein_precursor_pick()=="Peptide":
            peptidedict=dict()
            if input.upset_condition_or_run()=="condition":
                if numconditions==1:
                    for run in runnames:
                        peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"EG.ModifiedPeptide":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                else:
                    for condition in sampleconditions:
                        peptidelist=searchoutput[searchoutput["R.Condition"]==condition][["R.Condition","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["R.Condition"])
                        peptidelist.rename(columns={"EG.ModifiedPeptide":condition},inplace=True)
                        peptidedict[condition]=peptidelist[condition].tolist()
                peptides=from_contents(peptidedict)
            if input.upset_condition_or_run()=="individual":
                for run in runnames:
                    peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                    peptidelist.rename(columns={"EG.ModifiedPeptide":run},inplace=True)
                    peptidedict[run]=peptidelist[run].tolist()
                peptides=from_contents(peptidedict)
            if input.upset_condition_or_run()=="specific_condition":
                df=searchoutput[searchoutput["R.Condition"]==input.specific_condition_pick()]
                for run in df["Cond_Rep"].drop_duplicates():
                    peptidelist=df[df["Cond_Rep"]==run][["Cond_Rep","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                    peptidelist.rename(columns={"EG.ModifiedPeptide":run},inplace=True)
                    peptidedict[run]=peptidelist[run].tolist()
                peptides=from_contents(peptidedict)
            plottingdf=peptides
            titlemod="Peptide"
        elif input.protein_precursor_pick()=="Precursor":
            searchoutput["pep_charge"]=searchoutput["EG.ModifiedPeptide"]+searchoutput["FG.Charge"].astype(str)

            peptidedict=dict()
            if input.upset_condition_or_run()=="condition":
                if numconditions==1:
                    for run in runnames:
                        peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"pep_charge":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                else:
                    for condition in sampleconditions:
                        peptidelist=searchoutput[searchoutput["R.Condition"]==condition][["R.Condition","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["R.Condition"])
                        peptidelist.rename(columns={"pep_charge":condition},inplace=True)
                        peptidedict[condition]=peptidelist[condition].tolist()
                peptides=from_contents(peptidedict)
            if input.upset_condition_or_run()=="individual":
                for run in runnames:
                    peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                    peptidelist.rename(columns={"pep_charge":run},inplace=True)
                    peptidedict[run]=peptidelist[run].tolist()
                peptides=from_contents(peptidedict)
            if input.upset_condition_or_run()=="specific_condition":
                df=searchoutput[searchoutput["R.Condition"]==input.specific_condition_pick()]
                for run in df["Cond_Rep"].drop_duplicates():
                    peptidelist=df[df["Cond_Rep"]==run][["Cond_Rep","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                    peptidelist.rename(columns={"pep_charge":run},inplace=True)
                    peptidedict[run]=peptidelist[run].tolist()
                peptides=from_contents(peptidedict)
            plottingdf=peptides
            titlemod="Precursor"

        countsperrun=query(plottingdf).category_totals
        totalcounts=query(plottingdf,max_degree=1).total
        subset1_df=pd.DataFrame(query(plottingdf,max_degree=1).subset_sizes).reset_index()
        subset_nminus1_df=pd.DataFrame(query(plottingdf,min_degree=len(countsperrun)-1).subset_sizes).reset_index()

        str1="Total unique "+titlemod+"s across all runs"
        str2=titlemod+"s found in only 1 run"
        str3=titlemod+"s found in at least "+str(len(countsperrun)-1)+" out of "+str(len(countsperrun))+" runs"
        labellist=[str1,str2,str3]
        countlist=[totalcounts,sum(subset1_df["size"]),sum(subset_nminus1_df["size"])]
        df=pd.DataFrame({"Property":labellist,"Counts":countlist})

        return render.DataGrid(df,editable=False,width="600px",filters=False)
    #plot upset plot
    @reactive.effect
    def _():
        @render.plot(width=input.upsetplot_width(),height=input.upsetplot_height())
        def upsetplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            if input.protein_precursor_pick()=="Protein":
                proteindict=dict()
                if input.upset_condition_or_run()=="condition":
                    if numconditions==1:
                        for run in runnames:
                            proteinlist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","PG.ProteinGroups"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                            proteinlist.rename(columns={"PG.ProteinGroups":run},inplace=True)
                            proteindict[run]=proteinlist[run].tolist()
                    else:
                        for condition in sampleconditions:
                            proteinlist=searchoutput[searchoutput["R.Condition"]==condition][["R.Condition","PG.ProteinGroups"]].drop_duplicates().reset_index(drop=True).drop(columns=["R.Condition"])
                            proteinlist.rename(columns={"PG.ProteinGroups":condition},inplace=True)
                            proteindict[condition]=proteinlist[condition].tolist()
                    proteins=from_contents(proteindict)
                if input.upset_condition_or_run()=="individual":
                    for run in runnames:
                        proteinlist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","PG.ProteinGroups"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        proteinlist.rename(columns={"PG.ProteinGroups":run},inplace=True)
                        proteindict[run]=proteinlist[run].tolist()
                    proteins=from_contents(proteindict)
                if input.upset_condition_or_run()=="specific_condition":
                    df=searchoutput[searchoutput["R.Condition"]==input.specific_condition_pick()]
                    for run in df["Cond_Rep"].drop_duplicates():
                        proteinlist=df[df["Cond_Rep"]==run][["Cond_Rep","PG.ProteinGroups"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        proteinlist.rename(columns={"PG.ProteinGroups":run},inplace=True)
                        proteindict[run]=proteinlist[run].tolist()
                    proteins=from_contents(proteindict)
                plottingdf=proteins
                titlemod="Protein"
                min_degree=len(proteindict)-1

            elif input.protein_precursor_pick()=="Peptide":
                peptidedict=dict()
                if input.upset_condition_or_run()=="condition":
                    if numconditions==1:
                        for run in runnames:
                            peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                            peptidelist.rename(columns={"EG.ModifiedPeptide":run},inplace=True)
                            peptidedict[run]=peptidelist[run].tolist()
                    else:
                        for condition in sampleconditions:
                            peptidelist=searchoutput[searchoutput["R.Condition"]==condition][["R.Condition","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["R.Condition"])
                            peptidelist.rename(columns={"EG.ModifiedPeptide":condition},inplace=True)
                            peptidedict[condition]=peptidelist[condition].tolist()
                    peptides=from_contents(peptidedict)
                if input.upset_condition_or_run()=="individual":
                    for run in runnames:
                        peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"EG.ModifiedPeptide":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                    peptides=from_contents(peptidedict)
                if input.upset_condition_or_run()=="specific_condition":
                    df=searchoutput[searchoutput["R.Condition"]==input.specific_condition_pick()]
                    for run in df["Cond_Rep"].drop_duplicates():
                        peptidelist=df[df["Cond_Rep"]==run][["Cond_Rep","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"EG.ModifiedPeptide":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                    peptides=from_contents(peptidedict)
                plottingdf=peptides
                titlemod="Peptide"
                min_degree=len(peptidedict)-1

            elif input.protein_precursor_pick()=="Precursor":
                searchoutput["pep_charge"]=searchoutput["EG.ModifiedPeptide"]+searchoutput["FG.Charge"].astype(str)
                peptidedict=dict()
                if input.upset_condition_or_run()=="condition":
                    if numconditions==1:
                        for run in runnames:
                            peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                            peptidelist.rename(columns={"pep_charge":run},inplace=True)
                            peptidedict[run]=peptidelist[run].tolist()
                    else:
                        for condition in sampleconditions:
                            peptidelist=searchoutput[searchoutput["R.Condition"]==condition][["R.Condition","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["R.Condition"])
                            peptidelist.rename(columns={"pep_charge":condition},inplace=True)
                            peptidedict[condition]=peptidelist[condition].tolist()
                    peptides=from_contents(peptidedict)
                if input.upset_condition_or_run()=="individual":
                    for run in runnames:
                        peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"pep_charge":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                    peptides=from_contents(peptidedict)
                if input.upset_condition_or_run()=="specific_condition":
                    df=searchoutput[searchoutput["R.Condition"]==input.specific_condition_pick()]
                    for run in df["Cond_Rep"].drop_duplicates():
                        peptidelist=df[df["Cond_Rep"]==run][["Cond_Rep","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"pep_charge":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                    peptides=from_contents(peptidedict)
                plottingdf=peptides
                titlemod="Precursor"
                min_degree=len(peptidedict)-1

            fig=plt.figure(figsize=(input.upsetplot_width()*(1/plt.rcParams['figure.dpi']),input.upsetplot_height()*(1/plt.rcParams['figure.dpi'])))
            if input.upsetfilter()=="nofilter":
                upset=UpSet(plottingdf,show_counts=True,sort_by="cardinality").plot(fig)
            if input.upsetfilter()=="1run":
                upset=UpSet(plottingdf,show_counts=True,sort_by="cardinality",max_degree=1).plot(fig)
            if input.upsetfilter()=="n-1runs":
                upset=UpSet(plottingdf,show_counts=True,sort_by="cardinality",min_degree=min_degree).plot(fig)
            
            upset["totals"].set_title("# "+titlemod+"s")
            plt.ylabel(titlemod+" Intersections",fontsize=input.axisfont())
            imagedownload("upsetplot")

    # ====================================== UpSet Plot (stats)
    #render ui call for dropdown calling sample condition names
    @render.ui
    def upsetplotstats_conditionlist_ui():
        if input.upsetplotstats_whattoplot()=="condition" or input.upsetplotstats_whattoplot()=="specific_condition":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=sampleconditions
            return ui.input_selectize("upsetplotstats_conditionlist_pick","Pick sample condition:",choices=opts)
    @reactive.effect
    def _():
        @render.plot(width=input.upsetplotstats_width(),height=input.upsetplotstats_height())
        def upsetplotstats_singlehitIDplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.upsetplotstats_peptide_precursor()=="Peptide":
                peptidedict=dict()
                if input.upsetplotstats_whattoplot()=="individual":
                    for run in runnames:
                        peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"EG.ModifiedPeptide":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                    peptides=from_contents(peptidedict)
                if input.upsetplotstats_whattoplot()=="condition" or input.upsetplotstats_whattoplot()=="specific_condition":
                    df=searchoutput[searchoutput["R.Condition"]==input.upsetplotstats_conditionlist_pick()]
                    for run in df["Cond_Rep"].drop_duplicates():
                        peptidelist=df[df["Cond_Rep"]==run][["Cond_Rep","EG.ModifiedPeptide"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"EG.ModifiedPeptide":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                    peptides=from_contents(peptidedict)
                plottingdf=peptides
                setindex="EG.ModifiedPeptide"
            if input.upsetplotstats_peptide_precursor()=="Precursor":
                searchoutput["pep_charge"]=searchoutput["EG.ModifiedPeptide"]+searchoutput["FG.Charge"].astype(str)
                peptidedict=dict()
                if input.upsetplotstats_whattoplot()=="individual":
                    for run in runnames:
                        peptidelist=searchoutput[searchoutput["Cond_Rep"]==run][["Cond_Rep","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"pep_charge":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                    peptides=from_contents(peptidedict)
                if input.upsetplotstats_whattoplot()=="condition" or input.upsetplotstats_whattoplot()=="specific_condition":
                    df=searchoutput[searchoutput["R.Condition"]==input.upsetplotstats_conditionlist_pick()]
                    for run in df["Cond_Rep"].drop_duplicates():
                        peptidelist=df[df["Cond_Rep"]==run][["Cond_Rep","pep_charge"]].drop_duplicates().reset_index(drop=True).drop(columns=["Cond_Rep"])
                        peptidelist.rename(columns={"pep_charge":run},inplace=True)
                        peptidedict[run]=peptidelist[run].tolist()
                    peptides=from_contents(peptidedict)
                plottingdf=peptides
                setindex="pep_charge"
            
            IDs_1run=pd.DataFrame(query(plottingdf,max_degree=1).data).reset_index()
            if input.upsetplotstats_whattoplot()=="condition" or input.upsetplotstats_whattoplot()=="specific_condition":
                searchoutput_pepindex=searchoutput[searchoutput["R.Condition"]==input.upsetplotstats_conditionlist_pick()].set_index(setindex)
            else:
                searchoutput_pepindex=searchoutput.set_index(setindex)
            IDs_1run_searchoutput=searchoutput_pepindex.loc[IDs_1run["id"]]

            if input.upsetplotstats_whattoplot()=="individual" or input.upsetplotstats_whattoplot()=="condition":
                #pull mz and IM from every ID
                remainderIDs=searchoutput[["FG.PrecMz","EG.IonMobility"]].drop_duplicates()
                scatterlabel="All IDs"
            if input.upsetplotstats_whattoplot()=="specific_condition":
                #pull mz and IM from every ID from specified condition
                remainderIDs=searchoutput[searchoutput["R.Condition"]==input.upsetplotstats_conditionlist_pick()][["FG.PrecMz","EG.IonMobility"]].drop_duplicates()
                scatterlabel="All IDs from "+input.upsetplotstats_conditionlist_pick()
            if input.upsetplotstats_plottype()=="scatter":
                fig,ax=plt.subplots(figsize=(input.upsetplotstats_width()*(1/plt.rcParams['figure.dpi']),input.upsetplotstats_height()*(1/plt.rcParams['figure.dpi'])))
                #scatter of all IDs
                ax.scatter(remainderIDs["FG.PrecMz"],remainderIDs["EG.IonMobility"],zorder=1,s=2,label=scatterlabel)
                #scatter of just the single run hits
                ax.scatter(IDs_1run_searchoutput["FG.PrecMz"],IDs_1run_searchoutput["EG.IonMobility"],zorder=2,s=2,label="IDs in only 1 Run")
                ax.set_xlabel("m/z",fontsize=axisfont)
                ax.set_ylabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
                ax.legend(markerscale=5,prop={'size':legendfont})
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)

            if input.upsetplotstats_plottype()=="2dhist":
                numbins=[100,100]
                cmap=matplotlib.colors.LinearSegmentedColormap.from_list("",["white","blue","red"])
                
                fig,ax=plt.subplots(ncols=2,sharey=True,sharex=True,figsize=(input.upsetplotstats_width()*(1/plt.rcParams['figure.dpi']),input.upsetplotstats_height()*(1/plt.rcParams['figure.dpi'])))
                
                j=ax[0].hist2d(remainderIDs["FG.PrecMz"],remainderIDs["EG.IonMobility"],bins=numbins,cmap=cmap)
                k=ax[1].hist2d(IDs_1run_searchoutput["FG.PrecMz"],IDs_1run_searchoutput["EG.IonMobility"],bins=numbins,cmap=cmap)

                ax[0].set_xlabel("m/z",fontsize=axisfont)
                ax[1].set_xlabel("m/z",fontsize=axisfont)
                ax[0].set_ylabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
                ax[1].tick_params("y",labelleft=True)
                fig.colorbar(j[3],ax=ax[0])
                fig.colorbar(k[3],ax=ax[1])
                fig.set_tight_layout(True)
                ax[0].set_title("All IDs",fontsize=titlefont)
                ax[1].set_title("Unique IDs",fontsize=titlefont)
                ax[0].tick_params(axis="both",labelsize=axisfont_labels)
                ax[1].tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("upsetplotstats")

    # ====================================== Tracker
    #render table of detected proteins and their average PG.MS2Quantity
    @render.data_frame
    def protein_df():
        df_styles={"location":"body",
            "style":{"column-width":"25px",
                    "overflow":"hidden"
                    }
            }
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        df=searchoutput[[input.tracker_proteingrouping(),"PG.MS2Quantity"]].groupby(input.tracker_proteingrouping()).mean().reset_index().rename(columns={"PG.MS2Quantity":"Mean_PG.MS2Quantity"})
        return render.DataGrid(df,width="100%",selection_mode="row",editable=False,styles=df_styles)
    #render table of peptides identified for picked protein from above table
    @render.data_frame
    def pickedprotein_df():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        if len(protein_df.data_view(selected=True)[input.tracker_proteingrouping()].tolist())==0:
            df=pd.DataFrame()
            return render.DataGrid(df)
        else:
            selectedprotein=protein_df.data_view(selected=True)[input.tracker_proteingrouping()].tolist()[0]
            #df=searchoutput[searchoutput["PG.ProteinGroups"]==selectedprotein][["Cond_Rep","PEP.StrippedSequence","EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]].drop_duplicates()
            df=searchoutput[searchoutput[input.tracker_proteingrouping()]==selectedprotein][["PEP.StrippedSequence"]].drop_duplicates().sort_values("PEP.StrippedSequence")
            return render.DataGrid(df,width="100%",selection_mode="row")  
    #line plot of either protein or peptide signal depending on what is selected in the above two tables
    @reactive.effect
    def _():
        @render.plot(width=input.tracker_width(),height=input.tracker_height())
        def tracker_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            #show empty graph since nothing was selected
            if len(protein_df.data_view(selected=True)[input.tracker_proteingrouping()].tolist())==0:
                fig,ax=plt.subplots()
            #if just protien is selected, show protein intensities across runs
            elif len(pickedprotein_df.data_view(selected=True)["PEP.StrippedSequence"].tolist())==0:
                fig,ax=plt.subplots(figsize=(input.tracker_width()*(1/plt.rcParams['figure.dpi']),input.tracker_height()*(1/plt.rcParams['figure.dpi'])))
                selectedprotein=protein_df.data_view(selected=True)[input.tracker_proteingrouping()].tolist()[0]
                plottingdf=searchoutput[searchoutput[input.tracker_proteingrouping()]==selectedprotein][["Cond_Rep","PG.MS2Quantity"]].drop_duplicates().fillna(0)
                if len(plottingdf)<len(runnames):
                    expectedrows=pd.DataFrame({"Cond_Rep":runnames})
                    plottingdf=expectedrows.merge(plottingdf,how="left",left_on="Cond_Rep",right_on="Cond_Rep").fillna(0)
                if input.tracker_logscale()==True:
                    plottingdf["PG.MS2Quantity"]=np.log10(plottingdf["PG.MS2Quantity"])
                else:
                    plottingdf=plottingdf
                ax.plot(plottingdf["Cond_Rep"],plottingdf["PG.MS2Quantity"],marker="o")
                if str(selectedprotein).count(";")>=1:
                    selectedprotein=str(selectedprotein).split(";")[0]
                ax.set_title(str(selectedprotein),fontsize=titlefont)
            #show intensity for selected peptide from selected protein across runs
            else:
                selectedprotein=protein_df.data_view(selected=True)[input.tracker_proteingrouping()].tolist()[0]
                proteindf=searchoutput[searchoutput[input.tracker_proteingrouping()]==selectedprotein][["Cond_Rep","PEP.StrippedSequence","EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]].drop_duplicates()

                selectedpeptide=pickedprotein_df.data_view(selected=True)["PEP.StrippedSequence"].tolist()[0]
                peptidedf=proteindf[proteindf["PEP.StrippedSequence"]==selectedpeptide]

                modpeps=peptidedf["EG.ModifiedPeptide"].drop_duplicates().tolist()

                fig,ax=plt.subplots(figsize=(input.tracker_width()*(1/plt.rcParams['figure.dpi']),input.tracker_height()*(1/plt.rcParams['figure.dpi'])))
                for pep in modpeps:
                    chargelist=peptidedf[peptidedf["EG.ModifiedPeptide"]==pep]["FG.Charge"].drop_duplicates().tolist()
                    for charge in chargelist:
                        plottingdf=peptidedf[(peptidedf["EG.ModifiedPeptide"]==pep)&(peptidedf["FG.Charge"]==charge)]
                        if len(plottingdf)<len(runnames):
                            expectedrows=pd.DataFrame({"Cond_Rep":runnames})
                            plottingdf=expectedrows.merge(plottingdf,how="left",left_on="Cond_Rep",right_on="Cond_Rep").fillna(0)
                        else:
                            pass
                        if input.tracker_logscale()==True:
                            plottingdf["FG.MS2Quantity"]=np.log10(plottingdf["FG.MS2Quantity"])
                        else:
                            plottingdf=plottingdf
                        ax.plot(plottingdf["Cond_Rep"],plottingdf["FG.MS2Quantity"],marker="o",label=pep.strip("_")+"_"+str(charge)+"+")

                ax.legend(loc='center left',bbox_to_anchor=(1,0.5),prop={'size':legendfont})
                if str(selectedprotein).count(";")>=1:
                    selectedprotein=str(selectedprotein).split(";")[0]
                ax.set_title(str(selectedprotein)+"_"+str(selectedpeptide),fontsize=titlefont)
            ax.tick_params(axis="x",rotation=x_label_rotation)
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            ax.set_xlabel("Condition",fontsize=axisfont)
            if input.tracker_logscale()==True:
                ax.set_ylabel("MS2 Intensity (log10)",fontsize=axisfont)
            else:
                ax.set_ylabel("MS2 Intensity",fontsize=axisfont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            fig.set_tight_layout(True)
            if input.tracker_yaxiszero()==True:
                ax.set_ylim(bottom=0)
                top=ax.get_ylim()[1]
                ax.set_ylim(top=top+(top*0.1))
            if len(protein_df.data_view(selected=True)[input.tracker_proteingrouping()].tolist())!=0:
                imagedownload("signaltracker")
   
    # ====================================== Individual Protein ID Counts
    @render.data_frame
    def protein_ID_df():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        df=searchoutput[["PG.ProteinGroups","PG.ProteinNames","PG.Genes","PG.MS2Quantity"]].groupby(["PG.ProteinGroups","PG.ProteinNames","PG.Genes"]).mean().reset_index().rename(columns={"PG.MS2Quantity":"Mean_PG.MS2Quantity"})
        df_styles={"location":"body",
                   "style":{"column-width":"25px",
                            "overflow":"hidden"
                           }
                  }
        return render.DataGrid(df,selection_mode="row",editable=False,styles=df_styles)
    @reactive.effect
    def _():
        @render.plot(width=input.protein_id_plot_width(),height=input.protein_id_plot_height())
        def protein_id_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if len(protein_ID_df.data_view(selected=True)["PG.ProteinGroups"].tolist())==0:
                fig,ax=plt.subplots()
            else:
                selectedprotein=protein_ID_df.data_view(selected=True)["PG.ProteinGroups"].tolist()[0]
                df=searchoutput[searchoutput["PG.ProteinGroups"]==selectedprotein]

                if ";" in selectedprotein:
                    proteintitle=selectedprotein.split(";")[0]
                else:
                    proteintitle=selectedprotein

                proteinresultdf=pd.DataFrame(df[["Cond_Rep","R.FileName","R.Condition","R.Replicate"]].drop_duplicates()).reset_index(drop=True)
                sampleconditions=df["R.Condition"].drop_duplicates().tolist()
                maxreplicatelist=[]
                for i in sampleconditions:
                    samplegroup=pd.DataFrame(df[df["R.Condition"]==i])
                    maxreplicates=len(samplegroup["R.Replicate"].drop_duplicates().tolist())
                    maxreplicatelist.append(maxreplicates)
                proteinaveragedf=pd.DataFrame({"R.Condition":sampleconditions,"N.Replicates":maxreplicatelist})

                numpeptides=[]
                numstrippedpeptides=[]
                numprecursors=[]

                for i in df["Cond_Rep"].drop_duplicates().tolist():
                    replicatedata=df[df["Cond_Rep"]==i]
                    if replicatedata.empty:
                        continue
                    numpeptides.append(replicatedata["EG.ModifiedPeptide"].nunique())
                    numstrippedpeptides.append(replicatedata["PEP.StrippedSequence"].nunique())
                    numprecursors.append(len(replicatedata[["EG.ModifiedPeptide","FG.Charge"]].drop_duplicates()))
                proteinresultdf["peptides"]=numpeptides
                proteinresultdf["strippedpeptides"]=numstrippedpeptides
                proteinresultdf["precursors"]=numprecursors

                columnlist=proteinresultdf.columns.values.tolist()
                for i in columnlist:
                    if i=="R.FileName" or i=="Cond_Rep" or i=="R.Condition" or i=="R.Replicate":
                        continue
                    avglist=[]
                    stdevlist=[]
                    for j in sampleconditions:
                        samplecondition=proteinresultdf[proteinresultdf["R.Condition"]==j]
                        avglist.append(round(np.average(samplecondition[i].to_numpy())))
                        stdevlist.append(np.std(samplecondition[i].to_numpy()))
                    proteinaveragedf[i+"_avg"]=avglist
                    proteinaveragedf[i+"_stdev"]=stdevlist

                fig,ax=plt.subplots(ncols=3,sharey=True,figsize=(input.protein_id_plot_width()*(1/plt.rcParams['figure.dpi']),input.protein_id_plot_height()*(1/plt.rcParams['figure.dpi'])))

                if input.protein_id_individual_average()=="individual":
                    idmetricscolor=replicatecolors()
                    ax[0].bar(proteinresultdf["Cond_Rep"],proteinresultdf["peptides"],width=0.8,edgecolor="k",color=idmetricscolor)
                    ax[1].bar(proteinresultdf["Cond_Rep"],proteinresultdf["strippedpeptides"],width=0.8,edgecolor="k",color=idmetricscolor)
                    ax[2].bar(proteinresultdf["Cond_Rep"],proteinresultdf["precursors"],width=0.8,edgecolor="k",color=idmetricscolor)

                    ax[0].bar_label(ax[0].containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax[1].bar_label(ax[1].containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax[2].bar_label(ax[2].containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)

                    ax[0].set_ylim(top=max(proteinresultdf["peptides"].tolist())+y_padding*max(proteinresultdf["peptides"].tolist()))
                    ax[1].set_ylim(top=max(proteinresultdf["strippedpeptides"].tolist())+y_padding*max(proteinresultdf["strippedpeptides"].tolist()))
                    ax[2].set_ylim(top=max(proteinresultdf["precursors"].tolist())+y_padding*max(proteinresultdf["precursors"].tolist()))
                elif input.protein_id_individual_average()=="average":
                    idmetricscolor=colorpicker()
                    bars1=ax[0].bar(proteinaveragedf["R.Condition"],proteinaveragedf["peptides_avg"],yerr=proteinaveragedf["peptides_stdev"],capsize=10,width=0.8,edgecolor="k",color=idmetricscolor)
                    bars2=ax[1].bar(proteinaveragedf["R.Condition"],proteinaveragedf["strippedpeptides_avg"],yerr=proteinaveragedf["strippedpeptides_stdev"],capsize=10,width=0.8,edgecolor="k",color=idmetricscolor)
                    bars3=ax[2].bar(proteinaveragedf["R.Condition"],proteinaveragedf["precursors_avg"],yerr=proteinaveragedf["precursors_stdev"],capsize=10,width=0.8,edgecolor="k",color=idmetricscolor)

                    ax[0].bar_label(bars1,label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax[1].bar_label(bars2,label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax[2].bar_label(bars3,label_type="edge",rotation=90,padding=5,fontsize=labelfont)

                    ax[0].set_ylim(top=max(proteinaveragedf["peptides_avg"].tolist())+y_padding*max(proteinaveragedf["peptides_avg"].tolist()))
                    ax[1].set_ylim(top=max(proteinaveragedf["strippedpeptides_avg"].tolist())+y_padding*max(proteinaveragedf["strippedpeptides_avg"].tolist()))
                    ax[2].set_ylim(top=max(proteinaveragedf["precursors_avg"].tolist())+y_padding*max(proteinaveragedf["precursors_avg"].tolist()))

                ax[0].set_axisbelow(True)
                ax[0].grid(linestyle="--")
                ax[0].tick_params(axis="both",labelsize=axisfont_labels)
                ax[0].tick_params(axis="x",rotation=x_label_rotation)
                ax[0].set_title("Modified Peptides",fontsize=titlefont)
                ax[0].set_ylabel("Counts",rotation="vertical",fontsize=axisfont)
                ax[1].set_axisbelow(True)
                ax[1].grid(linestyle="--")
                ax[1].tick_params(axis="both",labelsize=axisfont_labels)
                ax[1].tick_params(axis="x",rotation=x_label_rotation)
                ax[1].set_title("Stripped Peptides",fontsize=titlefont)
                ax[2].set_axisbelow(True)
                ax[2].grid(linestyle="--")
                ax[2].tick_params(axis="both",labelsize=axisfont_labels)
                ax[2].tick_params(axis="x",rotation=x_label_rotation)
                ax[2].set_title("Precursors",fontsize=titlefont)

                fig.suptitle("ID Counts for: "+proteintitle,fontsize=titlefont)
                fig.set_tight_layout(True)
            if len(protein_ID_df.data_view(selected=True)["PG.ProteinGroups"].tolist())!=0:
                imagedownload("proteinids_"+proteintitle)

#endregion

# ============================================================================= Metrics
#region
    # ====================================== Charge State
    #plot charge states
    @render.ui
    def chargestate_peplength_slider_ui():
        if input.chargestate_peplength()==True:
            return ui.input_slider("chargestate_peplength_slider_pick","Pick peptide length to plot:",min=7,max=25,value=9,step=1,ticks=True)
    @render.ui
    def chargestate_averages_ui():
        if input.chargestate_condition_or_run()=="condition" and input.chargestate_stacked()==False:
            return ui.input_switch("chargestate_averages_switch","Show bar plots as averages with error bars",value=False)
    @reactive.effect
    def _():
        @render.plot(width=input.chargestate_width(),height=input.chargestate_height())
        def chargestateplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            df=searchoutput[["Cond_Rep","R.Condition","EG.ModifiedPeptide","PEP.StrippedSequence","FG.Charge","Peptide Length"]]
            charges=list(set(df["FG.Charge"]))
            chargedf=pd.DataFrame()
            if input.chargestate_peplength()==True:
                chargedf_ref=df[df["Peptide Length"]==input.chargestate_peplength_slider_pick()][["Cond_Rep","EG.ModifiedPeptide","FG.Charge"]]
            else:
                chargedf_ref=df[["Cond_Rep","EG.ModifiedPeptide","FG.Charge"]]
            if input.chargestate_condition_or_run()=="individual":
                for charge in charges:
                    chargecounts=[]
                    for run in chargedf_ref["Cond_Rep"].drop_duplicates().tolist():
                        chargecounts.append(len(chargedf_ref[(chargedf_ref["Cond_Rep"]==run)&(chargedf_ref["FG.Charge"]==charge)].drop_duplicates()))
                    chargedf[charge]=chargecounts
                chargedf["Run"]=chargedf_ref["Cond_Rep"].drop_duplicates().tolist()
                chargedf=chargedf.set_index("Run").transpose()
            elif input.chargestate_condition_or_run()=="condition":
                if input.chargestate_averages_switch()==True:
                    for charge in charges:
                        chargecounts=[]
                        for run in chargedf_ref["Cond_Rep"].drop_duplicates().tolist():
                            chargecounts.append(len(chargedf_ref[(chargedf_ref["Cond_Rep"]==run)&(chargedf_ref["FG.Charge"]==charge)].drop_duplicates()))
                        chargedf[charge]=chargecounts
                    chargedf["Run"]=chargedf_ref["Cond_Rep"].drop_duplicates().tolist()
                    chargedf=chargedf.set_index("Run").transpose()
                    #calculate averages of the counts or percents
                    chargedf_transpose=chargedf.transpose().reset_index()
                    chargedf_transpose["R.Condition"]=chargedf_transpose["Run"].str.rsplit("_",n=1,expand=True)[0]
                    chargedf_transpose=chargedf_transpose.drop(columns="Run")
                    chargeresultdf_avg=pd.DataFrame()
                    chargeresultdf_stdev=pd.DataFrame()
                    chargeresultdf_avg["R.Condition"]=chargedf_transpose["R.Condition"].drop_duplicates().tolist()
                    chargeresultdf_stdev["R.Condition"]=chargedf_transpose["R.Condition"].drop_duplicates().tolist()
                    columns=chargedf_transpose.columns.tolist()[:-1]
                    for charge in columns:
                        avg=[]
                        stdev=[]
                        for condition in chargedf_transpose["R.Condition"].drop_duplicates().tolist():
                            avg.append(round(chargedf_transpose.groupby("R.Condition").get_group(condition)[charge].mean(),2))
                            stdev.append(chargedf_transpose.groupby("R.Condition").get_group(condition)[charge].std())
                        chargeresultdf_avg[charge]=avg
                        chargeresultdf_stdev[charge]=stdev
                    chargeresultdf_avg=chargeresultdf_avg.set_index("R.Condition").transpose()
                    chargeresultdf_stdev=chargeresultdf_stdev.set_index("R.Condition").transpose()
                    chargedf=chargeresultdf_avg
                else:
                    #count number of unique precursors per charge state in each condition
                    if input.chargestate_peplength()==True:
                        chargedf_ref=df[df["Peptide Length"]==input.chargestate_peplength_slider_pick()][["R.Condition","EG.ModifiedPeptide","FG.Charge"]]
                    else:
                        chargedf_ref=df[["R.Condition","EG.ModifiedPeptide","FG.Charge"]]
                    for charge in charges:
                        chargecounts=[]
                        for run in chargedf_ref["R.Condition"].drop_duplicates().tolist():
                            chargecounts.append(len(chargedf_ref[(chargedf_ref["R.Condition"]==run)&(chargedf_ref["FG.Charge"]==charge)].drop_duplicates()))
                        chargedf[charge]=chargecounts
                    chargedf["Run"]=chargedf_ref["R.Condition"].drop_duplicates().tolist()
                    chargedf=chargedf.set_index("Run").transpose()
            #use to show counts calculated above as percents
            if input.chargestate_counts_percent()=="Percent":
                chargedf_percents=pd.DataFrame()
                for run in chargedf.columns:
                    total=chargedf[run].sum()
                    percentlist=[]
                    for value in chargedf[run]:
                        percentlist.append(round((value/total*100),1))
                    chargedf_percents[run]=percentlist
                chargedf_percents=chargedf_percents.set_index(pd.Index(charges))
                chargedf_percents=chargedf_percents.transpose().reset_index().rename(columns={"index":"Run"}).set_index("Run").transpose()
                chargedf=chargedf_percents

            if input.chargestate_counts_percent()=="Percent":
                ylabel="Frequency (%)"
            if input.chargestate_counts_percent()=="Counts":
                ylabel="Counts"
            if input.chargestate_condition_or_run()=="individual":
                chargestatecolor=replicatecolors()
                stackxlabel="Run"
            if input.chargestate_condition_or_run()=="condition":
                chargestatecolor=colorpicker()
                stackxlabel="Condition"

            #show as stacked bar graphs
            if input.chargestate_stacked()==True:
                fig,ax=plt.subplots(figsize=(input.chargestate_width()*(1/plt.rcParams['figure.dpi']),input.chargestate_height()*(1/plt.rcParams['figure.dpi'])))
                matplottabcolors=list(mcolors.TABLEAU_COLORS)
                x=np.arange(len(chargedf.columns))
                for i,column in enumerate(chargedf.columns):
                    bottom=0
                    for j,charge in enumerate(charges):
                        ax.bar(x[i],chargedf[column][charge],width=0.75,bottom=bottom,color=matplottabcolors[j])
                        bottom+=chargedf[column][charge]
                ax.legend(charges,loc="center left",bbox_to_anchor=(1, 0.5),prop={'size':legendfont})
                if input.chargestate_counts_percent()=="Percent":
                    ax.set_ylim(top=105)
                ax.set_xticks(x,labels=chargedf.columns,rotation=x_label_rotation)
                ax.set_ylabel(ylabel,fontsize=axisfont)
                ax.set_xlabel(stackxlabel,fontsize=axisfont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)
            #show as separate bar graphs
            else:
                x=charges
                if len(chargedf.columns.tolist())==1:
                    fig,ax=plt.subplots(figsize=(input.chargestate_width()*(1/plt.rcParams['figure.dpi']),input.chargestate_height()*(1/plt.rcParams['figure.dpi'])))
                    column=chargedf.columns.tolist()[0]
                    if input.chargestate_averages_switch()==True:
                        ax.bar(x,chargedf[column],yerr=chargeresultdf_stdev[column],capsize=8,edgecolor="k",color=chargestatecolor)
                        ax.bar_label(ax.containers[1],label_type="edge",padding=10,rotation=90,fontsize=labelfont)
                    else:
                        ax.bar(x,chargedf[column],edgecolor="k",color=chargestatecolor)
                        ax.bar_label(ax.containers[0],label_type="edge",padding=10,rotation=90,fontsize=labelfont)
                    ax.set_xticks(np.arange(1,max(x)+1,1))
                    ax.set_title(column,fontsize=titlefont)
                    ax.set_xlabel("Charge State",fontsize=axisfont)
                    ax.set_ylabel(ylabel,fontsize=axisfont)
                    ax.set_ylim(bottom=0-(max(chargedf[column])*0.05),top=max(chargedf[column])+y_padding*max(chargedf[column]))

                    ax.set_axisbelow(True)
                    ax.grid(linestyle="--")
                    ax.tick_params(axis="both",labelsize=axisfont_labels)
                else:
                    fig,ax=plt.subplots(nrows=1,ncols=len(chargedf.columns),sharey=True,figsize=(input.chargestate_width()*(1/plt.rcParams['figure.dpi']),input.chargestate_height()*(1/plt.rcParams['figure.dpi'])))
                    for i,run in enumerate(chargedf.columns.tolist()):
                        if input.chargestate_averages_switch()==True:
                            ax[i].bar(x,chargedf[run].tolist(),yerr=chargeresultdf_stdev[run],capsize=8,color=chargestatecolor[i],edgecolor="k")
                            ax[i].bar_label(ax[i].containers[1],label_type="edge",padding=10,rotation=90,fontsize=labelfont)
                        else:
                            if numconditions==1:
                                ax[i].bar(x,chargedf[run].tolist(),color=chargestatecolor,edgecolor="k")
                                ax[i].bar_label(ax[i].containers[0],label_type="edge",padding=10,rotation=90,fontsize=labelfont)
                            else:
                                ax[i].bar(x,chargedf[run].tolist(),color=chargestatecolor[i],edgecolor="k")
                                ax[i].bar_label(ax[i].containers[0],label_type="edge",padding=10,rotation=90,fontsize=labelfont)
                        ax[i].set_title(chargedf.columns[i],fontsize=titlefont)
                        ax[i].set_xticks(np.arange(1,max(x)+1,1))
                        
                        #ax[i].set_ylim(bottom=0-(max(chargedf[run])*0.05),top=max(chargedf[run])+y_padding*max(chargedf[run]))
                        ax[i].set_xticks(np.arange(1,max(x)+1,1))
                        ax[i].set_xlabel("Charge State",fontsize=axisfont)
                        
                        ax[i].set_axisbelow(True)
                        ax[i].grid(linestyle="--")
                        ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                    ax[0].set_ylim(bottom=0-(max(chargedf.max().tolist())*0.05),top=max(chargedf.max().tolist())+(y_padding*max(chargedf.max().tolist())))
                    ax[0].set_ylabel(ylabel,fontsize=axisfont)

            fig.set_tight_layout(True)
            imagedownload("chargestates")

    # ====================================== Peptide Length
    #ui call for dropdown for marking peptide lengths
    @render.ui
    def lengthmark_ui():
        if input.peplengthinput()=="barplot":
            minlength=7
            maxlength=30
            opts=[item for item in range(minlength,maxlength+1)]
            opts.insert(0,0)
            return ui.column(3,ui.input_switch("hide_lengthmark","Hide peptide length marker"),ui.input_selectize("lengthmark_pick","Pick peptide length to mark on bar plot (use 0 for maximum)",choices=opts))
    #plot peptide legnths
    @reactive.effect
    def _():
        @render.plot(width=input.peptidelength_width(),height=input.peptidelength_height())
        def peptidelengthplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            colors=colorpicker()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            peptidelengths_condition,peptidelengths_run=peptidelengths()

            if input.peptidelengths_condition_or_run()=="condition":
                plottingdf=peptidelengths_condition
                colors=colorpicker()
            if input.peptidelengths_condition_or_run()=="individual":
                plottingdf=peptidelengths_run
                colors=replicatecolors()

            if input.peplengthinput()=="lineplot":
                legendlist=[]
                fig,ax=plt.subplots(figsize=(input.peptidelength_width()*(1/plt.rcParams['figure.dpi']),input.peptidelength_height()*(1/plt.rcParams['figure.dpi'])))
                for i in range(len(plottingdf)):
                    if len(plottingdf)==1:
                        x=list(set(plottingdf["Peptide Lengths"][0]))
                        frequencies=[len(list(group)) for key, group in groupby(sorted(plottingdf["Peptide Lengths"][0]))]
                        ax.plot(x,frequencies,color=colors,linewidth=2)
                        legendlist.append(plottingdf["Sample Names"][0])
                    else:
                        x=list(set(plottingdf["Peptide Lengths"][i]))
                        frequencies=[len(list(group)) for key, group in groupby(sorted(plottingdf["Peptide Lengths"][i]))]
                        if numconditions==1:
                            ax.plot(x,frequencies,color=colors,linewidth=2)
                        else:
                            if input.peptidelengths_condition_or_run()=="individual":
                                colors_light=[]
                                for j in range(len(colors)):
                                    if j==0 or colors[j-1]!=colors[j]:
                                        colors_light.append(colors[j])
                                    else:
                                        c=colorsys.rgb_to_hls(*mcolors.to_rgb(colors_light[j-1]))
                                        colors_light.append(colorsys.hls_to_rgb(c[0],1-0.5*(1-c[1]),c[2]))
                                colors=colors_light
                            ax.plot(x,frequencies,color=colors[i],linewidth=2)
                        legendlist.append(plottingdf["Sample Names"][i])
                ax.set_xlabel("Peptide Length",fontsize=axisfont)
                ax.set_ylabel("Counts",fontsize=axisfont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                ax.legend(legendlist,prop={'size':legendfont})

            if input.peplengthinput()=="barplot":
                lengthmark=int(input.lengthmark_pick())
                if len(plottingdf)==1:
                    fig,ax=plt.subplots(figsize=(input.peptidelength_width()*(1/plt.rcParams['figure.dpi']),input.peptidelength_height()*(1/plt.rcParams['figure.dpi'])))
                    x=list(set(plottingdf["Peptide Lengths"][0]))
                    frequencies=[len(list(group)) for key, group in groupby(sorted(plottingdf["Peptide Lengths"][0]))]
                    ax.bar(x,frequencies,color=colors,edgecolor="k")
                    ax.set_title(plottingdf["Sample Names"][0],fontsize=titlefont)
                    ax.set_axisbelow(True)
                    ax.grid(linestyle="--")
                    if input.hide_lengthmark()==False:
                        if lengthmark!=0:
                            ax.vlines(x=x[lengthmark-min(x)],ymin=frequencies[lengthmark-min(x)],ymax=frequencies[lengthmark-min(x)]+0.2*frequencies[lengthmark-min(x)],color="k")
                            ax.text(x=x[lengthmark-min(x)],y=frequencies[lengthmark-min(x)]+0.2*frequencies[lengthmark-min(x)],s=str(x[lengthmark-min(x)])+": "+str(frequencies[lengthmark-min(x)]),fontsize=labelfont)
                        if lengthmark==0:
                            ax.vlines(x=x[np.argmax(frequencies)],ymin=max(frequencies),ymax=max(frequencies)+0.2*max(frequencies),color="k")
                            ax.text(x=x[np.argmax(frequencies)],y=max(frequencies)+0.2*max(frequencies),s=str(x[np.argmax(frequencies)])+": "+str(max(frequencies)),fontsize=labelfont)
                            ax.set_ylim(top=max(frequencies)+y_padding*max(frequencies))
                    ax.set_xlabel("Peptide Length",fontsize=axisfont)
                    ax.set_ylabel("Counts",fontsize=axisfont)
                    ax.xaxis.set_minor_locator(MultipleLocator(1))
                    ax.tick_params(axis="both",labelsize=axisfont_labels)
                else:
                    fig,ax=plt.subplots(nrows=1,ncols=len(plottingdf),sharey=True,figsize=(input.peptidelength_width()*(1/plt.rcParams['figure.dpi']),input.peptidelength_height()*(1/plt.rcParams['figure.dpi'])))
                    for i in range(len(plottingdf)):
                        x=list(set(plottingdf["Peptide Lengths"][i]))
                        frequencies=[len(list(group)) for key, group in groupby(sorted(plottingdf["Peptide Lengths"][i]))]
                        if numconditions==1:
                            ax[i].bar(x,frequencies,color=colors,edgecolor="k")
                        else:
                            ax[i].bar(x,frequencies,color=colors[i],edgecolor="k")
                        ax[i].set_title(plottingdf["Sample Names"][i],fontsize=titlefont)
                        ax[i].set_axisbelow(True)
                        ax[i].grid(linestyle="--")
                        if input.hide_lengthmark()==False:
                            if lengthmark!=0:
                                ax[i].vlines(x=x[lengthmark-min(x)],ymin=frequencies[lengthmark-min(x)],ymax=frequencies[lengthmark-min(x)]+0.2*frequencies[lengthmark-min(x)],color="k")
                                ax[i].text(x=x[lengthmark-min(x)],y=frequencies[lengthmark-min(x)]+0.2*frequencies[lengthmark-min(x)],s=str(x[lengthmark-min(x)])+": "+str(frequencies[lengthmark-min(x)]),fontsize=labelfont)
                            if lengthmark==0:
                                ax[i].vlines(x=x[np.argmax(frequencies)],ymin=max(frequencies),ymax=max(frequencies)+0.2*max(frequencies),color="k")
                                ax[i].text(x=x[np.argmax(frequencies)],y=max(frequencies)+0.2*max(frequencies),s=str(x[np.argmax(frequencies)])+": "+str(max(frequencies)),fontsize=labelfont)
                                #ax[i].set_ylim(top=max(frequencies)+y_padding*max(frequencies))
                        ax[i].set_xlabel("Peptide Length",fontsize=axisfont)
                        ax[i].xaxis.set_minor_locator(MultipleLocator(1))
                        ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                    ymax_list=[]
                    for i in range(len(ax)):
                        ymax_list.append(ax[i].get_ylim()[1])
                    ax[0].set_ylim(top=max(ymax_list)+(y_padding*max(ymax_list)))
                    ax[0].set_ylabel("Counts",fontsize=axisfont)
            fig.set_tight_layout(True)
            imagedownload("peptidelengths")
    
    # ====================================== Peptides per Protein
    #plot peptides per protein
    @reactive.effect
    def _():
        @render.plot(width=input.pepsperprotein_width(),height=input.pepsperprotein_height())
        def pepsperproteinplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            pepsperprotein_condition,pepsperprotein_run=pepsperprotein()

            if input.pepsperprotein_condition_or_run()=="condition":
                plottingdf=pepsperprotein_condition
                colors=colorpicker()
            if input.pepsperprotein_condition_or_run()=="individual":
                plottingdf=pepsperprotein_run
                colors=replicatecolors()

            if input.pepsperproteininput()=="lineplot":
                legendlist=[]
                fig,ax=plt.subplots(figsize=(input.pepsperprotein_width()*(1/plt.rcParams['figure.dpi']),input.pepsperprotein_height()*(1/plt.rcParams['figure.dpi'])))

                if len(plottingdf)==1:
                    x=sorted(list(set(plottingdf["Peptides per Protein"][0])))
                    frequencies=[len(list(group)) for key, group in groupby(sorted(plottingdf["Peptides per Protein"][0]))]
                    ax.plot(x,frequencies,color=colors,linewidth=2)
                    legendlist.append(plottingdf["Sample Names"][0])
                else:
                    for i in range(len(plottingdf)):
                        x=sorted(list(set(plottingdf["Peptides per Protein"][i])))
                        frequencies=[len(list(group)) for key, group in groupby(sorted(plottingdf["Peptides per Protein"][i]))]
                        if numconditions==1:
                            ax.plot(x,frequencies,color=colors,linewidth=2)
                        else:
                            if input.pepsperprotein_condition_or_run()=="individual":
                                colors_light=[]
                                for j in range(len(colors)):
                                    if j==0 or colors[j-1]!=colors[j]:
                                        colors_light.append(colors[j])
                                    else:
                                        c=colorsys.rgb_to_hls(*mcolors.to_rgb(colors_light[j-1]))
                                        colors_light.append(colorsys.hls_to_rgb(c[0],1-0.5*(1-c[1]),c[2]))
                                colors=colors_light
                            ax.plot(x,frequencies,color=colors[i],linewidth=2)
                        legendlist.append(plottingdf["Sample Names"][i])
                ax.set_xlim(left=0,right=input.pepsperprotein_xrange())
                ax.set_xlabel("Peptides per Protein",fontsize=axisfont)
                ax.set_ylabel("Counts",fontsize=axisfont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.legend(legendlist,prop={'size':legendfont})
                ax.tick_params(axis="both",labelsize=axisfont_labels)

            if input.pepsperproteininput()=="barplot":
                if len(plottingdf)==1:
                    fig,ax=plt.subplots(figsize=(input.pepsperprotein_width()*(1/plt.rcParams['figure.dpi']),input.pepsperprotein_height()*(1/plt.rcParams['figure.dpi'])))
                    for i in range(len(plottingdf)):
                        x=sorted(list(set(plottingdf["Peptides per Protein"][0])))
                        frequencies=[len(list(group)) for key, group in groupby(sorted(plottingdf["Peptides per Protein"][0]))]

                        ax.bar(x,frequencies,color=colors,width=0.025)
                        ax.set_title(plottingdf["Sample Names"][0],fontsize=titlefont)
                        ax.set_axisbelow(True)
                        ax.grid(linestyle="--")

                        ax.set_xticks(np.arange(0,max(x)+1,25))
                        ax.set_xlabel("# Peptides",fontsize=axisfont)
                        ax.set_ylabel("Counts",fontsize=axisfont)
                        ax.set_xlim(left=0,right=input.pepsperprotein_xrange())
                        ax.tick_params(axis="both",labelsize=axisfont_labels)
                else:
                    fig,ax=plt.subplots(nrows=1,ncols=len(plottingdf),figsize=(input.pepsperprotein_width()*(1/plt.rcParams['figure.dpi']),input.pepsperprotein_height()*(1/plt.rcParams['figure.dpi'])))
                    for i in range(len(plottingdf)):
                        x=list(set(plottingdf["Peptides per Protein"][i]))
                        frequencies=[len(list(group)) for key, group in groupby(sorted(plottingdf["Peptides per Protein"][i]))]
                        if numconditions==1:
                            ax[i].bar(x,frequencies,color=colors)
                        else:
                            ax[i].bar(x,frequencies,color=colors[i])
                        ax[i].set_title(plottingdf["Sample Names"][i],fontsize=titlefont)
                        ax[i].set_axisbelow(True)
                        ax[i].grid(linestyle="--")

                        ax[i].set_xticks(np.arange(0,max(x)+1,25))
                        ax[i].set_xlabel("# Peptides",fontsize=axisfont)
                        ax[i].set_xlim(left=0,right=input.pepsperprotein_xrange())
                        ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                    ax[0].set_ylabel("Counts",fontsize=axisfont)
                fig.set_tight_layout(True)
                imagedownload("peptidesperprotein")
    
    # ====================================== Dynamic Range
    #plot dynamic range
    @reactive.effect
    def _():
        @render.plot(width=input.dynamicrange_width(),height=input.dynamicrange_height())
        def dynamicrangeplot():
            conditioninput=input.conditionname()
            propertyinput=input.meanmedian()
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            markersize=25
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if propertyinput=="mean":
                intensitydf=searchoutput[searchoutput["R.Condition"]==conditioninput][["PG.ProteinGroups","PG.MS2Quantity"]].drop_duplicates().groupby("PG.ProteinGroups").mean().reset_index(drop=True)

            elif propertyinput=="median":
                intensitydf=searchoutput[searchoutput["R.Condition"]==conditioninput][["PG.ProteinGroups","PG.MS2Quantity"]].drop_duplicates().groupby("PG.ProteinGroups").median().reset_index(drop=True)

            fig,ax=plt.subplots(nrows=2,ncols=1,sharex=True,gridspec_kw={"height_ratios":[1,3]},figsize=(input.dynamicrange_width()*(1/plt.rcParams['figure.dpi']),input.dynamicrange_height()*(1/plt.rcParams['figure.dpi'])))
            ax1=ax[0]
            ax2=ax[1]

            maxintensity=intensitydf.max()
            relative_fraction=(1-(intensitydf/maxintensity)).sort_values(by="PG.MS2Quantity").reset_index(drop=True)
            n_25=relative_fraction[relative_fraction["PG.MS2Quantity"]<0.25].shape[0]
            n_50=relative_fraction[relative_fraction["PG.MS2Quantity"]<0.50].shape[0]
            n_75=relative_fraction[relative_fraction["PG.MS2Quantity"]<0.75].shape[0]

            ax1.scatter(relative_fraction.index,relative_fraction["PG.MS2Quantity"],marker=".",s=markersize)
            ax1.set_ylabel("Relative Fraction of\nProtein Signal",fontsize=axisfont)
            ax1.text(0,0.2,"- - - - - - - "+str(n_25)+" Protein groups",fontsize=labelfont)
            ax1.text(0,0.45,"- - - - - - - "+str(n_50)+" Protein groups",fontsize=labelfont)
            ax1.text(0,0.7,"- - - - - - - "+str(n_75)+" Protein groups",fontsize=labelfont)

            log10df=np.log10(intensitydf).sort_values(by="PG.MS2Quantity",ascending=False).reset_index(drop=True)
            dynamicrange=round(max(log10df["PG.MS2Quantity"])-min(log10df[log10df["PG.MS2Quantity"]!=float("-inf")]["PG.MS2Quantity"]),1)

            ax2.scatter(log10df.index,log10df["PG.MS2Quantity"],marker=".",s=markersize)
            ax2.set_ylabel("Log10(Protein Signal)",fontsize=axisfont)
            ax2.text(max(log10df.index)-0.6*(max(log10df.index)),max(log10df["PG.MS2Quantity"])-0.15*(max(log10df["PG.MS2Quantity"])),str(dynamicrange)+" log",fontsize=titlefont)

            plt.xlabel("Rank",fontsize=axisfont)
            plt.suptitle(conditioninput+" ("+propertyinput+"_PG)",x=0.13,horizontalalignment="left",fontsize=titlefont)
            ax1.set_axisbelow(True)
            ax2.set_axisbelow(True)
            ax1.grid(linestyle="--")
            ax2.grid(linestyle="--")
            ax1.tick_params(axis="both",labelsize=axisfont_labels)
            ax2.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("dynamicrange")
    #get ranked proteins based on signal
    @render.data_frame
    def dynamicrange_proteinrank():
        conditioninput=input.conditionname()
        propertyinput=input.meanmedian()

        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

        if propertyinput=="mean":
            intensitydf=searchoutput[searchoutput["R.Condition"]==conditioninput][["PG.ProteinGroups","PG.MS2Quantity"]].drop_duplicates().groupby("PG.ProteinGroups").mean().reset_index()
        elif propertyinput=="median":
            intensitydf=searchoutput[searchoutput["R.Condition"]==conditioninput][["PG.ProteinGroups","PG.MS2Quantity"]].drop_duplicates().groupby("PG.ProteinGroups").median().reset_index()
        intensitydf=intensitydf.sort_values("PG.MS2Quantity",ascending=False).reset_index(drop=True)
        df_styles={"location":"body",
                   "style":{"column-width":"200px",
                            "overflow":"hidden"
                           }
                  }
        return render.DataGrid(intensitydf.iloc[:input.top_n()],editable=False,styles=df_styles)

    # ====================================== Mass Accuracy
    @render.ui
    def massaccuracy_bins_ui():
        if input.massaccuracy_violin_hist()=="histogram":
            return ui.input_slider("massaccuracy_hist_bins","Number of bins",min=10,max=200,value=100,step=10,ticks=True,width="300px")
    #plot mass accuracy
    @reactive.effect
    def _():
        @render.plot(width=input.massaccuracy_width(),height=input.massaccuracy_height())
        def massaccuracy_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            massaccuracy_df=pd.DataFrame()
            massaccuracy_df["Cond_Rep"]=runnames
            massaccuracy=[]
            for run in runnames:
                massaccuracylist=searchoutput[searchoutput["Cond_Rep"]==run]["FG.CalibratedMassAccuracy (PPM)"].dropna().tolist()
                massaccuracy.append(massaccuracylist)

            massaccuracy_df["Mass Accuracy"]=massaccuracy

            violincolors=replicatecolors()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.massaccuracy_violin_hist()=="violin":
                fig,ax=plt.subplots(figsize=(input.massaccuracy_width()*(1/plt.rcParams['figure.dpi']),input.massaccuracy_height()*(1/plt.rcParams['figure.dpi'])))
                medianlineprops=dict(linestyle="--",color="black")
                flierprops=dict(markersize=3)
                x=np.arange(len(massaccuracy_df["Cond_Rep"].tolist()))
                plot=ax.violinplot(massaccuracy_df["Mass Accuracy"],showextrema=False)
                ax.boxplot(massaccuracy_df["Mass Accuracy"],medianprops=medianlineprops,flierprops=flierprops)
                ax.set_ylabel("Mass Accuracy (ppm)",fontsize=axisfont)
                ax.set_xlabel("Run",fontsize=axisfont)
                ax.set_xticks(x+1,labels=massaccuracy_df["Cond_Rep"].tolist(),rotation=x_label_rotation)
                ax.grid(linestyle="--")
                ax.set_axisbelow(True)
                ax.tick_params(axis="both",labelsize=axisfont_labels)

                if numconditions==1:
                    for z in plot["bodies"]:
                        z.set_facecolor(violincolors)
                        z.set_edgecolor("black")
                        z.set_alpha(0.7)
                else:
                    for z,color in zip(plot["bodies"],violincolors):
                        z.set_facecolor(color)
                        z.set_edgecolor("black")
                        z.set_alpha(0.7)

            if input.massaccuracy_violin_hist()=="histogram":
                if numsamples==1:
                    fig,ax=plt.subplots(figsize=(input.massaccuracy_width()*(1/plt.rcParams['figure.dpi']),input.massaccuracy_height()*(1/plt.rcParams['figure.dpi'])))
                    ax.hist(massaccuracy_df["Mass Accuracy"],bins=input.massaccuracy_hist_bins(),color=violincolors)
                    ax.set_xlabel("Mass Accuracy (ppm)",fontsize=axisfont)
                    ax.set_ylabel("Frequency",fontsize=axisfont)
                    ax.set_title(run,fontsize=titlefont)
                    ax.grid(linestyle="--")
                    ax.set_axisbelow(True)
                    ax.tick_params(axis="both",labelsize=axisfont_labels)
                else:
                    fig,ax=plt.subplots(ncols=len(massaccuracy_df["Cond_Rep"]),figsize=(input.massaccuracy_width()*(1/plt.rcParams['figure.dpi']),input.massaccuracy_height()*(1/plt.rcParams['figure.dpi'])))
                    for i,run in enumerate(massaccuracy_df["Cond_Rep"]):
                        if numconditions==1:
                            ax[i].hist(massaccuracy_df["Mass Accuracy"][i],bins=input.massaccuracy_hist_bins(),color=violincolors)
                        else:
                            ax[i].hist(massaccuracy_df["Mass Accuracy"][i],bins=input.massaccuracy_hist_bins(),color=violincolors[i])
                        ax[i].set_xlabel("Mass Accuracy (ppm)",fontsize=axisfont)
                        ax[i].set_title(run,fontsize=titlefont)
                        ax[i].grid(linestyle="--")
                        ax[i].set_axisbelow(True)
                        ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                    ax[0].set_ylabel("Frequency",fontsize=axisfont)
            fig.set_tight_layout(True)
            imagedownload("massaccuracy")

    # ====================================== Data Completeness
    #render ui call for dropdown calling sample condition names
    @render.ui
    def datacompleteness_sampleconditions_ui():
        if input.datacompleteness_sampleconditions_switch()==True:
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=sampleconditions
            return ui.input_selectize("datacompleteness_sampleconditions_pick","Pick sample condition",choices=opts)
    #plot data completeness
    @reactive.effect
    def _():
        @render.plot(width=input.datacompleteness_width(),height=input.datacompleteness_height())
        def datacompletenessplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()
            labelpadding=1

            color1="tab:blue"
            color2="black"

            if input.datacompleteness_sampleconditions_switch()==True:
                plottingdf=searchoutput[searchoutput["R.Condition"]==input.datacompleteness_sampleconditions_pick()]
                if input.protein_peptide()=="proteins":
                    proteincounts=[len(list(group)) for key, group in groupby(sorted(plottingdf[["R.Condition","R.Replicate","PG.ProteinGroups"]].drop_duplicates().drop(columns=["R.Condition","R.Replicate"]).reset_index(drop=True).groupby(["PG.ProteinGroups"]).size().tolist()))]

                elif input.protein_peptide()=="peptides":
                    proteincounts=[len(list(group)) for key, group in groupby(sorted(plottingdf[["R.Condition","R.Replicate","EG.ModifiedPeptide"]].drop_duplicates().drop(columns=["R.Condition","R.Replicate"]).reset_index(drop=True).groupby(["EG.ModifiedPeptide"]).size().tolist()))]

                elif input.protein_peptide()=="strippedpeptides":
                    proteincounts=[len(list(group)) for key, group in groupby(sorted(plottingdf[["R.Condition","R.Replicate","PEP.StrippedSequence"]].drop_duplicates().drop(columns=["R.Condition","R.Replicate"]).reset_index(drop=True).groupby(["PEP.StrippedSequence"]).size().tolist()))]

            else:
                if input.protein_peptide()=="proteins":
                    proteincounts=[len(list(group)) for key, group in groupby(sorted(searchoutput[["R.Condition","R.Replicate","PG.ProteinGroups"]].drop_duplicates().drop(columns=["R.Condition","R.Replicate"]).reset_index(drop=True).groupby(["PG.ProteinGroups"]).size().tolist()))]

                elif input.protein_peptide()=="peptides":
                    proteincounts=[len(list(group)) for key, group in groupby(sorted(searchoutput[["R.Condition","R.Replicate","EG.ModifiedPeptide"]].drop_duplicates().drop(columns=["R.Condition","R.Replicate"]).reset_index(drop=True).groupby(["EG.ModifiedPeptide"]).size().tolist()))]

                elif input.protein_peptide()=="strippedpeptides":
                    proteincounts=[len(list(group)) for key, group in groupby(sorted(searchoutput[["R.Condition","R.Replicate","PEP.StrippedSequence"]].drop_duplicates().drop(columns=["R.Condition","R.Replicate"]).reset_index(drop=True).groupby(["PEP.StrippedSequence"]).size().tolist()))]

            totalproteins=sum(proteincounts)

            proteinfrequencies=[]
            for i in proteincounts:
                proteinfrequencies.append(i/totalproteins*100)

            xaxis=np.arange(1,len(proteinfrequencies)+1,1).tolist()

            y1=proteincounts
            y2=proteinfrequencies

            fig,ax1=plt.subplots(figsize=(input.datacompleteness_width()*(1/plt.rcParams['figure.dpi']),input.datacompleteness_height()*(1/plt.rcParams['figure.dpi'])))

            ax2=ax1.twinx()
            ax1.bar(xaxis,y1,edgecolor="k")
            ax2.plot(xaxis,y2,"-o",color=color2)

            ax1.set_xlabel('Observed in X Runs',fontsize=axisfont)
            if input.protein_peptide()=="proteins":
                ax1.set_ylabel('# Proteins',color=color1,fontsize=axisfont)
            elif input.protein_peptide()=="peptides":
                ax1.set_ylabel('# Peptides',color=color1,fontsize=axisfont)
            ax2.set_ylabel('% of IDs',color=color2,fontsize=axisfont)
            ax1.tick_params(axis="y",colors=color1)
            ax2.tick_params(axis="y",colors=color2)
            ax1.tick_params(axis="both",labelsize=axisfont_labels)
            ax2.tick_params(axis="both",labelsize=axisfont_labels)

            ax1.bar_label(ax1.containers[0],label_type="edge",padding=35,color=color1,fontsize=labelfont)
            ax1.set_ylim(top=max(proteincounts)+y_padding*max(proteincounts))
            ax2.set_ylim(top=max(proteinfrequencies)+y_padding*max(proteinfrequencies))

            for x,y in enumerate(proteinfrequencies):
                ax2.text(xaxis[x],proteinfrequencies[x]+labelpadding,str(round(y,1))+"%",horizontalalignment="center",verticalalignment="bottom",color=color2,fontsize=labelfont)

            ax1.set_axisbelow(True)
            ax1.grid(linestyle="--")
            plt.xticks(range(1,len(xaxis)+1))
            plt.xlim(0.5,len(xaxis)+1)
            fig.set_tight_layout(True)
            imagedownload("datacompleteness")

    # ====================================== Peak Widths
    #plot peak widths
    @reactive.effect
    def _():
        @render.plot(width=input.peakwidth_width(),height=input.peakwidth_height())
        def peakwidthplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            violincolors=replicatecolors()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            fwhm_df=peakwidths()

            x=np.arange(len(fwhm_df["Cond_Rep"].tolist()))
            fig,ax=plt.subplots(figsize=(input.peakwidth_width()*(1/plt.rcParams['figure.dpi']),input.peakwidth_height()*(1/plt.rcParams['figure.dpi'])))
            medianlineprops=dict(linestyle="--",color="black")
            flierprops=dict(markersize=3)
            if input.peakwidth_removetop5percent()==False:
                plot=ax.violinplot(fwhm_df["Peak Width"],showextrema=False)
                bplot=ax.boxplot(fwhm_df["Peak Width"],medianprops=medianlineprops,flierprops=flierprops)
            elif input.peakwidth_removetop5percent()==True:
                plot=ax.violinplot(fwhm_df["Peak Width 95%"],showextrema=False)
                bplot=ax.boxplot(fwhm_df["Peak Width 95%"],medianprops=medianlineprops,flierprops=flierprops)

            ax.set_ylabel("Peak Width (s)",fontsize=axisfont)
            ax.set_xlabel("Run",fontsize=axisfont)
            ax.set_xticks(x+1,labels=fwhm_df["Cond_Rep"].tolist(),rotation=input.xaxis_label_rotation())
            ax.grid(linestyle="--")
            ax.set_axisbelow(True)
            ax.tick_params(axis="both",labelsize=axisfont_labels)

            if numconditions==1:
                for z in plot["bodies"]:
                    z.set_facecolor(violincolors)
                    z.set_edgecolor("black")
                    z.set_alpha(0.7)
            else:
                for z,color in zip(plot["bodies"],violincolors):
                    z.set_facecolor(color)
                    z.set_edgecolor("black")
                    z.set_alpha(0.7)
            fig.set_tight_layout(True)
            imagedownload("peakwidths")
    #show a table of mean/median peak widths per run 
    @render.table
    def peakwidth_table():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            fwhm_df=peakwidths()
            peakwidthtable=pd.DataFrame()

            meanlist=[]
            medianlist=[]
            meanlist95=[]
            medianlist95=[]
            for run in fwhm_df["Cond_Rep"]:
                meanlist.append(np.mean(fwhm_df[fwhm_df["Cond_Rep"]==run]["Peak Width"].tolist()))
                meanlist95.append(np.mean(fwhm_df[fwhm_df["Cond_Rep"]==run]["Peak Width 95%"].tolist()))
                medianlist.append(np.median(fwhm_df[fwhm_df["Cond_Rep"]==run]["Peak Width"].tolist()))
                medianlist95.append(np.median(fwhm_df[fwhm_df["Cond_Rep"]==run]["Peak Width 95%"].tolist()))

            peakwidthtable["Run"]=fwhm_df["Cond_Rep"]
            peakwidthtable["Mean Peak Width (s)"]=meanlist
            peakwidthtable["Median Peak Width (s)"]=medianlist
            peakwidthtable["Mean Peak Width 95% (s)"]=meanlist95
            peakwidthtable["Median Peak Width 95% (s)"]=medianlist95

            if input.peakwidth_removetop5percent()==True:
                return peakwidthtable[["Run","Mean Peak Width 95% (s)","Median Peak Width 95% (s)"]]
            elif input.peakwidth_removetop5percent()==False:
                return peakwidthtable[["Run","Mean Peak Width (s)","Median Peak Width (s)"]]

    # ====================================== Missed Cleavages
    @reactive.effect
    def _():
        @render.plot(width=input.missedcleavages_width(),height=input.missedcleavages_height())
        def missedcleavages_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            enzyme_rules=input.enzyme_rules()
            missedcleavages=[]
            for pep in searchoutput["PEP.StrippedSequence"]:
                #check if tryptic
                if enzyme_rules=="trypsin":
                    if pep[-1]=="K" or pep[-1]=="R":
                        missedcleavages.append(pep[:-1].count("K")+pep[:-1].count("R"))
                    #-1 if peptide is nontryptic
                    else:
                        missedcleavages.append(-1)
            searchoutput["Missed Cleavages"]=missedcleavages

            missedcleavages_df=pd.DataFrame()
            for run in runnames:
                df=pd.DataFrame(searchoutput[searchoutput["Cond_Rep"]==run][["PEP.StrippedSequence","Missed Cleavages"]].drop_duplicates().reset_index(drop=True)["Missed Cleavages"].value_counts()).sort_index().transpose()
                missedcleavages_df=pd.concat([missedcleavages_df,df],axis=0)
            missedcleavages_df=missedcleavages_df.reset_index(drop=True)
            missedcleavages_df["Cond_Rep"]=runnames
            missedcleavages_df=missedcleavages_df.set_index("Cond_Rep")

            maxvalue=max(missedcleavages_df.select_dtypes(include=[np.number]).max().tolist())
            x=np.arange(len(missedcleavages_df.index.tolist()))

            width=input.missedcleavages_barwidth()
            y_padding=input.ypadding()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            x_label_rotation=input.xaxis_label_rotation()

            fig,ax=plt.subplots(figsize=(input.missedcleavages_width()*(1/plt.rcParams['figure.dpi']),input.missedcleavages_height()*(1/plt.rcParams['figure.dpi'])))
            legend_patches=[]
            for i,column in enumerate(missedcleavages_df.columns):
                ax.bar(x+(i*width),missedcleavages_df[column],width=width,edgecolor="k")
                ax.bar_label(ax.containers[i],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                if column==-1:
                    label="Non-Tryptic"
                else:
                    label=column
                legend_patches.append(mpatches.Patch(color=list(mcolors.TABLEAU_COLORS.keys())[i],label=label))
            ax.set_ylim(top=(maxvalue+(y_padding*maxvalue)),bottom=-(0.025*maxvalue))
            ax.set_xticks(x+width,missedcleavages_df.index.tolist(),rotation=x_label_rotation)
            ax.legend(handles=legend_patches,loc='center left', bbox_to_anchor=(1,0.5),prop={'size':legendfont})
            ax.set_title("Missed Cleavages",fontsize=titlefont)
            ax.set_ylabel("Counts",fontsize=axisfont)
            ax.set_xlabel("Condition",fontsize=axisfont)
            ax.grid(linestyle="--")
            ax.set_axisbelow(True)
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("missedcleavages")

#endregion

# ============================================================================= PTMs 
#region
    #store ptmidmetrics call in a reactive.calc
    @reactive.calc
    def ptmidmetrics_calc():
        ptmresultdf,ptmaveragedf=ptmcounts(metadata_update(),input.foundptms())
        return ptmresultdf,ptmaveragedf
    @reactive.calc
    def ptmcvs_calc():
        cvcalc_df=cvmetrics(metadata_update(),input.foundptms())
        return cvcalc_df
    
    #function for finding the PTMs in the data
    @reactive.calc
    def find_ptms(): 
        searchoutput=metadata_update()
        peplist=searchoutput["EG.ModifiedPeptide"]
        ptmlist=[]
        for i in peplist:
            ptmlist.append(re.findall(r"[^[]*\[([^]]*)\]",i))
        searchoutput["PTMs"]=ptmlist

        #convert ProteinPTMLocations column to lists instead of strings
        #split ProteinPTMLocations column by ; delimiter, only use first protein that is ID'd
        if "EG.ProteinPTMLocations" in searchoutput.columns:
            if math.isnan(searchoutput["EG.ProteinPTMLocations"].drop_duplicates()[0])==True:
                pass
            else:
                searchoutput["EG.ProteinPTMLocations"]=searchoutput["EG.ProteinPTMLocations"].str.split(";").str[0]
                ptmproteinlocations=[]
                for ele in searchoutput["EG.ProteinPTMLocations"]:
                    if type(ele)==float:
                        ptmproteinlocations.append([])
                    else:
                        ptmproteinlocations.append(re.split("[( )]",ele)[1].split(","))
                searchoutput["PTM Protein Locations"]=ptmproteinlocations
        return(list(OrderedDict.fromkeys(itertools.chain(*ptmlist))))
    #text showing what PTMs were detected
    @render.text
    def ptmlist_text():
        listofptms=find_ptms()
        return listofptms
    #generate list to pull from to pick PTMs
    @render.ui
    def ptmlist_ui():
        listofptms=find_ptms()
        return ui.input_selectize("foundptms","Pick PTM to plot data for (Counts per Condition, CVs, and Mass Accuracy):",choices=listofptms,selected=listofptms[0])

    # ====================================== Counts per Condition
    #plot PTM ID metrics
    @reactive.effect
    def _():
        @render.plot(width=input.ptmidmetrics_width(),height=input.ptmidmetrics_height())
        def ptmidmetricsplot():
            plotinput=input.ptmidplotinput()
            ptmresultdf,ptmaveragedf=ptmidmetrics_calc()
            ptm=input.foundptms()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.ptm_counts_vs_enrich()=="counts":
                y1="proteins"
                y2="proteins2pepts"
                y3="peptides"
                y4="precursors"
                titlemod="ID Counts for PTM: "
                ylabel="Counts"
            if input.ptm_counts_vs_enrich()=="percent":
                y1="proteins_enrich%"
                y2="proteins2pepts_enrich%"
                y3="peptides_enrich%"
                y4="precursors_enrich%"
                titlemod="% of IDs for PTM: "
                ylabel="% of IDs"

            if plotinput=="all":
                fig,ax=plt.subplots(nrows=2,ncols=2,sharex=True,figsize=(input.ptmidmetrics_width()*(1/plt.rcParams['figure.dpi']),input.ptmidmetrics_height()*(1/plt.rcParams['figure.dpi'])))
                ax1=ax[0,0]
                ax2=ax[0,1]
                ax3=ax[1,0]
                ax4=ax[1,1]

                if input.ptmidmetrics_individual_average()=="individual":
                    idmetricscolor=replicatecolors()
                    titlemod1=""
                    ax1.bar(ptmresultdf["Cond_Rep"],ptmresultdf[y1],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax2.bar(ptmresultdf["Cond_Rep"],ptmresultdf[y2],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax3.bar(ptmresultdf["Cond_Rep"],ptmresultdf[y3],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax4.bar(ptmresultdf["Cond_Rep"],ptmresultdf[y4],width=0.8,color=idmetricscolor,edgecolor="k")

                    ax1.bar_label(ax1.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax2.bar_label(ax2.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax3.bar_label(ax3.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax4.bar_label(ax4.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)

                    ax1.set_ylim(top=max(ptmresultdf[y1].tolist())+y_padding*max(ptmresultdf[y1].tolist()))
                    ax2.set_ylim(top=max(ptmresultdf[y2].tolist())+y_padding*max(ptmresultdf[y2].tolist()))
                    ax3.set_ylim(top=max(ptmresultdf[y3].tolist())+(y_padding+0.1)*max(ptmresultdf[y3].tolist()))
                    ax4.set_ylim(top=max(ptmresultdf[y4].tolist())+(y_padding+0.1)*max(ptmresultdf[y4].tolist()))
                elif input.ptmidmetrics_individual_average()=="average":
                    idmetricscolor=colorpicker()
                    titlemod1="Average "
                    bars1=ax1.bar(ptmaveragedf["R.Condition"],ptmaveragedf[y1+"_avg"],yerr=ptmaveragedf[y1+"_stdev"],width=0.8,capsize=10,color=idmetricscolor,edgecolor="k")
                    bars2=ax2.bar(ptmaveragedf["R.Condition"],ptmaveragedf[y2+"_avg"],yerr=ptmaveragedf[y2+"_stdev"],width=0.8,capsize=10,color=idmetricscolor,edgecolor="k")
                    bars3=ax3.bar(ptmaveragedf["R.Condition"],ptmaveragedf[y3+"_avg"],yerr=ptmaveragedf[y3+"_stdev"],width=0.8,capsize=10,color=idmetricscolor,edgecolor="k")
                    bars4=ax4.bar(ptmaveragedf["R.Condition"],ptmaveragedf[y4+"_avg"],yerr=ptmaveragedf[y4+"_stdev"],width=0.8,capsize=10,color=idmetricscolor,edgecolor="k")

                    ax1.bar_label(bars1,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                    ax2.bar_label(bars2,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                    ax3.bar_label(bars3,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                    ax4.bar_label(bars4,label_type="edge",rotation=90,padding=10,fontsize=labelfont)

                    ax1.set_ylim(top=max(ptmaveragedf[y1+"_avg"].tolist())+y_padding*max(ptmaveragedf[y1+"_avg"].tolist()))
                    ax2.set_ylim(top=max(ptmaveragedf[y2+"_avg"].tolist())+y_padding*max(ptmaveragedf[y2+"_avg"].tolist()))
                    ax3.set_ylim(top=max(ptmaveragedf[y3+"_avg"].tolist())+(y_padding+0.1)*max(ptmaveragedf[y3+"_avg"].tolist()))
                    ax4.set_ylim(top=max(ptmaveragedf[y4+"_avg"].tolist())+(y_padding+0.1)*max(ptmaveragedf[y4+"_avg"].tolist()))

                fig.text(0, 0.6,ylabel,ha="left",va="center",rotation="vertical",fontsize=axisfont)
                plt.suptitle(titlemod1+titlemod+ptm,y=1,fontsize=titlefont)
                ax1.set_title("Protein Groups",fontsize=titlefont)
                ax2.set_title("Protein Groups with >2 Peptides",fontsize=titlefont)
                ax3.set_title("Peptides",fontsize=titlefont)
                ax3.set_xlabel("Condition",fontsize=axisfont)
                ax3.set_ylabel("  ",fontsize=axisfont)
                ax3.tick_params(axis="x",rotation=x_label_rotation)
                ax4.set_title("Precursors",fontsize=titlefont)
                ax4.set_xlabel("Condition",fontsize=axisfont)
                ax4.tick_params(axis="x",rotation=x_label_rotation)

                ax1.set_axisbelow(True)
                ax1.grid(linestyle="--")
                ax1.tick_params(axis="both",labelsize=axisfont_labels)
                ax2.set_axisbelow(True)
                ax2.grid(linestyle="--")
                ax2.tick_params(axis="both",labelsize=axisfont_labels)
                ax3.set_axisbelow(True)
                ax3.grid(linestyle="--")
                ax3.tick_params(axis="both",labelsize=axisfont_labels)
                ax4.set_axisbelow(True)
                ax4.grid(linestyle="--")
                ax4.tick_params(axis="both",labelsize=axisfont_labels)

            else:
                plotinput=input.ptmidplotinput()
                if plotinput=="proteins":
                    titleprop="Protein Groups"
                if plotinput=="proteins2pepts":
                    titleprop="Protein Groups with >2 Peptides"
                if plotinput=="peptides":
                    titleprop="Peptides"
                if plotinput=="precursors":
                    titleprop="Precursors"

                if input.ptm_counts_vs_enrich()=="counts":
                    plot_y=plotinput
                    titlemod="ID Counts for PTM: "
                    y_mod="Counts"
                if input.ptm_counts_vs_enrich()=="percent":
                    plot_y=plotinput+"_enrich%"
                    titlemod="% of IDs for PTM: "
                    y_mod="% of IDs"

                fig,ax=plt.subplots(figsize=(input.ptmidmetrics_width()*(1/plt.rcParams['figure.dpi']),input.ptmidmetrics_height()*(1/plt.rcParams['figure.dpi'])))
                if input.ptmidmetrics_individual_average()=="individual":
                    idmetricscolor=replicatecolors()
                    ax.bar(ptmresultdf["Cond_Rep"],ptmresultdf[plot_y],width=0.8,color=idmetricscolor,edgecolor="k")
                    ax.bar_label(ax.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                    ax.set_ylim(top=max(ptmresultdf[plot_y].tolist())+y_padding*max(ptmresultdf[plot_y].tolist()))
                    titlemod1=""
                    
                elif input.ptmidmetrics_individual_average()=="average":
                    idmetricscolor=colorpicker()
                    bars=ax.bar(ptmaveragedf["R.Condition"],ptmaveragedf[plot_y+"_avg"],yerr=ptmaveragedf[plot_y+"_stdev"],width=0.8,capsize=10,color=idmetricscolor,edgecolor="k")
                    ax.bar_label(bars,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                    ax.set_ylim(top=max(ptmaveragedf[plot_y+"_avg"].tolist())+y_padding*max(ptmaveragedf[plot_y+"_avg"].tolist()))
                    titlemod1="Average "
                
                ax.set_ylabel(y_mod,fontsize=axisfont)
                ax.set_xlabel("Condition",fontsize=axisfont)
                ax.set_title(titlemod1+titleprop+" "+titlemod+ptm,fontsize=titlefont)
                ax.tick_params(axis="x",rotation=x_label_rotation)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("ptmidcounts_"+ptm)

    # ====================================== CV Plots
    @render.ui
    def ptm_cvplot_histogram_bins_ui():
        if input.ptm_cvplot_histogram_bins_switch()==True:
            return ui.input_slider("ptm_cvplot_histogram_bins_slider","Number of bins:",min=10,max=250,value=100,step=10,ticks=True)
    #plot PTM CV violin plots
    @reactive.effect
    def _():
        @render.plot(width=input.ptmcvplot_width(),height=input.ptmcvplot_height())
        def ptm_cvplot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            ptm=input.foundptms()
            cvcalc_df=ptmcvs_calc()

            colors=colorpicker()
            cvplotinput=input.ptm_proteins_precursors()
            cutoff95=input.ptm_removetop5percent()
            
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.ptm_cvplot_histogram_bins_switch()==True:
                numbins=input.ptm_cvplot_histogram_bins_slider()
                fig,ax=plt.subplots(ncols=len(cvcalc_df),figsize=(input.ptmcvplot_width()*(1/plt.rcParams['figure.dpi']),input.ptmcvplot_height()*(1/plt.rcParams['figure.dpi'])))
                colors=colorpicker()
                if numconditions==1:
                    if cutoff95==True:
                        ax.hist(cvcalc_df[cvplotinput+" 95% CVs"],bins=numbins,color=colors)
                    elif cutoff95==False:
                        ax.hist(cvcalc_df[cvplotinput+" CVs"],bins=numbins,color=colors)
                    ax.set_xlabel(cvplotinput+" % CV",fontsize=axisfont)
                    ax.axvline(x=20,color="black",linestyle="--")
                    ax.set_ylabel("Counts",fontsize=axisfont)
                    ax.set_title(cvcalc_df["R.Condition"].tolist()[0],fontsize=titlefont)
                    ax.grid(linestyle="--")
                    ax.set_axisbelow(True)
                else:
                    for i in range(len(cvcalc_df)):
                        if cutoff95==True:
                            ax[i].hist(cvcalc_df[cvplotinput+" 95% CVs"][i],bins=numbins,color=colors[i])
                        elif cutoff95==False:
                            ax[i].hist(cvcalc_df[cvplotinput+" CVs"][i],bins=numbins,color=colors[i])
                        ax[i].set_xlabel(cvplotinput+" % CV",fontsize=axisfont)
                        ax[i].set_title(cvcalc_df["R.Condition"].tolist()[i],fontsize=titlefont)
                        ax[i].axvline(x=20,color="black",linestyle="--")
                        ax[i].grid(linestyle="--")
                        ax[i].set_axisbelow(True)
                    ax[0].set_ylabel("Counts",fontsize=axisfont)
            else:
                fig,ax=plt.subplots(figsize=(input.ptmcvplot_width()*(1/plt.rcParams['figure.dpi']),input.ptmcvplot_height()*(1/plt.rcParams['figure.dpi'])))
                x=np.arange(len(cvcalc_df["R.Condition"]))
                lineprops=dict(linestyle="--",color="black")
                flierprops=dict(markersize=3)
                if cutoff95==True:
                    bplot=ax.boxplot(cvcalc_df[cvplotinput+" 95% CVs"],medianprops=lineprops,flierprops=flierprops)
                    plot=ax.violinplot(cvcalc_df[cvplotinput+" 95% CVs"],showextrema=False)
                    ax.set_title(cvplotinput+" CVs, 95% Cutoff",fontsize=titlefont)
                elif cutoff95==False:
                    bplot=ax.boxplot(cvcalc_df[cvplotinput+" CVs"],medianprops=lineprops,flierprops=flierprops)
                    plot=ax.violinplot(cvcalc_df[cvplotinput+" CVs"],showextrema=False)
                    ax.set_title(cvplotinput+" CVs",fontsize=titlefont)
                ax.set_xticks(x+1,labels=cvcalc_df["R.Condition"],rotation=x_label_rotation)
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                ax.set_ylabel("CV%",fontsize=axisfont)
                ax.set_xlabel("Condition",fontsize=axisfont)
                ax.grid(linestyle="--")
                ax.set_axisbelow(True)
                ax.axhline(y=20,color="black",linestyle="--")

                if numconditions==1:
                    for z in plot["bodies"]:
                        z.set_facecolor(colors)
                        z.set_edgecolor("black")
                        z.set_alpha(0.7)
                else:
                    for z,color in zip(plot["bodies"],colors):
                        z.set_facecolor(color)
                        z.set_edgecolor("black")
                        z.set_alpha(0.7)
            fig.set_tight_layout(True)
            imagedownload("ptmcvplot_"+ptm)
    #show a table of mean/median CV values per condition for selected PTM
    @render.table
    def ptm_cvtable():
        ptmcvs=ptmcvs_calc()

        ptmcvs_summary=pd.DataFrame()
        proteinptmcv_mean=[]
        proteinptmcv95_mean=[]
        precursorptmcv_mean=[]
        precursorptmcv95_mean=[]
        peptideptmcv_mean=[]
        peptideptmcv95_mean=[]

        proteinptmcv_median=[]
        proteinptmcv95_median=[]
        precursorptmcv_median=[]
        precursorptmcv95_median=[]
        peptideptmcv_median=[]
        peptideptmcv95_median=[]
        for i,run in enumerate(ptmcvs["R.Condition"].tolist()):
            proteinptmcv_mean.append(np.mean(ptmcvs[ptmcvs["R.Condition"]==run]["Protein CVs"][i]))
            proteinptmcv95_mean.append(np.mean(ptmcvs[ptmcvs["R.Condition"]==run]["Protein 95% CVs"][i]))
            
            precursorptmcv_mean.append(np.mean(ptmcvs[ptmcvs["R.Condition"]==run]["Precursor CVs"][i]))
            precursorptmcv95_mean.append(np.mean(ptmcvs[ptmcvs["R.Condition"]==run]["Precursor 95% CVs"][i]))

            peptideptmcv_mean.append(np.mean(ptmcvs[ptmcvs["R.Condition"]==run]["Peptide CVs"][i]))
            peptideptmcv95_mean.append(np.mean(ptmcvs[ptmcvs["R.Condition"]==run]["Peptide 95% CVs"][i]))
            
            proteinptmcv_median.append(np.median(ptmcvs[ptmcvs["R.Condition"]==run]["Protein CVs"][i]))
            proteinptmcv95_median.append(np.median(ptmcvs[ptmcvs["R.Condition"]==run]["Protein 95% CVs"][i]))
            
            precursorptmcv_median.append(np.median(ptmcvs[ptmcvs["R.Condition"]==run]["Precursor CVs"][i]))
            precursorptmcv95_median.append((np.median(ptmcvs[ptmcvs["R.Condition"]==run]["Precursor 95% CVs"][i])))
            
            peptideptmcv_median.append(np.median(ptmcvs[ptmcvs["R.Condition"]==run]["Peptide CVs"][i]))
            peptideptmcv95_median.append((np.median(ptmcvs[ptmcvs["R.Condition"]==run]["Peptide 95% CVs"][i])))
        
        ptmcvs_summary["R.Condition"]=ptmcvs["R.Condition"]
        ptmcvs_summary["Protein Mean CVs"]=proteinptmcv_mean
        ptmcvs_summary["Protein Median CVs"]=proteinptmcv_median
        ptmcvs_summary["Protein Mean CVs 95%"]=proteinptmcv95_mean
        ptmcvs_summary["Protein Median CVs 95%"]=proteinptmcv95_median

        ptmcvs_summary["Precursor Mean CVs"]=precursorptmcv_mean
        ptmcvs_summary["Precursor Median CVs"]=precursorptmcv_median
        ptmcvs_summary["Precursor Mean CVs 95%"]=precursorptmcv95_mean
        ptmcvs_summary["Precursor Median CVs 95%"]=precursorptmcv95_median

        ptmcvs_summary["Peptide Mean CVs"]=peptideptmcv_mean
        ptmcvs_summary["Peptide Median CVs"]=peptideptmcv_median
        ptmcvs_summary["Peptide Mean CVs 95%"]=peptideptmcv95_mean
        ptmcvs_summary["Peptide Median CVs 95%"]=peptideptmcv95_median

        cvplotinput=input.ptm_proteins_precursors()
        cutoff95=input.ptm_removetop5percent()
        if cvplotinput=="Protein":
            if cutoff95==True:
                return ptmcvs_summary[["R.Condition","Protein Mean CVs 95%","Protein Median CVs 95%"]]
            if cutoff95==False:
                return ptmcvs_summary[["R.Condition","Protein Mean CVs","Protein Median CVs"]]
        if cvplotinput=="Precursor":
            if cutoff95==True:
                return ptmcvs_summary[["R.Condition","Precursor Mean CVs 95%","Precursor Median CVs 95%"]]
            if cutoff95==False:
                return ptmcvs_summary[["R.Condition","Precursor Mean CVs","Precursor Median CVs"]]
        if cvplotinput=="Peptide":
            if cutoff95==True:
                return ptmcvs_summary[["R.Condition","Peptide Mean CVs 95%","Peptide Median CVs 95%"]]
            if cutoff95==False:
                return ptmcvs_summary[["R.Condition","Peptide Mean CVs","Peptide Median CVs"]]

    # ====================================== PTMs per Precursor
    #generate list to pull from to pick PTMs
    @render.ui
    def ptmsperprecursor_ptmlist_ui():
        if input.ptmsperprecursor_specific()==True:
            listofptms=find_ptms()
            return ui.input_selectize("ptmsperprecursor_foundptms","Pick PTM to plot data for:",choices=listofptms,selected=listofptms[0])
    #plot PTMs per precursor
    @reactive.effect
    def _():
        @render.plot(width=input.ptmsperprecursor_width(),height=input.ptmsperprecursor_height())
        def ptmsperprecursor():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            colors=colorpicker()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            width=input.barwidth()

            fig,ax=plt.subplots(figsize=(input.ptmsperprecursor_width()*(1/plt.rcParams['figure.dpi']),input.ptmsperprecursor_height()*(1/plt.rcParams['figure.dpi'])))
            ptmdf=pd.DataFrame()

            for j,condition in enumerate(sampleconditions):
                df=searchoutput[searchoutput["R.Condition"]==condition][["EG.ModifiedPeptide","FG.Charge"]].drop_duplicates().reset_index(drop=True)
                dfptmlist=[]
                numptms=[]
                for i in df["EG.ModifiedPeptide"]:
                    if input.ptmsperprecursor_specific()==True:
                        ptm=input.ptmsperprecursor_foundptms().split(" ")[0]
                        foundptms=re.findall(ptm,i)
                        titlemod=ptm+" "
                    else:
                        foundptms=re.findall(r"[^[]*\[([^]]*)\]",i)
                        titlemod=""
                    dfptmlist.append(foundptms)
                    numptms.append(len(foundptms))
                dfptmlist=pd.Series(dfptmlist).value_counts().to_frame().reset_index().rename(columns={"index":condition,"count":condition+"_count"})
                ptmdf=pd.concat([ptmdf,dfptmlist],axis=1)
                
                x=np.arange(0,max(numptms)+1,1)
                frequencies=pd.Series(numptms).value_counts().sort_index()
                if numconditions==1:
                    ax.bar(x+(j*width),frequencies,width=width,color=colors,edgecolor="black")
                else:
                    ax.bar(x+(j*width),frequencies,width=width,color=colors[j],edgecolor="black")
                ax.bar_label(ax.containers[j],label_type="edge",rotation=90,padding=5,fontsize=labelfont)

            ax.legend(sampleconditions,loc="upper right",fontsize=legendfont)
            ax.set_ylim(bottom=-ax.get_ylim()[1]+0.9*ax.get_ylim()[1],top=ax.get_ylim()[1]+y_padding*ax.get_ylim()[1])
            ax.set_xticks(x+((numconditions-1)/2)*width,x)
            ax.set_ylabel("Counts",fontsize=axisfont)
            ax.set_xlabel("# of PTMs",fontsize=axisfont)
            ax.set_title("# of "+titlemod+"PTMs per Precursor",fontsize=titlefont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("ptmsperprecursor")

    # ====================================== Mass Accuracy
    @render.ui
    def ptm_massaccuracy_bins_ui():
        if input.ptm_massaccuracy_violin_hist()=="histogram":
            return ui.input_slider("ptm_massaccuracy_hist_bins","Number of bins",min=10,max=200,value=100,step=10,ticks=True,width="300px")
    #plot mass accuracy
    @reactive.effect
    def _():
        @render.plot(width=input.ptm_massaccuracy_width(),height=input.ptm_massaccuracy_height())
        def ptm_massaccuracy_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            ptm=input.foundptms()

            massaccuracy_df=pd.DataFrame()
            massaccuracy_df["Cond_Rep"]=runnames
            massaccuracy=[]
            for run in runnames:
                massaccuracylist=searchoutput[(searchoutput["Cond_Rep"]==run)&(searchoutput["EG.ModifiedPeptide"].str.contains(ptm,regex=False))]["FG.CalibratedMassAccuracy (PPM)"].dropna().tolist()
                massaccuracy.append(massaccuracylist)

            massaccuracy_df["Mass Accuracy"]=massaccuracy

            violincolors=replicatecolors()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.ptm_massaccuracy_violin_hist()=="violin":
                fig,ax=plt.subplots(figsize=(input.ptm_massaccuracy_width()*(1/plt.rcParams['figure.dpi']),input.ptm_massaccuracy_height()*(1/plt.rcParams['figure.dpi'])))
                medianlineprops=dict(linestyle="--",color="black")
                flierprops=dict(markersize=3)
                x=np.arange(len(massaccuracy_df["Cond_Rep"].tolist()))
                plot=ax.violinplot(massaccuracy_df["Mass Accuracy"],showextrema=False)
                ax.boxplot(massaccuracy_df["Mass Accuracy"],medianprops=medianlineprops,flierprops=flierprops)
                ax.set_ylabel("Mass Accuracy (ppm)",fontsize=axisfont)
                ax.set_xlabel("Run",fontsize=axisfont)
                ax.set_title(ptm+" Precursor Mass Accuracy",fontsize=titlefont)
                ax.set_xticks(x+1,labels=massaccuracy_df["Cond_Rep"].tolist(),rotation=x_label_rotation)
                ax.grid(linestyle="--")
                ax.set_axisbelow(True)
                ax.tick_params(axis="both",labelsize=axisfont_labels)

                if numconditions==1:
                    for z in plot["bodies"]:
                        z.set_facecolor(violincolors)
                        z.set_edgecolor("black")
                        z.set_alpha(0.7)
                else:
                    for z,color in zip(plot["bodies"],violincolors):
                        z.set_facecolor(color)
                        z.set_edgecolor("black")
                        z.set_alpha(0.7)

            if input.ptm_massaccuracy_violin_hist()=="histogram":
                if numsamples==1:
                    fig,ax=plt.subplots(figsize=(input.ptm_massaccuracy_width()*(1/plt.rcParams['figure.dpi']),input.ptm_massaccuracy_height()*(1/plt.rcParams['figure.dpi'])))
                    ax.hist(massaccuracy_df["Mass Accuracy"],bins=input.ptm_massaccuracy_hist_bins(),color=violincolors)
                    ax.set_xlabel("Mass Accuracy (ppm)",fontsize=axisfont)
                    ax.set_ylabel("Frequency",fontsize=axisfont)
                    ax.set_title(ptm+"Precursor Mass Accuracy",fontsize=titlefont)
                    ax.grid(linestyle="--")
                    ax.set_axisbelow(True)
                    ax.tick_params(axis="both",labelsize=axisfont_labels)
                else:
                    fig,ax=plt.subplots(ncols=len(massaccuracy_df["Cond_Rep"]),figsize=(input.ptm_massaccuracy_width()*(1/plt.rcParams['figure.dpi']),input.ptm_massaccuracy_height()*(1/plt.rcParams['figure.dpi'])))
                    for i,run in enumerate(massaccuracy_df["Cond_Rep"]):
                        if numconditions==1:
                            ax[i].hist(massaccuracy_df["Mass Accuracy"][i],bins=input.ptm_massaccuracy_hist_bins(),color=violincolors)
                        else:
                            ax[i].hist(massaccuracy_df["Mass Accuracy"][i],bins=input.ptm_massaccuracy_hist_bins(),color=violincolors[i])
                        ax[i].set_xlabel("Mass Accuracy (ppm)",fontsize=axisfont)
                        ax[i].set_title(run,fontsize=titlefont)
                        ax[i].grid(linestyle="--")
                        ax[i].set_axisbelow(True)
                        ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                    ax[0].set_ylabel("Frequency",fontsize=axisfont)
                    plt.suptitle(ptm+"Precursor Mass Accuracy",fontsize=titlefont)
            fig.set_tight_layout(True)
            imagedownload("ptmmassaccuracy_"+ptm)

    # ====================================== Unique PTM Sites
    @render.data_frame
    def proteinptms_table():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        
        #calculate number of unique modifications found across all runs for each protein
        proteinptmdf=searchoutput[["PG.ProteinNames","PTM Protein Locations"]]
        proteinptmdf=proteinptmdf.explode("PTM Protein Locations").dropna().drop_duplicates().sort_values("PG.ProteinNames").reset_index(drop=True)
        proteinptmdf_grouped=proteinptmdf.groupby(["PG.ProteinNames"]).size().reset_index().rename(columns={0:"Unique PTMs"}).sort_values("Unique PTMs",ascending=False).reset_index(drop=True)

        df_styles={"location":"body",
                   "style":{"column-width":"200px",
                            "overflow":"hidden"
                           }
                  }
        return render.DataGrid(proteinptmdf_grouped,editable=False,selection_mode="row",styles=df_styles)
    @reactive.calc
    def modificationdfs():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        modlist=find_ptms()

        #sorting functions to sort the PTM Protein Locations column correctly
        #https://stackoverflow.com/questions/5967500/how-to-correctly-sort-a-string-with-a-number-inside
        def atoi(text):
            return int(text) if text.isdigit() else text
        def natural_keys(text):
            '''
            alist.sort(key=natural_keys) sorts in human order
            http://nedbatchelder.com/blog/200712/human_sorting.html
            (See Toothy's implementation in the comments)
            '''
            return [atoi(c) for c in re.split(r'(\d+)',text)]

        modificationdf=pd.DataFrame(columns=modlist)
        modificationcountdf=pd.DataFrame(columns=modlist)

        for run in runnames:
            if input.ptmsites_pickprotein()==True:
                if len(proteinptms_table.data_view(selected=True)["PG.ProteinNames"].tolist())==0:
                    df=searchoutput[searchoutput["Cond_Rep"]==run]
                else:
                    selectedprotein=proteinptms_table.data_view(selected=True)["PG.ProteinNames"].tolist()[0]
                    df=searchoutput[(searchoutput["Cond_Rep"]==run)&(searchoutput["PG.ProteinNames"]==selectedprotein)]
            else:
                df=searchoutput[searchoutput["Cond_Rep"]==run]

            runmodificationlist=[]
            runmodcountlist=[]
            
            for mod in modlist:
                #df for modifications from specific reagent/source (e.g. just Carbaidomethylations or Oxidations)
                #values in the resulting dfs are most similar to # of peptides since the dfs are based on the whole searchoutput
                #precursors of the same modification site and different charge would be categorized under the same unique modification site
                df_mod=df[df["PTMs"].str.contains(mod,regex=False)].reset_index(drop=True)
                modifiedresiduelist=[]
                for i in range(len(df_mod)):
                    if len(df_mod["PTM Protein Locations"][i])==1:
                        modifiedresiduelist.append([df_mod["PTM Protein Locations"][i][df_mod["PTMs"][i].index(mod)]])
                    else:
                        #adjust multiple modifications to a list of strings instead of a long string
                        multimods=[]
                        for j in range(len(df_mod["PTM Protein Locations"][i])):
                            if mod in df_mod["PTMs"][i][j]:
                                multimods.append(df_mod["PTM Protein Locations"][i][j])
                        modifiedresiduelist.append(multimods)
                df_mod["Reagent Specific Mods"]=modifiedresiduelist
                
                #append protein groups and modification locations
                #this should account for when there are multiple modifications ID'd at the same sequence position of different proteins
                #this would otherwise cause multiple mods to be counted as one if this were not accounted for 
                protein_mod=[]
                for i in range(len(df_mod)):
                    if len(df_mod["Reagent Specific Mods"])==1:
                        protein_mod.append(df_mod["PG.ProteinGroups"][i]+";"+df_mod["Reagent Specific Mods"][i][0])
                    else:
                        for j in range(len(df_mod["Reagent Specific Mods"][i])):
                            protein_mod.append(df_mod["PG.ProteinGroups"][i]+";"+df_mod["Reagent Specific Mods"][i][j])
                #list of the actual identifications
                runmodcountlist.append(len(list(set(protein_mod))))
                #list of the number of unique modification sites
                runmodificationlist.append(sorted(list(set(protein_mod)),key=natural_keys))
            modificationdf.loc[run]=runmodificationlist
            modificationcountdf.loc[run]=runmodcountlist

        modificationcountdf["R.Condition"]=sampleconditions
        modificationcountdf_for_group=modificationcountdf.reset_index(drop=True)
        #mean and stdev numbers of modifications per reagent
        modmean=modificationcountdf_for_group.groupby(["R.Condition"]).mean().round(1)
        modstdev=modificationcountdf_for_group.groupby(["R.Condition"]).std()
        modificationcountdf=modificationcountdf.drop(columns="R.Condition")

        return modificationdf,modificationcountdf,modmean,modstdev
    @reactive.effect
    def _():
        @render.plot(width=input.ptmsites_width(),height=input.ptmsites_height())
        def ptmsites_plot():
            modificationdf,modificationcountdf,modmean,modstdev=modificationdfs()
            modlist=find_ptms()

            if input.ptmsites_pickprotein()==True:
                if len(proteinptms_table.data_view(selected=True)["PG.ProteinNames"].tolist())==0:
                    titlemod=""
                else:
                    selectedprotein=proteinptms_table.data_view(selected=True)["PG.ProteinNames"].tolist()[0]
                    titlemod=", "+selectedprotein
            else:
                titlemod=""

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            width=input.ptmsites_barwidth()
            if input.ptmsites_individual_average()=="individual":
                x=np.arange(len(modificationcountdf))
                fig,ax=plt.subplots(figsize=(input.ptmsites_width()*(1/plt.rcParams['figure.dpi']),input.ptmsites_height()*(1/plt.rcParams['figure.dpi'])))
                ymax=max(modificationcountdf.max(axis=1).tolist())
                for i in range(len(modificationcountdf.columns)):
                    ax.bar(x+(i*width),modificationcountdf[modlist[i]],width=width,label=modlist[i])
                    ax.bar_label(ax.containers[i],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                ax.set_xticks(x+(len(modlist)-1)/2*width,modificationcountdf.index,rotation=x_label_rotation)

            if input.ptmsites_individual_average()=="average":
                x=np.arange(len(modmean.index))
                fig,ax=plt.subplots(figsize=(input.ptmsites_width()*(1/plt.rcParams['figure.dpi']),input.ptmsites_height()*(1/plt.rcParams['figure.dpi'])))
                ymax=max(modmean.max(axis=1).tolist())
                for i in range(len(modmean.columns)):
                    bars=ax.bar(x+(i*width),modmean[modlist[i]],yerr=modstdev[modlist[i]],capsize=4,width=width,label=modlist[i])
                    ax.bar_label(bars,label_type="edge",rotation=90,padding=10,fontsize=labelfont)
                ax.set_xticks(x+(len(modlist)-1)/2*width,modmean.index,rotation=x_label_rotation)

            ax.legend(loc="center left",bbox_to_anchor=(1,0.5),fontsize=legendfont)
            ax.set_title("Unique PTM Sites"+titlemod,fontsize=titlefont)
            ax.set_ylabel("Counts",fontsize=axisfont)
            ax.set_xlabel("Run",fontsize=axisfont)
            ax.set_ylim(top=ymax+ymax*y_padding)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            fig.set_tight_layout(True)
            imagedownload("uniqueptmsites")

    # ====================================== Per-Residue PTM Sites
    @render.data_frame
    def proteinptms_table2():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        
        #calculate number of unique modifications found across all runs for each protein
        proteinptmdf=searchoutput[["PG.ProteinNames","PTM Protein Locations"]]
        proteinptmdf=proteinptmdf.explode("PTM Protein Locations").dropna().drop_duplicates().sort_values("PG.ProteinNames").reset_index(drop=True)
        proteinptmdf_grouped=proteinptmdf.groupby(["PG.ProteinNames"]).size().reset_index().rename(columns={0:"Unique PTMs"}).sort_values("Unique PTMs",ascending=False).reset_index(drop=True)

        df_styles={"location":"body",
                   "style":{"column-width":"200px",
                            "overflow":"hidden"
                           }
                  }
        return render.DataGrid(proteinptmdf_grouped,editable=False,selection_mode="row",styles=df_styles)
    @reactive.calc
    def modificationdfs2():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        modlist=find_ptms()

        #sorting functions to sort the PTM Protein Locations column correctly
        #https://stackoverflow.com/questions/5967500/how-to-correctly-sort-a-string-with-a-number-inside
        def atoi(text):
            return int(text) if text.isdigit() else text
        def natural_keys(text):
            '''
            alist.sort(key=natural_keys) sorts in human order
            http://nedbatchelder.com/blog/200712/human_sorting.html
            (See Toothy's implementation in the comments)
            '''
            return [atoi(c) for c in re.split(r'(\d+)',text)]

        modificationdf=pd.DataFrame(columns=modlist)
        modificationcountdf=pd.DataFrame(columns=modlist)

        for run in runnames:
            if input.multiresidueptms_pickprotein()==True:
                if len(proteinptms_table2.data_view(selected=True)["PG.ProteinNames"].tolist())==0:
                    df=searchoutput[searchoutput["Cond_Rep"]==run]
                else:
                    selectedprotein=proteinptms_table2.data_view(selected=True)["PG.ProteinNames"].tolist()[0]
                    df=searchoutput[(searchoutput["Cond_Rep"]==run)&(searchoutput["PG.ProteinNames"]==selectedprotein)]
            else:
                df=searchoutput[searchoutput["Cond_Rep"]==run]

            runmodificationlist=[]
            runmodcountlist=[]
            
            for mod in modlist:
                #df for modifications from specific reagent/source (e.g. just Carbaidomethylations or Oxidations)
                #values in the resulting dfs are most similar to # of peptides since the dfs are based on the whole searchoutput
                #precursors of the same modification site and different charge would be categorized under the same unique modification site
                df_mod=df[df["PTMs"].str.contains(mod,regex=False)].reset_index(drop=True)
                modifiedresiduelist=[]
                for i in range(len(df_mod)):
                    if len(df_mod["PTM Protein Locations"][i])==1:
                        modifiedresiduelist.append([df_mod["PTM Protein Locations"][i][df_mod["PTMs"][i].index(mod)]])
                    else:
                        #adjust multiple modifications to a list of strings instead of a long string
                        multimods=[]
                        for j in range(len(df_mod["PTM Protein Locations"][i])):
                            if mod in df_mod["PTMs"][i][j]:
                                multimods.append(df_mod["PTM Protein Locations"][i][j])
                        modifiedresiduelist.append(multimods)
                df_mod["Reagent Specific Mods"]=modifiedresiduelist
                
                #append protein groups and modification locations
                #this should account for when there are multiple modifications ID'd at the same sequence position of different proteins
                #this would otherwise cause multiple mods to be counted as one if this were not accounted for 
                protein_mod=[]
                for i in range(len(df_mod)):
                    if len(df_mod["Reagent Specific Mods"])==1:
                        protein_mod.append(df_mod["PG.ProteinGroups"][i]+";"+df_mod["Reagent Specific Mods"][i][0])
                    else:
                        for j in range(len(df_mod["Reagent Specific Mods"][i])):
                            protein_mod.append(df_mod["PG.ProteinGroups"][i]+";"+df_mod["Reagent Specific Mods"][i][j])
                #list of the actual identifications
                runmodcountlist.append(len(list(set(protein_mod))))
                #list of the number of unique modification sites
                runmodificationlist.append(sorted(list(set(protein_mod)),key=natural_keys))
            modificationdf.loc[run]=runmodificationlist
            modificationcountdf.loc[run]=runmodcountlist

        modificationcountdf["R.Condition"]=sampleconditions
        modificationcountdf_for_group=modificationcountdf.reset_index(drop=True)
        #mean and stdev numbers of modifications per reagent
        modmean=modificationcountdf_for_group.groupby(["R.Condition"]).mean().round(1)
        modstdev=modificationcountdf_for_group.groupby(["R.Condition"]).std()
        modificationcountdf=modificationcountdf.drop(columns="R.Condition")

        return modificationdf,modificationcountdf,modmean,modstdev
    @render.ui
    def multiresidueptms_ui():
        modlist=find_ptms()

        multiresidueptms=[]
        for i in range(len(modlist)):
            modifiedresidues=modlist[i][modlist[i].find("(")+1:modlist[i].find(")")]
            if modifiedresidues!="Protein N-term" and len(modifiedresidues)>1:
                multiresidueptms.append(modlist[i])
            else:
                pass
        return ui.input_selectize("multiresidueptms_list","Multi-Residue PTMs",choices=multiresidueptms)
    @reactive.effect
    def _():
        @render.plot(width=input.multiresidueptms_width(),height=input.multiresidueptms_height())
        def multiresidueptms_plot():
            mod=input.multiresidueptms_list()
            modificationdf,modificationcountdf,modmean,modstdev=modificationdfs2()
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.multiresidueptms_pickprotein()==True:
                if len(proteinptms_table2.data_view(selected=True)["PG.ProteinNames"].tolist())==0:
                    titlemod=""
                else:
                    selectedprotein=proteinptms_table2.data_view(selected=True)["PG.ProteinNames"].tolist()[0]
                    titlemod=", "+selectedprotein
            else:
                titlemod=""

            fig,ax=plt.subplots(figsize=(input.multiresidueptms_width()*(1/plt.rcParams['figure.dpi']),input.multiresidueptms_height()*(1/plt.rcParams['figure.dpi'])))
            colors=list(mcolors.TABLEAU_COLORS.values())
            x=0
            for run in runnames:
                modifiedresidueslist=pd.Series(modificationdf[mod][run]).str.split(";").str[1]
                residuetargets=[i for i in mod.split("(")[1].split(")")[0]]
                #for fragpipe, the index needs to be 2, for spectronaut it should be 0
                if input.software()=="spectronaut" and "EG.ProteinPTMLocations" in searchoutput.columns:
                    index=0
                #for when spectronaut is selected but it's from a fragpipe initial file
                elif input.software()=="spectronaut" and "EG.ProteinPTMLocations" not in searchoutput.columns:
                    index=2
                if input.software()=="fragpipe":
                    index=2
                residuesonly=[re.split(r"(\d+)",modifiedresidueslist[j])[index] for j in range(len(modifiedresidueslist))]
                mod_dict=dict([k,residuesonly.count(k)] for k in residuetargets)
                
                bottom=0
                for i,key in enumerate(mod_dict):
                    ax.bar(x,mod_dict[key],bottom=bottom,label=key,color=colors[i])
                    bottom+=mod_dict[key]
                x+=1

            handles,labels=ax.get_legend_handles_labels()
            legend=OrderedDict(zip(labels, handles))
            ax.legend(legend.values(),legend.keys(),loc="center left",bbox_to_anchor=(1,0.5),fontsize=legendfont)

            ax.set_ylabel("Counts",fontsize=axisfont)
            ax.set_xlabel("Run",fontsize=axisfont)
            ax.set_xticks(np.arange(0,len(runnames),1),runnames,rotation=x_label_rotation)
            ax.set_title("Per-Residue PTM Sites, "+mod+titlemod,fontsize=titlefont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            fig.set_tight_layout(True)
            imagedownload("multiresidueptms_"+mod)

#endregion

# ============================================================================= Heatmaps
#region
    # ====================================== RT, m/z, IM Heatmaps
    @render.ui
    def cond_rep_list_heatmap():
        if input.conditiontype()=="replicate":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=runnames
            return ui.input_selectize("cond_rep_heatmap","Pick run to show:",choices=opts)               
        elif input.conditiontype()=="condition":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=sampleconditions
            return ui.input_selectize("cond_rep_heatmap","Pick condition to show:",choices=opts)   
    #plot 2D heatmaps for RT, m/z, mobility
    @reactive.effect
    def _():
        @render.plot(width=input.heatmap_width(),height=input.heatmap_height())
        def replicate_heatmap():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            numbins_x=input.heatmap_numbins_x()
            numbins_y=input.heatmap_numbins_y()
            numbins=(numbins_x,numbins_y)

            conditioninput=input.cond_rep_heatmap()
            if input.conditiontype()=="replicate":
                his2dsample=searchoutput[searchoutput["Cond_Rep"]==conditioninput][["Cond_Rep","EG.IonMobility","EG.ApexRT","FG.PrecMz","FG.MS2Quantity"]].sort_values(by="EG.ApexRT").reset_index(drop=True)
            elif input.conditiontype()=="condition":
                his2dsample=searchoutput[searchoutput["R.Condition"]==conditioninput][["R.Condition","EG.IonMobility","EG.ApexRT","FG.PrecMz","FG.MS2Quantity"]].sort_values(by="EG.ApexRT").reset_index(drop=True)

            samplename=conditioninput
            if input.heatmap_cmap()=="default":
                cmap=matplotlib.colors.LinearSegmentedColormap.from_list("",["white","blue","red"])
            else:
                cmap=input.heatmap_cmap()

            fig,ax=plt.subplots(nrows=2,ncols=2,figsize=(input.heatmap_width()*(1/plt.rcParams['figure.dpi']),input.heatmap_height()*(1/plt.rcParams['figure.dpi'])))

            i=ax[0,0].hist2d(his2dsample["EG.ApexRT"],his2dsample["FG.PrecMz"],bins=numbins,cmap=cmap)
            ax[0,0].set_title("RT vs m/z",fontsize=titlefont)
            ax[0,0].set_xlabel("Retention Time (min)",fontsize=axisfont)
            ax[0,0].set_ylabel("m/z",fontsize=axisfont)
            ax[0,0].tick_params(axis="both",labelsize=axisfont_labels)
            fig.colorbar(i[3],ax=ax[0,0])

            j=ax[0,1].hist2d(his2dsample["FG.PrecMz"],his2dsample["EG.IonMobility"],bins=numbins,cmap=cmap)
            ax[0,1].set_title("m/z vs Mobility",fontsize=titlefont)
            ax[0,1].set_xlabel("m/z",fontsize=axisfont)
            ax[0,1].set_ylabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
            ax[0,1].tick_params(axis="both",labelsize=axisfont_labels)
            fig.colorbar(j[3],ax=ax[0,1])

            if len(his2dsample["FG.MS2Quantity"].drop_duplicates())==1:
                ax[1,0].remove()
            else:
                ax[1,0].plot(his2dsample["EG.ApexRT"],his2dsample["FG.MS2Quantity"],color="blue",linewidth=0.5)
                ax[1,0].set_title("RT vs Intensity (line plot)",fontsize=titlefont)
                ax[1,0].set_xlabel("Retention Time (min)",fontsize=axisfont)
                ax[1,0].set_ylabel("Intensity",fontsize=axisfont)
                ax[1,0].tick_params(axis="both",labelsize=axisfont_labels)

            k=ax[1,1].hist2d(his2dsample["EG.ApexRT"].sort_values(),his2dsample["EG.IonMobility"],bins=numbins,cmap=cmap)
            ax[1,1].set_title("RT vs Mobility",fontsize=titlefont)
            ax[1,1].set_xlabel("Retention Time (min)",fontsize=axisfont)
            ax[1,1].set_ylabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
            ax[1,1].tick_params(axis="both",labelsize=axisfont_labels)
            fig.colorbar(k[3],ax=ax[1,1])
            fig.set_tight_layout(True)
            plt.suptitle("Histograms of Identified Precursors"+", "+samplename,y=1,fontsize=titlefont)
            imagedownload("2dhistograms_RT_MZ_IM")

    # ====================================== Charge/PTM Precursor Heatmap
    #download windows template
    @render.download(filename="dia_windows_template.csv")
    def diawindows_template():
        template_df=pd.DataFrame(columns=[
            "#MS Type",
            "Cycle Id",
            "Start IM [1/K0]",
            "End IM [1/K0]",
            "Start Mass [m/z]",
            "End Mass [m/z]",
            "CE [eV]"
            ])
        with io.BytesIO() as buf:
            template_df.to_csv(buf,index=False)
            yield buf.getvalue()
    #imported DIA windows
    def diawindows_import():
        if input.diawindow_upload() is None:
            return pd.DataFrame()
        if input.windows_choice()=="diagonal":
            diagonal=pd.read_csv(input.diawindow_upload()[0]["datapath"])
            diagonal=diagonal.drop(index=0).reset_index(drop=True).drop(columns=["type"])
            columnlist=[]
            for column in diagonal.columns.tolist():
                diagonal[column]=diagonal[column].astype(float)
                columnlist.append(column.strip())
            diagonal=diagonal.set_axis(columnlist,axis=1)
            diagonalspacing=list(set(diagonal["mass pos.1 end [m/z]"]-diagonal["mass pos.1 start [m/z]"].tolist()))[0]
            diagonal["mass pos.2 end [m/z]"]=diagonal["mass pos.2 start [m/z]"]+diagonalspacing
            diawindows=diagonal
        else:
            diawindows=pd.read_csv(input.diawindow_upload()[0]["datapath"])
            diawindows=diawindows.drop(index=0).reset_index(drop=True)
            startcoords=[]
            for i in range(len(diawindows)):
                startcorner=float(diawindows["Start Mass [m/z]"][i]),float(diawindows["Start IM [1/K0]"][i])
                startcoords.append(startcorner)
            diawindows["W"]=diawindows["End Mass [m/z]"].astype(float)-diawindows["Start Mass [m/z]"].astype(float)
            diawindows["H"]=diawindows["End IM [1/K0]"].astype(float)-diawindows["Start IM [1/K0]"].astype(float)
            diawindows["xy"]=startcoords
        return diawindows
    #Lubeck DIA windows
    def lubeckdiawindow():
        lubeckdia=pd.DataFrame({
            "#MS Type":['PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF'],
            "Cycle Id":[1, 1, 1, 2, 2, 2, 3, 3, 3, 4, 4, 4, 5, 5, 5, 6, 6, 6, 7, 7, 7, 8, 8, 8, 9, 9, 9, 10, 10, 10, 11, 11, 11, 12, 12, 12, 13, 13, 13, 14, 14, 14, 15, 15, 15, 16, 16, 16, 17, 17, 17, 18, 18, 18, 19, 19, 19, 20, 20, 20],
            "Start IM [1/K0]":[0.965, 0.805, 0.6, 0.977, 0.87, 0.6, 0.986, 0.89, 0.6, 0.995, 0.9, 0.6, 1.006, 0.906, 0.6, 1.025, 0.91, 0.6, 1.045, 0.917, 0.6, 1.064, 0.927, 0.6, 1.085, 0.94, 0.6, 1.105, 0.955, 0.6, 0.97, 0.85, 0.6, 0.982, 0.884, 0.6, 0.99, 0.895, 0.6, 1.0, 0.903, 0.6, 1.015, 0.908, 0.6, 1.034, 0.914, 0.6, 1.054, 0.92, 0.6, 1.075, 0.934, 0.6, 1.095, 0.947, 0.6, 1.11, 0.96, 0.6],
            "End IM [1/K0]":[1.12, 0.965, 0.805, 1.15, 0.977, 0.87, 1.19, 0.986, 0.89, 1.23, 0.995, 0.9, 1.27, 1.006, 0.906, 1.31, 1.025, 0.91, 1.35, 1.045, 0.917, 1.39, 1.064, 0.927, 1.43, 1.085, 0.94, 1.45, 1.105, 0.955, 1.13, 0.97, 0.85, 1.17, 0.982, 0.884, 1.21, 0.99, 0.895, 1.25, 1.0, 0.903, 1.29, 1.015, 0.908, 1.33, 1.034, 0.914, 1.37, 1.054, 0.92, 1.41, 1.075, 0.934, 1.44, 1.095, 0.947, 1.45, 1.11, 0.96],
            "Start Mass [m/z]":[725.13, 559.8, 350.68, 746.34, 574.21, 412.39, 769.41, 588.81, 437.42, 794.87, 603.8, 456.05, 821.43, 619.33, 472.76, 851.93, 635.35, 488.1, 886.49, 651.86, 502.77, 928.86, 668.83, 517.08, 982.67, 686.35, 531.29, 1059.54, 704.89, 545.77, 735.74, 567.01, 381.54, 757.88, 581.51, 424.9, 782.14, 596.31, 446.74, 808.15, 611.57, 464.4, 836.68, 627.34, 480.43, 869.21, 643.61, 495.43, 907.67, 660.35, 509.92, 955.76, 677.59, 524.18, 1021.11, 695.62, 538.53, 1154.84, 715.01, 552.79],
            "End Mass [m/z]":[736.24, 567.5, 382.04, 758.38, 582.01, 425.4, 782.64, 596.81, 447.23, 808.65, 612.07, 464.9, 837.18, 627.84, 480.93, 869.71, 644.1, 495.93, 908.17, 660.85, 510.43, 956.26, 678.09, 524.68, 1021.61, 696.12, 539.03, 1155.34, 715.51, 553.28, 746.84, 574.71, 412.89, 769.91, 589.31, 437.92, 795.37, 604.3, 456.55, 821.93, 619.83, 473.26, 852.43, 635.85, 488.6, 886.99, 652.36, 503.27, 929.36, 669.33, 517.58, 983.17, 686.85, 531.79, 1060.04, 705.39, 546.27, 1250.64, 725.63, 560.3],
            "CE [eV]":['-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-'],
            "W":[11.110000000000014, 7.7000000000000455, 31.360000000000014, 12.039999999999964, 7.7999999999999545, 13.009999999999991, 13.230000000000018, 8.0, 9.810000000000002, 13.779999999999973, 8.270000000000095, 8.849999999999966, 15.75, 8.509999999999991, 8.170000000000016, 17.780000000000086, 8.75, 7.829999999999984, 21.67999999999995, 8.990000000000009, 7.660000000000025, 27.399999999999977, 9.259999999999991, 7.599999999999909, 38.940000000000055, 9.769999999999982, 7.740000000000009, 95.79999999999995, 10.620000000000005, 7.509999999999991, 11.100000000000023, 7.7000000000000455, 31.349999999999966, 12.029999999999973, 7.7999999999999545, 13.020000000000039, 13.230000000000018, 7.990000000000009, 9.810000000000002, 13.779999999999973, 8.259999999999991, 8.860000000000014, 15.75, 8.509999999999991, 8.170000000000016, 17.779999999999973, 8.75, 7.839999999999975, 21.690000000000055, 8.980000000000018, 7.660000000000025, 27.409999999999968, 9.259999999999991, 7.610000000000014, 38.92999999999995, 9.769999999999982, 7.740000000000009, 95.80000000000018, 10.620000000000005, 7.509999999999991],
            "H":[0.15500000000000014, 0.15999999999999992, 0.20500000000000007, 0.17299999999999993, 0.10699999999999998, 0.27, 0.20399999999999996, 0.09599999999999997, 0.29000000000000004, 0.235, 0.09499999999999997, 0.30000000000000004, 0.264, 0.09999999999999998, 0.30600000000000005, 0.28500000000000014, 0.11499999999999988, 0.31000000000000005, 0.30500000000000016, 0.1279999999999999, 0.31700000000000006, 0.32599999999999985, 0.137, 0.32700000000000007, 0.345, 0.14500000000000002, 0.33999999999999997, 0.345, 0.15000000000000002, 0.355, 0.15999999999999992, 0.12, 0.25, 0.18799999999999994, 0.09799999999999998, 0.28400000000000003, 0.21999999999999997, 0.09499999999999997, 0.29500000000000004, 0.25, 0.09699999999999998, 0.30300000000000005, 0.27500000000000013, 0.10699999999999987, 0.30800000000000005, 0.29600000000000004, 0.12, 0.31400000000000006, 0.31600000000000006, 0.134, 0.32000000000000006, 0.33499999999999996, 0.1409999999999999, 0.3340000000000001, 0.345, 0.14800000000000002, 0.347, 0.33999999999999986, 0.15000000000000013, 0.36],
            "xy":[(725.13, 0.965), (559.8, 0.805), (350.68, 0.6), (746.34, 0.977), (574.21, 0.87), (412.39, 0.6), (769.41, 0.986), (588.81, 0.89), (437.42, 0.6), (794.87, 0.995), (603.8, 0.9), (456.05, 0.6), (821.43, 1.006), (619.33, 0.906), (472.76, 0.6), (851.93, 1.025), (635.35, 0.91), (488.1, 0.6), (886.49, 1.045), (651.86, 0.917), (502.77, 0.6), (928.86, 1.064), (668.83, 0.927), (517.08, 0.6), (982.67, 1.085), (686.35, 0.94), (531.29, 0.6), (1059.54, 1.105), (704.89, 0.955), (545.77, 0.6), (735.74, 0.97), (567.01, 0.85), (381.54, 0.6), (757.88, 0.982), (581.51, 0.884), (424.9, 0.6), (782.14, 0.99), (596.31, 0.895), (446.74, 0.6), (808.15, 1.0), (611.57, 0.903), (464.4, 0.6), (836.68, 1.015), (627.34, 0.908), (480.43, 0.6), (869.21, 1.034), (643.61, 0.914), (495.43, 0.6), (907.67, 1.054), (660.35, 0.92), (509.92, 0.6), (955.76, 1.075), (677.59, 0.934), (524.18, 0.6), (1021.11, 1.095), (695.62, 0.947), (538.53, 0.6), (1154.84, 1.11), (715.01, 0.96), (552.79, 0.6)]
        })
        return lubeckdia
    #phospho DIA windows
    def phosphodiawindow():
        phosphodia=pd.DataFrame({
            "#MS Type":['PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF', 'PASEF'],
            "Cycle Id":[1, 1, 2, 2, 3, 3, 4, 4, 5, 5, 6, 6, 7, 7, 8, 8, 9, 9, 10, 11, 12, 13, 14, 15, 16, 17, 18, 19],
            "Start IM [1/K0]":[0.8691, 0.6811, 0.8834, 0.691, 0.912, 0.701, 0.9264, 0.7109, 0.9407, 0.7208, 0.9694, 0.7308, 0.9837, 0.7407, 1.0123, 0.7544, 1.0267, 0.7688, 0.7831, 0.7974, 0.8117, 0.8261, 0.8404, 0.8547, 0.8977, 0.955, 0.998, 1.041],
            "End IM [1/K0]":[1.1616, 0.8579, 1.1791, 0.8786, 1.2142, 0.8993, 1.2317, 0.92, 1.2492, 0.9407, 1.2842, 0.9614, 1.3017, 0.9821, 1.3368, 1.0028, 1.3543, 1.0235, 1.0442, 1.0649, 1.0856, 1.1063, 1.1266, 1.1441, 1.1966, 1.2667, 1.3192, 1.3718],
            "Start Mass [m/z]":[839.43, 419.43, 867.43, 447.43, 923.43, 475.43, 951.43, 503.43, 979.43, 531.43, 1035.43, 559.43, 1063.43, 587.43, 1119.43, 615.43, 1147.43, 643.43, 671.43, 699.43, 727.43, 755.43, 783.43, 811.43, 895.43, 1007.43, 1091.43, 1175.43],
            "End Mass [m/z]":[867.43, 447.43, 895.43, 475.43, 951.43, 503.43, 979.43, 531.43, 1007.43, 559.43, 1063.43, 587.43, 1091.43, 615.43, 1147.43, 643.43, 1175.43, 671.43, 699.43, 727.43, 755.43, 783.43, 811.43, 839.43, 923.43, 1035.43, 1119.43, 1203.43],
            "CE [eV]":['-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-', '-'],
            "W":[28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 27.999999999999943, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.0, 28.000000000000114, 28.0, 28.0],
            "H":[0.2925, 0.17679999999999996, 0.2957000000000001, 0.1876000000000001, 0.3021999999999999, 0.19830000000000003, 0.3053, 0.20910000000000006, 0.3085000000000001, 0.21989999999999998, 0.31479999999999997, 0.23060000000000003, 0.31800000000000006, 0.24139999999999995, 0.3245, 0.24839999999999995, 0.3276000000000001, 0.25470000000000004, 0.2611, 0.26749999999999996, 0.2738999999999999, 0.2802000000000001, 0.2862, 0.2893999999999999, 0.29890000000000005, 0.3117, 0.32119999999999993, 0.3308],
            "xy":[(839.43, 0.8691), (419.43, 0.6811), (867.43, 0.8834), (447.43, 0.691), (923.43, 0.912), (475.43, 0.701), (951.43, 0.9264), (503.43, 0.7109), (979.43, 0.9407), (531.43, 0.7208), (1035.43, 0.9694), (559.43, 0.7308), (1063.43, 0.9837), (587.43, 0.7407), (1119.43, 1.0123), (615.43, 0.7544), (1147.43, 1.0267), (643.43, 0.7688), (671.43, 0.7831), (699.43, 0.7974), (727.43, 0.8117), (755.43, 0.8261), (783.43, 0.8404), (811.43, 0.8547), (895.43, 0.8977), (1007.43, 0.955), (1091.43, 0.998), (1175.43, 1.041)],
        })
        return phosphodia
    #bremen DIA windows
    def bremendiawindow():
        bremendia=pd.DataFrame({
            "#MS Type":['PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF','PASEF'],
            "Cycle Id":[1,1,1,2,2,2,3,3,3,4,4,4,5,5,5,6,6,6,7,7,7,8,8,8],
            "Start IM [1/K0]":[0.64,0.83,1.01,0.64,0.85,1.04,0.64,0.87,1.06,0.64,0.9,1.09,0.64,0.92,1.11,0.64,0.94,1.13,0.64,0.97,1.16,0.64,0.99,1.18],
            "End IM [1/K0]":[0.83,1.01,1.37,0.85,1.04,1.37,0.87,1.06,1.37,0.9,1.09,1.37,0.92,1.11,1.37,0.94,1.13,1.37,0.97,1.16,1.37,0.99,1.18,1.37],
            "Start Mass [m/z]":[400,600,800,425,625,825,450,650,850,475,675,875,500,700,900,525,725,925,550,750,950,575,775,975],
            "End Mass [m/z]":[425,625,825,450,650,850,475,675,875,500,700,900,525,725,925,550,750,950,575,775,975,600,800,1000],
            "CE [eV]":['-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-','-'],
            "W":[25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25,25],
            "H":[0.19,0.18,0.36,0.21,0.19,0.33,0.23,0.19,0.31,0.26,0.19,0.28,0.28,0.19,0.26,0.3,0.19,0.24,0.33,0.19,0.21,0.35,0.19,0.19],
            "xy":[(400,0.64),(600,0.83),(800,1.01),(425,0.64),(625,0.85),(825,1.04),(450,0.64),(650,0.87),(850,1.06),(475,0.64),(675,0.9),(875,1.09),(500,0.64),(700,0.92),(900,1.11),(525,0.64),(725,0.94),(925,1.13),(550,0.64),(750,0.97),(950,1.16),(575,0.64),(775,0.99),(975,1.18)],
        })
        return bremendia

    @render.ui
    def chargeptm_cond_rep_ui():
        if input.chargeptm_conditiontype()=="replicate":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=runnames
            return ui.input_selectize("chargeptm_cond_rep","Pick run to show:",choices=opts)               
        elif input.chargeptm_conditiontype()=="condition":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=sampleconditions
            return ui.input_selectize("chargeptm_cond_rep","Pick condition to show:",choices=opts)   
    #render ui call for dropdown calling charge states that were detected
    @render.ui
    def chargestates_chargeptmheatmap_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        mincharge=min(searchoutput["FG.Charge"])
        maxcharge=max(searchoutput["FG.Charge"])
        opts=[item for item in range(mincharge,maxcharge+1)]
        opts.insert(0,0)
        return ui.input_selectize("chargestates_chargeptmheatmap_list","Pick charge to plot for (use 0 for all):",choices=opts)
    #render ui call for dropdown calling PTMs that were detected
    @render.ui
    def ptm_chargeptmheatmap_ui():
        listofptms=find_ptms()
        listofptms.insert(0,"None")
        return ui.input_selectize("ptm_chargeptmheatmap_list","Pick PTM to plot for (use None for all):",choices=listofptms,selected="None")
    #Charge/PTM precursor heatmap
    @reactive.effect
    def _():
        @render.plot(width=input.chargeptmheatmap_width(),height=input.chargeptmheatmap_height())
        def chargeptmheatmap():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

            charge=input.chargestates_chargeptmheatmap_list()
            ptm=input.ptm_chargeptmheatmap_list()

            numbins_x=input.chargeptm_numbins_x()
            numbins_y=input.chargeptm_numbins_y()
            numbins=[numbins_x,numbins_y]

            if input.chargeptmheatmap_cmap()=="default":
                cmap=matplotlib.colors.LinearSegmentedColormap.from_list("",["white","blue","red"])
            else:
                cmap=input.chargeptmheatmap_cmap()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()

            if input.chargeptm_conditiontype()=="replicate":
                df=searchoutput[searchoutput["Cond_Rep"]==input.chargeptm_cond_rep()]
            elif input.chargeptm_conditiontype()=="condition":
                df=searchoutput[searchoutput["R.Condition"]==input.chargeptm_cond_rep()]

            if ptm=="None":
                if charge=="0":
                    #all precursors
                    his2dsample=df[["R.Condition","R.Replicate","EG.IonMobility","FG.PrecMz"]]
                    title="Precursors"
                    savetitle="All Precursor IDs Heatmap_"
                elif charge!="0":
                    #all precursors of specific charge
                    his2dsample=df[df["FG.Charge"]==int(charge)][["R.Condition","R.Replicate","EG.IonMobility","FG.PrecMz"]]
                    title=str(charge)+"+ Precursors"
                    savetitle=str(charge)+"+_"+"_Precursor IDs Heatmap_"   
            if ptm!="None":
                if charge=="0":
                    #all modified precursors
                    his2dsample=df[df["EG.ModifiedPeptide"].str.contains(ptm,regex=False)][["R.Condition","R.Replicate","EG.IonMobility","FG.PrecMz"]]
                    title=ptm+" Precursors"
                    savetitle=ptm+"_Precursor IDs Heatmap_"   
                elif charge!="0":
                    #modified precursors of specific charge
                    his2dsample=df[(df["FG.Charge"]==int(charge))&(df["EG.ModifiedPeptide"].str.contains(ptm,regex=False))][["R.Condition","R.Replicate","EG.IonMobility","FG.PrecMz"]]
                    title=ptm+" "+str(charge)+"+ Precursors"
                    savetitle=ptm+"_"+str(charge)+"+_"+"_Precursor IDs Heatmap_"

            if input.chargeptmheatmap_axishistogram()==True:
                fig,ax=plt.subplot_mosaic([["histx","."],["hist2d","histy"]],width_ratios=(4,1),height_ratios=(1,4),layout="constrained",figsize=(input.chargeptmheatmap_width()*(1/plt.rcParams['figure.dpi']),input.chargeptmheatmap_height()*(1/plt.rcParams['figure.dpi'])))

                j=ax["hist2d"].hist2d(his2dsample["FG.PrecMz"],his2dsample["EG.IonMobility"],bins=numbins,cmap=cmap)
                ax["hist2d"].set_xlabel("m/z",fontsize=axisfont)
                ax["hist2d"].set_ylabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)

                ax["histx"].hist(his2dsample["FG.PrecMz"],bins=numbins_x)
                ax["histx"].tick_params(axis="x",labelbottom=False)
                ax["histx"].set_ylabel("Counts",fontsize=axisfont)

                ax["histy"].hist(his2dsample["EG.IonMobility"],bins=numbins_y,orientation="horizontal")
                ax["histy"].tick_params(axis="y",labelleft=False)
                ax["histy"].set_xlabel("Counts",fontsize=axisfont)
                fig.suptitle(input.chargeptm_cond_rep()+", "+title,fontsize=titlefont)
            else:
                fig,ax=plt.subplots(figsize=(input.chargeptmheatmap_width()*(1/plt.rcParams['figure.dpi']),input.chargeptmheatmap_height()*(1/plt.rcParams['figure.dpi'])))
                j=ax.hist2d(his2dsample["FG.PrecMz"],his2dsample["EG.IonMobility"],bins=numbins,cmap=cmap)
                ax.set_title(input.chargeptm_cond_rep()+", "+title,fontsize=titlefont)
                ax.set_xlabel("m/z",fontsize=axisfont)
                ax.set_ylabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                fig.colorbar(j[3],ax=ax)
                
            if input.windows_choice()!="None":
                if input.windows_choice()=="diagonal":
                    diawindows=diawindows_import()
                    coordlist=[]
                    for i in range(len(diawindows)):
                        polycoords=[(diawindows["mass pos.1 start [m/z]"][i],diawindows["mobility pos.1 [1/K0]"][i]),
                                    (diawindows["mass pos.1 end [m/z]"][i],diawindows["mobility pos.1 [1/K0]"][i]),
                                    (diawindows["mass pos.2 end [m/z]"][i],diawindows["mobility pos.2 [1/K0]"][i]),
                                    (diawindows["mass pos.2 start [m/z]"][i],diawindows["mobility pos.2 [1/K0]"][i])
                                ]
                        coordlist.append(polycoords)
                    diawindows["coords"]=coordlist
                    for i in range(len(diawindows)):
                        y=diawindows["coords"][i]

                        p=matplotlib.patches.Polygon(y,facecolor="red",alpha=0.1,edgecolor="black")
                        ax.add_patch(p)
                else:
                    if input.windows_choice()=="lubeck":
                        diawindows=lubeckdiawindow()
                    elif input.windows_choice()=="phospho":
                        diawindows=phosphodiawindow()
                    elif input.windows_choice()=="bremen":
                        diawindows=bremendiawindow()
                    elif input.windows_choice()=="imported":
                        diawindows=diawindows_import()

                    for i in range(len(diawindows)):
                        rect=matplotlib.patches.Rectangle(xy=diawindows["xy"][i],width=diawindows["W"][i],height=diawindows["H"][i],facecolor="red",alpha=0.1,edgecolor="grey")
                        if input.chargeptmheatmap_axishistogram()==True:
                            ax["hist2d"].add_patch(rect)
                        else:    
                            ax.add_patch(rect)
            fig.set_tight_layout(True)
            imagedownload("precursorheatmap")

    # ====================================== Charge/PTM Precursor Scatter
    #render ui call for dropdown calling Cond_Rep column
    @render.ui
    def chargeptmscatter_cond_rep():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_selectize("chargeptmscatter_cond_rep_pick","Pick run:",choices=opts)
    #render ui call for dropdown calling PTMs that were detected
    @render.ui
    def ptm_chargeptmscatter_ui():
        listofptms=find_ptms()
        listofptms.insert(0,"None")
        return ui.input_selectize("ptm_chargeptmscatter_list","Pick PTM to highlight (use None for all):",choices=listofptms,selected="None")
    #render ui call for dropdown calling charge states that were detected
    @render.ui
    def chargestates_chargeptmscatter_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        mincharge=min(searchoutput["FG.Charge"])
        maxcharge=max(searchoutput["FG.Charge"])
        opts=[item for item in range(mincharge,maxcharge+1)]
        return ui.input_checkbox_group("chargestates_chargeptmscatter_list","Pick charge to highlight:",choices=opts)
    #scatterplot of picked PTM or charge against the rest of the detected precursors (better for DDA to show charge groups in the heatmap)
    @reactive.effect
    def _():
        @render.plot(width=input.chargeptmscatter_width(),height=input.chargeptmscatter_height())
        def chargeptmscatter():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()

            ptm=input.ptm_chargeptmscatter_list()
            charge=input.chargestates_chargeptmscatter_list()
            cond_rep_pick=input.chargeptmscatter_cond_rep_pick()

            colorlist=["tab:blue","tab:orange","tab:green","tab:red","tab:purple","tab:brown","tab:pink","tab:gray","tab:olive","tab:cyan"]
            precursor_other=searchoutput[searchoutput["Cond_Rep"]==cond_rep_pick]

            fig,ax=plt.subplots(figsize=(input.chargeptmscatter_width()*(1/plt.rcParams['figure.dpi']),input.chargeptmscatter_height()*(1/plt.rcParams['figure.dpi'])))
                
            if ptm=="None":
                if len(charge)==0:
                    ax.scatter(x=precursor_other["FG.PrecMz"],y=precursor_other["EG.IonMobility"],s=2,label="All Precursors",color="dimgray")
                else:
                    ax.scatter(x=precursor_other["FG.PrecMz"],y=precursor_other["EG.IonMobility"],s=2,label="All Other Precursors",color="dimgray")
                    for i in range(len(charge)):
                        x=searchoutput[(searchoutput["FG.Charge"]==int(charge[i]))&(searchoutput["Cond_Rep"]==cond_rep_pick)]["FG.PrecMz"]
                        y=searchoutput[(searchoutput["FG.Charge"]==int(charge[i]))&(searchoutput["Cond_Rep"]==cond_rep_pick)]["EG.IonMobility"]
                        ax.scatter(x=x,y=y,s=2,label=str(charge[i])+"+",color=colorlist[i])
            if ptm!="None":
                ax.scatter(x=precursor_other["FG.PrecMz"],y=precursor_other["EG.IonMobility"],s=2,label="All Other Precursors",color="dimgray")
                if len(charge)==0:
                        x=searchoutput[(searchoutput["EG.ModifiedPeptide"].str.contains(ptm,regex=False))&(searchoutput["Cond_Rep"]==cond_rep_pick)]["FG.PrecMz"]
                        y=searchoutput[(searchoutput["EG.ModifiedPeptide"].str.contains(ptm,regex=False))&(searchoutput["Cond_Rep"]==cond_rep_pick)]["EG.IonMobility"]
                        ax.scatter(x=x,y=y,s=2,label=ptm+" Precursors",color="tab:cyan")
                else:
                    for i in range(len(charge)):
                        x=searchoutput[(searchoutput["FG.Charge"]==int(charge[i]))&(searchoutput["EG.ModifiedPeptide"].str.contains(ptm,regex=False))&(searchoutput["Cond_Rep"]==cond_rep_pick)]["FG.PrecMz"]
                        y=searchoutput[(searchoutput["FG.Charge"]==int(charge[i]))&(searchoutput["EG.ModifiedPeptide"].str.contains(ptm,regex=False))&(searchoutput["Cond_Rep"]==cond_rep_pick)]["EG.IonMobility"]
                        ax.scatter(x=x,y=y,s=2,label=ptm+" "+str(charge[i])+"+",color=colorlist[i])

            ax.set_xlabel("m/z",fontsize=axisfont)
            ax.set_ylabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
            ax.legend(loc="upper left",fontsize=legendfont,markerscale=5)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("precursorscatter")
    #show table of # of precursors selected in charge-ptm scatterplot
    @render.table
    def chargeptmscatter_table():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        ptm=input.ptm_chargeptmscatter_list()
        cond_rep_pick=input.chargeptmscatter_cond_rep_pick()

        if ptm=="None":
            #all ptms
            precursor_pick=searchoutput[(searchoutput["Cond_Rep"]==cond_rep_pick)]
        if ptm!="None":
            #specific ptm
            precursor_pick=searchoutput[(searchoutput["EG.ModifiedPeptide"].str.contains(ptm,regex=False)==True)&(searchoutput["Cond_Rep"]==cond_rep_pick)]

        sortedchargelist=sorted(precursor_pick["FG.Charge"].drop_duplicates().tolist())
        charge_populationlist=[]
        for charge in sortedchargelist:
            charge_populationlist.append(len(precursor_pick[precursor_pick["FG.Charge"]==charge]))

        sortedchargelist.insert(0,"Sum")
        charge_populationlist.insert(0,sum(charge_populationlist))

        df=pd.DataFrame({"Charge":sortedchargelist,"# Precursors":charge_populationlist})
        
        return df

    # ====================================== IDs vs RT
    @render.ui
    def binslider_ui():
        return ui.input_slider("binslider","Number of RT bins:",min=100,max=1000,step=50,value=500,ticks=True)
    #render ui for checkboxes to plot specific runs
    @render.ui
    def ids_vs_rt_checkbox():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        keys=runnames
        opts=dict()
        for x in keys:
            opts[x]=x
        return ui.input_checkbox_group("ids_vs_rt_checkbox_pick","Pick runs to plot data for:",choices=opts)
    #plot # of IDs vs RT for each run
    @reactive.effect
    def _():
        @render.plot(width=input.idsvsrt_width(),height=input.idsvsrt_height())
        def ids_vs_rt():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            
            rtmax=float(math.ceil(max(searchoutput["EG.ApexRT"]))) #needs to be a float
            numbins=input.binslider()
            runlist=input.ids_vs_rt_checkbox_pick()

            bintime=rtmax/numbins*60

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()

            fig,ax=plt.subplots(figsize=(input.idsvsrt_width()*(1/plt.rcParams['figure.dpi']),input.idsvsrt_height()*(1/plt.rcParams['figure.dpi'])))

            for run in runlist:
                rt_run=searchoutput[searchoutput["Cond_Rep"]==run]["EG.ApexRT"]
                if rt_run.empty:
                    continue
                hist=np.histogram(rt_run,bins=numbins,range=(0.0,rtmax))
                ax.plot(np.delete(hist[1],0),hist[0],linewidth=0.5,label=run)

            ax.set_ylabel("# of IDs",fontsize=axisfont)
            ax.set_xlabel("RT (min)",fontsize=axisfont)
            ax.text(0,(ax.get_ylim()[1]-(0.1*ax.get_ylim()[1])),"~"+str(round(bintime,2))+" s per bin",fontsize=axisfont)
            legend=ax.legend(loc='center left', bbox_to_anchor=(1, 0.5),prop={'size':legendfont})
            for i in legend.legend_handles:
                i.set_linewidth(5)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.set_title("# of Precursor IDs vs RT",fontsize=titlefont)
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            if runlist:
                imagedownload("ids_vs_rt")

    # ====================================== Histogram
    #render ui call for dropdown calling Cond_Rep column
    @render.ui
    def histogram_cond_rep_list():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_checkbox_group("histogram_cond_rep_pick","Pick runs to plot:",choices=opts)
    @reactive.effect
    def _():
        @render.plot(width=input.histogram_width(),height=input.histogram_height())
        def histogram_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            runlist=input.histogram_cond_rep_pick()
            bins=input.histogram_numbins()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()

            fig,ax=plt.subplots(figsize=(input.histogram_width()*(1/plt.rcParams['figure.dpi']),input.histogram_height()*(1/plt.rcParams['figure.dpi'])))
            ax.set_ylabel("Counts",fontsize=axisfont)
            for run in runlist:
                df=searchoutput[searchoutput["Cond_Rep"]==run]
                if input.histogram_pick()=="ionmobility":
                    ax.hist(df["EG.IonMobility"],bins=bins,label=run,alpha=0.75)
                    ax.set_xlabel("EG.IonMobility",fontsize=axisfont)
                if input.histogram_pick()=="precursormz":
                    ax.hist(df["FG.PrecMz"],bins=bins,label=run,alpha=0.75)
                    ax.set_xlabel("Precursor m/z",fontsize=axisfont)
                if input.histogram_pick()=="precursorintensity":
                    df=df[df["FG.MS2Quantity"]!=0]
                    ax.hist(np.log10(df["FG.MS2Quantity"]),bins=bins,label=run,alpha=0.75)
                    ax.set_xlabel("log10(Precursor Intensity)",fontsize=axisfont)
                if input.histogram_pick()=="proteinintensity":
                    df=df[df["PG.MS2Quantity"]!=0]
                    ax.hist(np.log10(df["PG.MS2Quantity"]),bins=bins,label=run,alpha=0.75)
                    ax.set_xlabel("log10(Protein Intensity)",fontsize=axisfont)
            ax.legend(prop={'size':legendfont})
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            if runlist:
                imagedownload("histogram_"+input.histogram_pick())

#endregion

# ============================================================================= Statistics
#region
    # ====================================== Volcano Plot
    #render ui for picking conditions
    @render.ui
    def volcano_condition1():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=sampleconditions
        return ui.input_selectize("control_condition","Pick control condition:",choices=opts)
    @render.ui
    def volcano_condition2():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=sampleconditions
        return ui.input_selectize("test_condition","Pick test condition:",choices=opts,selected=sampleconditions[1])
    #calculation for fold change and p value
    @reactive.calc
    def volcano_calc():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

        control=input.control_condition()
        test=input.test_condition()

        controldf=searchoutput[searchoutput["R.Condition"]==control][["PG.ProteinGroups","PG.Genes","PG.ProteinNames","PG.MS2Quantity"]].drop_duplicates().groupby(["PG.ProteinGroups","PG.Genes","PG.ProteinNames"])
        testdf=searchoutput[searchoutput["R.Condition"]==test][["PG.ProteinGroups","PG.Genes","PG.ProteinNames","PG.MS2Quantity"]].drop_duplicates().groupby(["PG.ProteinGroups","PG.Genes","PG.ProteinNames"])
        controldf_means=controldf.mean()
        testdf_means=testdf.mean()
        merged=controldf_means.merge(testdf_means,on=["PG.ProteinGroups","PG.Genes","PG.ProteinNames"],suffixes=("_control_mean","_test_mean")).reset_index().dropna()

        proteinlist=[]
        controllist=[]
        testlist=[]
        for protein,gene,name in zip(merged["PG.ProteinGroups"],merged["PG.Genes"],merged["PG.ProteinNames"]):
            proteinlist.append(protein)
            controllist.append(list(controldf.get_group((protein,gene,name))["PG.MS2Quantity"]))
            testlist.append(list(testdf.get_group((protein,gene,name))["PG.MS2Quantity"]))

        merged["Control"]=controllist
        merged["Test"]=testlist
        merged["log2_FoldChange"]=np.log2(merged["PG.MS2Quantity_test_mean"]/merged["PG.MS2Quantity_control_mean"])

        merged=merged.set_index("PG.ProteinGroups")

        pvalue=[]
        for protein in merged.index:
            pvalue.append(-np.log10(stats.ttest_ind(a=merged.loc[protein]["Control"],b=merged.loc[protein]["Test"])[1]))
        merged["-log10_pvalue"]=pvalue

        pvalue_cutoff=input.volcano_pvalue()
        foldchange_cutoff=input.volcano_foldchange()

        coordlist=[]
        colorlist=[]
        for protein in merged.index:
            if merged.loc[protein]["log2_FoldChange"] >= foldchange_cutoff and np.absolute(merged.loc[protein]["-log10_pvalue"]) >= pvalue_cutoff:
                colorlist.append("r")
                coordlist.append((merged.loc[protein]["log2_FoldChange"],merged.loc[protein]["-log10_pvalue"]))
            elif merged.loc[protein]["log2_FoldChange"] <= -foldchange_cutoff and np.absolute(merged.loc[protein]["-log10_pvalue"]) >= pvalue_cutoff:
                colorlist.append("b")
                coordlist.append((merged.loc[protein]["log2_FoldChange"],merged.loc[protein]["-log10_pvalue"]))
            else:
                colorlist.append("grey")
                coordlist.append("")

        merged["log2_FoldChange_absolute"]=merged["log2_FoldChange"].abs()
        merged["color"]=colorlist
        merged["label"]=coordlist

        return merged
    #download a table with values from the volcano plot
    @render.download(filename="volcanoplot_values.csv")
    def volcano_download():
        merged=volcano_calc()
        merged_download=merged[["PG.Genes","PG.ProteinNames","log2_FoldChange","-log10_pvalue","log2_FoldChange_absolute"]].reset_index()
        merged_download=merged_download.sort_values("log2_FoldChange_absolute",ascending=False)
        with io.BytesIO() as buf:
            merged_download.to_csv(buf,index=False)
            yield buf.getvalue()
    #volcano plot
    @reactive.effect
    def _():
        @render.plot(width=input.volcano_width(),height=input.volcano_height())
        def volcanoplot():
            import textalloc as ta
            merged=volcano_calc()

            pvalue_cutoff=input.volcano_pvalue()
            foldchange_cutoff=input.volcano_foldchange()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()

            fig,ax=plt.subplots(figsize=(input.volcano_width()*(1/plt.rcParams['figure.dpi']),input.volcano_height()*(1/plt.rcParams['figure.dpi'])))
            ax.scatter(merged["log2_FoldChange"],merged["-log10_pvalue"],s=1,c=merged["color"])
            if input.volcano_h_v_lines()==True:
                ax.axhline(y=pvalue_cutoff,color="dimgrey",linestyle="-",linewidth=0.75)
                ax.axvline(x=foldchange_cutoff,color="dimgrey",linestyle="-",linewidth=0.75)
                ax.axvline(x=-foldchange_cutoff,color="dimgrey",linestyle="-",linewidth=0.75)
            ax.set_xlabel("log2 Fold Change",fontsize=axisfont)
            ax.set_ylabel("-log10 p value",fontsize=axisfont)
            ax.set_title("Control: "+input.control_condition()+", Test: "+input.test_condition(),fontsize=titlefont)
            
            ax.set_axisbelow(True)
            ax.grid(linestyle="--",alpha=0.75)
            ax.tick_params(axis="both",labelsize=axisfont_labels)

            if input.volcano_plotrange_switch()==True:
                ax.set_xlim(input.volcano_xplotrange()[0],input.volcano_xplotrange()[1])
                ax.set_ylim(input.volcano_yplotrange()[0],input.volcano_yplotrange()[1])

            if input.show_labels()==True:
                text_list=[]
                inc_text_x=[]
                inc_text_y=[]
                dec_text_x=[]
                dec_text_y=[]
                for protein in merged.index:
                    if merged.loc[protein]["log2_FoldChange"] >= foldchange_cutoff and np.absolute(merged.loc[protein]["-log10_pvalue"]) >= pvalue_cutoff:
                        text_list.append(protein.split(";")[0])
                        inc_text_x.append(merged.loc[protein]["log2_FoldChange"])
                        inc_text_y.append(merged.loc[protein]["-log10_pvalue"])
                    elif merged.loc[protein]["log2_FoldChange"] <= -foldchange_cutoff and np.absolute(merged.loc[protein]["-log10_pvalue"]) >= pvalue_cutoff:
                        text_list.append(protein.split(";")[0])
                        dec_text_x.append(merged.loc[protein]["log2_FoldChange"])
                        dec_text_y.append(merged.loc[protein]["-log10_pvalue"])
                ta.allocate(ax,x=inc_text_x,y=inc_text_y,text_list=text_list,
                            x_scatter=merged["log2_FoldChange"],y_scatter=merged["-log10_pvalue"],
                            textsize=input.label_fontsize(),linewidth=0.25,linecolor="r",direction="east",
                            max_distance=0.1,
                            avoid_label_lines_overlap=True,avoid_crossing_label_lines=True)
                ta.allocate(ax,x=dec_text_x,y=dec_text_y,text_list=text_list,
                            x_scatter=merged["log2_FoldChange"],y_scatter=merged["-log10_pvalue"],
                            textsize=input.label_fontsize(),linewidth=0.25,linecolor="b",direction="west",
                            max_distance=0.1,
                            avoid_label_lines_overlap=True,avoid_crossing_label_lines=True)
            fig.set_tight_layout(True)
            imagedownload("volcanoplot_"+input.control_condition()+"_vs_"+input.test_condition())

                # for protein in merged.index:
                #     if np.absolute(merged.loc[protein]["log2_FoldChange"]) >= foldchange_cutoff and np.absolute(merged.loc[protein]["-log10_pvalue"]) >= pvalue_cutoff:
                #         ax.annotate(protein,merged.loc[protein]["label"],fontsize=input.label_fontsize())
                #     else:
                #         pass

    # ====================================== Volcano Plot - Feature Plot
    #selectable table for corresponding plot
    @render.data_frame
    def feature_table():
        merged=volcano_calc()
        merged_table=merged[["PG.Genes","PG.ProteinNames","log2_FoldChange","-log10_pvalue","log2_FoldChange_absolute"]].reset_index()
        df_styles={"location":"body",
                   "style":{"column-width":"200px",
                            "overflow":"hidden"
                           }
                  }
        return render.DataGrid(merged_table,width="100%",selection_mode="rows",editable=False,styles=df_styles)
    #box plot for abundances of selected proteins
    @reactive.effect
    def _():
        @render.plot(width=input.volcano_feature_width(),height=input.volcano_feature_height())
        def feature_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            if len(feature_table.data_view(selected=True)["PG.ProteinGroups"].tolist())==0:
                fig,ax=plt.subplots()
            else:
                pickedproteins=feature_table.data_view(selected=True)["PG.ProteinGroups"].tolist()

                list_pairs=[searchoutput[searchoutput["PG.ProteinGroups"]==val][["R.Condition","PG.ProteinGroups","PG.MS2Quantity"]].drop_duplicates().reset_index(drop=True) for val in pickedproteins]
                conditions=list_pairs[0]["R.Condition"].drop_duplicates().tolist()

                colorlist=[]
                legend_patches=[]
                for i in range(len(conditions)):
                    colorlist.append(list(mcolors.TABLEAU_COLORS.keys())[i])
                    legend_patches.append(mpatches.Patch(color=list(mcolors.TABLEAU_COLORS.keys())[i],label=conditions[i]))

                lineprops=dict(linestyle="--",color="black")
                fig,ax=plt.subplots(figsize=(input.volcano_feature_width()*(1/plt.rcParams['figure.dpi']),input.volcano_feature_height()*(1/plt.rcParams['figure.dpi'])))
                positions=list(np.arange(len(conditions)))
                for i in range(len(list_pairs)):
                    plottinglist=[]
                    #make a list of lists of the signals for the given protein in the given condition
                    for condition in conditions:
                        plottinglist.append(list_pairs[i].groupby("R.Condition").get_group(condition)["PG.MS2Quantity"].tolist())
                    bplot=ax.boxplot(plottinglist,positions=positions,widths=0.75,patch_artist=True,medianprops=lineprops)
                    #change box color corresponding to condition
                    for x in range(len(bplot["boxes"])):
                        color=colorlist[x]
                        for item in ["boxes"]:
                            plt.setp(bplot[item][x],color=color)
                    #adjust positions for the next places to plot
                    for y in range(len(positions)):
                        positions[y]+=len(positions)

                xticks=[]
                x=1/(len(conditions))
                for i in range(len(list_pairs)):
                    xticks.append(x)
                    ax.axvline(x=x+1,color="lightgrey",linestyle="--")
                    x+=len(conditions)
                pickedproteins_shortened=[]
                for protein in pickedproteins:
                    if protein.count(";")>=1:
                        pickedproteins_shortened.append(protein.split(";")[0])
                    else:
                        pickedproteins_shortened.append(protein)
                ax.set_xticks(xticks)
                ax.set_xticklabels(pickedproteins_shortened)
                ax.legend(handles=legend_patches)

            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()

            ax.set_ylabel("Protein Group Intensity",fontsize=axisfont)
            ax.set_xlabel("Protein Group",fontsize=axisfont)
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            #ax.grid(axis="y",linestyle="--")
            fig.set_tight_layout(True)
            if len(feature_table.data_view(selected=True)["PG.ProteinGroups"].tolist())!=0:
                imagedownload("volcano_featureplot")

    # ====================================== Volcano Plot - Up/Down Regulation
    #volcano plot up/down regulation protein list
    @reactive.effect
    def _():
        @render.plot(width=input.volcano_regulation_width(),height=input.volcano_regulation_height())
        def volcano_updownregulation_plot():
            merged=volcano_calc()
            plot_up_or_down=input.regulation_upordown()
            top_n=int(input.regulation_topN())
            foldchange_cutoff=input.volcano_foldchange()
            pvalue_cutoff=input.volcano_pvalue()
            plotproperty=input.regulation_p_fold()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()

            if plot_up_or_down=="up":
                if plotproperty=="pvalue":
                    merged_sort=merged[merged["log2_FoldChange"]>=foldchange_cutoff]["-log10_pvalue"].dropna().sort_values(axis=0,ascending=False).reset_index()
                    ylabel="-log10 pvalue"
                    plotcolumn="-log10_pvalue"
                if plotproperty=="foldchange":
                    merged_sort=merged[merged["-log10_pvalue"]>=pvalue_cutoff]["log2_FoldChange"].dropna().sort_values(axis=0,ascending=False).reset_index()
                    ylabel="log2 Fold Change"
                    plotcolumn="log2_FoldChange"
                color="r"
                title="Upregulated Protiens"
            if plot_up_or_down=="down":
                if plotproperty=="pvalue":
                    merged_sort=merged[merged["log2_FoldChange"]<=-foldchange_cutoff]["-log10_pvalue"].dropna().sort_values(axis=0,ascending=False).reset_index()
                    ylabel="-log10 pvalue"
                    plotcolumn="-log10_pvalue"
                if plotproperty=="foldchange":
                    merged_sort=merged[merged["-log10_pvalue"]>=pvalue_cutoff]["log2_FoldChange"].dropna().sort_values(axis=0,ascending=True).reset_index()
                    ylabel="log2 Fold Change"
                    plotcolumn="log2_FoldChange"
                color="b"
                title="Downregulated Proteins"

            proteinlist=merged_sort["PG.ProteinGroups"].str.split(";").tolist()
            proteinlist_simplified=[]
            for protein in proteinlist:
                proteinlist_simplified.append(protein[0])
            fig,ax=plt.subplots(figsize=(input.volcano_regulation_width()*(1/plt.rcParams['figure.dpi']),input.volcano_regulation_height()*(1/plt.rcParams['figure.dpi'])))
            y=np.flip(np.arange(len(merged_sort)))
            ax.barh(y[:top_n],merged_sort[plotcolumn][:top_n],color=color)
            ax.set_yticks(y[:top_n],labels=proteinlist_simplified[:top_n])
            ax.set_xlabel(ylabel,fontsize=axisfont)
            ax.set_ylabel("Protein Group Name",fontsize=axisfont)
            ax.set_title(title,fontsize=titlefont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            ax.margins(0.02)
            fig.set_tight_layout(True)
            imagedownload("volcano_regulationplot_"+plot_up_or_down+"regulated")

    # ====================================== Correlations
    @render.ui
    def correlations_sampleconditions_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=sampleconditions
        return ui.input_selectize("correlations_sampleconditions","Pick sample condition:",choices=opts)
    #plot correlations of log10(FG.MS2Quantity) between replicates of the same conditions
    @reactive.effect
    def _():
        @render.plot(width=input.correlations_width(),height=input.correlations_height())
        def correlations_plot():
            from itertools import permutations
            
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            condition=input.correlations_sampleconditions()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()

            reps=maxreplicatelist[sampleconditions.index(condition)]
            gridrange=np.arange(reps)

            rep_permutations=list(permutations(gridrange,2))

            #generate unique cells in the matrix to plot in, list of cells to ignore
            to_keep=[]
            for combo in rep_permutations:
                if combo[0]>combo[1]:
                    to_keep.append(combo)
            ignore_plotcoords=[[i for i in coord] for coord in rep_permutations if coord not in to_keep]
            rep_permutations_plot=[[i for i in coord] for coord in to_keep]

            fig,axes=plt.subplots(sharex=True,nrows=reps,ncols=reps)
            for i in range(len(rep_permutations_plot)):
                run1=condition+"_"+str(rep_permutations_plot[i][0]+1)
                run2=condition+"_"+str(rep_permutations_plot[i][1]+1)

                r1=searchoutput[searchoutput["Cond_Rep"]==run1][["Cond_Rep","EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]].drop_duplicates().sort_values("EG.ModifiedPeptide").reset_index(drop=True).drop(columns="Cond_Rep")
                r2=searchoutput[searchoutput["Cond_Rep"]==run2][["Cond_Rep","EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]].drop_duplicates().sort_values("EG.ModifiedPeptide").reset_index(drop=True).drop(columns="Cond_Rep")
                merged=r1.merge(r2,on=["EG.ModifiedPeptide","FG.Charge"])

                x=np.log10(merged["FG.MS2Quantity_x"])
                y=np.log10(merged["FG.MS2Quantity_y"])

                slope,intercept,r_value,p_value,std_err=scipy.stats.linregress(x,y)
                z=np.polyfit(x,y,1)
                p=np.poly1d(z)

                axis=rep_permutations_plot[i]
                axes[axis[0],axis[1]].scatter(x,y,s=2)
                axes[axis[0],axis[1]].plot(x,p(x),"tab:orange",linestyle=":",linewidth=2)

                ymin,ymax=axes[axis[1],axis[0]].get_ylim()
                xmin,xmax=axes[axis[1],axis[0]].get_xlim()
                axes[axis[1],axis[0]].text(x=0.25*xmax,y=0.5*ymax,s="R$^2$="+str(round(r_value,3)),fontsize=labelfont)

            for i in range(reps):
                run=condition+"_"+str(i+1)
                x=searchoutput[searchoutput["Cond_Rep"]==run][["EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]].drop_duplicates().reset_index(drop=True)
                axes[i,i].hist(np.log10(x["FG.MS2Quantity"]),bins=50,density=True)
                
            for i in range(len(ignore_plotcoords)):
                x=ignore_plotcoords[i][0]
                y=ignore_plotcoords[i][1]
                axes[x,y].set_yticks([])

            runs=searchoutput[searchoutput["Cond_Rep"].str.contains(condition)]["Cond_Rep"].drop_duplicates().tolist()
            rows=runs
            for ax,run in zip(axes[0],runs):
                ax.set_title(run,fontsize=titlefont)
            for ax,row in zip(axes[:,0],rows):
                ax.set_ylabel(row,rotation=90,fontsize=axisfont)
            fig.align_ylabels(axes[:])
            fig.set_tight_layout(True)

    # ====================================== Dendrogram/Protein Signal
    @reactive.effect
    def _():
        @render.plot(width=input.dendrogram_width(),height=input.dendrogram_height())
        def dendrogram_heatmap():
            from scipy.cluster.hierarchy import dendrogram, linkage

            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            df=searchoutput[["Cond_Rep","PG.ProteinGroups","PG.MS2Quantity"]].drop_duplicates().reset_index(drop=True)

            data=pd.DataFrame()
            for run in runnames:
                placeholder=df[df["Cond_Rep"]==run][["PG.ProteinGroups","PG.MS2Quantity"]].set_index("PG.ProteinGroups").sort_values("PG.MS2Quantity",ascending=False)
                placeholder["PG.MS2Quantity"]=np.log10(placeholder["PG.MS2Quantity"])
                data=pd.concat([data,placeholder],axis=1).rename(columns={"PG.MS2Quantity":run})
            data=data.fillna(0)
            data_transposed=data.transpose()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            #make dendrogram linking the runs
            clusteringlist=[]
            for column in data.columns:
                clusteringlist.append(data[column].tolist())
            fig,ax=plt.subplots(nrows=2,ncols=2,height_ratios=(1,input.dendrogram_scaling()),width_ratios=(5,1),figsize=(input.dendrogram_width()*(1/plt.rcParams['figure.dpi']),input.dendrogram_height()*(1/plt.rcParams['figure.dpi'])))
            datalinkage=linkage(clusteringlist)
            dendrogram(datalinkage,labels=data.columns,ax=ax[0,0])

            ax[0,0].set_ylabel("Distance",fontsize=axisfont)
            ax[0,0].set_title("Run Linkage",fontsize=titlefont)
            ax[0,0].tick_params(axis="both",labelsize=axisfont_labels)
            #ax[0].tick_params(axis="x",rotation=90)

            #get tick labels in the order they're shown in the dendrogram
            ticklabels=ax[0,0].get_xticklabels(which="major")
            labels=[]
            for ele in ticklabels:
                labels.append(str(ele).split(", ")[2].split(")")[0].strip("''"))
            ax[0,0].xaxis.set_ticklabels([])

            #map the dendrogram labels to the df used for the heatmap
            data_transposed=data_transposed.reindex(labels)

            plottinglist=[]
            for column in data_transposed.columns:
                plottinglist.append(data_transposed[column].tolist())
            heatmap=ax[1,0].imshow(plottinglist,cmap=input.dendrogram_cmap(),aspect="auto")
            ax[1,0].set_xlabel("Run",fontsize=axisfont)
            ax[1,0].set_ylabel("Protein",fontsize=axisfont)           
            ax[1,0].set_title("Protein Intensity Heatmap (log10)",fontsize=titlefont)
            ax[1,0].tick_params(axis="x",rotation=x_label_rotation)
            ax[1,0].set_xticks(np.arange(0,6),labels)
            ax[1,0].tick_params(axis="both",labelsize=axisfont_labels)
            plt.colorbar(heatmap,ax=ax[1,1],location="left")

            fig.set_tight_layout(True)
            ax[0,1].remove()
            ax[1,1].axes.set_axis_off()
            imagedownload("dendrogram")

            # #dendrogram linking proteins
            # datalinkage=linkage(plottinglist)
            # fig,ax=plt.subplots(figsize=(10,20))
            # dendrogram(datalinkage,ax=ax,orientation="right")#,truncate_mode="lastp",p=100)
            # ax.set_xlabel("Distance")
            # ax.set_ylabel("Protein Group")
            # ax.set_title("Protein Group Linkage")

    # ====================================== PCA
    @render.ui
    def pca_3d_ui():
        if input.pca_3d_switch()==True:
            return ui.row(ui.input_slider("pca_azimuth","Azimuth angle",min=0,max=360,value=0,step=5,ticks=True),ui.input_slider("pca_elevation","Elevation angle",min=0,max=90,value=0,step=5,ticks=True))
    #compute PCA and plot principal components
    @reactive.effect
    def _():
        @render.plot(width=input.pca_width(),height=input.pca_height())
        def pca_plot():
            from sklearn.decomposition import PCA
            from sklearn.pipeline import Pipeline
            from sklearn.preprocessing import StandardScaler

            #https://www.youtube.com/watch?v=WPRysPAhG5Q&ab_channel=CompuFlair
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()

            samplelist=runnames
            intensitydict=dict()
            for run in samplelist:
                intensitydict[run]=searchoutput[searchoutput["Cond_Rep"]==run][["PG.ProteinGroups","PG.MS2Quantity"]].drop_duplicates().reset_index(drop=True).rename(columns={"PG.MS2Quantity":run}).set_index(keys="PG.ProteinGroups")
            df_intensity=[]
            for run in intensitydict.keys():
                temp_df=intensitydict[run]
                df_intensity.append(temp_df)
            concatenated_proteins=pd.concat(df_intensity,axis=1).dropna()

            X=np.array(concatenated_proteins).T
            pip=Pipeline([("scaler",StandardScaler()),("pca",PCA())]).fit(X)
            X_trans=pip.transform(X)
            #each row is a sample, each element of each row is a principal component

            colors=colorpicker()

            firstindex=0
            secondindex=0

            if input.pca_3d_switch()==True:
                fig=plt.figure(figsize=(input.pca_width()*(1/plt.rcParams['figure.dpi']),input.pca_height()*(1/plt.rcParams['figure.dpi'])))
                ax1=fig.add_subplot(1,3,(1,2),projection='3d')
                ax2=fig.add_subplot(1,3,3)
                for i in range(numconditions):
                    if i==0:
                        firstindex=0
                    else:
                        firstindex+=maxreplicatelist[i-1]
                    secondindex=firstindex+maxreplicatelist[i]
                    if numconditions==1:
                        ax1.scatter(X_trans[firstindex:secondindex,0],X_trans[firstindex:secondindex,1],X_trans[firstindex:secondindex,2],label=sampleconditions[i],color=colors)
                    else:
                        ax1.scatter(X_trans[firstindex:secondindex,0],X_trans[firstindex:secondindex,1],X_trans[firstindex:secondindex,2],label=sampleconditions[i],color=colors[i])
                ax1.set_zlabel("PC3"+" ("+str(round(pip.named_steps.pca.explained_variance_ratio_[2]*100,1))+"%)",fontsize=axisfont)
                ax1.view_init(azim=input.pca_azimuth(),elev=input.pca_elevation())

            else:
                fig,ax=plt.subplots(ncols=2,gridspec_kw={"width_ratios":[10,5]},figsize=(input.pca_width()*(1/plt.rcParams['figure.dpi']),input.pca_height()*(1/plt.rcParams['figure.dpi'])))
                ax1=ax[0]
                ax2=ax[1]
                for i in range(numconditions):
                    if i==0:
                        firstindex=0
                    else:
                        firstindex+=maxreplicatelist[i-1]
                    secondindex=firstindex+maxreplicatelist[i]
                    if numconditions==1:
                        ax1.scatter(X_trans[firstindex:secondindex,0],X_trans[firstindex:secondindex,1],label=sampleconditions[i],color=colors)
                    else:
                        ax1.scatter(X_trans[firstindex:secondindex,0],X_trans[firstindex:secondindex,1],label=sampleconditions[i],color=colors[i])

                ax1.legend(loc="upper left",bbox_to_anchor=[0,-0.1],fontsize=legendfont)
                ax1.spines['bottom'].set_position('zero')
                ax1.spines['top'].set_color('none')
                ax1.spines['right'].set_color('none')
                ax1.spines['left'].set_position('zero')
                ax1.set_axisbelow(True)
                ax1.grid(linestyle="--")
                ax1.tick_params(axis="both",labelsize=axisfont_labels)
                ax1.xaxis.set_label_coords(x=0.5,y=-0.02)
                ax1.yaxis.set_label_coords(x=-0.02,y=0.45)

            ax1.set_xlabel("PC1"+" ("+str(round(pip.named_steps.pca.explained_variance_ratio_[0]*100,1))+"%)",fontsize=axisfont)
            ax1.set_ylabel("PC2"+" ("+str(round(pip.named_steps.pca.explained_variance_ratio_[1]*100,1))+"%)",fontsize=axisfont)

            #ax2.bar(np.arange(1,len(pip.named_steps.pca.explained_variance_ratio_)+1),pip.named_steps.pca.explained_variance_ratio_*100,edgecolor="k")
            ax2.bar(np.arange(1,5),pip.named_steps.pca.explained_variance_ratio_[:4]*100,edgecolor="k")
            ax2.set_xlabel("Principal Component",fontsize=axisfont)
            ax2.set_ylabel("Total % Variance Explained",fontsize=axisfont)
            #ax2.set_xticks(np.arange(1,len(pip.named_steps.pca.explained_variance_ratio_)+1))
            ax2.set_xticks(np.arange(1,5))
            ax2.set_axisbelow(True)
            ax2.grid(linestyle="--")
            ax2.tick_params(axis="both",labelsize=axisfont_labels)

            fig.set_tight_layout(True)
            imagedownload("pca")

#endregion

# ============================================================================= Immunopeptidomics
#region
    # ====================================== Sequence Motifs
    @render.ui
    def seqmotif_run_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_selectize("seqmotif_run_pick","Pick run:",choices=opts)
    #sequence motif plot
    @reactive.effect
    def _():
        @render.plot(width=input.seqmotif_width(),height=input.seqmotif_height())
        def seqmotif_plot():
            import logomaker as lm
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            seq_df=searchoutput[searchoutput["Cond_Rep"]==input.seqmotif_run_pick()][["PEP.StrippedSequence","Peptide Length"]].drop_duplicates()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()

            seq=seq_df[seq_df["Peptide Length"]==input.seqmotif_peplengths()].drop(columns=["Peptide Length"])["PEP.StrippedSequence"].tolist()

            matrix=lm.alignment_to_matrix(seq)
            if input.seqmotif_plottype()=="counts":
                ylabel="Counts"
            if input.seqmotif_plottype()=="information":
                matrix=lm.transform_matrix(matrix,from_type="counts",to_type="information")
                ylabel="Information (bits)"
            fig,ax=plt.subplots(figsize=(input.seqmotif_width()*(1/plt.rcParams['figure.dpi']),input.seqmotif_height()*(1/plt.rcParams['figure.dpi'])))
            logo=lm.Logo(matrix,ax=ax,vsep=0.01,vpad=0.01,color_scheme="weblogo_protein")
            ax.set_xlabel("Position",fontsize=axisfont)
            ax.set_ylabel(ylabel,fontsize=axisfont)
            ax.set_title(input.seqmotif_run_pick()+": "+str(input.seqmotif_peplengths())+"mers",fontsize=titlefont)
            ax.set_ylim(bottom=0)
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            logo.fig.tight_layout()
            imagedownload("sequencemotifs_"+input.seqmotif_run_pick()+"_"+str(input.seqmotif_peplengths())+"mers")

    # ====================================== Charge States (Bar)
    @reactive.calc
    def ipep_charge_peplength():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        fillvalue="0"

        #run function by individual run
        dict_chargecountdf_run=dict()
        dict_peplengthcountdf_run=dict()

        for run in runnames:
            df=searchoutput[searchoutput["Cond_Rep"]==run][["EG.ModifiedPeptide","FG.Charge","EG.ApexRT"]]
            charges=sorted(df["FG.Charge"].drop_duplicates(ignore_index=True))
            df_pivot=pd.pivot_table(df,index=["EG.ModifiedPeptide"],columns=["FG.Charge"],fill_value=fillvalue)

            chargedf=pd.DataFrame()
            chargedf["EG.ModifiedPeptide"]=df_pivot.index
            chargedf.set_index(["EG.ModifiedPeptide"],inplace=True)

            for charge in charges:
                #pivot table of just peptides of specific charge
                df_pivot_column=pd.DataFrame(df_pivot[('EG.ApexRT', charge)][df_pivot[('EG.ApexRT', charge)]!=fillvalue])
                #make an extra column of just the charge value
                df_pivot_column["Charge"]=[str(charge)]*len(df_pivot_column)
                #make an extra column in our final df to update from the initial pivot table
                chargedf["Charge"]=[fillvalue]*len(chargedf)
                #update the final df from the pivot table along same indices
                chargedf.update(df_pivot_column["Charge"])
                #change the name of the column so we can make the next one in the next loop
                chargedf=chargedf.rename(columns={"Charge":charge})

            #get list of charges present for peptides as a list
            chargelist=[]
            for pep in chargedf.index:
                listperpep=[]
                for item in chargedf.loc[pep].tolist():
                    if item!="0":
                        listperpep.append(item)
                    else:
                        pass
                chargelist.append(listperpep)
            #unpack each list as strings for better visualization and downstream use
            unpackedcharges=[]
            for ele in chargelist:
                if len(ele)==1:
                    unpackedcharges.append(ele[0])
                else:
                    placeholder=""
                    for i,string in enumerate(ele):
                        if i==0:
                            placeholder=placeholder+string
                        else:
                            placeholder=placeholder+","+string
                    unpackedcharges.append(placeholder)

            chargedf["list"]=unpackedcharges
            chargedf=chargedf.reset_index()
            chargegroups=sorted(chargedf["list"].drop_duplicates())

            #add in the StrippedSequence column
            strippedseq=searchoutput[searchoutput["Cond_Rep"]==run][["EG.ModifiedPeptide","PEP.StrippedSequence"]].drop_duplicates().sort_values("EG.ModifiedPeptide")
            chargedf["PEP.StrippedSequence"]=strippedseq["PEP.StrippedSequence"].reset_index(drop=True)

            #calculate peptide lengths
            lengthlist=[]
            for pep in chargedf["PEP.StrippedSequence"]:
                lengthlist.append(len(pep))
            chargedf["Peptide Length"]=lengthlist

            #df for plotting charge states only
            chargecountdf=pd.DataFrame(chargedf["list"].value_counts()).reset_index().sort_values("list").reset_index(drop=True)
            totals=sum(chargecountdf["count"])
            frequencies=[]
            for ele in chargecountdf["count"]:
                frequencies.append(round((ele/totals)*100,5))
            chargecountdf["frequency%"]=frequencies
            
            #df for plotting intersection of peptide length and charge state
            charge_length_df=chargedf[["list","Peptide Length"]]
            peplengths=list(set(charge_length_df["Peptide Length"]))
            peplengthcountdf=pd.DataFrame()
            peplengthcountdf["Peptide Length"]=peplengths
            for charge in chargecountdf["list"]:
                countpercharge=[]
                for peplength in peplengths:
                    val_count=charge_length_df[charge_length_df["Peptide Length"]==peplength]["list"].value_counts().get(charge)
                    if val_count is None:
                        countpercharge.append(0)
                    else:
                        countpercharge.append(val_count)
                peplengthcountdf[charge]=countpercharge
            
            #dump the chargedf and chargecountdf generated in the loop into a dict
            dict_chargecountdf_run[run]=chargecountdf
            dict_peplengthcountdf_run[run]=peplengthcountdf
            
        #run function by condition
        dict_chargecountdf_condition=dict()
        dict_peplengthcountdf_condition=dict()

        for run in sampleconditions:
            df=searchoutput[searchoutput["R.Condition"]==run][["EG.ModifiedPeptide","FG.Charge","EG.ApexRT"]]
            charges=sorted(df["FG.Charge"].drop_duplicates(ignore_index=True))
            df_pivot=pd.pivot_table(df,index=["EG.ModifiedPeptide"],columns=["FG.Charge"],fill_value=fillvalue)

            chargedf=pd.DataFrame()
            chargedf["EG.ModifiedPeptide"]=df_pivot.index
            chargedf.set_index(["EG.ModifiedPeptide"],inplace=True)

            for charge in charges:
                #pivot table of just peptides of specific charge
                df_pivot_column=pd.DataFrame(df_pivot[('EG.ApexRT', charge)][df_pivot[('EG.ApexRT', charge)]!=fillvalue])
                #make an extra column of just the charge value
                df_pivot_column["Charge"]=[str(charge)]*len(df_pivot_column)
                #make an extra column in our final df to update from the initial pivot table
                chargedf["Charge"]=[fillvalue]*len(chargedf)
                #update the final df from the pivot table along same indices
                chargedf.update(df_pivot_column["Charge"])
                #change the name of the column so we can make the next one in the next loop
                chargedf=chargedf.rename(columns={"Charge":charge})

            #get list of charges present for peptides as a list
            chargelist=[]
            for pep in chargedf.index:
                listperpep=[]
                for item in chargedf.loc[pep].tolist():
                    if item!="0":
                        listperpep.append(item)
                    else:
                        pass
                chargelist.append(listperpep)
            #unpack each list as strings for better visualization and downstream use
            unpackedcharges=[]
            for ele in chargelist:
                if len(ele)==1:
                    unpackedcharges.append(ele[0])
                else:
                    placeholder=""
                    for i,string in enumerate(ele):
                        if i==0:
                            placeholder=placeholder+string
                        else:
                            placeholder=placeholder+","+string
                    unpackedcharges.append(placeholder)

            chargedf["list"]=unpackedcharges
            chargedf=chargedf.reset_index()
            chargegroups=sorted(chargedf["list"].drop_duplicates())

            #add in the StrippedSequence column
            strippedseq=searchoutput[searchoutput["R.Condition"]==run][["EG.ModifiedPeptide","PEP.StrippedSequence"]].drop_duplicates().sort_values("EG.ModifiedPeptide")
            chargedf["PEP.StrippedSequence"]=strippedseq["PEP.StrippedSequence"].reset_index(drop=True)
            #calculate peptide lengths
            lengthlist=[]
            for pep in chargedf["PEP.StrippedSequence"]:
                lengthlist.append(len(pep))
            chargedf["Peptide Length"]=lengthlist

            #df for plotting charge states only
            chargecountdf=pd.DataFrame(chargedf["list"].value_counts()).reset_index().sort_values("list").reset_index(drop=True)
            totals=sum(chargecountdf["count"])
            frequencies=[]
            for ele in chargecountdf["count"]:
                frequencies.append(round((ele/totals)*100,5))
            chargecountdf["frequency%"]=frequencies
            
            #df for plotting intersection of peptide length and charge state
            charge_length_df=chargedf[["list","Peptide Length"]]
            peplengths=list(set(charge_length_df["Peptide Length"]))
            peplengthcountdf=pd.DataFrame()
            peplengthcountdf["Peptide Length"]=peplengths
            for charge in chargecountdf["list"]:
                countpercharge=[]
                for peplength in peplengths:
                    val_count=charge_length_df[charge_length_df["Peptide Length"]==peplength]["list"].value_counts().get(charge)
                    if val_count is None:
                        countpercharge.append(0)
                    else:
                        countpercharge.append(val_count)
                peplengthcountdf[charge]=countpercharge

            #dump the chargedf and chargecountdf generated in the loop into a dict
            dict_chargecountdf_condition[run]=chargecountdf
            dict_peplengthcountdf_condition[run]=peplengthcountdf

        return dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition
    @render.ui
    def chargestate_charges_ui():
        dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition=ipep_charge_peplength()
        if input.chargestate_peplength_condition_or_run()=="condition":
            plottingdf=dict_peplengthcountdf_condition
        if input.chargestate_peplength_condition_or_run()=="individual":
            plottingdf=dict_peplengthcountdf_run
        chargelist=[]
        for key in plottingdf.keys():
            chargelist.append(plottingdf[key].columns.tolist())
        chargelist=sorted(list(set(itertools.chain(*chargelist))))
        chargelist.remove("Peptide Length")

        return ui.input_checkbox_group("chargestate_charges","Charges to plot:",choices=chargelist)
    #bar chart of charge state counts
    @reactive.effect
    def _():
        @render.plot(width=input.charge_barchart_width(),height=input.charge_barchart_height())
        def charge_barchart():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition=ipep_charge_peplength()
            if input.chargestate_bar_condition_or_run()=="condition":
                plottingdf=dict_chargecountdf_condition
                colors=colorpicker()
            if input.chargestate_bar_condition_or_run()=="individual":
                plottingdf=dict_chargecountdf_run
                colors=replicatecolors()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            x_label_rotation=input.xaxis_label_rotation()

            if input.chargestate_charges_usepickedcharges()==True:
                chargelist=list(input.chargestate_charges())
                plottingdf_picked=dict()
                for key in plottingdf.keys():
                    plottingdf_picked[key]=plottingdf[key].set_index("list").loc[chargelist].reset_index()
                plottingdf=plottingdf_picked

            if len(plottingdf)==1:
                fig,ax=plt.subplots(figsize=(input.charge_barchart_width()*(1/plt.rcParams['figure.dpi']),input.charge_barchart_height()*(1/plt.rcParams['figure.dpi'])))
                key=list(plottingdf.keys())[0]
                x=np.arange(1,len(plottingdf[key]["list"])+1)
                ax.bar(x,plottingdf[key]["count"],edgecolor="k")
                ax.set_xticks(x,plottingdf[key]["list"],rotation=x_label_rotation)
                ax.set_ylabel("Counts (%)",fontsize=axisfont)
                ax.set_xlabel("Charge(s)",fontsize=axisfont)
                ax.set_title(key,fontsize=titlefont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)

            else:
                fig,ax=plt.subplots(ncols=len(plottingdf),figsize=(input.charge_barchart_width()*(1/plt.rcParams['figure.dpi']),input.charge_barchart_height()*(1/plt.rcParams['figure.dpi'])))

                for i,key in enumerate(plottingdf.keys()):
                    x=np.arange(1,len(plottingdf[key]["list"])+1)
                    if numconditions==1:
                        ax[i].bar(x,plottingdf[key]["count"],edgecolor="k",color=colorpicker())
                    else:
                        ax[i].bar(x,plottingdf[key]["count"],edgecolor="k",color=colors[i])
                    ax[i].set_xticks(x,plottingdf[key]["list"],rotation=x_label_rotation)
                    ax[i].set_title(key,fontsize=titlefont)
                    ax[i].set_xlabel("Charge(s)",fontsize=axisfont)
                    ax[i].set_axisbelow(True)
                    ax[i].grid(linestyle="--")
                    ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                ax[0].set_ylabel("Counts",fontsize=axisfont)
            fig.set_tight_layout(True)
            imagedownload("chargestate_barchart")

    # ====================================== Charge State (Stacked)
    #stacked bar chart of charge frequencies
    @reactive.effect
    def _():
        @render.plot(width=input.charge_stackedbarchart_width(),height=input.charge_stackedbarchart_height())
        def charge_stacked_barchart():
            dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition=ipep_charge_peplength()
            if input.chargestate_stacked_condition_or_run()=="condition":
                plottingdf=dict_chargecountdf_condition
            if input.chargestate_stacked_condition_or_run()=="individual":
                plottingdf=dict_chargecountdf_run

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()
            x_label_rotation=input.xaxis_label_rotation()

            fig,ax=plt.subplots(figsize=(input.charge_stackedbarchart_width()*(1/plt.rcParams['figure.dpi']),input.charge_stackedbarchart_height()*(1/plt.rcParams['figure.dpi'])))
            for i,key in enumerate(plottingdf.keys()):
                bottom=np.zeros(1)
                plottingdf[key]=plottingdf[key].sort_values("frequency%",ascending=False).reset_index(drop=True)
                for x in range(len(plottingdf[key])):
                    #generate a color list based on the length of each df in the dict
                    matplottabcolors=list(mcolors.TABLEAU_COLORS)
                    plotcolors=[]
                    if len(plottingdf[key]) > len(matplottabcolors):
                        dif=len(plottingdf[key])-len(matplottabcolors)
                        plotcolors=matplottabcolors
                        for y in range(dif):
                            plotcolors.append(matplottabcolors[y])
                    else:
                        plotcolors=matplottabcolors
                    ax.bar(i,plottingdf[key]["frequency%"][x],bottom=bottom,label=plottingdf[key]["list"][x],color=plotcolors[x])
                    bottom+=plottingdf[key]["frequency%"][x]
            ax.set_xticks(np.arange(0,len(plottingdf)),list(plottingdf.keys()),rotation=x_label_rotation)
            ax.set_ylabel("Frequency (%)",fontsize=axisfont)
            ax.set_xlabel("Condition",fontsize=axisfont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            handles,labels=plt.gca().get_legend_handles_labels()
            by_label=OrderedDict(zip(labels,handles))
            ax.legend(by_label.values(), by_label.keys(),loc="center",bbox_to_anchor=(1.1,0.5),prop={'size':legendfont})
            fig.set_tight_layout(True)
            imagedownload("chargestate_stackedchart")
    #table of frequencies for the stacked bar chart
    @render.data_frame
    def charge_stacked_table():
        dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition=ipep_charge_peplength()
        if input.chargestate_stacked_condition_or_run()=="condition":
            plottingdf=dict_chargecountdf_condition
        if input.chargestate_stacked_condition_or_run()=="individual":
            plottingdf=dict_chargecountdf_run
        displaydf=pd.DataFrame()
        for i,key in enumerate(plottingdf.keys()):
            displaydf=pd.concat([displaydf,plottingdf[key][["list","frequency%"]].sort_values("frequency%",ascending=False).rename(columns={"list":key,"frequency%":"frequency%_"+str(i)})],axis=1)

        return render.DataGrid(displaydf,editable=False)

    # ====================================== Charge State/Peptide Length
    @render.ui
    def chargestate_peplength_plotrange_ui():
        dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition=ipep_charge_peplength()
        if input.chargestate_peplength_condition_or_run()=="condition":
            plottingdf=dict_peplengthcountdf_condition
        if input.chargestate_peplength_condition_or_run()=="individual":
            plottingdf=dict_peplengthcountdf_run

        min=7
        max=30
        
        return ui.input_slider("chargestate_peplength_plotrange","Plot Range",min=1,max=50,value=[min,max],step=1,ticks=True,width="300px",drag_range=True)
    @render.ui
    def chargestate_peplength_charges_ui():
        dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition=ipep_charge_peplength()
        if input.chargestate_peplength_condition_or_run()=="condition":
            plottingdf=dict_peplengthcountdf_condition
        if input.chargestate_peplength_condition_or_run()=="individual":
            plottingdf=dict_peplengthcountdf_run
        chargelist=[]
        for key in plottingdf.keys():
            chargelist.append(plottingdf[key].columns.tolist())
        chargelist=sorted(list(set(itertools.chain(*chargelist))))
        chargelist.remove("Peptide Length")

        return ui.input_checkbox_group("chargestate_peplength_charges","Charges to plot:",choices=chargelist)
    @render.ui
    def chargestate_peplength_download_ui():
        if input.chargestate_peplength_condition_or_run()=="condition":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=sampleconditions
            return ui.input_selectize("chargestate_peplength_download_pick","Pick condition for download:",choices=opts,width="300px")
        if input.chargestate_peplength_condition_or_run()=="individual":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=runnames
            return ui.input_selectize("chargestate_peplength_download_pick","Pick run for download:",choices=opts,width="300px")
    @render.download(filename=lambda: f"{input.chargestate_peplength_download_pick()}_chargestates-peptidelengths.csv")
    def chargestate_peplength_download():
        dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition=ipep_charge_peplength()
        if input.chargestate_peplength_condition_or_run()=="condition":
            plottingdf=dict_peplengthcountdf_condition
        if input.chargestate_peplength_condition_or_run()=="individual":
            plottingdf=dict_peplengthcountdf_run
        exportdf=plottingdf[input.chargestate_peplength_download_pick()]

        with io.BytesIO() as buf:
            exportdf.to_csv(buf,index=False)
            yield buf.getvalue()
    #plot charge states per peptide length
    @reactive.effect
    def _():
        @render.plot(width=input.chargestate_peplength_width(),height=input.chargestate_peplength_height())
        def chargestate_peplength_plot():
            dict_chargecountdf_run,dict_peplengthcountdf_run,dict_chargecountdf_condition,dict_peplengthcountdf_condition=ipep_charge_peplength()
            if input.chargestate_peplength_condition_or_run()=="condition":
                plottingdf=dict_peplengthcountdf_condition
            if input.chargestate_peplength_condition_or_run()=="individual":
                plottingdf=dict_peplengthcountdf_run

            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            titlefont=input.titlefont()
            legendfont=input.legendfont()
            y_padding=0.05

            if len(plottingdf)==1:
                fig,ax=plt.subplots(figsize=(input.chargestate_peplength_width()*(1/plt.rcParams['figure.dpi']),input.chargestate_peplength_height()*(1/plt.rcParams['figure.dpi'])))
                key=list(plottingdf.keys())[0]
                plotdf=plottingdf[key].set_index("Peptide Length")
                if input.usepickedcharges()==True:
                    columns=input.chargestate_peplength_charges()
                else:
                    columns=plotdf.columns.tolist()
                x=plotdf.index.tolist()
                bottom=np.zeros(len(plotdf))
                for col in columns:
                    ax.bar(x,plotdf[col],label=col,bottom=bottom)
                    bottom+=plotdf[col].tolist()
                ax.set_xlabel("Peptide Length",fontsize=axisfont)
                ax.set_title(key,fontsize=titlefont)
                ax.legend(loc="upper right",prop={'size':legendfont})
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.set_ylabel("Counts",fontsize=axisfont)
                ax.set_xlim(left=input.chargestate_peplength_plotrange()[0]-0.75,right=input.chargestate_peplength_plotrange()[1]+0.5)
                ax.set_ylim(top=ax.get_ylim()[1]+(y_padding*ax.get_ylim()[1]))
                ax.xaxis.set_major_locator(MaxNLocator(integer=True))
                ax.xaxis.set_minor_locator(MultipleLocator(1))
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                fig.set_tight_layout(True)

            else:
                fig,ax=plt.subplots(ncols=len(plottingdf),sharey=True,figsize=(input.chargestate_peplength_width()*(1/plt.rcParams['figure.dpi']),input.chargestate_peplength_height()*(1/plt.rcParams['figure.dpi'])))
                maxlist=[]
                for i,key in enumerate(plottingdf.keys()):
                    plotdf=plottingdf[key].set_index("Peptide Length")
                    maxlist.append(max(plotdf.max().tolist()))
                    if input.usepickedcharges()==True:
                        columns=input.chargestate_peplength_charges()
                    else:
                        columns=plotdf.columns.tolist()
                    x=plotdf.index.tolist()
                    bottom=np.zeros(len(plotdf))
                    for col in columns:
                        ax[i].bar(x,plotdf[col],label=col,bottom=bottom)
                        bottom+=plotdf[col].tolist()

                    ax[i].set_xlabel("Peptide Length",fontsize=axisfont)
                    ax[i].set_title(key,fontsize=titlefont)
                    ax[i].legend(loc="upper right",prop={'size':legendfont})
                    ax[i].set_axisbelow(True)
                    ax[i].grid(linestyle="--")
                    ax[i].set_xlim(left=input.chargestate_peplength_plotrange()[0]-0.75,right=input.chargestate_peplength_plotrange()[1]+0.5)
                    ax[i].xaxis.set_major_locator(MaxNLocator(integer=True))
                    ax[i].xaxis.set_minor_locator(MultipleLocator(1))
                    ax[i].tick_params(axis="both",labelsize=axisfont_labels)

                ax[maxlist.index(max(maxlist))].set_ylim(top=ax[i].get_ylim()[1]+(y_padding*ax[i].get_ylim()[1]))
                ax[0].set_ylabel("Counts",fontsize=axisfont)
                fig.set_tight_layout(True)
            imagedownload("chargestate_peplength")

#endregion

# ============================================================================= Mixed Proteome
#region
    # ====================================== Info
    #show a table of the detected organisms and an order column to reorder them 
    @render.data_frame
    def organismtable():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        organismcolumn=[]
        for entry in searchoutput["PG.ProteinNames"].str.split(";"):
            organismcolumn.append(entry[0].split("_")[-1])
        searchoutput["Organism"]=organismcolumn
        organismlist=searchoutput["Organism"].drop_duplicates().tolist()

        organism_table=pd.DataFrame()
        organism_table["Organism"]=organismlist
        organism_table["Order"]=np.arange(1,len(organism_table["Organism"])+1).astype(str)
        organism_table["Remove?"]=[""]*len(organism_table["Organism"])
        for column in searchoutput["R.Condition"].drop_duplicates().reset_index(drop=True):
            organism_table[column+"_Quant Ratio"]=""
        return render.DataGrid(organism_table,editable=True,width="100%")
    #take the view of organismtable and generate organismlist in the order specified
    @reactive.calc
    def organism_list_from_table():
        organism_table_view=organismtable.data_view()
        organismlist=list(organism_table_view.sort_values("Order")["Organism"])
        organismlist_adjusted=organism_table_view[organism_table_view["Remove?"]!="x"]["Organism"].tolist()
        organismlist=organismlist_adjusted
        return organismlist
    @render.text
    def organisms():
        organismlist=organism_list_from_table()
        return organismlist

    #generate dfs for ID counts and summed intensities per organism
    @reactive.calc
    def mixedproteomestats():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        organismlist=organism_list_from_table()

        mixedproteomecounts=pd.DataFrame()
        mixedproteomecounts["Cond_Rep"]=runnames
        mixedproteomeintensity=pd.DataFrame()
        mixedproteomeintensity["Cond_Rep"]=runnames

        for organism in organismlist:
            proteincountlist=[]
            peptidecountlist=[]
            precursorcountlist=[]
            summedproteinintensitylist=[]
            summedprecursorintensitylist=[]
            for run in runnames:
                df=searchoutput[searchoutput["Cond_Rep"]==run]
                summedproteinintensitylist.append(df[df["Organism"]==organism]["PG.MS2Quantity"].drop_duplicates().reset_index(drop=True).sum())
                summedprecursorintensitylist.append(df[df["Organism"]==organism]["FG.MS2Quantity"].drop_duplicates().reset_index(drop=True).sum())
                proteincountlist.append(df[(df["Organism"]==organism)&(df["PG.MS2Quantity"]>0)]["PG.ProteinNames"].drop_duplicates().reset_index(drop=True).count())
                peptidecountlist.append(df[(df["Organism"]==organism)&(df["FG.MS2Quantity"]>0)]["EG.ModifiedPeptide"].drop_duplicates().reset_index(drop=True).count())
                precursorcountlist.append(len(df[(df["Organism"]==organism)&(df["FG.MS2Quantity"]>0)][["EG.ModifiedPeptide","FG.Charge"]].drop_duplicates().reset_index(drop=True)))
            mixedproteomeintensity[organism+"_summedproteinintensity"]=summedproteinintensitylist
            mixedproteomeintensity[organism+"_summedprecursorintensity"]=summedprecursorintensitylist
            mixedproteomecounts[organism+"_proteins"]=proteincountlist
            mixedproteomecounts[organism+"_peptides"]=peptidecountlist
            mixedproteomecounts[organism+"_precursors"]=precursorcountlist
        
        return mixedproteomecounts,mixedproteomeintensity

    # ====================================== Counts per Organism
    #counts per organism
    @reactive.effect
    def _():
        @render.plot(width=input.countsperorganism_width(),height=input.countsperorganism_height())
        def countsperorganism():
            mixedproteomecounts,mixedproteomeintensity=mixedproteomestats()
            organismlist=organism_list_from_table()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            matplottabcolors=list(mcolors.TABLEAU_COLORS)
            bluegray_colors=["#054169","#737373","#0071BC","#DCDCDC","#313331","#008DEB","#A0A0A0","#E5F1F8"]
            if input.coloroptions_sumint()=="matplot":
                colors=matplottabcolors
            elif input.coloroptions_sumint()=="bluegray":
                colors=bluegray_colors
            #extend color list if it's shorter than organismlist
            if len(colors)<len(organismlist):
                for i in colors:
                    if len(colors)<len(organismlist):
                        colors.append(i)

            if input.countsplotinput()=="proteins":
                titleprop="Protein"
            if input.countsplotinput()=="peptides":
                titleprop="Peptide"
            if input.countsplotinput()=="precursors":
                titleprop="Precursor"

            x=np.arange(len(mixedproteomecounts["Cond_Rep"].tolist()))
            width=input.countsperorganism_barwidth()
            relevantcolumns=[col for col in mixedproteomecounts if input.countsplotinput() in col]
            maxvalue=mixedproteomecounts[relevantcolumns].values.max()

            fig,ax=plt.subplots(figsize=(input.countsperorganism_width()*(1/plt.rcParams['figure.dpi']),input.countsperorganism_height()*(1/plt.rcParams['figure.dpi'])))
            for i in range(len(organismlist)):
                ax.bar(x+(i*width),mixedproteomecounts[organismlist[i]+"_"+input.countsplotinput()],width=width,label=organismlist[i],color=colors[i])
                ax.bar_label(ax.containers[i],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
            if len(organismlist)==2:
                ax.set_xticks(x+width/2,mixedproteomecounts["Cond_Rep"].tolist(),rotation=x_label_rotation)
            else:
                ax.set_xticks(x+width,mixedproteomecounts["Cond_Rep"].tolist(),rotation=x_label_rotation)
            ax.set_ylim(top=maxvalue+(y_padding*maxvalue))
            ax.legend(loc='center left',bbox_to_anchor=(1, 0.5),prop={'size':legendfont})
            ax.set_ylabel("Counts",fontsize=axisfont)
            ax.set_xlabel("Condition",fontsize=axisfont)
            ax.set_title(titleprop+" Counts per Organism",fontsize=titlefont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("countsperorganism")

    # ====================================== Summed Intensities
    #summed intensities per organism per run
    @reactive.effect
    def _():
        @render.plot(width=input.summedintensities_width(),height=input.summedintensities_height())
        def summedintensities():
            mixedproteomecounts,mixedproteomeintensity=mixedproteomestats()
            organismlist=organism_list_from_table()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            x_label_rotation=input.xaxis_label_rotation()

            matplottabcolors=list(mcolors.TABLEAU_COLORS)
            bluegray_colors=["#054169","#737373","#0071BC","#DCDCDC","#313331","#008DEB","#A0A0A0","#E5F1F8"]
            if input.coloroptions_sumint()=="matplot":
                colors=matplottabcolors
            elif input.coloroptions_sumint()=="bluegray":
                colors=bluegray_colors
            #extend color list if it's shorter than organismlist
            if len(colors)<len(organismlist):
                for i in colors:
                    if len(colors)<len(organismlist):
                        colors.append(i)

            if input.summedintensities_pick()=="protein":
                intensitycolumn="_summedproteinintensity"
                intensitypick="Protein"
            if input.summedintensities_pick()=="precursor":
                intensitycolumn="_summedprecursorintensity"
                intensitypick="Precursor"

            x=np.arange(len(mixedproteomeintensity["Cond_Rep"].tolist()))
            bottom=np.zeros(len(mixedproteomeintensity["Cond_Rep"].tolist()))
            fig,ax=plt.subplots(figsize=(input.summedintensities_width()*(1/plt.rcParams['figure.dpi']),input.summedintensities_height()*(1/plt.rcParams['figure.dpi'])))
            for i in range(len(organismlist)):
                ax.bar(x,mixedproteomeintensity[organismlist[i]+intensitycolumn],bottom=bottom,label=organismlist[i],color=colors[i])
                bottom+=mixedproteomeintensity[organismlist[i]+intensitycolumn].tolist()

            ax.set_xticks(x,labels=mixedproteomecounts["Cond_Rep"].tolist(),rotation=x_label_rotation)
            ax.set_ylabel("Total Intensity",fontsize=axisfont)
            ax.set_xlabel("Condition",fontsize=axisfont)
            ax.legend(loc="center left",bbox_to_anchor=(1, 0.5),fontsize=legendfont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.set_title("Total "+intensitypick+" Intensity per Organism per Run",fontsize=titlefont)
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("summedintensities")

    # ====================================== Quant Ratios
    #render ui call for dropdown calling sample condition names
    @render.ui
    def referencecondition():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=sampleconditions
        return ui.input_selectize("referencecondition_list","Pick reference condition:",choices=opts,selected=opts[0])
    #render ui call for dropdown calling sample condition names
    @render.ui
    def testcondition():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=sampleconditions
        return ui.input_selectize("testcondition_list","Pick test condition:",choices=opts,selected=opts[1])
    #quant ratios
    @reactive.effect
    def _():
        @render.plot(width=input.quantratios_width(),height=input.quantratios_height())
        def quantratios():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            organismlist=organism_list_from_table()
            organism_table=organismtable.data_view()

            organism_table=organism_table[organism_table["Remove?"]!="x"]
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            x_label_rotation=input.xaxis_label_rotation()

            matplottabcolors=list(mcolors.TABLEAU_COLORS)
            bluegray_colors=["#054169","#737373","#0071BC","#DCDCDC","#313331","#008DEB","#A0A0A0","#E5F1F8"]
            if input.coloroptions_sumint()=="matplot":
                colors=matplottabcolors
            elif input.coloroptions_sumint()=="bluegray":
                colors=bluegray_colors
            #extend color list if it's shorter than organismlist
            if len(colors)<len(organismlist):
                for i in colors:
                    if len(colors)<len(organismlist):
                        colors.append(i)

            referencecondition=input.referencecondition_list()
            testcondition=input.testcondition_list()

            testcolumn=organism_table[[col for col in organism_table if testcondition in col]].columns[0]
            referencecolumn=organism_table[[col for col in organism_table if referencecondition in col]].columns[0]

            testratios=organism_table.sort_values("Order")[testcolumn].astype(float).tolist()
            referenceratios=organism_table.sort_values("Order")[referencecolumn].astype(float).tolist()
            
            organism_merged=dict()
            ratio_average=pd.DataFrame()
            ratio_average["Organism"]=organismlist

            if input.y_log_scale()=="log2":
                ratio_average["Theoretical_Ratio"]=[np.log2(i/j) for i,j in zip(testratios,referenceratios)]
            if input.y_log_scale()=="log10":
                ratio_average["Theoretical_Ratio"]=[np.log10(i/j) for i,j in zip(testratios,referenceratios)]

            if input.quantratios_IDpick()=="protein":
                group_key=["PG.ProteinNames"]
                intensitycolumn="PG.MS2Quantity"
                df_columns=["PG.ProteinNames","PG.MS2Quantity"]
            if input.quantratios_IDpick()=="precursor":
                group_key=["EG.ModifiedPeptide","FG.Charge"]
                intensitycolumn="FG.MS2Quantity"
                df_columns=["EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]

            averagelist=[]
            for organism in organismlist:
                df=searchoutput[(searchoutput["Cond_Rep"].str.contains(referencecondition))&(searchoutput["Organism"]==organism)][df_columns].drop_duplicates().reset_index(drop=True)
                if input.quantratios_mean_median()=="mean":
                    df_reference=df.groupby(group_key).mean().reset_index().rename(columns={intensitycolumn:referencecondition})
                if input.quantratios_mean_median()=="median":
                    df_reference=df.groupby(group_key).median().reset_index().rename(columns={intensitycolumn:referencecondition})

                df_reference[referencecondition+"_stdev"]=df.groupby(group_key).std().reset_index(drop=True)
                df_reference[referencecondition+"_CV"]=df_reference[referencecondition+"_stdev"]/df_reference[referencecondition]*100

                df=searchoutput[(searchoutput["Cond_Rep"].str.contains(testcondition))&(searchoutput["Organism"]==organism)][df_columns].drop_duplicates().reset_index(drop=True)
                if input.quantratios_mean_median()=="mean":
                    df_test=df.groupby(group_key).mean().reset_index().rename(columns={intensitycolumn:testcondition})
                if input.quantratios_mean_median()=="median":
                    df_test=df.groupby(group_key).median().reset_index().rename(columns={intensitycolumn:testcondition})

                df_test[testcondition+"_stdev"]=df.groupby(group_key).std().reset_index(drop=True)
                df_test[testcondition+"_CV"]=df_test[testcondition+"_stdev"]/df_test[testcondition]*100
                
                if 1 in maxreplicatelist:
                    merged=df_reference.merge(df_test,how="inner")
                else:
                    merged=df_reference.merge(df_test,how="inner").dropna()

                if input.x_log_scale()=="log2":
                    merged["reference"]=np.log2(merged[referencecondition])
                if input.x_log_scale()=="log10":
                    merged["reference"]=np.log10(merged[referencecondition])
                if input.y_log_scale()=="log2":
                    merged["ratio"]=np.log2(merged[testcondition]/merged[referencecondition])
                if input.y_log_scale()=="log10":
                    merged["ratio"]=np.log10(merged[testcondition]/merged[referencecondition])

                averagelist.append(np.average(merged["ratio"].dropna()))

                organism_merged[organism]=merged

                if input.cvcutoff_switch()==True:
                    cv_cutoff=input.cvcutofflevel()
                    organism_merged[organism]=organism_merged[organism][(organism_merged[organism][referencecondition+"_CV"]<cv_cutoff)&(organism_merged[organism][testcondition+"_CV"]<cv_cutoff)]

            ratio_average["Experimental_Ratio"]=averagelist

            fig,ax=plt.subplots(nrows=1,ncols=3,gridspec_kw={"width_ratios":[2,5,2]},figsize=(input.quantratios_width()*(1/plt.rcParams['figure.dpi']),input.quantratios_height()*(1/plt.rcParams['figure.dpi'])))
            fig.set_tight_layout(True)
            for x,organism in enumerate(organismlist):
                ax[0].bar(x,len(organism_merged[organism]),color=colors[x])
                ax[0].bar_label(ax[0].containers[x],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
                ax[1].scatter(organism_merged[organism]["reference"],organism_merged[organism]["ratio"],alpha=0.25,color=colors[x])
                ax[2].hist(organism_merged[organism]["ratio"],bins=100,orientation=u"horizontal",alpha=0.5,density=True,color=colors[x])

            for x in range(len(organismlist)):
                if input.mixedproteome_showexperimentalratios()==True:
                    ax[1].axhline(y=ratio_average["Experimental_Ratio"][x],linestyle="dashed",color=colors[x])
                    ax[2].axhline(y=ratio_average["Experimental_Ratio"][x],linestyle="dashed",color=colors[x])
                if input.mixedproteome_showtheoreticalratios()==True:
                    ax[1].axhline(y=ratio_average["Theoretical_Ratio"][x],color=colors[x])
                    ax[2].axhline(y=ratio_average["Theoretical_Ratio"][x],color=colors[x])

            if input.plotrange_switch()==True:
                ymin=input.plotrange()[0]
                ymax=input.plotrange()[1]
                ax[1].set_ylim(ymin,ymax)
                ax[2].set_ylim(ymin,ymax)

            ax[0].set_xticks(np.arange(len(organismlist)),organismlist,rotation=x_label_rotation)
            ax[0].set_ylabel("Number of Common IDs",fontsize=axisfont)
            bottom,top=ax[0].get_ylim()
            ax[0].set_ylim(top=top+(0.15*top))

            leg=ax[1].legend(organismlist,loc="upper right",prop={'size':legendfont})
            for tag in leg.legend_handles:
                tag.set_alpha(1)
            ax[1].set_xlabel(input.x_log_scale()+" Intensity, Reference",fontsize=axisfont)
            ax[1].set_ylabel(input.y_log_scale()+" Ratio, Test/Reference",fontsize=axisfont)
            ax[1].set_title("Reference: "+referencecondition+", Test: "+testcondition,pad=10,fontsize=titlefont)

            ax[2].set_xlabel("Density",fontsize=axisfont)
            ax[2].set_ylabel(input.y_log_scale()+" Ratio, Test/Reference",fontsize=axisfont)
            ax[0].tick_params(axis="both",labelsize=axisfont_labels)
            ax[1].tick_params(axis="both",labelsize=axisfont_labels)
            ax[2].tick_params(axis="both",labelsize=axisfont_labels)
            ax[0].set_axisbelow(True)
            ax[0].grid(linestyle="--")
            ax[1].set_axisbelow(True)
            ax[1].grid(linestyle="--")
            ax[2].set_axisbelow(True)
            ax[2].grid(linestyle="--")

            imagedownload("quantratios")
    #show table of quant ratios
    @render.table
    def quantratios_table():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        organismlist=organism_list_from_table()
        organism_table=organismtable.data_view()

        organism_table=organism_table[organism_table["Remove?"]!="x"]

        referencecondition=input.referencecondition_list()
        testcondition=input.testcondition_list()

        testcolumn=organism_table[[col for col in organism_table if testcondition in col]].columns[0]
        referencecolumn=organism_table[[col for col in organism_table if referencecondition in col]].columns[0]

        testratios=organism_table.sort_values("Order")[testcolumn].astype(float).tolist()
        referenceratios=organism_table.sort_values("Order")[referencecolumn].astype(float).tolist()
        
        organism_merged=dict()
        ratio_average=pd.DataFrame()
        ratio_average["Organism"]=organismlist
        
        if input.y_log_scale()=="log2":
            ratio_average["Theoretical_Ratio (log2)"]=[np.log2(i/j) for i,j in zip(testratios,referenceratios)]
        if input.y_log_scale()=="log10":
            ratio_average["Theoretical_Ratio (log10)"]=[np.log10(i/j) for i,j in zip(testratios,referenceratios)]

        if input.quantratios_IDpick()=="protein":
            group_key=["PG.ProteinNames"]
            intensitycolumn="PG.MS2Quantity"
            df_columns=["PG.ProteinNames","PG.MS2Quantity"]
        if input.quantratios_IDpick()=="precursor":
            group_key=["EG.ModifiedPeptide","FG.Charge"]
            intensitycolumn="FG.MS2Quantity"
            df_columns=["EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]

        averagelist=[]
        for organism in organismlist:
            df=searchoutput[(searchoutput["Cond_Rep"].str.contains(referencecondition))&(searchoutput["Organism"]==organism)][df_columns].drop_duplicates().reset_index(drop=True)
            if input.quantratios_mean_median()=="mean":
                df_reference=df.groupby(group_key).mean().reset_index().rename(columns={intensitycolumn:referencecondition})
            if input.quantratios_mean_median()=="median":
                df_reference=df.groupby(group_key).median().reset_index().rename(columns={intensitycolumn:referencecondition})

            df_reference[referencecondition+"_stdev"]=df.groupby(group_key).std().reset_index(drop=True)
            df_reference[referencecondition+"_CV"]=df_reference[referencecondition+"_stdev"]/df_reference[referencecondition]*100

            df=searchoutput[(searchoutput["Cond_Rep"].str.contains(testcondition))&(searchoutput["Organism"]==organism)][df_columns].drop_duplicates().reset_index(drop=True)
            if input.quantratios_mean_median()=="mean":
                df_test=df.groupby(group_key).mean().reset_index().rename(columns={intensitycolumn:testcondition})
            if input.quantratios_mean_median()=="median":
                df_test=df.groupby(group_key).median().reset_index().rename(columns={intensitycolumn:testcondition})

            df_test[testcondition+"_stdev"]=df.groupby(group_key).std().reset_index(drop=True)
            df_test[testcondition+"_CV"]=df_test[testcondition+"_stdev"]/df_test[testcondition]*100
            
            if 1 in maxreplicatelist:
                merged=df_reference.merge(df_test,how="inner")
            else:
                merged=df_reference.merge(df_test,how="inner").dropna()

            if input.y_log_scale()=="log2":
                merged["ratio"]=np.log2(merged[testcondition]/merged[referencecondition])
            if input.y_log_scale()=="log10":
                merged["ratio"]=np.log10(merged[testcondition]/merged[referencecondition])
            averagelist.append(np.average(merged["ratio"].dropna()))

            organism_merged[organism]=merged

            if input.cvcutoff_switch()==True:
                cv_cutoff=input.cvcutofflevel()
                organism_merged[organism]=organism_merged[organism][(organism_merged[organism][referencecondition+"_CV"]<cv_cutoff)&(organism_merged[organism][testcondition+"_CV"]<cv_cutoff)]
            
        ratio_average["Experimental_Ratio ("+input.y_log_scale()+")"]=averagelist
        return ratio_average

#endregion

# ============================================================================= PRM
#region
    # ====================================== PRM List
    #download PRM table template
    @render.download(filename="prm_template.csv")
    def prm_template():
        template_df=pd.DataFrame(columns=[
            "PG.ProteinGroups",
            "EG.ModifiedPeptide"
            ])
        with io.BytesIO() as buf:
            template_df.to_csv(buf,index=False)
            yield buf.getvalue()
    #import prm list and generate a searchoutput-like table for just the prm peptides
    @reactive.calc
    def prm_import():
        if input.prm_list() is None:
            return pd.DataFrame()
        prm_list=pd.read_csv(input.prm_list()[0]["datapath"])
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        df_list=[]
        for peptide in prm_list["EG.ModifiedPeptide"]:
            prm_peptide=searchoutput[searchoutput["EG.ModifiedPeptide"]==peptide]
            df_list.append(prm_peptide)
        searchoutput_prmpepts=pd.concat(df_list).reset_index(drop=True)
        if "Concentration" in searchoutput.columns:
            searchoutput_prmpepts.sort_values("Concentration")
        return prm_list,searchoutput_prmpepts

    # ====================================== PRM Table
    #generate prm table to be exported to timscontrol
    @reactive.calc
    def prm_list_import():
        if input.prm_list() is None:
            return pd.DataFrame()
        prm_list=pd.read_csv(input.prm_list()[0]["datapath"])
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        try:
            isolationwidth=float(input.isolationwidth_input())
        except:
            isolationwidth=0
        try:
            rtwindow=float(input.rtwindow_input())
        except:
            rtwindow=0
        try:
            imwindow=float(input.imwindow_input())
        except:
            imwindow=0
        df_list=[]
        for peptide in prm_list["EG.ModifiedPeptide"]:
            prm_peptide=searchoutput[searchoutput["EG.ModifiedPeptide"]==peptide][["PG.ProteinGroups","EG.ModifiedPeptide","FG.PrecMz","FG.Charge","EG.ApexRT","EG.IonMobility"]].groupby(["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge"]).mean().reset_index()
            df_list.append(prm_peptide)
        searchoutput_prm=pd.concat(df_list).reset_index(drop=True)
        searchoutput_prm["EG.ApexRT"]=searchoutput_prm["EG.ApexRT"]*60

        searchoutput_prm.rename(columns={"FG.PrecMz":"Mass [m/z]","FG.Charge":"Charge","EG.ApexRT":"RT [s]"},inplace=True)

        mzisolationwidth=[]
        RTrange=[]
        startIM=[]
        endIM=[]
        CE=[]
        externalID=[]
        description=[]
        for i in range(len(searchoutput_prm)):
            mzisolationwidth.append(isolationwidth)
            RTrange.append(rtwindow)
            startIM.append(searchoutput_prm["EG.IonMobility"][i]-imwindow)
            endIM.append(searchoutput_prm["EG.IonMobility"][i]+imwindow)
            CE.append("")
            externalID.append(searchoutput_prm["EG.ModifiedPeptide"][i])
            description.append("")

        searchoutput_prm["Isolation Width [m/z]"]=mzisolationwidth
        searchoutput_prm["RT Range [s]"]=RTrange
        searchoutput_prm["Start IM [1/k0]"]=startIM
        searchoutput_prm["End IM [1/k0]"]=endIM
        searchoutput_prm["CE [eV]"]=CE
        searchoutput_prm["External ID"]=externalID
        searchoutput_prm["Description"]=description

        searchoutput_prm=searchoutput_prm[["Mass [m/z]","Charge","Isolation Width [m/z]","RT [s]","RT Range [s]","Start IM [1/k0]","End IM [1/k0]","CE [eV]","External ID","Description","PG.ProteinGroups","EG.ModifiedPeptide","EG.IonMobility"]]

        searchoutput_prm.drop(columns=["PG.ProteinGroups","EG.ModifiedPeptide","EG.IonMobility"],inplace=True)

        return searchoutput_prm
    
    #show prm list in window
    @render.data_frame
    def prm_table():
        searchoutput_prm=prm_list_import()
        return render.DataGrid(searchoutput_prm,width="100%",editable=True)

    #download prm list that's been edited in the window
    @render.download(filename="prm_peptide_list.csv")
    def prm_table_download():
        prm_table_view=prm_table.data_view()
        
        yield prm_table_view.to_csv(index=False)

    # ====================================== Individual Tracker
    #prm selectize peptide list
    @render.ui
    def prmpeptracker_pick():
        prm_list,searchoutput_prmpepts=prm_import()
        opts=prm_list["EG.ModifiedPeptide"]
        return ui.input_selectize("prmpeptracker_picklist","Pick PRM peptide to plot data for:",choices=opts,width="600px")

    #plot intensity across runs, number of replicates, and CVs of selected peptide from prm list
    @reactive.effect
    def _():
        @render.plot(width=input.prmpeptracker_width(),height=input.prmpeptracker_height())
        def prmpeptracker_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            prm_list,searchoutput_prmpepts=prm_import()

            peplist=prm_list["EG.ModifiedPeptide"]
            peptide=peplist[int(input.prmpeptracker_picklist())]

            pepdf=searchoutput_prmpepts[searchoutput_prmpepts["EG.ModifiedPeptide"]==peptide]
            chargelist=pepdf["FG.Charge"].drop_duplicates().tolist()
            #make sure the chargelist is sorted so we're plotting in charge order
            chargelist.sort()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            x_label_rotation=input.xaxis_label_rotation()
            width=0.25

            #make a table of expected values for the concentration ratios
            df=searchoutput[["R.Condition","Concentration"]].drop_duplicates().reset_index(drop=True)
            expectedratio=[]
            for i in range(len(df)):
                conc_min=min(df["Concentration"])
                expectedratio.append(df["Concentration"][i]/conc_min)
            df["Expected Ratio"]=expectedratio

            fig,ax=plt.subplots(ncols=2,nrows=2,figsize=(input.prmpeptracker_width()*(1/plt.rcParams['figure.dpi']),input.prmpeptracker_height()*(1/plt.rcParams['figure.dpi'])))
            for i,charge in enumerate(chargelist):
                meandf=pepdf[pepdf["FG.Charge"]==charge][["R.Condition","FG.MS2Quantity"]].groupby("R.Condition",sort=False).mean()
                stdev=pepdf[pepdf["FG.Charge"]==charge][["R.Condition","FG.MS2Quantity"]].groupby("R.Condition",sort=False).std()
                cv=stdev/meandf*100
                meandf=meandf.reset_index()
                stdev=stdev.reset_index()
                cv=cv.reset_index()
                detectedinreps=pd.DataFrame({"Count":pepdf[pepdf["FG.Charge"]==charge].groupby("R.Condition",sort=False).size()}).reset_index()

                if len(meandf)<len(searchoutput["R.Condition"].drop_duplicates().reset_index(drop=True)):
                    expectedrows=pd.DataFrame({"R.Condition":searchoutput["R.Condition"].drop_duplicates().reset_index(drop=True)})
                    meandf=expectedrows.merge(meandf,how="left",left_on="R.Condition",right_on="R.Condition").fillna(0)
                    stdev=expectedrows.merge(stdev,how="left",left_on="R.Condition",right_on="R.Condition").fillna(0)
                    cv=expectedrows.merge(cv,how="left",left_on="R.Condition",right_on="R.Condition").fillna(0)
                    detectedinreps=expectedrows.merge(detectedinreps,how="left",left_on="R.Condition",right_on="R.Condition").fillna(0)

                x=np.arange(0,len(meandf["R.Condition"]),1)
                y=meandf["FG.MS2Quantity"].tolist()
                fit=np.poly1d(np.polyfit(x,y,1))
                ax[0,0].errorbar(x,y,yerr=stdev["FG.MS2Quantity"],marker="o",linestyle="None")
                ax[0,0].plot(x,fit(x),linestyle="--",color="black")
                ax[1,0].plot(meandf["R.Condition"],cv["FG.MS2Quantity"],marker="o")
                ax[0,1].bar(x+i*width,detectedinreps["Count"],width=width,label=str(charge)+"+")
                if "Concentration" in searchoutput.columns:
                    #use a single dataframe to plot the expected and measured ratios for the selected peptides
                    merged=df.merge(meandf)
                    min_conc_signal=merged.loc[merged["Concentration"]==min(merged["Concentration"]),"FG.MS2Quantity"].values[0]
                    measuredratio=[]
                    #if there's no signal for the lowest concentration condition, set measured ratio values to zero since the numbers won't be informative
                    if min_conc_signal!=0:
                        for conc in merged["Concentration"]:
                            measuredratio.append(merged.loc[merged["Concentration"]==conc,"FG.MS2Quantity"].values[0]/min_conc_signal)
                    else:
                        measuredratio.append([0.0]*len(merged["Concentration"]))
                        measuredratio=list(itertools.chain(*measuredratio))
                    merged["Measured Ratio"]=measuredratio
                    #plotting average FG.MS2Quantity over each condition
                    ax[1,1].scatter(merged["Expected Ratio"],merged["Measured Ratio"])
                else:
                    ax[1,1].set_visible(False)
                
            ax[0,0].set_ylabel("FG.MS2Quantity",fontsize=axisfont)
            ax[0,0].set_xticks(x)
            ax[0,0].set_xticklabels(searchoutput["R.Condition"].drop_duplicates().reset_index(drop=True),rotation=x_label_rotation)
            ax[0,0].set_ylim(bottom=-(ax[0,0].get_ylim()[1])/10)
            ax[0,0].set_title("Intensity Across Runs",fontsize=titlefont)
            ax[0,0].set_axisbelow(True)
            ax[0,0].grid(linestyle="--")
            ax[0,0].tick_params(axis="both",labelsize=axisfont_labels)

            ax[1,0].axhline(y=20,linestyle="--",color="black")
            ax[1,0].set_xticks(x)
            ax[1,0].set_xticklabels(searchoutput["R.Condition"].drop_duplicates().reset_index(drop=True),rotation=x_label_rotation)
            ax[1,0].set_ylabel("CV (%)",fontsize=axisfont)
            ax[1,0].set_title("CVs",fontsize=titlefont)
            ax[1,0].set_axisbelow(True)
            ax[1,0].grid(linestyle="--")
            ax[1,0].tick_params(axis="both",labelsize=axisfont_labels)

            ax[0,1].set_ylabel("Number of Replicates",fontsize=axisfont)
            ax[0,1].set_xticks(x,meandf["R.Condition"],rotation=x_label_rotation)
            ax[0,1].set_title("Number of Replicates Observed",fontsize=titlefont)
            ax[0,1].set_axisbelow(True)
            ax[0,1].grid(linestyle="--")
            ax[0,1].tick_params(axis="both",labelsize=axisfont_labels)

            ax[1,1].set_xlabel("Expected Ratio",fontsize=axisfont)
            ax[1,1].set_ylabel("Measured Ratio",fontsize=axisfont)
            lims=[np.min([ax[1,1].get_xlim(),ax[1,1].get_ylim()]),np.max([ax[1,1].get_xlim(),ax[1,1].get_ylim()])]
            ax[1,1].plot(lims,lims,color="k",linestyle="--",alpha=0.5)
            ax[1,1].set_title("Dilution Curve",fontsize=titlefont)
            ax[1,1].set_axisbelow(True)
            ax[1,1].grid(linestyle="--")
            ax[1,1].tick_params(axis="both",labelsize=axisfont_labels)

            fig.legend(loc="lower right",bbox_to_anchor=(0.99,0.9))
            fig.suptitle(peptide.strip("_"),fontsize=titlefont)
            fig.set_tight_layout(True)
            imagedownload("prmpeptidetracker_"+peptide)

    # ====================================== Intensity Across Runs
    #plot intensity of all prm peptides across runs
    @reactive.effect
    def _():
        @render.plot(width=input.prmpepintensity_width(),height=input.prmpepintensity_height())
        def prmpepintensity_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            prm_list,searchoutput_prmpepts=prm_import()

            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()
            x_label_rotation=input.xaxis_label_rotation()

            fig,ax=plt.subplots(figsize=(input.prmpepintensity_width()*(1/plt.rcParams['figure.dpi']),input.prmpepintensity_height()*(1/plt.rcParams['figure.dpi'])))
            for peptide in prm_list["EG.ModifiedPeptide"]:
                pepdf=searchoutput_prmpepts[searchoutput_prmpepts["EG.ModifiedPeptide"]==peptide]
                chargelist=pepdf["FG.Charge"].drop_duplicates().tolist()

                for charge in chargelist:
                    plottingdf=pepdf[pepdf["FG.Charge"]==charge][["Cond_Rep","FG.Charge","FG.MS2Quantity"]]

                    if len(plottingdf)<len(runnames):
                        expectedrows=pd.DataFrame({"Cond_Rep":runnames})
                        plottingdf=expectedrows.merge(plottingdf,how="left",left_on="Cond_Rep",right_on="Cond_Rep").fillna(1)
                    else:
                        pass
                    ax.plot(plottingdf["Cond_Rep"],np.log10(plottingdf["FG.MS2Quantity"]),marker="o",label=peptide.strip("_")+"_"+str(charge)+"+")

            ax.legend(loc='center left', bbox_to_anchor=(1,0.5),prop={'size':legendfont})
            ax.tick_params(axis="x",rotation=x_label_rotation)
            ax.set_xlabel("Condition",fontsize=axisfont)
            ax.set_ylabel("log10(FG.MS2Quantity)",fontsize=axisfont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("prmpeptideintensity")

#endregion

# ============================================================================= Dilution Series
#region
    #ui call to pick normalizing condition for dilution series calculations
    @render.ui
    def normalizingcondition():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=sampleconditions
        return ui.input_selectize("normalizingcondition_pick","Pick normalizing condition:",choices=opts)

    #dilution series calculations for plotting
    @reactive.calc
    def dilutionseries_calc():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        sortedconditions=[]
        concentrations=searchoutput["Concentration"].drop_duplicates().tolist()
        sortedconcentrations=sorted(concentrations)
        for i in sortedconcentrations:
            sortedconditions.append(searchoutput[searchoutput["Concentration"]==i]["R.Condition"].values[0])

        normalizingcondition=input.normalizingcondition_pick()

        dilutionseries=[]
        theoreticalratio=[]
        norm=searchoutput[searchoutput["R.Condition"]==normalizingcondition][["PG.ProteinGroups","PG.MS2Quantity"]].groupby("PG.ProteinGroups").mean().reset_index()

        for condition in sortedconditions:
            test=searchoutput[searchoutput["R.Condition"]==condition][["PG.ProteinGroups","PG.MS2Quantity"]].groupby("PG.ProteinGroups").mean().reset_index()
            merge=norm.merge(test,how="left",on="PG.ProteinGroups",suffixes=("_norm","_test"))
            merge["Ratio"]=merge["PG.MS2Quantity_test"]/merge["PG.MS2Quantity_norm"]
            merge=merge.dropna()
            dilutionseries.append(merge["Ratio"])
            
            norm_conc=searchoutput[searchoutput["R.Condition"]==normalizingcondition]["Concentration"].drop_duplicates().reset_index(drop=True)[0]
            conc=searchoutput[searchoutput["R.Condition"]==condition]["Concentration"].drop_duplicates().reset_index(drop=True)[0]
            theoreticalratio.append(conc/norm_conc)
        
        return sortedconcentrations,sortedconditions,dilutionseries,theoreticalratio

    #plot dilution ratios
    @reactive.effect
    def _():
        @render.plot(width=input.dilutionseries_width(),height=input.dilutionseries_height())
        def dilutionseries_plot(width=input.dilutionseries_width(),height=input.dilutionseries_height()):
            sortedconcentrations,sortedconditions,dilutionseries,theoreticalratio=dilutionseries_calc()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()

            fig,ax=plt.subplots(figsize=(input.dilutionseries_width()*(1/plt.rcParams['figure.dpi']),input.dilutionseries_height()*(1/plt.rcParams['figure.dpi'])))

            medianlineprops=dict(linestyle="--",color="black")
            flierprops=dict(markersize=3)

            widths=[]
            for conc in sortedconcentrations:
                widths.append(0.5*conc)

            bplot=ax.boxplot(dilutionseries,medianprops=medianlineprops,flierprops=flierprops,positions=sortedconcentrations,widths=widths)
            plot=ax.violinplot(dilutionseries,showextrema=False,positions=sortedconcentrations,widths=widths)
            ax.plot(sortedconcentrations,theoreticalratio,zorder=2.5,marker="o",color="k",label="Theoretical Ratio")

            ax.set_yscale("log")
            ax.set_xscale("log")
            #ax.set_xticks(np.arange(1,len(sortedconditions)+1,1),labels=sortedconditions,rotation=input.xaxis_label_rotation())
            ax.set_xlabel("Concentration",fontsize=axisfont)
            ax.set_ylabel("Ratio",fontsize=axisfont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.legend(loc="upper left",prop={'size':legendfont})
            ax.tick_params(axis="both",labelsize=axisfont_labels)

            colors=colorpicker()
            for z,color in zip(plot["bodies"],colors):
                z.set_facecolor(color)
                z.set_edgecolor("black")
                z.set_alpha(0.7)
            fig.set_tight_layout(True)
            imagedownload("dilutionseries")

#endregion

# ============================================================================= Glycoproteomics
#region
    #generate dfs and variables for the glyco results
    @reactive.calc
    def glyco_variables():
        searchoutput=metadata_update()
        searchoutput["R.Condition"]=searchoutput["R.Condition"].apply(str)
        if "Cond_Rep" not in searchoutput.columns:
            searchoutput.insert(0,"Cond_Rep",searchoutput["R.Condition"]+"_"+searchoutput["R.Replicate"].apply(str))
        elif "Cond_Rep" in searchoutput.columns:
            searchoutput["Cond_Rep"]=searchoutput["R.Condition"]+"_"+searchoutput["R.Replicate"].apply(str)
        
        resultdf_glyco=pd.DataFrame(searchoutput[["Cond_Rep","R.Condition","R.Replicate"]].drop_duplicates()).reset_index(drop=True)
        num_glycoproteins=[]
        num_glycopeptides=[]
        num_glycoPSMs=[]
        dict_glycoproteins=dict()
        dict_glycopeptides=dict()
        dict_glycoPSMs=dict()
        for run in searchoutput["Cond_Rep"].drop_duplicates():
            df=searchoutput[searchoutput["Cond_Rep"]==run]

            #filtered PSMs with a q value cutoff (used for the dfs generated here)
            if input.software()=="fragpipe_glyco":
                df=df[(df["Glycan q-value"].isnull()==False)&(df["Glycan q-value"]<=0.01)]
            if input.software()=="glycoscape":
                df=df[df["Total Glycan Composition"].isnull()==False]

            #Assign truth value to whether number after Hex mod is 5<x<10
            highmannose=[]
            for i in range(len(df["Total Glycan Composition"].tolist())):
                glycanstring=str(df["Total Glycan Composition"].tolist()[i])
                if "Hex(" in glycanstring:
                    hexvalue=df["Total Glycan Composition"].tolist()[i][df["Total Glycan Composition"].tolist()[i].find("Hex(")+len("Hex(")]
                    if int(hexvalue)<=10 and int(hexvalue)>=5:
                        highmannose.append("True")
                    elif int(hexvalue)<5:
                        highmannose.append("False")
                else:
                    highmannose.append("False")
            df["High Mannose"]=highmannose
            
            #all unique glycopeptides (charge agnostic)
            glycopeptide=df[["PEP.StrippedSequence","EG.ModifiedPeptide","Total Glycan Composition","PG.ProteinGroups"]].drop_duplicates()
            #all unique glycoproteins
            glycoprotein=df["PG.ProteinGroups"].drop_duplicates()
            
            num_glycoproteins.append(len(glycoprotein))
            num_glycopeptides.append(len(glycopeptide))
            num_glycoPSMs.append(len(df))
            
            dict_glycoproteins[run]=glycoprotein.reset_index(drop=True)
            dict_glycopeptides[run]=glycopeptide.reset_index(drop=True)
            dict_glycoPSMs[run]=df.reset_index(drop=True)

        resultdf_glyco["Cond_Rep"]=searchoutput["Cond_Rep"].drop_duplicates().reset_index(drop=True)
        resultdf_glyco["glycoproteins"]=num_glycoproteins
        resultdf_glyco["glycopeptides"]=num_glycopeptides
        resultdf_glyco["glycoPSM"]=num_glycoPSMs

        glycoproteins_df=pd.DataFrame()
        glycopeptides_df=pd.DataFrame()
        glycoPSMs_df=pd.DataFrame()
        for key in dict_glycoproteins.keys():
            protein_key=pd.DataFrame({"Cond_Rep":[key]*len(dict_glycoproteins[key]),"PG.ProteinGroups":dict_glycoproteins[key]})
            glycoproteins_df=pd.concat([glycoproteins_df,protein_key]).reset_index(drop=True)
            
            peptide_key=pd.DataFrame({"Cond_Rep":[key]*len(dict_glycopeptides[key])})
            peptide_key_df=pd.concat([peptide_key,dict_glycopeptides[key]],axis=1)
            glycopeptides_df=pd.concat([glycopeptides_df,peptide_key_df]).reset_index(drop=True)

            glycoPSMs_df=pd.concat([glycoPSMs_df,dict_glycoPSMs[key]]).reset_index(drop=True)

        return resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df

    # ====================================== Glyco ID Metrics
    #plot only glycosylated IDs
    @reactive.effect
    def _():
        @render.plot(width=input.glycoIDsplot_width(),height=input.glycoIDsplot_height())
        def glycoIDsplot():
            resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
            color=replicatecolors()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            y_padding=input.ypadding()

            fig,ax=plt.subplots(ncols=3,sharex=True,figsize=(input.glycoIDsplot_width()*(1/plt.rcParams['figure.dpi']),input.glycoIDsplot_height()*(1/plt.rcParams['figure.dpi'])))
            ax1=ax[0]
            ax2=ax[1]
            ax3=ax[2]

            resultdf_glyco.plot.bar(ax=ax1,x="Cond_Rep",y="glycoproteins",legend=False,width=0.8,color=color,edgecolor="k")
            ax1.bar_label(ax1.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
            ax1.set_ylim(top=max(resultdf_glyco["glycoproteins"].tolist())+y_padding*max(resultdf_glyco["glycoproteins"].tolist()))
            ax1.set_ylabel("Counts",fontsize=axisfont)
            ax1.set_xlabel("Condition",fontsize=axisfont)
            ax1.set_title("Glycoproteins",fontsize=titlefont)
            ax1.tick_params(axis="both",labelsize=axisfont_labels)

            resultdf_glyco.plot.bar(ax=ax2,x="Cond_Rep",y="glycopeptides",legend=False,width=0.8,color=color,edgecolor="k")
            ax2.bar_label(ax2.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
            ax2.set_ylim(top=max(resultdf_glyco["glycopeptides"].tolist())+y_padding*max(resultdf_glyco["glycopeptides"].tolist()))
            ax2.set_xlabel("Condition",fontsize=axisfont)
            ax2.set_title("Glycopeptides",fontsize=titlefont)
            ax2.tick_params(axis="both",labelsize=axisfont_labels)

            resultdf_glyco.plot.bar(ax=ax3,x="Cond_Rep",y="glycoPSM",legend=False,width=0.8,color=color,edgecolor="k")
            ax3.bar_label(ax3.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
            ax3.set_ylim(top=max(resultdf_glyco["glycoPSM"].tolist())+(y_padding+0.1)*max(resultdf_glyco["glycoPSM"].tolist()))
            ax3.set_xlabel("Condition",fontsize=axisfont)
            ax3.set_title("Glyco-PSMs",fontsize=titlefont)
            ax3.tick_params(axis="both",labelsize=axisfont_labels)

            ax1.set_axisbelow(True)
            ax1.grid(linestyle="--")
            ax2.set_axisbelow(True)
            ax2.grid(linestyle="--")
            ax3.set_axisbelow(True)
            ax3.grid(linestyle="--")
            fig.set_tight_layout(True)
            imagedownload("glycoidcounts")

    # ====================================== Venn Diagram
    @render.ui
    def glyco_venn_run1_ui():
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
        if input.glyco_venn_conditionorrun()=="condition":
            opts=resultdf_glyco["R.Condition"].drop_duplicates().tolist()
            return ui.input_selectize("glyco_venn_run1_list","Condition 1:",choices=opts)
        if input.glyco_venn_conditionorrun()=="individual":
            opts=resultdf_glyco["Cond_Rep"].tolist()
            return ui.input_selectize("glyco_venn_run1_list","Run 1:",choices=opts)   
    @render.ui
    def glyco_venn_run2_ui():
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
        if input.glyco_venn_conditionorrun()=="condition":
            opts=resultdf_glyco["R.Condition"].drop_duplicates().tolist()
            return ui.input_selectize("glyco_venn_run2_list","Condition 2:",choices=opts)
        if input.glyco_venn_conditionorrun()=="individual":
            opts=resultdf_glyco["Cond_Rep"].tolist()
            return ui.input_selectize("glyco_venn_run2_list","Run 2:",choices=opts)
    @render.ui
    def glyco_venn_run3_ui():
        if input.glyco_venn_numcircles()=="3":
            resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
            if input.glyco_venn_conditionorrun()=="condition":
                opts=resultdf_glyco["R.Condition"].drop_duplicates().tolist()
                return ui.input_selectize("glyco_venn_run3_list","Condition 3:",choices=opts)
            if input.glyco_venn_conditionorrun()=="individual":
                opts=resultdf_glyco["Cond_Rep"].tolist()
                return ui.input_selectize("glyco_venn_run3_list","Run 3:",choices=opts)   
    #plot Venn Diagram
    @reactive.effect
    def _():
        @render.plot(width=input.glyco_venn_width(),height=input.glyco_venn_height())
        def glyco_venn_plot():
            from matplotlib_venn import venn2,venn2_circles,venn3,venn3_circles
            
            resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
            if input.glyco_venn_conditionorrun()=="condition":
                A=glycoPSMs_df[glycoPSMs_df["R.Condition"]==input.glyco_venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
                B=glycoPSMs_df[glycoPSMs_df["R.Condition"]==input.glyco_venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
                if input.glyco_venn_numcircles()=="3":
                    C=glycoPSMs_df[glycoPSMs_df["R.Condition"]==input.glyco_venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
            if input.glyco_venn_conditionorrun()=="individual":
                A=glycoPSMs_df[glycoPSMs_df["Cond_Rep"]==input.glyco_venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
                B=glycoPSMs_df[glycoPSMs_df["Cond_Rep"]==input.glyco_venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
                if input.glyco_venn_numcircles()=="3":
                    C=glycoPSMs_df[glycoPSMs_df["Cond_Rep"]==input.glyco_venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
            
            A["modpep_glycan"]=A["EG.ModifiedPeptide"]+"_"+A["Total Glycan Composition"]
            B["modpep_glycan"]=B["EG.ModifiedPeptide"]+"_"+B["Total Glycan Composition"]

            A["modpep_glycan_charge"]=A["EG.ModifiedPeptide"]+"_"+A["Total Glycan Composition"]+"_"+A["FG.Charge"].astype(str)
            B["modpep_glycan_charge"]=B["EG.ModifiedPeptide"]+"_"+B["Total Glycan Composition"]+"_"+B["FG.Charge"].astype(str)

            if input.glyco_venn_numcircles()=="3":
                C["modpep_glycan"]=C["EG.ModifiedPeptide"]+"_"+C["Total Glycan Composition"]
                C["modpep_glycan_charge"]=C["EG.ModifiedPeptide"]+"_"+C["Total Glycan Composition"]+"_"+C["FG.Charge"].astype(str)

            if input.glyco_venn_plotproperty()=="glycoproteins":
                a=set(A["PG.ProteinGroups"])
                b=set(B["PG.ProteinGroups"])
                if input.glyco_venn_numcircles()=="3":
                    c=set(C["PG.ProteinGroups"])
                titlemod="Glycoproteins"
            if input.glyco_venn_plotproperty()=="glycopeptides":
                a=set(A["modpep_glycan"])
                b=set(B["modpep_glycan"])
                if input.glyco_venn_numcircles()=="3":
                    c=set(C["modpep_glycan"])
                titlemod="Glycopeptides"
            if input.glyco_venn_plotproperty()=="glycoPSMs":
                a=set(A["modpep_glycan_charge"])
                b=set(B["modpep_glycan_charge"])
                if input.glyco_venn_numcircles()=="3":
                    c=set(C["modpep_glycan_charge"])
                titlemod="GlycoPSMs"

            fig,ax=plt.subplots(figsize=(input.glyco_venn_width()*(1/plt.rcParams['figure.dpi']),input.glyco_venn_height()*(1/plt.rcParams['figure.dpi'])))
            if input.glyco_venn_numcircles()=="2":
                Ab=len(a-b)
                aB=len(b-a)
                AB=len(a&b)
                venn2(subsets=(Ab,aB,AB),set_labels=(input.glyco_venn_run1_list(),input.glyco_venn_run2_list()),set_colors=("tab:blue","tab:orange"),ax=ax)
                venn2_circles(subsets=(Ab,aB,AB),linestyle="dashed",linewidth=0.5)
                plt.title("Venn Diagram for "+titlemod)
            if input.glyco_venn_numcircles()=="3":
                Abc=len(a-b-c)
                aBc=len(b-a-c)
                ABc=len((a&b)-c)
                abC=len(c-a-b)
                AbC=len((a&c)-b)
                aBC=len((b&c)-a)
                ABC=len(a&b&c)
                venn3(subsets=(Abc,aBc,ABc,abC,AbC,aBC,ABC),set_labels=(input.glyco_venn_run1_list(),input.glyco_venn_run2_list(),input.glyco_venn_run3_list()),set_colors=("tab:blue","tab:orange","tab:green"),ax=ax)
                venn3_circles(subsets=(Abc,aBc,ABc,abC,AbC,aBC,ABC),linestyle="dashed",linewidth=0.5)
                plt.title("Venn Diagram for "+titlemod)
            fig.set_tight_layout(True)
            imagedownload("glycovenn_"+input.glyco_venn_plotproperty())
    #download table of Venn Diagram intersections
    @reactive.effect
    def _():
        if input.glyco_venn_numcircles()=="2":
            filename=lambda: f"VennList_{input.glyco_venn_run1_list()}_vs_{input.glyco_venn_run2_list()}_{input.glyco_venn_plotproperty()}.csv"
        if input.glyco_venn_numcircles()=="3":
            filename=lambda: f"VennList_A-{input.glyco_venn_run1_list()}_vs_B-{input.glyco_venn_run2_list()}_vs_C-{input.glyco_venn_run3_list()}_{input.glyco_venn_plotproperty()}.csv"
        @render.download(filename=filename)
        def glyco_venn_download():
            resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
            if input.glyco_venn_conditionorrun()=="condition":
                A=glycoPSMs_df[glycoPSMs_df["R.Condition"]==input.glyco_venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
                B=glycoPSMs_df[glycoPSMs_df["R.Condition"]==input.glyco_venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
                if input.glyco_venn_numcircles()=="3":
                    C=glycoPSMs_df[glycoPSMs_df["R.Condition"]==input.glyco_venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
            if input.glyco_venn_conditionorrun()=="individual":
                A=glycoPSMs_df[glycoPSMs_df["Cond_Rep"]==input.glyco_venn_run1_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
                B=glycoPSMs_df[glycoPSMs_df["Cond_Rep"]==input.glyco_venn_run2_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
                if input.glyco_venn_numcircles()=="3":
                    C=glycoPSMs_df[glycoPSMs_df["Cond_Rep"]==input.glyco_venn_run3_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Total Glycan Composition"]].drop_duplicates().reset_index(drop=True)
            
            A["modpep_glycan"]=A["EG.ModifiedPeptide"]+"_"+A["Total Glycan Composition"]
            B["modpep_glycan"]=B["EG.ModifiedPeptide"]+"_"+B["Total Glycan Composition"]

            A["modpep_glycan_charge"]=A["EG.ModifiedPeptide"]+"_"+A["Total Glycan Composition"]+"_"+A["FG.Charge"].astype(str)
            B["modpep_glycan_charge"]=B["EG.ModifiedPeptide"]+"_"+B["Total Glycan Composition"]+"_"+B["FG.Charge"].astype(str)

            if input.glyco_venn_numcircles()=="3":
                C["modpep_glycan"]=C["EG.ModifiedPeptide"]+"_"+C["Total Glycan Composition"]
                C["modpep_glycan_charge"]=C["EG.ModifiedPeptide"]+"_"+C["Total Glycan Composition"]+"_"+C["FG.Charge"].astype(str)

            if input.glyco_venn_plotproperty()=="glycoproteins":
                a=set(A["PG.ProteinGroups"])
                b=set(B["PG.ProteinGroups"])
                if input.glyco_venn_numcircles()=="3":
                    c=set(C["PG.ProteinGroups"])
                titlemod="Glycoproteins"
            if input.glyco_venn_plotproperty()=="glycopeptides":
                a=set(A["modpep_glycan"])
                b=set(B["modpep_glycan"])
                if input.glyco_venn_numcircles()=="3":
                    c=set(C["modpep_glycan"])
                titlemod="Glycopeptides"
            if input.glyco_venn_plotproperty()=="glycoPSMs":
                a=set(A["modpep_glycan_charge"])
                b=set(B["modpep_glycan_charge"])
                if input.glyco_venn_numcircles()=="3":
                    c=set(C["modpep_glycan_charge"])
                titlemod="GlycoPSMs"

            df=pd.DataFrame()
            if input.glyco_venn_numcircles()=="2":
                Ab=list(a-b)
                aB=list(b-a)
                AB=list(a&b)
                df=pd.concat([df,pd.Series(Ab,name=input.glyco_venn_run1_list())],axis=1)
                df=pd.concat([df,pd.Series(aB,name=input.glyco_venn_run2_list())],axis=1)
                df=pd.concat([df,pd.Series(AB,name="Both")],axis=1)
            if input.glyco_venn_numcircles()=="3":
                Abc=list(a-b-c)
                aBc=list(b-a-c)
                ABc=list((a&b)-c)
                abC=list(c-a-b)
                AbC=list((a&c)-b)
                aBC=list((b&c)-a)
                ABC=list(a&b&c)
                df=pd.concat([df,pd.Series(Abc,name="A only")],axis=1)
                df=pd.concat([df,pd.Series(aBc,name="B only")],axis=1)
                df=pd.concat([df,pd.Series(ABc,name="A and B, not C")],axis=1)
                df=pd.concat([df,pd.Series(abC,name="C only")],axis=1)
                df=pd.concat([df,pd.Series(AbC,name="A and C, not B")],axis=1)
                df=pd.concat([df,pd.Series(aBC,name="B and C, not A")],axis=1)
                df=pd.concat([df,pd.Series(ABC,name="ABC")],axis=1)
            with io.BytesIO() as buf:
                df.to_csv(buf,index=False)
                yield buf.getvalue()            

    # ====================================== Glyco ID Tables
    #show data grid for glycoproteins
    @render.data_frame
    def glycoproteins_df_view():
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
        return render.DataGrid(glycoproteins_df,editable=False,height="600px")
    #download shown data grid for glycoproteins
    @render.download(filename=lambda: f"glycoproteins_{input.searchreport()[0]['name']}.csv")
    def glycoproteins_download():
        glycoprotein_table=glycoproteins_df_view.data_view()
        with io.BytesIO() as buf:
            glycoprotein_table.to_csv(buf,index=False)
            yield buf.getvalue()

    #show data grid for glycopeptides
    @render.data_frame
    def glycopeptides_df_view():
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
        return render.DataGrid(glycopeptides_df,editable=False,height="600px")
    #download shown data grid for glycopeptides
    @render.download(filename=lambda: f"glycopeptides_{input.searchreport()[0]['name']}.csv")
    def glycopeptides_download():
        glycopeptide_table=glycopeptides_df_view.data_view()
        with io.BytesIO() as buf:
            glycopeptide_table.to_csv(buf,index=False)
            yield buf.getvalue()

    #show data grid for glycoPSMs
    @render.data_frame
    def glycoPSMs_df_view():
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
        return render.DataGrid(glycoPSMs_df,editable=False,height="600px")
    #download shown data grid for glyco PSMs
    @render.download(filename=lambda: f"glycoPSMs_{input.searchreport()[0]['name']}.csv")
    def glycoPSMs_download():
        glycoPSMs_table=glycoPSMs_df_view.data_view()
        with io.BytesIO() as buf:
            glycoPSMs_table.to_csv(buf,index=False)
            yield buf.getvalue()

    # ====================================== Glycan Tracker
    #generate selectize list of detected glyco mods
    @render.ui
    def glycomodlist_ui():
        #glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_dataframes()
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()

        raw_modlist=[]
        modlist=[]
        allmods=glycoPSMs_df["Total Glycan Composition"].drop_duplicates().tolist()
        for i in range(len(allmods)):
            raw_modlist.append(allmods[i].split("%")[0])
            modstring=raw_modlist[i].replace("(","").replace(")",",").replace(", ","")
            
            modlist.append("".join(i for i in modstring if not i.isdigit()))
        set_modlist=list(set(modlist))
        unique_mods=[]
        for item in set_modlist:
            if item.count(",")>=1:
                unique_mods.append(item.split(","))
            else:
                unique_mods.append([item])
        unique_mods=list(set(list(itertools.chain(*unique_mods))))
        mods_dict=dict()
        for item in unique_mods:
            mods_dict[item]=item

        return ui.input_selectize("found_glycomods","Pick glycan mod to plot data for",choices=mods_dict)
    @render.ui
    def high_mannose_ui():
        if input.found_glycomods()=="Hex":
            return ui.input_switch("high_mannose","Show only high mannose?")
    #plot ID bar graph for selected glyco mod, option to show % instead of counts
    @reactive.effect
    def _():
        @render.plot(width=input.glycomodIDsplot_width(),height=input.glycomodIDsplot_height())
        def glycomod_IDs():
            resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
            picked_mod=input.found_glycomods()

            color=replicatecolors()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            y_padding=input.ypadding()

            resultdf_glycomod=pd.DataFrame()

            num_glycoproteins_mod=[]
            num_glycopeptides_mod=[]
            num_glycoPSMs_mod=[]

            for run in glycoPSMs_df["Cond_Rep"].drop_duplicates():
                if picked_mod!="Hex":
                    num_glycoproteins_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))]["PG.ProteinGroups"].drop_duplicates()))
                    num_glycopeptides_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))]["EG.ModifiedPeptide"].drop_duplicates()))
                    num_glycoPSMs_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))]))
                else:
                    if input.high_mannose()==True:
                        num_glycoproteins_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))&(glycoPSMs_df["High Mannose"]=="True")]["PG.ProteinGroups"].drop_duplicates()))
                        num_glycopeptides_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))&(glycoPSMs_df["High Mannose"]=="True")]["EG.ModifiedPeptide"].drop_duplicates()))
                        num_glycoPSMs_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))&(glycoPSMs_df["High Mannose"]=="True")]))
                    else:
                        num_glycoproteins_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))]["PG.ProteinGroups"].drop_duplicates()))
                        num_glycopeptides_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))]["EG.ModifiedPeptide"].drop_duplicates()))
                        num_glycoPSMs_mod.append(len(glycoPSMs_df[(glycoPSMs_df["Cond_Rep"]==run)&(glycoPSMs_df["Total Glycan Composition"].str.contains(picked_mod))]))

            resultdf_glycomod["Cond_Rep"]=glycoPSMs_df["Cond_Rep"].drop_duplicates().reset_index(drop=True)
            resultdf_glycomod["glycoproteins"]=num_glycoproteins_mod
            resultdf_glycomod["glycopeptides"]=num_glycopeptides_mod
            resultdf_glycomod["glycoPSM"]=num_glycoPSMs_mod

            if input.counts_vs_enrich()=="counts":
                plottingdf=resultdf_glycomod
                ylabel="Counts"
                titlemod="ID Counts for "
            if input.counts_vs_enrich()=="percent":
                resultdf_glyco_enrich=round(resultdf_glycomod.drop(columns=["Cond_Rep"])/resultdf_glyco.drop(columns=["Cond_Rep"])*100,1)
                resultdf_glyco_enrich["Cond_Rep"]=glycoPSMs_df["Cond_Rep"].drop_duplicates().reset_index(drop=True)
                plottingdf=resultdf_glyco_enrich
                ylabel="% of IDs"
                titlemod="% of IDs for "

            fig,ax=plt.subplots(ncols=3,sharex=True,figsize=(input.glycomodIDsplot_width()*(1/plt.rcParams['figure.dpi']),input.glycomodIDsplot_height()*(1/plt.rcParams['figure.dpi'])))
            ax1=ax[0]
            ax2=ax[1]
            ax3=ax[2]

            plottingdf.plot.bar(ax=ax1,x="Cond_Rep",y="glycoproteins",legend=False,width=0.8,color=color,edgecolor="k")
            ax1.bar_label(ax1.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
            ax1.set_ylim(top=max(plottingdf["glycoproteins"].tolist())+y_padding*max(plottingdf["glycoproteins"].tolist()))
            ax1.set_ylabel(ylabel,fontsize=axisfont)
            ax1.set_xlabel("Condition",fontsize=axisfont)
            ax1.set_title("Glycoproteins",fontsize=titlefont)
            ax1.tick_params(axis="both",labelsize=axisfont_labels)

            plottingdf.plot.bar(ax=ax2,x="Cond_Rep",y="glycopeptides",legend=False,width=0.8,color=color,edgecolor="k")
            ax2.bar_label(ax2.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
            ax2.set_ylim(top=max(plottingdf["glycopeptides"].tolist())+y_padding*max(plottingdf["glycopeptides"].tolist()))
            ax2.set_xlabel("Condition",fontsize=axisfont)
            ax2.set_title("Glycopeptides",fontsize=titlefont)
            ax2.tick_params(axis="both",labelsize=axisfont_labels)

            plottingdf.plot.bar(ax=ax3,x="Cond_Rep",y="glycoPSM",legend=False,width=0.8,color=color,edgecolor="k")
            ax3.bar_label(ax3.containers[0],label_type="edge",rotation=90,padding=5,fontsize=labelfont)
            ax3.set_ylim(top=max(plottingdf["glycoPSM"].tolist())+(y_padding+0.1)*max(plottingdf["glycoPSM"].tolist()))
            ax3.set_xlabel("Condition",fontsize=axisfont)
            ax3.set_title("Glyco-PSMs",fontsize=titlefont)
            ax3.tick_params(axis="both",labelsize=axisfont_labels)

            ax1.set_axisbelow(True)
            ax1.grid(linestyle="--")
            ax2.set_axisbelow(True)
            ax2.grid(linestyle="--")
            ax3.set_axisbelow(True)
            ax3.grid(linestyle="--")
            plt.suptitle(titlemod+picked_mod,y=1,fontsize=titlefont)
            fig.set_tight_layout(True)
            imagedownload("glycoids_"+picked_mod)

    # ====================================== Peptide Tracker
    #generate selectize list of stripped peptide sequences
    @render.ui
    def glyco_peplist():
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
        peplist=glycoPSMs_df["PEP.StrippedSequence"].drop_duplicates().sort_values().reset_index(drop=True).tolist()
        pep_dict=dict()
        for pep in peplist:
            pep_dict[pep]=pep
        return ui.input_selectize("glyco_peplist_pick","Pick stripped peptide sequence:",choices=pep_dict)
    #show glycoPSMs for selected stripped peptide sequence 
    @render.data_frame
    def selected_glyco_peplist():
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
        selectedpep=input.glyco_peplist_pick()
        selectedpep_df=glycoPSMs_df[glycoPSMs_df["PEP.StrippedSequence"]==selectedpep]
        return render.DataGrid(selectedpep_df,editable=False,height="600px")
    #download shown table of glycoPSMs for selected stripped peptide sequence
    @render.download(filename=lambda: f"glycoPSMs_{input.glyco_peplist_pick()}.csv")
    def selected_glyco_download():
        resultdf_glyco,glycoproteins_df,glycopeptides_df,glycoPSMs_df=glyco_variables()
        glycopep_table=selected_glyco_peplist.data_view()
        with io.BytesIO() as buf:
            glycopep_table.to_csv(buf,index=False)
            yield buf.getvalue()

    # ====================================== Precursor Scatterplot
    @render.ui
    def glycoscatter_ui():
        searchoutput=metadata_update()
        opts=searchoutput["Cond_Rep"].drop_duplicates().tolist()
        return ui.input_selectize("glycoscatter_pick","Pick run:",choices=opts)
    #scatterplot like the charge/PTM scatter for glycosylated precursors
    @reactive.effect
    def _():
        @render.plot(width=input.glycoscatter_width(),height=input.glycoscatter_height())
        def glycoscatter():
            searchoutput=metadata_update()

            if input.software()=="fragpipe_glyco":
                searchoutput_nonglyco=searchoutput[(searchoutput["Glycan q-value"].isnull()==True)&(searchoutput["Cond_Rep"]==input.glycoscatter_pick())]
                searchoutput_glyco=searchoutput[(searchoutput["Glycan q-value"].isnull()==False)&(searchoutput["Cond_Rep"]==input.glycoscatter_pick())]
            if input.software()=="glycoscape":
                searchoutput_nonglyco=searchoutput[(searchoutput["Total Glycan Composition"].isnull()==True)&(searchoutput["Cond_Rep"]==input.glycoscatter_pick())]
                searchoutput_glyco=searchoutput[(searchoutput["Total Glycan Composition"].isnull()==False)&(searchoutput["Cond_Rep"]==input.glycoscatter_pick())]

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()

            fig,ax=plt.subplots(figsize=(input.glycoscatter_width()*(1/plt.rcParams['figure.dpi']),input.glycoscatter_height()*(1/plt.rcParams['figure.dpi'])))
            ax.scatter(x=searchoutput_nonglyco["FG.PrecMz"],y=searchoutput_nonglyco["EG.IonMobility"],s=2,label="All Other Precursors")
            ax.scatter(x=searchoutput_glyco["FG.PrecMz"],y=searchoutput_glyco["EG.IonMobility"],s=2,label="Glycosylated Precursors")
            ax.set_xlabel("m/z",fontsize=axisfont)
            ax.set_ylabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
            ax.legend(loc="upper left",fontsize=legendfont,markerscale=5)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("glycosylated_precursors")

#endregion

# ============================================================================= MOMA
#region
    #choose whether to import data from individual file names or from a directory
    @render.ui
    def moma_rawfile_input_ui():
        if input.moma_file_or_folder()=="individual":
            return ui.input_text_area("moma_rawfile_input","Paste the path for each .d file you want to upload (note: do not leave whitespace at the end):",width="1500px",autoresize=True,placeholder="ex - C:\\Users\\Data\\K562_500ng_1_Slot1-49_1_3838.d")
        if input.moma_file_or_folder()=="directory":
            return ui.input_text_area("moma_rawfile_input","Paste the path for the directory containing the raw files to upload (note: do not leave whitespace at the end):",width="1500px",autoresize=True,placeholder="ex - C:\\Users\\Data")

    #list to choose raw data file to plot EIM for
    @render.ui
    def moma_rawfile_buttons_ui():
        os_var=os_check()
        if input.moma_file_or_folder()=="individual":
            filelist=list(input.moma_rawfile_input().split("\n"))
            samplenames=[]
            for run in filelist:
                samplenames.append(run.split(os_var)[-1])
        if input.moma_file_or_folder()=="directory":
            try:
                os.chdir(input.moma_rawfile_input())
                cwd=os.getcwd()
                filelist=[]
                for file in os.listdir():
                    if ".d" in file:
                        filelist.append(cwd+os_var+file)
                samplenames=[]
                for run in filelist:
                    samplenames.append(run.split(os_var)[-1])
            except:
                samplenames=[""]
        return ui.input_radio_buttons("moma_rawfile_buttons_pick","Pick raw file for EIM:",choices=samplenames,width="100%")

    #render ui call for dropdown calling Cond_Rep column
    @render.ui
    def moma_cond_rep_list():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_selectize("moma_cond_rep","Pick run from search file:",choices=opts)
    
    #store imported raw file as a reactive value upon button press
    @reactive.calc
    @reactive.event(input.moma_load_rawfile)
    def rawfile_import():
        import alphatims.bruker as atb
        import alphatims.plotting as atp

        os_var=os_check()
        if input.moma_file_or_folder()=="individual":
            run=input.moma_rawfile_input()
        if input.moma_file_or_folder()=="directory":
            directory=input.moma_rawfile_input()
            selectedrun=input.moma_rawfile_buttons_pick()
            run=directory+os_var+selectedrun
        rawfile=atb.TimsTOF(run)
        return rawfile

    #text popup for when the input file is loaded
    @render.text
    def file_uploaded():
        rawfile=rawfile_import()
        if rawfile is None:
            return ""
        else:
            return "Raw File Loaded"

    #table of MOMA events from search file
    @render.data_frame
    def moma_events():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        mztolerance=input.moma_mztolerance()
        rttolerance=input.moma_rttolerance()
        imtolerance=input.moma_imtolerance()
        sample=input.moma_cond_rep()

        columns=["EG.ModifiedPeptide","FG.Charge","EG.IonMobility","EG.ApexRT","FG.PrecMz"]
        df=searchoutput[searchoutput["Cond_Rep"]==sample][["EG.ModifiedPeptide","FG.Charge","EG.IonMobility","EG.ApexRT","FG.PrecMz"]].sort_values(["EG.ApexRT"]).reset_index(drop=True)
        coelutingpeptides=pd.DataFrame(columns=columns)
        for i in range(len(df)):
            if i+1 not in range(len(df)):
                break
            #rtpercentdiff=(abs(df["EG.ApexRT"][i]-df["EG.ApexRT"][i+1])/df["EG.ApexRT"][i])*100
            rtdiff=abs(df["EG.ApexRT"][i]-df["EG.ApexRT"][i+1])
            mzdiff=abs(df["FG.PrecMz"][i]-df["FG.PrecMz"][i+1])
            imdiff=abs(df["EG.IonMobility"][i]-df["EG.IonMobility"][i+1])
            if rtdiff <= rttolerance and mzdiff <= mztolerance and imdiff >= imtolerance:
                coelutingpeptides.loc[len(coelutingpeptides)]=df.iloc[i].tolist()
                coelutingpeptides.loc[len(coelutingpeptides)]=df.iloc[i+1].tolist()

        #adding a column for a rough group number for each group of peptides detected
        for i in range(len(coelutingpeptides)):
            if i+1 not in range(len(coelutingpeptides)):
                break
            #rtpercentdiff=(abs(coelutingpeptides["EG.ApexRT"][i]-coelutingpeptides["EG.ApexRT"][i+1])/coelutingpeptides["EG.ApexRT"][i])*100
            rtdiff=abs(coelutingpeptides["EG.ApexRT"][i]-coelutingpeptides["EG.ApexRT"][i+1])
            mzdiff=abs(coelutingpeptides["FG.PrecMz"][i]-coelutingpeptides["FG.PrecMz"][i+1])
            imdiff=abs(coelutingpeptides["EG.IonMobility"][i]-coelutingpeptides["EG.IonMobility"][i+1])
            if rtdiff <= rttolerance and mzdiff <= mztolerance and imdiff >= imtolerance:
                coelutingpeptides.loc[coelutingpeptides.index[i],"Group"]=i
        #"Group" not in index tells us that there aren't any MOMA events
        try:
            coelutingpeptides=coelutingpeptides[["Group","EG.ModifiedPeptide","FG.Charge","FG.PrecMz","EG.ApexRT","EG.IonMobility"]]
        except KeyError:
            raise KeyError("No MOMA events found") from None

        return render.DataGrid(coelutingpeptides,editable=False)
    
    @render.plot(width=600,height=600)
    def moma_eim():
        rawfile=rawfile_import()
        try:
            mz=float(input.moma_mz())
            rt=float(input.moma_rt())
            mz_window=input.moma_mztolerance()
            rt_window=input.moma_rttolerance()

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            low_mz=mz-mz_window
            high_mz=mz+mz_window
            low_rt=(rt-rt_window)*60
            #(rt-(rt*(rt_window/100)))*60
            high_rt=(rt+rt_window)*60
            #(rt+(rt*(rt_window/100)))*60

            eim_df=rawfile[low_rt: high_rt,:,0,low_mz: high_mz].sort_values("mobility_values")

            fig,ax=plt.subplots()
            ax.plot(eim_df["mobility_values"],eim_df["intensity_values"])
            ax.set_xlabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
            ax.set_ylabel("Intensity",fontsize=axisfont)
            ax.xaxis.set_minor_locator(MultipleLocator(0.025))
            ax.set_title("EIM",fontsize=titlefont)
            ax.tick_params(axis="both",labelsize=axisfont_labels)
        except:
            fig,ax=plt.subplots()

    @render.download(filename=lambda: f"{input.moma_cond_rep()}_MOMA_events.csv")
    def momatable_download():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        mztolerance=input.moma_mztolerance()
        rttolerance=input.moma_rttolerance()
        imtolerance=input.moma_imtolerance()
        sample=input.moma_cond_rep()

        columns=["EG.ModifiedPeptide","FG.Charge","EG.IonMobility","EG.ApexRT","FG.PrecMz","FG.MS2Quantity","PG.ProteinGroups","PG.ProteinNames"]
        df=searchoutput[searchoutput["Cond_Rep"]==sample][["EG.ModifiedPeptide","FG.Charge","EG.IonMobility","EG.ApexRT","FG.PrecMz","FG.MS2Quantity","PG.ProteinGroups","PG.ProteinNames"]].sort_values(["EG.ApexRT"]).reset_index(drop=True)
        coelutingpeptides=pd.DataFrame(columns=columns)
        for i in range(len(df)):
            if i+1 not in range(len(df)):
                break
            #rtpercentdiff=(abs(df["EG.ApexRT"][i]-df["EG.ApexRT"][i+1])/df["EG.ApexRT"][i])*100
            rtdiff=abs(df["EG.ApexRT"][i]-df["EG.ApexRT"][i+1])
            mzdiff=abs(df["FG.PrecMz"][i]-df["FG.PrecMz"][i+1])
            imdiff=abs(df["EG.IonMobility"][i]-df["EG.IonMobility"][i+1])
            if rtdiff <= rttolerance and mzdiff <= mztolerance and imdiff >= imtolerance:
                coelutingpeptides.loc[len(coelutingpeptides)]=df.iloc[i].tolist()
                coelutingpeptides.loc[len(coelutingpeptides)]=df.iloc[i+1].tolist()

        #adding a column for a rough group number for each group of peptides detected
        for i in range(len(coelutingpeptides)):
            if i+1 not in range(len(coelutingpeptides)):
                break
            #rtpercentdiff=(abs(coelutingpeptides["EG.ApexRT"][i]-coelutingpeptides["EG.ApexRT"][i+1])/coelutingpeptides["EG.ApexRT"][i])*100
            rtdiff=abs(coelutingpeptides["EG.ApexRT"][i]-coelutingpeptides["EG.ApexRT"][i+1])
            mzdiff=abs(coelutingpeptides["FG.PrecMz"][i]-coelutingpeptides["FG.PrecMz"][i+1])
            imdiff=abs(coelutingpeptides["EG.IonMobility"][i]-coelutingpeptides["EG.IonMobility"][i+1])
            if rtdiff <= rttolerance and mzdiff <= mztolerance and imdiff >= imtolerance:
                coelutingpeptides.loc[coelutingpeptides.index[i],"Group"]=i
        coelutingpeptides=coelutingpeptides[["Group","EG.ModifiedPeptide","FG.Charge","FG.PrecMz","EG.ApexRT","EG.IonMobility","FG.MS2Quantity","PG.ProteinGroups","PG.ProteinNames"]]
        with io.BytesIO() as buf:
            coelutingpeptides.to_csv(buf,index=False)
            yield buf.getvalue()

    @render.download(filename=lambda: f"{input.moma_rawfile_buttons_pick()}_MOMA_precursors.csv")
    def precursortable_download():
        rawfile=rawfile_import()
        mz=float(input.moma_mz())
        rt=float(input.moma_rt())
        mz_window=input.moma_mztolerance()
        rt_window=input.moma_rttolerance()

        low_mz=mz-mz_window
        high_mz=mz+mz_window
        low_rt=(rt-rt_window)*60
        #(rt-(rt*(rt_window/100)))*60
        high_rt=(rt+rt_window)*60
        #(rt+(rt*(rt_window/100)))*60

        eim_df=rawfile[low_rt: high_rt,:,0,low_mz: high_mz].sort_values("mobility_values")
        with io.BytesIO() as buf:
            eim_df.to_csv(buf,index=False)
            yield buf.getvalue()

#endregion

# ============================================================================= De Novo
#region
    # ====================================== Secondary File Import
    #region

    #software choices to simplify the interface
    @render.ui
    def software_secondary_ui():
        software=input.software_secondary_general()
        if software=="spectronaut":
            opts={"spectronaut":"directDIA / library-based search",
                  "ddalibrary":"DDA Library"}
        if software=="diann":
            opts={"diann":"DIA-NN pre 2.0",
                  "diann2.0":"DIA-NN 2.0"}
        if software=="fragpipe":
            opts={"fragpipe":"FragPipe",
                  "fragpipe_glyco":"FragPipe Glyco",
                  "fragpipe_combined_ion":"FragPipe Quant"}
        if software=="bps":
            opts={"bps_timsrescore":"tims-rescore",
                  "bps_timsdiann":"tims-DIANN",
                  "bps_spectronaut":"Spectronaut",
                  "bps_pulsar":"Pulsar",
                  "bps_denovo":"BPS Novor",
                  "bps_sage":"Sage",
                  "glycoscape":"Glycoscape"}
        if software=="spectromine":
            opts={"spectromine":"Spectromine"}
        if software=="peaks":
            opts={"peaks":"PEAKS"}
        if software=="sage":
            opts={"sage":"Sage",
                  "sage_lfq":"Sage LFQ"}
        return ui.input_radio_buttons("software_secondary","",choices=opts)
    
    # ====================================== UI and calc functions
    #give a reminder for what to do with search reports from different software
    @render.text
    def metadata_reminder_secondary():
        text=metadata_reminder_text(input.software_secondary())
        return text
    #store uploaded file as a reactive value so we don't need to call inputfile function each time we need to access the uploaded file
    @reactive.calc
    def fileupload_secondary():
        searchoutput=inputfile(input.searchreport_secondary(),input.software_secondary(),input.software_secondary_bps_report_type(),"off",input.searchreport_secondary_reupload())
        return searchoutput
    #store generated metadata df as a reactive value
    @reactive.calc
    def metadata_secondary_calc():
        metadata=metadata_gen(input.use_uploaded_metadata_secondary(),input.metadata_upload_secondary(),fileupload_secondary())
        return metadata
    #render metadata table
    @render.data_frame
    def metadata_table_secondary():
        metadata=metadata_secondary_calc()
        metadata=metadata[["R.FileName","R.Condition","R.Replicate","remove"]]
        if len(metadata)==0:
            #return non-editable metadata if there's nothing there
            return render.DataGrid(metadata,width="100%")
        else:
            return render.DataGrid(metadata,editable=True,width="100%")
    #store metadata_condition df as a reactive value
    @reactive.calc
    def metadata_condition_secondary_calc():
        metadata_condition=metadata_condition_gen(metadata_table_secondary.data_view(),metadata_secondary_calc(),input.use_uploaded_metadata_secondary(),input.remove_secondary())
        return metadata_condition
    #render metadata_condition table
    @render.data_frame
    def metadata_condition_table_secondary():
        metadata_condition=metadata_condition_secondary_calc()
        if len(metadata_condition)==0:
            return render.DataGrid(metadata_condition,width="100%")
        else:
            return render.DataGrid(metadata_condition,editable=True,width="100%")
    #download metadata table as shown
    @render.download(filename="metadata_"+str(date.today())+".csv")
    def metadata_download_secondary():
        metadata_download=metadata_download_gen(metadata_table_secondary.data_view(),metadata_condition_table_secondary.data_view())
        with io.BytesIO() as buf:
            metadata_download.to_csv(buf,index=False)
            yield buf.getvalue()
    #update the searchoutput df to match how we edited the metadata sheet
    @reactive.calc
    @reactive.event(input.rerun_metadata_secondary,ignore_none=False)
    def metadata_update_secondary():
        searchoutput=metadata_update_gen(fileupload_secondary(),metadata_table_secondary.data_view(),metadata_condition_table_secondary.data_view(),input.remove_secondary(),input.cRAP_filter_secondary())
        return searchoutput
    #checkmark popup to show that the metadata has been applied to the search report
    @render.text
    @reactive.event(input.rerun_metadata_secondary)
    def denovo_metadata_applied_ui():
        return '\u2705'

    #endregion
    
    #calculations for setting up dfs for comparison
    @reactive.calc
    def software_comparison():
        #de novo
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        #secondary software
        searchoutput_secondary=metadata_update_secondary()

        #make dfs to use in later calculations, drop duplicates to only have unique sequences to each run
        bps_df=searchoutput[["Cond_Rep","PEP.StrippedSequence","found_in_fasta"]].drop_duplicates().reset_index(drop=True)
        secondary_df=searchoutput_secondary[["Cond_Rep","PEP.StrippedSequence"]].drop_duplicates().reset_index(drop=True)

        #convert Ile to Leu in secondary software (BPS can't differentiate, only writes Leu)
        peplist_subItoL=[]
        for pep in secondary_df["PEP.StrippedSequence"]:
            peplist_subItoL.append(pep.replace("I","L"))
        secondary_df["PEP.StrippedSequence"]=peplist_subItoL

        #calculate and add peptide lengths column to both dfs
        bps_peplen=[]
        for pep in bps_df["PEP.StrippedSequence"]:
            bps_peplen.append(len(pep))
        bps_df["Peptide Length"]=bps_peplen

        secondary_peplen=[]
        for pep in secondary_df["PEP.StrippedSequence"]:
            secondary_peplen.append(len(pep))
        secondary_df["Peptide Length"]=secondary_peplen

        return bps_df,secondary_df

    # ====================================== Compare - Peptide Lengths
    #render ui call for dropdown calling Cond_Rep column
    @render.ui
    def compare_len_samplelist():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_selectize("compare_len_samplelist_pick","Pick run:",choices=opts)
    #bar graph for comparing lengths of detected peptides between BPS Novor and other software
    @reactive.effect
    def _():
        @render.plot(width=input.peplength_compare_width(),height=input.peplength_compare_height())
        def peplength_compare_plot():
            bps_df,secondary_df=software_comparison()

            run=input.compare_len_samplelist_pick()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            titlefont=input.titlefont()
            legendfont=input.legendfont()

            name_mod_dict=software_dict()
            name_mod=name_mod_dict[input.software_secondary()]

            bps_df=bps_df[bps_df["Cond_Rep"]==run]
            secondary_df=secondary_df[secondary_df["Cond_Rep"]==run]

            #secondary software lengths and counts for plotting IDs and unique de novo IDs
            secondary_x=list(set(itertools.chain(*[list(group) for key, group in groupby(sorted(secondary_df["Peptide Length"]))])))
            secondary_y=[len(list(group)) for key, group in groupby(sorted(secondary_df["Peptide Length"]))]
            secondary_len_counts=pd.DataFrame()
            secondary_len_counts["Peptide Length"]=secondary_x
            secondary_len_counts["Secondary"]=secondary_y
            secondary_len_counts=secondary_len_counts.set_index("Peptide Length")

            #intersections of the two software
            A=set(secondary_df["PEP.StrippedSequence"])
            B=set(bps_df["PEP.StrippedSequence"])

            AvsB=list(A-B)
            BvsA=list(B-A)

            AnotB=len(A-B)
            BnotA=len(B-A)
            bothAB=len(A&B)

            #build a df of just the peptides unique to denovo
            unique_to_denovo=pd.DataFrame()
            unique_to_denovo=pd.concat([unique_to_denovo,pd.Series(BvsA,name="PEP.StrippedSequence")],axis=1)
            unique_peplen=[]
            for pep in unique_to_denovo["PEP.StrippedSequence"]:
                unique_peplen.append(len(pep))
            unique_to_denovo["Peptide Length"]=unique_peplen
            unique_to_denovo_x=list(set(itertools.chain(*[list(group) for key, group in groupby(sorted(unique_to_denovo["Peptide Length"]))])))
            unique_to_denovo_y=[len(list(group)) for key, group in groupby(sorted(unique_to_denovo["Peptide Length"]))]
            unique_to_denovo_len_counts=pd.DataFrame()
            unique_to_denovo_len_counts["Peptide Length"]=unique_to_denovo_x
            unique_to_denovo_len_counts["Denovo Unique"]=unique_to_denovo_y
            unique_to_denovo_len_counts=unique_to_denovo_len_counts.set_index("Peptide Length")

            peplengths_combined=pd.concat([secondary_len_counts,unique_to_denovo_len_counts],axis=1).sort_values("Peptide Length").fillna(0)
            peplengths_combined
            fig,ax=plt.subplots(figsize=(input.peplength_compare_width()*(1/plt.rcParams['figure.dpi']),input.peplength_compare_height()*(1/plt.rcParams['figure.dpi'])))
            x=list(peplengths_combined.index)
            ax.bar(x,peplengths_combined["Secondary"],label=name_mod)
            ax.bar(x,peplengths_combined["Denovo Unique"],label="BPS Novor (Unique)",bottom=peplengths_combined["Secondary"])
            ax.legend(loc="upper right",fontsize=legendfont)
            ax.xaxis.set_major_locator(MultipleLocator(2))
            ax.set_title(run,fontsize=titlefont)
            ax.set_xlabel("Peptide Length",fontsize=axisfont)
            ax.set_ylabel("Counts",fontsize=axisfont)
            ax.set_ylim(bottom=-(0.025*ax.get_ylim()[1]))
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("compare_peptidelength")

    # ====================================== Compare - Stripped Peptide IDs
    #render ui call for dropdown calling Cond_Rep column
    @render.ui
    def denovocompare_venn_samplelist():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_selectize("denovocompare_venn_samplelist_pick","Pick run:",choices=opts)
    @render.ui
    def denovocompare_specific_length_ui():
        if input.denovocompare_specific_length()==True:
            return ui.input_slider("denovocompare_specific_length_pick","Pick specific peptide length to compare:",min=3,max=30,value=9,step=1,ticks=True)
    #plot venn diagram of stripped peptide IDs between the two software
    @reactive.effect
    def _():
        @render.plot(width=input.denovocompare_venn_width(),height=input.denovocompare_venn_height())
        def denovocompare_venn_plot():
            from matplotlib_venn import venn2,venn2_circles,venn3,venn3_circles

            bps_df,secondary_df=software_comparison()

            run=input.denovocompare_venn_samplelist_pick()
            titlefont=input.titlefont()

            name_mod_dict=software_dict()
            name_mod=name_mod_dict[input.software_secondary()]

            if input.denovocompare_specific_length()==True:
                bps_df=bps_df[(bps_df["Cond_Rep"]==run)&(bps_df["Peptide Length"]==int(input.denovocompare_specific_length_pick()))]
                secondary_df=secondary_df[(secondary_df["Cond_Rep"]==run)&(secondary_df["Peptide Length"]==int(input.denovocompare_specific_length_pick()))]
                titlemod=" "+str(input.denovocompare_specific_length_pick())+"mers"
            else:
                bps_df=bps_df[bps_df["Cond_Rep"]==run]
                secondary_df=secondary_df[secondary_df["Cond_Rep"]==run]
                titlemod=""

            if input.denovocompare_peptidecore()==True:
                bps_df_coreAA=[]
                secondary_df_coreAA=[]
                for pep in bps_df["PEP.StrippedSequence"].tolist():
                    bps_df_coreAA.append(pep[2:-2])
                for pep in secondary_df["PEP.StrippedSequence"].tolist():
                    secondary_df_coreAA.append(pep[2:-2])
                A=set(secondary_df_coreAA)
                B=set(bps_df_coreAA)

            else:
                #intersections of the two software
                A=set(secondary_df["PEP.StrippedSequence"])
                B=set(bps_df["PEP.StrippedSequence"])

            AvsB=list(A-B)
            BvsA=list(B-A)

            AnotB=len(A-B)
            BnotA=len(B-A)
            bothAB=len(A&B)

            fig,ax=plt.subplots(figsize=(input.denovocompare_venn_width()*(1/plt.rcParams['figure.dpi']),input.denovocompare_venn_height()*(1/plt.rcParams['figure.dpi'])))
            vennlist=[AnotB,BnotA,bothAB]
            venn2(subsets=vennlist,set_labels=(name_mod,"BPS Novor"),set_colors=("tab:blue","tab:orange"),ax=ax)
            venn2_circles(subsets=vennlist,linestyle="dashed",linewidth=0.5)
            plt.title(run+" Stripped Peptide IDs"+titlemod,fontsize=titlefont)
            fig.set_tight_layout(True)
            imagedownload("compare_strippedsequences_venn")
    #download list of common and unique IDs
    @render.download(filename=lambda: f'{input.denovocompare_venn_samplelist_pick()}_BPSNovor_Venn.csv')
    def denovocompare_venn_download():
        bps_df,secondary_df=software_comparison()

        run=input.denovocompare_venn_samplelist_pick()
        titlefont=input.titlefont()

        name_mod_dict=software_dict()
        name_mod=name_mod_dict[input.software_secondary()]

        if input.denovocompare_specific_length()==True:
            bps_df=bps_df[(bps_df["Cond_Rep"]==run)&(bps_df["Peptide Length"]==int(input.denovocompare_specific_length_pick()))]
            secondary_df=secondary_df[(secondary_df["Cond_Rep"]==run)&(secondary_df["Peptide Length"]==int(input.denovocompare_specific_length_pick()))]
        else:
            bps_df=bps_df[bps_df["Cond_Rep"]==run]
            secondary_df=secondary_df[secondary_df["Cond_Rep"]==run]

        if input.denovocompare_peptidecore()==True:
            bps_df_coreAA=[]
            secondary_df_coreAA=[]
            for pep in bps_df["PEP.StrippedSequence"].tolist():
                bps_df_coreAA.append(pep[2:-2])
            for pep in secondary_df["PEP.StrippedSequence"].tolist():
                secondary_df_coreAA.append(pep[2:-2])
            A=set(secondary_df_coreAA)
            B=set(bps_df_coreAA)

        else:
            #intersections of the two software
            A=set(secondary_df["PEP.StrippedSequence"])
            B=set(bps_df["PEP.StrippedSequence"])

        AvsB=list(A-B)
        BvsA=list(B-A)
        ABcommon=list(A&B)

        AnotB=len(A-B)
        BnotA=len(B-A)
        bothAB=len(A&B)

        unique_ids_df=pd.DataFrame()
        unique_ids_df=pd.concat([unique_ids_df,pd.Series(ABcommon,name="Common IDs"),pd.Series(AvsB,name=name_mod+" Unique"),pd.Series(BvsA,name="BPS Novor Unique")],axis=1)

        with io.BytesIO() as buf:
            unique_ids_df.to_csv(buf,index=False)
            yield buf.getvalue()

    # ====================================== Compare - Sequence Motifs
    @render.ui
    def seqmotif_compare_run_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_selectize("seqmotif_compare_run_pick","Pick run:",choices=opts)
    #BPS sequence motif plot
    @render.plot(width=800,height=400)
    def seqmotif_compare_plot1():
        import logomaker as lm
        bps_df,secondary_df=software_comparison()
        titlefont=input.titlefont()
        axisfont=input.axisfont()
        axisfont_labels=input.axisfont_labels()

        if input.seqmotif_compare_onlyunique()==False:
            seq=bps_df[(bps_df["Cond_Rep"]==input.seqmotif_compare_run_pick())&(bps_df["Peptide Length"]==input.seqmotif_compare_peplengths())]["PEP.StrippedSequence"].drop_duplicates().tolist()
            titlemod=""
        if input.seqmotif_compare_onlyunique()==True:
            titlemod=" Only Unique IDs"
            bps_df=bps_df[(bps_df["Cond_Rep"]==input.seqmotif_compare_run_pick())&(bps_df["Peptide Length"]==int(input.seqmotif_compare_peplengths()))]
            secondary_df=secondary_df[(secondary_df["Cond_Rep"]==input.seqmotif_compare_run_pick())&(secondary_df["Peptide Length"]==int(input.seqmotif_compare_peplengths()))]
            A=set(secondary_df["PEP.StrippedSequence"])
            B=set(bps_df["PEP.StrippedSequence"])
            seq=list(A-B)

        matrix=lm.alignment_to_matrix(seq)
        if input.seqmotif_compare_plottype()=="counts":
            ylabel="Counts"
        if input.seqmotif_compare_plottype()=="information":
            matrix=lm.transform_matrix(matrix,from_type="counts",to_type="information")
            ylabel="Information (bits)"
        fig,ax=plt.subplots(figsize=(800*(1/plt.rcParams['figure.dpi']),400*(1/plt.rcParams['figure.dpi'])))
        logo=lm.Logo(matrix,ax=ax,vsep=0.01,vpad=0.01,color_scheme="weblogo_protein")
        logo.ax.set_xlabel("Position",fontsize=axisfont)
        logo.ax.set_ylabel(ylabel,fontsize=axisfont)
        logo.ax.set_title("BPS "+input.seqmotif_compare_run_pick()+": "+str(input.seqmotif_peplengths())+"mers"+titlemod,fontsize=titlefont)
        logo.ax.set_ylim(bottom=0)
        logo.ax.tick_params(axis="both",labelsize=axisfont_labels)
        logo.fig.tight_layout()
        imagedownload("sequencemotifs_BPSNovor_"+str(input.seqmotif_peplengths())+"mers")
    #secondary sequence motif plot
    @render.plot(width=800,height=400)
    def seqmotif_compare_plot2():
        import logomaker as lm
        bps_df,secondary_df=software_comparison()

        seq2=secondary_df[(secondary_df["Cond_Rep"]==input.seqmotif_compare_run_pick())&(secondary_df["Peptide Length"]==input.seqmotif_compare_peplengths())]["PEP.StrippedSequence"].drop_duplicates().tolist()
        titlefont=input.titlefont()
        axisfont=input.axisfont()
        axisfont_labels=input.axisfont_labels()

        name_mod_dict=software_dict()
        name_mod=name_mod_dict[input.software_secondary()]

        matrix2=lm.alignment_to_matrix(seq2)
        if input.seqmotif_compare_plottype()=="counts":
            ylabel="Counts"
        if input.seqmotif_compare_plottype()=="information":
            matrix2=lm.transform_matrix(matrix2,from_type="counts",to_type="information")
            ylabel="Information (bits)"
        fig,ax=plt.subplots(figsize=(800*(1/plt.rcParams['figure.dpi']),400*(1/plt.rcParams['figure.dpi'])))
        logo2=lm.Logo(matrix2,ax=ax,vsep=0.01,vpad=0.01,color_scheme="weblogo_protein")
        logo2.ax.set_xlabel("Position",fontsize=axisfont)
        logo2.ax.set_ylabel(ylabel,fontsize=axisfont)
        logo2.ax.set_title(name_mod+" "+input.seqmotif_compare_run_pick()+": "+str(input.seqmotif_peplengths())+"mers",fontsize=titlefont)
        logo2.ax.set_ylim(bottom=0)
        logo2.ax.tick_params(axis="both",labelsize=axisfont_labels)
        logo2.fig.tight_layout()
        imagedownload("sequencemotifs_"+name_mod+"_"+str(input.seqmotif_peplengths())+"mers")

    # ====================================== Compare - Sequence Motifs (stats)
    @render.ui
    def seqmotif_compare_run_ui2():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=runnames
        return ui.input_selectize("seqmotif_compare_run_pick2","Pick run:",choices=opts)
    #PCA plot
    @render.plot(height=600)
    def seqmotif_pca():
        import logomaker as lm
        from sklearn.decomposition import PCA
        from sklearn.pipeline import Pipeline
        from sklearn.preprocessing import StandardScaler

        bps_df,secondary_df=software_comparison()

        titlefont=input.titlefont()
        axisfont=input.axisfont()
        axisfont_labels=input.axisfont_labels()
        labelfont=input.labelfont()
        legendfont=input.legendfont()
        y_padding=input.ypadding()

        seq_BPS=bps_df[(bps_df["Cond_Rep"]==input.seqmotif_compare_run_pick2())&(bps_df["Peptide Length"]==int(input.seqmotif_compare_peplengths2()))]["PEP.StrippedSequence"].drop_duplicates().tolist()
        seq_secondary=secondary_df[(secondary_df["Cond_Rep"]==input.seqmotif_compare_run_pick2())&(secondary_df["Peptide Length"]==int(input.seqmotif_compare_peplengths2()))]["PEP.StrippedSequence"].drop_duplicates().tolist()

        name_mod_dict=software_dict()
        name_mod=name_mod_dict[input.software_secondary()]

        if input.seqmotif_compare_onlyunique2()==True:
            a=set(seq_secondary)
            b=set(seq_BPS)
            seq_BPS=list(a-b)
            B=lm.alignment_to_matrix(seq_BPS,to_type="information")
            titlemod=" Only Unique BPS IDs"
        if input.seqmotif_compare_onlyunique2()==False:
            B=lm.alignment_to_matrix(seq_BPS,to_type="information")
            titlemod=""

        A=lm.alignment_to_matrix(seq_secondary,to_type="information")
        for col in A.columns.tolist():
            A=A.rename(columns={col:str(col)+"_secondary"})
        for col in B.columns.tolist():
            B=B.rename(columns={col:str(col)+"_BPS"})
        concatenated=pd.concat([A,B],axis=1)

        X=np.array(concatenated).T
        pip=Pipeline([("scaler",StandardScaler()),("pca",PCA())]).fit(X)
        X_trans=pip.transform(X)

        fig,ax=plt.subplots()

        colors=["firebrick","red","lightcoral","orange","gold","goldenrod","green","lime","darkturquoise","cyan",
                "dodgerblue","lightsteelblue","blue","blueviolet","purple","fuchsia","violet","pink","slategrey"]
        for i in range(len(concatenated.columns.tolist())):
            if i>18:
                z=i-19
                label="B"
            else:
                z=i
                label="S"
            ax.scatter(X_trans[i,0],X_trans[i,1],color=colors[z],s=45,label=concatenated.columns.tolist()[i].split("_")[0])
            ax.annotate(label,(X_trans[i,0]+0.05,X_trans[i,1]+0.05),fontsize=8)
        ax.set_xlabel("PC1"+" ("+str(round(pip.named_steps.pca.explained_variance_ratio_[0]*100,1))+"%)",fontsize=axisfont)
        ax.set_ylabel("PC2"+" ("+str(round(pip.named_steps.pca.explained_variance_ratio_[1]*100,1))+"%)",fontsize=axisfont)
        handles,labels=ax.get_legend_handles_labels()
        handle_list,label_list =[],[]
        for handle,label in zip(handles,labels):
            if label not in label_list:
                handle_list.append(handle)
                label_list.append(label)
        ax.legend(handle_list,label_list,loc="center",bbox_to_anchor=(1.1,0.5),markerscale=1,prop={'size':legendfont})
        ax.set_title("BPS vs. "+name_mod+": "+input.seqmotif_compare_run_pick2()+" "+str(input.seqmotif_compare_peplengths2())+"mers"+titlemod,fontsize=titlefont)
        ax.tick_params(axis="both",labelsize=axisfont_labels)
        fig.set_tight_layout(True)
    #3D plot
    @render.plot(height=600)
    def seqmotif_3d():
        import logomaker as lm
        bps_df,secondary_df=software_comparison()

        titlefont=input.titlefont()
        axisfont=input.axisfont()
        axisfont_labels=input.axisfont_labels()
        labelfont=input.labelfont()
        legendfont=input.legendfont()
        y_padding=input.ypadding()

        name_mod_dict=software_dict()
        name_mod=name_mod_dict[input.software_secondary()]

        seq_BPS=bps_df[(bps_df["Cond_Rep"]==input.seqmotif_compare_run_pick2())&(bps_df["Peptide Length"]==int(input.seqmotif_compare_peplengths2()))]["PEP.StrippedSequence"].drop_duplicates().tolist()
        seq_secondary=secondary_df[(secondary_df["Cond_Rep"]==input.seqmotif_compare_run_pick2())&(secondary_df["Peptide Length"]==int(input.seqmotif_compare_peplengths2()))]["PEP.StrippedSequence"].drop_duplicates().tolist()

        if input.seqmotif_compare_onlyunique2()==True:
            a=set(seq_secondary)
            b=set(seq_BPS)
            seq_BPS=list(a-b)
            B=lm.alignment_to_matrix(seq_BPS,to_type="information")
            titlemod=" Only Unique BPS IDs"
        if input.seqmotif_compare_onlyunique2()==False:
            B=lm.alignment_to_matrix(seq_BPS,to_type="information")
            titlemod=""

        A=lm.alignment_to_matrix(seq_secondary,to_type="information")
        
        xdata=np.arange(1,20).tolist()
        ydata=np.arange(1,1+int(input.seqmotif_compare_peplengths2())).tolist()
        columns=A.columns.tolist()
        colors=["firebrick","red","lightcoral","orange","gold","goldenrod","green","lime","darkturquoise","cyan",
                "dodgerblue","lightsteelblue","blue","blueviolet","purple","fuchsia","violet","pink","slategrey"]
        fig=plt.figure()
        ax=fig.add_subplot(111,projection='3d')
        for i in range(len(columns)):
            zdata_BPS=B[columns[i]]
            zdata_fragger=A[columns[i]]
            if i>18:
                z=i-19
            else:
                z=i
            ax.scatter(xdata[i],ydata,zdata_BPS,color=colors[z],label=columns[i])
            ax.scatter(xdata[i],ydata,zdata_fragger,color=colors[z])
        ax.view_init(azim=input.seqmotif_3d_azimuth(),elev=input.seqmotif_3d_elevation())
        ax.set_xticks(xdata,columns)
        ax.set_xlabel("Residue",fontsize=axisfont)
        ax.set_ylabel("Position",fontsize=axisfont)
        ax.legend(loc="center",bbox_to_anchor=(1.1,0.5),prop={"size":legendfont})
        ax.set_title("BPS vs. "+name_mod+": "+input.seqmotif_compare_run_pick2()+" "+str(input.seqmotif_compare_peplengths2())+"mers"+titlemod,fontsize=titlefont)
        ax.tick_params(axis="both",labelsize=axisfont_labels)
        fig.set_tight_layout(True)

    # ====================================== IDs Found in Fasta
    @reactive.effect
    def _():
        @render.plot(width=input.fasta_width(),height=input.fasta_height())
        def fasta_plot():
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            cond_rep_list=runnames
            fasta_df=pd.DataFrame()
            fasta_true=[]
            fasta_false=[]
            for run in cond_rep_list:
                df=searchoutput[searchoutput["Cond_Rep"]==run]
                fasta_false.append(len(df[df["found_in_fasta"]==False]))
                fasta_true.append(len(df[df["found_in_fasta"]==True]))
            fasta_df["Cond_Rep"]=cond_rep_list
            fasta_df["Fasta_True"]=fasta_true
            fasta_df["Fasta_False"]=fasta_false

            fig,ax=plt.subplots(figsize=(input.fasta_width()*(1/plt.rcParams['figure.dpi']),input.fasta_height()*(1/plt.rcParams['figure.dpi'])))
            x=np.arange(len(fasta_df))
            width=0.4
            
            labelfont=input.labelfont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            legendfont=input.legendfont()     
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            maxvalue=max(fasta_df[["Fasta_True","Fasta_False"]].max().tolist())
            ax.bar(x,fasta_df["Fasta_True"],width=width,label="Fasta=True",edgecolor="k")
            ax.bar(x+width,fasta_df["Fasta_False"],width=width,label="Fasta=False",edgecolor="k")

            ax.bar_label(ax.containers[0],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
            ax.bar_label(ax.containers[1],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
            ax.set_xticks(x+(width/2),fasta_df["Cond_Rep"],rotation=x_label_rotation)
            ax.set_ylim(top=maxvalue+(y_padding*maxvalue))
            ax.margins(x=0.02)
            ax.set_xlabel("Run",fontsize=axisfont)
            ax.set_ylabel("Counts",fontsize=axisfont)
            ax.legend(fontsize=legendfont)
            ax.set_axisbelow(True)
            ax.grid(linestyle="--")
            ax.tick_params(axis="both",labelsize=axisfont_labels)
            fig.set_tight_layout(True)
            imagedownload("idsfoundinfasta")

    # ====================================== Position Confidence
    @render.ui
    def confidence_condition_ui():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        opts=sampleconditions
        return ui.input_selectize("confidence_condition_pick","Pick sample condition:",choices=opts)
    @reactive.effect
    def _():
        @render.plot(width=input.confidence_width(),height=input.confidence_height())
        def confidence_plot():
            peplen=input.confidence_lengthslider()
            axisfont=input.axisfont()
            titlefont=input.titlefont()
            axisfont_labels=input.axisfont_labels()

            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            listoflengths=[]
            for pep in searchoutput["PEP.StrippedSequence"]:
                listoflengths.append(len(pep))
            searchoutput_len=searchoutput
            searchoutput_len["Peptide Length"]=listoflengths

            condition_pick=input.confidence_condition_pick()
            confidence_plot_df=searchoutput_len[searchoutput_len["R.Condition"]==condition_pick]

            fig,ax=plt.subplots(nrows=len(confidence_plot_df["Cond_Rep"].drop_duplicates()),sharex=True,figsize=(input.confidence_width()*(1/plt.rcParams['figure.dpi']),input.confidence_height()*(1/plt.rcParams['figure.dpi'])))

            for i,run in enumerate(confidence_plot_df["Cond_Rep"].drop_duplicates().tolist()):

                df=confidence_plot_df[(confidence_plot_df["Cond_Rep"]==run)&(confidence_plot_df["Peptide Length"]==peplen)]
                columns=list(np.arange(peplen).astype(str))
                denovo_confidence_df=pd.DataFrame(df["denovo_local_confidence"].tolist(),columns=columns)
                
                position_confidence=[]
                for col in denovo_confidence_df.columns:
                    position_confidence.append(denovo_confidence_df[col].tolist())
                plottingdf=pd.DataFrame()
                plottingdf["Position"]=columns
                plottingdf["Confidence"]=position_confidence

                medianlineprops=dict(linestyle="--",color="black")
                flierprops=dict(markersize=3)
                bplot=ax[i].boxplot(plottingdf["Confidence"],medianprops=medianlineprops,flierprops=flierprops,patch_artist=True)
                ax[i].set_ylabel("Confidence (%)",fontsize=axisfont)
                ax[i].set_title(run,fontsize=titlefont)
                
                ax[i].set_axisbelow(True)
                ax[i].grid(linestyle="--")
                ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                
                colors=mcolors.TABLEAU_COLORS
                for patch,color in zip(bplot["boxes"],colors):
                    patch.set_facecolor(color)
            plt.xlabel("Position",fontsize=axisfont)
            fig.set_tight_layout(True)
            imagedownload("positionconfidence")

#endregion

# ============================================================================= Two-Software Comparison
#region
    # ====================================== File Import
    @render.ui
    def software1_ui():
        software=input.software1_general()
        if software=="spectronaut":
            opts={"spectronaut":"directDIA / library-based search",
                  "ddalibrary":"DDA Library"}
        if software=="diann":
            opts={"diann":"DIA-NN pre 2.0",
                  "diann2.0":"DIA-NN 2.0"}
        if software=="fragpipe":
            opts={"fragpipe":"FragPipe",
                  "fragpipe_glyco":"FragPipe Glyco",
                  "fragpipe_combined_ion":"FragPipe Quant"}
        if software=="bps":
            opts={"bps_timsrescore":"tims-rescore",
                  "bps_timsdiann":"tims-DIANN",
                  "bps_spectronaut":"Spectronaut",
                  "bps_pulsar":"Pulsar",
                  "bps_denovo":"BPS Novor",
                  "bps_sage":"Sage",
                  "glycoscape":"Glycoscape"}
        if software=="spectromine":
            opts={"spectromine":"Spectromine"}
        if software=="peaks":
            opts={"peaks":"PEAKS"}
        if software=="sage":
            opts={"sage":"Sage",
                  "sage_lfq":"Sage LFQ"}
        return ui.input_radio_buttons("software1","",choices=opts)
    @render.ui
    def software2_ui():
        software=input.software2_general()
        if software=="spectronaut":
            opts={"spectronaut":"directDIA / library-based search",
                  "ddalibrary":"DDA Library"}
        if software=="diann":
            opts={"diann":"DIA-NN pre 2.0",
                  "diann2.0":"DIA-NN 2.0"}
        if software=="fragpipe":
            opts={"fragpipe":"FragPipe",
                  "fragpipe_glyco":"FragPipe Glyco",
                  "fragpipe_combined_ion":"FragPipe Quant"}
        if software=="bps":
            opts={"bps_timsrescore":"tims-rescore",
                  "bps_timsdiann":"tims-DIANN",
                  "bps_spectronaut":"Spectronaut",
                  "bps_pulsar":"Pulsar",
                  "bps_denovo":"BPS Novor",
                  "bps_sage":"Sage",
                  "glycoscape":"Glycoscape"}
        if software=="spectromine":
            opts={"spectromine":"Spectromine"}
        if software=="peaks":
            opts={"peaks":"PEAKS"}
        if software=="sage":
            opts={"sage":"Sage",
                  "sage_lfq":"Sage LFQ"}
        return ui.input_radio_buttons("software2","",choices=opts)
    #table for coloring options
    @render.data_frame
    def compare_colortable():
        name_mod_dict=software_dict()
        software1=name_mod_dict[input.software1()]
        software2=name_mod_dict[input.software2()]

        colortable=pd.DataFrame()
        colortable["Software"]=[software1,software2]
        colortable["Color"]=["tab:blue","tab:orange"]
        return render.DataGrid(colortable,editable=True,width="100%")
    #function for figure download
    def imagedownload_twosoftware(plotname):
        if input.twosoftware_download_path()!="":
            os_var=os_check()
            filename=input.software1()+"_vs_"+input.software2()+"_"+plotname
            plt.savefig(input.twosoftware_download_path()+os_var+filename+".png")
   
    # ====================================== Metadata functions (need new ones here because of need for two FileName columns)
    #generate metadata df
    def compare_metadata_gen(use_uploaded_metadata,metadata_upload,file1,file2):
        if use_uploaded_metadata==True:
            if metadata_upload is None:
                metadata=pd.DataFrame(columns=["R.FileName1","R.FileName2","R.Condition","R.Replicate","remove"])
            else:
                metadata=pd.read_csv(metadata_upload[0]["datapath"],sep=",")
                metadata["remove"]=metadata["remove"].fillna("")
        else:
            #fileupload outputs an empty df if input.searchreport() is None
            if file1.empty==True or file2.empty==True:
                metadata=pd.DataFrame(columns=["R.FileName1","R.FileName2","R.Condition","R.Replicate","remove"])
            else:
                metadata=pd.DataFrame(columns=["R.FileName1","R.FileName2","R.Condition","R.Replicate","remove"])
                metadata["R.FileName1"]=file1["R.FileName"].drop_duplicates().tolist()
                metadata["R.FileName2"]=file2["R.FileName"].drop_duplicates().tolist()
                metadata["R.Condition"]=[""]*len(metadata["R.FileName1"])
                metadata["R.Replicate"]=[""]*len(metadata["R.FileName1"])
                metadata["remove"]=[""]*len(metadata["R.FileName1"])
        return metadata
    #generate df for downloading metadata
    def compare_metadata_download_gen(metadata,metadata_condition):
        metadata_download=pd.DataFrame()
        metadata_download["R.FileName1"]=metadata["R.FileName1"]
        metadata_download["R.FileName2"]=metadata["R.FileName2"]
        metadata_download["R.Condition"]=metadata["R.Condition"]
        metadata_download["R.Replicate"]=metadata["R.Replicate"]
        metadata_download["remove"]=metadata["remove"]

        #check if conditions have been removed from the main metadata df
        orderlist=[]
        concentrationlist=[]
        if len(metadata["R.Condition"].drop_duplicates().tolist())!=len(metadata_condition["R.Condition"].tolist()):
            for condition in metadata["R.Condition"].drop_duplicates().tolist():
                if condition not in metadata_condition["R.Condition"].tolist():
                    concentrationlist.append([""]*len(metadata[metadata["R.Condition"]==condition]["R.Replicate"]))
                    orderlist.append([0]*len(metadata[metadata["R.Condition"]==condition]["R.Replicate"]))
                if condition in metadata_condition["R.Condition"].tolist():
                    orderlist.append([metadata_condition[metadata_condition["R.Condition"]==condition]["order"].values[0]]*len(metadata[metadata["R.Condition"]==condition]["R.Replicate"]))
                    concentrationlist.append([metadata_condition[metadata_condition["R.Condition"]==condition]["Concentration"].values[0]]*len(metadata[metadata["R.Condition"]==condition]["R.Replicate"]))
        else:
            for run in metadata_condition["R.Condition"]:
                fileindex=metadata_condition[metadata_condition["R.Condition"]==run].index.values[0]
                orderlist.append([metadata_condition["order"][fileindex]]*len(metadata[metadata["R.Condition"]==run]))
                concentrationlist.append([metadata_condition["Concentration"][fileindex]]*len(metadata[metadata["R.Condition"]==run]))
        metadata_download["order"]=list(itertools.chain(*orderlist))
        metadata_download["Concentration"]=list(itertools.chain(*concentrationlist))
        return metadata_download
    #update uploaded search file based on metadata and switches
    def compare_metadata_update_gen(file1,file2,metadata,metadata_condition,remove,cRAP_filter):
        #remove checked runs from search files
        if remove==True:
            metadata=metadata[metadata["remove"]!="x"].reset_index(drop=True)
            file1=file1.set_index("R.FileName").loc[metadata["R.FileName1"].tolist()].reset_index()
            file2=file2.set_index("R.FileName").loc[metadata["R.FileName2"].tolist()].reset_index()

        #update condition/replicate names/values from metadata. if not filled out, autofill with generic values
        RConditionlist=[]
        RReplicatelist=[]
        for i,run in enumerate(file1["R.FileName"].drop_duplicates().tolist()):
            fileindex=metadata[metadata["R.FileName1"]==run].index.values[0]
            if metadata["R.Condition"][fileindex]=="":
                RConditionlist.append(["Not Defined"]*len(file1.set_index("R.FileName").loc[run]))
            else:
                RConditionlist.append([metadata["R.Condition"][fileindex]]*len(file1.set_index("R.FileName").loc[run]))
            if metadata["R.Replicate"][fileindex]=="":
                RReplicatelist.append([i+1]*len(file1.set_index("R.FileName").loc[run]))
            else:
                RReplicatelist.append([metadata["R.Replicate"][fileindex]]*len(file1.set_index("R.FileName").loc[run]))
        file1["R.Condition"]=list(itertools.chain(*RConditionlist))
        file1["R.Replicate"]=list(itertools.chain(*RReplicatelist))
        file1["R.Replicate"]=file1["R.Replicate"].astype(int)
        #update condition/replicate names/values from metadata. if not filled out, autofill with generic values
        RConditionlist=[]
        RReplicatelist=[]
        for i,run in enumerate(file2["R.FileName"].drop_duplicates().tolist()):
            fileindex=metadata[metadata["R.FileName2"]==run].index.values[0]
            if metadata["R.Condition"][fileindex]=="":
                RConditionlist.append(["Not Defined"]*len(file2.set_index("R.FileName").loc[run]))
            else:
                RConditionlist.append([metadata["R.Condition"][fileindex]]*len(file2.set_index("R.FileName").loc[run]))
            if metadata["R.Replicate"][fileindex]=="":
                RReplicatelist.append([i+1]*len(file2.set_index("R.FileName").loc[run]))
            else:
                RReplicatelist.append([metadata["R.Replicate"][fileindex]]*len(file2.set_index("R.FileName").loc[run]))
        file2["R.Condition"]=list(itertools.chain(*RConditionlist))
        file2["R.Replicate"]=list(itertools.chain(*RReplicatelist))
        file2["R.Replicate"]=file2["R.Replicate"].astype(int)

        #filter protein names with cRAP database
        if cRAP_filter==True:
            cRAP=['ADH1_YEAST','ALBU_BOVIN','ALBU_HUMAN','ALDOA_RABIT','AMYS_HUMAN','ANT3_HUMAN','ANXA5_HUMAN','B2MG_HUMAN','BGAL_ECOLI','BID_HUMAN','CAH1_HUMAN','CAH2_BOVIN','CAH2_HUMAN',
                  'CAS1_BOVIN','CAS2_BOVIN','CASB_BOVIN','CASK_BOVIN','CATA_HUMAN','CATD_HUMAN','CATG_HUMAN','CO5_HUMAN','CRP_HUMAN','CTRA_BOVIN','CTRB_BOVIN','CYB5_HUMAN','CYC_HORSE',
                  'CYC_HUMAN','DHE3_BOVIN','EGF_HUMAN','FABPH_HUMAN','GAG_SCVLA','GELS_HUMAN','GFP_AEQVI','GSTA1_HUMAN','GSTP1_HUMAN','HBA_HUMAN','HBB_HUMAN','IGF2_HUMAN','IL8_HUMAN',
                  'K1C10_HUMAN','K1C15_SHEEP','K1C9_HUMAN','K1H1_HUMAN','K1H2_HUMAN','K1H4_HUMAN','K1H5_HUMAN','K1H6_HUMAN','K1H7_HUMAN','K1H8_HUMAN','K1HA_HUMAN','K1HB_HUMAN','K1M1_SHEEP',
                  'K1M2_SHEEP','K22E_HUMAN','K2C1_HUMAN','K2M1_SHEEP','K2M2_SHEEP','K2M3_SHEEP','KCRM_HUMAN','KRA3_SHEEP','KRA33_SHEEP','KRA34_SHEEP','KRA3A_SHEEP','KRA61_SHEEP','KRB2A_SHEEP',
                  'KRB2B_SHEEP','KRB2C_SHEEP','KRB2D_SHEEP','KRHB1_HUMAN','KRHB2_HUMAN','KRHB3_HUMAN','KRHB4_HUMAN','KRHB5_HUMAN','KRHB6_HUMAN','KRUC_SHEEP','LALBA_BOVIN','LALBA_HUMAN',
                  'LEP_HUMAN','LYSC_CHICK','LYSC_HUMAN','LYSC_LYSEN','MYG_HORSE','MYG_HUMAN','NEDD8_HUMAN','NQO1_HUMAN','NQO2_HUMAN','OVAL_CHICK','PDGFB_HUMAN','PEPA_BOVIN','PEPA_PIG',
                  'PEPB_PIG','PEPC_PIG','PLMP_GRIFR','PPIA_HUMAN','PRDX1_HUMAN','RASH_HUMAN','REF_HEVBR','RETBP_HUMAN','RS27A_HUMAN','SODC_HUMAN','SRPP_HEVBR','SSPA_STAAU','SUMO1_HUMAN',
                  'SYH_HUMAN','TAU_HUMAN','THIO_HUMAN','TNFA_HUMAN','TRFE_HUMAN','TRFL_HUMAN','TRY1_BOVIN','TRY2_BOVIN','TRYP_PIG','UB2E1_HUMAN','UBE2C_HUMAN','UBE2I_HUMAN']
            decoylist=[]
            for protein in file1["PG.ProteinNames"].drop_duplicates().tolist():
                if protein in cRAP:
                    decoylist.append(protein)
            for protein in file2["PG.ProteinNames"].drop_duplicates().tolist():
                if protein in cRAP:
                    decoylist.append(protein)
            file1=file1.set_index("PG.ProteinNames").drop(decoylist).reset_index()
            file2=file2.set_index("PG.ProteinNames").drop(decoylist).reset_index()

        #handling for if there are still blanks in metadata_condition R.Condition column
        if "" in metadata_condition["R.Condition"].tolist():
            metadata_condition.loc[metadata_condition.index[metadata_condition["R.Condition"]==""].values[0],"R.Condition"]="Not Defined"

        #add concentration values
        if metadata_condition["Concentration"].drop_duplicates().tolist()!=[""]:
            concentrationlist=[]
            for run in file1["R.Condition"].drop_duplicates().tolist():
                fileindex=metadata_condition[metadata_condition["R.Condition"]==run].index.values[0]
                concentrationlist.append([metadata_condition["Concentration"][fileindex]]*len(file1.set_index("R.Condition").loc[run]))
            if "Concentration" in file1.columns:
                file1["Concentration"]=list(itertools.chain(*concentrationlist))
                file1["Concentration"]=file1["Concentration"].astype(float)
            else:
                file1.insert(3,"Concentration",list(itertools.chain(*concentrationlist)))
                file1["Concentration"]=file1["Concentration"].astype(float)
            concentrationlist=[]
            for run in file2["R.Condition"].drop_duplicates().tolist():
                fileindex=metadata_condition[metadata_condition["R.Condition"]==run].index.values[0]
                concentrationlist.append([metadata_condition["Concentration"][fileindex]]*len(file2.set_index("R.Condition").loc[run]))
            if "Concentration" in file2.columns:
                file2["Concentration"]=list(itertools.chain(*concentrationlist))
                file2["Concentration"]=file2["Concentration"].astype(float)
            else:
                file2.insert(3,"Concentration",list(itertools.chain(*concentrationlist)))
                file2["Concentration"]=file2["Concentration"].astype(float)

        #generate Cond_Rep column
        file1["R.Condition"]=file1["R.Condition"].apply(str)
        if "Cond_Rep" not in file1.columns:
            file1.insert(0,"Cond_Rep",file1["R.Condition"]+"_"+file1["R.Replicate"].apply(str))
        elif "Cond_Rep" in file1.columns:
            file1["Cond_Rep"]=file1["R.Condition"]+"_"+file1["R.Replicate"].apply(str)
        file2["R.Condition"]=file2["R.Condition"].apply(str)
        if "Cond_Rep" not in file2.columns:
            file2.insert(0,"Cond_Rep",file2["R.Condition"]+"_"+file2["R.Replicate"].apply(str))
        elif "Cond_Rep" in file2.columns:
            file2["Cond_Rep"]=file2["R.Condition"]+"_"+file2["R.Replicate"].apply(str)
        
        #reorder
        metadata_condition["order"]=metadata_condition["order"].astype(int)
        if len(metadata_condition["order"])==1:
            pass
        else:
            sortedmetadata_bycondition=metadata_condition.sort_values(by="order").reset_index(drop=True)
            file1=file1.set_index("R.Condition").loc[sortedmetadata_bycondition["R.Condition"].tolist()].reset_index()
            file2=file2.set_index("R.Condition").loc[sortedmetadata_bycondition["R.Condition"].tolist()].reset_index()
        return file1,file2

    # ====================================== UI and calc functions
    #file1 import
    @reactive.calc
    def compare_inputfile1():
        searchoutput=inputfile(input.compare_searchreport1(),input.software1(),input.software1_bps_report_type(),"off",input.searchreport1_reupload())
        return searchoutput
    #file2 import
    @reactive.calc
    def compare_inputfile2():
        searchoutput=inputfile(input.compare_searchreport2(),input.software2(),input.software2_bps_report_type(),"off",input.searchreport2_reupload())
        return searchoutput
    #store generated metadata df as a reactive value
    @reactive.calc
    def compare_metadata_calc():
        metadata=compare_metadata_gen(input.compare_use_uploaded_metadata(),input.compare_metadata_upload(),compare_inputfile1(),compare_inputfile2())
        return metadata
    #render metadata table
    @render.data_frame
    def compare_metadata_table():
        metadata=compare_metadata_calc()
        metadata=metadata[["R.FileName1","R.FileName2","R.Condition","R.Replicate","remove"]]
        if len(metadata)==0:
            #return non-editable metadata if there's nothing there
            return render.DataGrid(metadata,width="100%")
        else:
            return render.DataGrid(metadata,editable=True,width="100%")
    #store metadata_condition df as a reactive value (don't need to make new generative function)
    @reactive.calc
    def compare_metadata_condition_calc():
        metadata_condition=metadata_condition_gen(compare_metadata_table.data_view(),compare_metadata_calc(),input.compare_use_uploaded_metadata(),input.compare_remove())
        return metadata_condition
    #render metadata_condition table
    @render.data_frame
    def compare_metadata_condition_table():
        metadata_condition=compare_metadata_condition_calc()
        if len(metadata_condition)==0:
            return render.DataGrid(metadata_condition,width="100%")
        else:
            return render.DataGrid(metadata_condition,editable=True,width="100%")
    #download metadata table as shown
    @render.download(filename="metadata_"+str(date.today())+".csv")
    def compare_metadata_download():
        metadata_download=compare_metadata_download_gen(compare_metadata_table.data_view(),compare_metadata_condition_table.data_view())
        with io.BytesIO() as buf:
            metadata_download.to_csv(buf,index=False)
            yield buf.getvalue()
    #update search files based on metadata
    @reactive.calc
    @reactive.event(input.compare_rerun_metadata,ignore_none=False)
    def compare_metadata_update():
        file1,file2=compare_metadata_update_gen(compare_inputfile1(),compare_inputfile2(),compare_metadata_table.data_view(),compare_metadata_condition_table.data_view(),input.compare_remove(),input.compare_cRAP_filter())
        return file1,file2

    #checkmark popup to show that the metadata has been applied to the search report
    @render.text
    @reactive.event(input.compare_rerun_metadata)
    def twosoftware_metadata_applied_ui():
        return '\u2705'
    #function for finding the PTMs in the data
    @reactive.calc
    def compare_find_ptms():
        file1,file2=compare_metadata_update()
        peplist=file1["EG.ModifiedPeptide"]
        ptmlist=[]
        for i in peplist:
            ptmlist.append(re.findall(r"[^[]*\[([^]]*)\]",i))
        return(list(OrderedDict.fromkeys(itertools.chain(*ptmlist))))

    #Calc functions for ID Counts and PTM ID Counts
    @reactive.calc
    def idmetrics_calc_1():
        file1,file2=compare_metadata_update()
        resultdf,averagedf,uniquecountsdf=idmetrics(file1)
        return resultdf,averagedf,uniquecountsdf
    @reactive.calc
    def idmetrics_calc_2():
        file1,file2=compare_metadata_update()
        resultdf,averagedf,uniquecountsdf=idmetrics(file2)
        return resultdf,averagedf,uniquecountsdf
    @reactive.calc
    def ptmidmetrics_calc_1():
        file1,file2=compare_metadata_update()
        ptmresultdf,ptmaveragedf=ptmcounts(file1,input.compare_idmetrics_ptm())
        return ptmresultdf,ptmaveragedf
    @reactive.calc
    def ptmidmetrics_calc_2():
        file1,file2=compare_metadata_update()
        ptmresultdf,ptmaveragedf=ptmcounts(file2,input.compare_idmetrics_ptm())
        return ptmresultdf,ptmaveragedf

    # ====================================== ID Counts
    @render.ui
    def compare_idmetrics_ptm_ui():
        listofptms=compare_find_ptms()
        listofptms.insert(0,"None")
        return ui.input_selectize("compare_idmetrics_ptm","Pick PTM to plot data for:",choices=listofptms,selected=listofptms[0])
    #ID Counts plot
    @reactive.effect
    def _():
        @render.plot(width=input.compare_id_counts_width(),height=input.compare_id_counts_height())
        def compare_id_counts():
            file1,file2=compare_metadata_update()
            resultdf1,averagedf1,uniquecountsdf1=idmetrics_calc_1()
            resultdf2,averagedf2,uniquecountsdf2=idmetrics_calc_2()
            ptmresultdf1,ptmaveragedf1=ptmidmetrics_calc_1()
            ptmresultdf2,ptmaveragedf2=ptmidmetrics_calc_2()

            width=0.4
            y_padding=input.ypadding()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            x_label_rotation=input.xaxis_label_rotation()

            colortable=compare_colortable.data_view()
            color1=colortable["Color"].tolist()[0]
            color2=colortable["Color"].tolist()[1]
            label1=colortable["Software"].tolist()[0]
            label2=colortable["Software"].tolist()[1]

            fig,ax=plt.subplots(nrows=2,ncols=2,sharex=True,figsize=(input.compare_id_counts_width()*(1/plt.rcParams['figure.dpi']),input.compare_id_counts_height()*(1/plt.rcParams['figure.dpi'])))
            ax1=ax[0,0]
            ax2=ax[0,1]
            ax3=ax[1,0]
            ax4=ax[1,1]

            if input.compare_idmetrics_peptides()=="modified":
                pepcolumn="peptides"
                peptitle="Peptides"
            elif input.compare_idmetrics_peptides()=="stripped":
                pepcolumn="strippedpeptides"
                peptitle="Peptide Sequences"
            
            if input.compare_idmetrics_ptm()=="None":
                plotresultdf1=resultdf1
                plotresultdf2=resultdf2
                plotaveragedf1=averagedf1
                plotaveragedf2=averagedf2
                ptm=""
            else:
                ptm=input.compare_idmetrics_ptm()
                plotresultdf1=ptmresultdf1
                plotresultdf2=ptmresultdf2
                plotaveragedf1=ptmaveragedf1
                plotaveragedf2=ptmaveragedf2

            if input.compare_idmetrics_individual_average()=="individual":
                runlist=plotresultdf1["Cond_Rep"]
                x=np.arange(len(runlist))

                ax1.bar(x,plotresultdf1["proteins"],width=width,label=label1,color=color1,edgecolor="k")
                ax1.bar(x+width,plotresultdf2["proteins"],width=width,label=label2,color=color2,edgecolor="k")
                ax1.bar_label(ax1.containers[0],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                ax1.bar_label(ax1.containers[1],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                maxvalue1=max([max(plotresultdf1["proteins"]),max(plotresultdf2["proteins"])])
                ax1.set_ylim(top=maxvalue1+(y_padding*maxvalue1))
                ax1.set_title("Protein Groups",fontsize=titlefont)
                ax1.tick_params(axis="both",labelsize=axisfont_labels)

                ax2.bar(x,plotresultdf1["proteins2pepts"],width=width,color=color1,edgecolor="k")
                ax2.bar(x+width,plotresultdf2["proteins2pepts"],width=width,color=color2,edgecolor="k")
                ax2.bar_label(ax2.containers[0],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                ax2.bar_label(ax2.containers[1],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                maxvalue2=max([max(plotresultdf1["proteins2pepts"]),max(plotresultdf2["proteins2pepts"])])
                ax2.set_ylim(top=maxvalue2+(y_padding*maxvalue2))
                ax2.set_title("Protein Groups with >2 Peptides",fontsize=titlefont)
                ax2.tick_params(axis="both",labelsize=axisfont_labels)

                ax3.bar(x,plotresultdf1[pepcolumn],width=width,color=color1,edgecolor="k")
                ax3.bar(x+width,plotresultdf2[pepcolumn],width=width,color=color2,edgecolor="k")
                ax3.bar_label(ax3.containers[0],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                ax3.bar_label(ax3.containers[1],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                maxvalue3=max([max(plotresultdf1[pepcolumn]),max(plotresultdf2[pepcolumn])])
                ax3.set_ylim(top=maxvalue3+(y_padding*maxvalue3))
                ax3.set_xlabel("Condition",fontsize=axisfont)
                ax3.set_xticks(x+width/2,runlist)
                ax3.tick_params(axis="x",rotation=x_label_rotation)
                ax3.set_title(peptitle,fontsize=titlefont)
                ax3.tick_params(axis="both",labelsize=axisfont_labels)

                ax4.bar(x,plotresultdf1["precursors"],width=width,color=color1,edgecolor="k")
                ax4.bar(x+width,plotresultdf2["precursors"],width=width,color=color2,edgecolor="k")
                ax4.bar_label(ax4.containers[0],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                ax4.bar_label(ax4.containers[1],label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                maxvalue4=max([max(plotresultdf1["precursors"]),max(plotresultdf2["precursors"])])
                ax4.set_ylim(top=maxvalue4+(y_padding*maxvalue4))
                ax4.set_xlabel("Condition",fontsize=axisfont)
                ax4.set_xticks(x+width/2,runlist)
                ax4.tick_params(axis="x",rotation=x_label_rotation)
                ax4.set_title("Precursors",fontsize=titlefont)
                ax4.tick_params(axis="both",labelsize=axisfont_labels)

            elif input.compare_idmetrics_individual_average()=="average":
                runlist=plotaveragedf1["R.Condition"]
                x=np.arange(len(runlist))

                bars1_1=ax1.bar(x,plotaveragedf1["proteins_avg"],yerr=plotaveragedf1["proteins_stdev"],width=width,label=label1,capsize=10,color=color1,edgecolor="k")
                bars1_2=ax1.bar(x+width,plotaveragedf2["proteins_avg"],yerr=plotaveragedf2["proteins_stdev"],width=width,label=label2,capsize=10,color=color2,edgecolor="k")
                ax1.bar_label(bars1_1,label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                ax1.bar_label(bars1_2,label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                maxvalue1=max([max(plotaveragedf1["proteins_avg"]),max(plotaveragedf2["proteins_avg"])])
                ax1.set_ylim(top=maxvalue1+(y_padding*maxvalue1))
                ax1.set_title("Protein Groups",fontsize=titlefont)
                ax1.tick_params(axis="both",labelsize=axisfont_labels)

                bars2_1=ax2.bar(x,plotaveragedf1["proteins2pepts_avg"],yerr=plotaveragedf1["proteins2pepts_stdev"],width=width,capsize=10,color=color1,edgecolor="k")
                bars2_2=ax2.bar(x+width,plotaveragedf2["proteins2pepts_avg"],yerr=plotaveragedf2["proteins2pepts_stdev"],width=width,capsize=10,color=color2,edgecolor="k")
                ax2.bar_label(bars2_1,label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                ax2.bar_label(bars2_2,label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                maxvalue2=max([max(plotaveragedf1["proteins2pepts_avg"]),max(plotaveragedf2["proteins2pepts_avg"])])
                ax2.set_ylim(top=maxvalue2+(y_padding*maxvalue2))
                ax2.set_title("Protein Groups with >2 Peptides",fontsize=titlefont)
                ax2.tick_params(axis="both",labelsize=axisfont_labels)

                bars3_1=ax3.bar(x,plotaveragedf1[pepcolumn+"_avg"],yerr=plotaveragedf1[pepcolumn+"_stdev"],width=width,capsize=10,color=color1,edgecolor="k")
                bars3_2=ax3.bar(x+width,plotaveragedf2[pepcolumn+"_avg"],yerr=plotaveragedf2[pepcolumn+"_stdev"],width=width,capsize=10,color=color2,edgecolor="k")
                ax3.bar_label(bars3_1,label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                ax3.bar_label(bars3_2,label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                maxvalue3=max([max(plotaveragedf1[pepcolumn+"_avg"]),max(plotaveragedf2[pepcolumn+"_avg"])])
                ax3.set_ylim(top=maxvalue3+(y_padding*maxvalue3))
                ax3.set_xlabel("Condition",fontsize=axisfont)
                ax3.set_xticks(x+width/2,runlist)
                ax3.tick_params(axis="x",rotation=x_label_rotation)
                ax3.set_title(peptitle,fontsize=titlefont)
                ax3.tick_params(axis="both",labelsize=axisfont_labels)

                bars4_1=ax4.bar(x,plotaveragedf1["precursors_avg"],yerr=plotaveragedf1["precursors_stdev"],width=width,capsize=10,color=color1,edgecolor="k")
                bars4_2=ax4.bar(x+width,plotaveragedf2["precursors_avg"],yerr=plotaveragedf2["precursors_stdev"],width=width,capsize=10,color=color2,edgecolor="k")
                ax4.bar_label(bars4_1,label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                ax4.bar_label(bars4_2,label_type="edge",padding=5,rotation=90,fontsize=labelfont)
                maxvalue4=max([max(plotaveragedf1["precursors_avg"]),max(plotaveragedf2["precursors_avg"])])
                ax4.set_ylim(top=maxvalue4+(y_padding*maxvalue4))
                ax4.set_xlabel("Condition",fontsize=axisfont)
                ax4.set_xticks(x+width/2,runlist)
                ax4.tick_params(axis="x",rotation=x_label_rotation)
                ax4.set_title("Precursors",fontsize=titlefont)
                ax4.tick_params(axis="both",labelsize=axisfont_labels)

            fig.legend(loc="upper left",bbox_to_anchor=(0,1),prop={'size':legendfont})
            fig.text(0, 0.6,"Counts",ha="left",va="center",rotation="vertical",fontsize=axisfont)
            
            if input.compare_idmetrics_ptm()!="None":
                fig.suptitle("ID Counts for PTM: "+ptm,y=1,fontsize=titlefont)

            ax1.set_axisbelow(True)
            ax1.grid(linestyle="--")
            ax2.set_axisbelow(True)
            ax2.grid(linestyle="--")
            ax3.set_axisbelow(True)
            ax3.grid(linestyle="--")
            ax4.set_axisbelow(True)
            ax4.grid(linestyle="--")

            fig.set_tight_layout(True)
            imagedownload_twosoftware("idcounts")

    # ====================================== Venn Diagram
    @render.ui
    def compare_venn_run_ui():
        file1,file2=compare_metadata_update()
        if input.compare_venn_conditionorrun()=="condition":
            opts=file1["R.Condition"].drop_duplicates().tolist()
            return ui.input_selectize("compare_venn_run_list","Pick condition to compare",choices=opts)
        if input.compare_venn_conditionorrun()=="individual":
            opts=file1["Cond_Rep"].tolist()
            return ui.input_selectize("compare_venn_run_list","Pick run to compare",choices=opts)   
    @render.ui
    def compare_venn_ptm_ui():
        if input.compare_venn_plotproperty()=="peptides" or input.compare_venn_plotproperty()=="precursors" or input.compare_venn_plotproperty()=="peptides_stripped":
            return ui.input_switch("compare_venn_ptm","Compare only for specific PTM?",value=False)
    @render.ui
    def compare_venn_ptmlist_ui():
        if input.compare_venn_plotproperty()=="peptides" or input.compare_venn_plotproperty()=="precursors" or input.compare_venn_plotproperty()=="peptides_stripped":
            if input.compare_venn_ptm()==True:
                listofptms=compare_find_ptms()
                listofptms.insert(0,"None")
                return ui.input_selectize("compare_venn_foundptms","Pick PTM to plot data for:",choices=listofptms,selected=listofptms[0])
    @render.ui
    def compare_venn_specific_length_ui():
        if input.compare_venn_plotproperty()=="peptides" or input.compare_venn_plotproperty()=="precursors" or input.compare_venn_plotproperty()=="peptides_stripped":
            return ui.input_switch("compare_venn_specific_length","Compare specific peptide length?",value=False,width="300px")
    @render.ui
    def compare_venn_peplength_ui():
        if input.compare_venn_specific_length()==True:
            return ui.input_slider("compare_venn_peplength_pick","Pick specific peptide length to compare:",min=3,max=30,value=9,step=1,ticks=True)
    #plot Venn Diagram
    @reactive.effect
    def _():
        @render.plot(width=input.compare_venn_width(),height=input.compare_venn_height())
        def compare_venn_plot():
            from matplotlib_venn import venn2,venn2_circles,venn3,venn3_circles
            file1,file2=compare_metadata_update()

            colortable=compare_colortable.data_view()
            color1=colortable["Color"].tolist()[0]
            color2=colortable["Color"].tolist()[1]
            label1=colortable["Software"].tolist()[0]
            label2=colortable["Software"].tolist()[1]

            if input.compare_venn_conditionorrun()=="condition":
                A=file1[file1["R.Condition"]==input.compare_venn_run_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                B=file2[file2["R.Condition"]==input.compare_venn_run_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            if input.compare_venn_conditionorrun()=="individual":
                A=file1[file1["Cond_Rep"]==input.compare_venn_run_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
                B=file2[file2["Cond_Rep"]==input.compare_venn_run_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)

            #some software return EG.ModifiedPeptide with _seq_, need to strip the underscores in order for them to be able to be compared
            A_strippedmodpeplist=[]
            B_strippedmodpeplist=[]
            for ele1 in A["EG.ModifiedPeptide"]:
                if ele1.find("_")==-1:
                    A_strippedmodpeplist.append(ele1)
                else:
                    A_strippedmodpeplist.append(ele1.split("_")[1]) 
            for ele2 in B["EG.ModifiedPeptide"]:
                if ele2.find("_")==-1:
                    B_strippedmodpeplist.append(ele2)
                else:
                    B_strippedmodpeplist.append(ele2.split("_")[1])
            A["EG.ModifiedPeptide"]=A_strippedmodpeplist
            B["EG.ModifiedPeptide"]=B_strippedmodpeplist

            A["pep_charge"]=A["EG.ModifiedPeptide"]+A["FG.Charge"].astype(str)
            B["pep_charge"]=B["EG.ModifiedPeptide"]+B["FG.Charge"].astype(str)

            titlemodlist=[]
            if input.compare_venn_plotproperty()=="proteingroups":
                a=set(A["PG.ProteinGroups"])
                b=set(B["PG.ProteinGroups"])
                titlemod="Protein Groups"
            if input.compare_venn_plotproperty()=="peptides":
                if input.compare_venn_ptm()==True:
                    ptm=input.compare_venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    titlemodlist.append(ptm)
                if input.compare_venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                    titlemodlist.append(str(input.compare_venn_peplength_pick())+"mers")
                else:
                    A=A
                    B=B
                a=set(A["EG.ModifiedPeptide"])
                b=set(B["EG.ModifiedPeptide"])
                titlemod="Peptides"
            if input.compare_venn_plotproperty()=="precursors":
                if input.compare_venn_ptm()==True:
                    ptm=input.compare_venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    titlemodlist.append(ptm)
                if input.compare_venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                    titlemodlist.append(str(input.compare_venn_peplength_pick())+"mers")
                else:
                    A=A
                    B=B
                a=set(A["pep_charge"])
                b=set(B["pep_charge"])
                titlemod="Precursors"
            if input.compare_venn_plotproperty()=="peptides_stripped":
                if input.compare_venn_ptm()==True:
                    ptm=input.compare_venn_foundptms()
                    A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                    titlemodlist.append(ptm)
                if input.compare_venn_specific_length()==True:
                    A=A[A["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                    B=B[B["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                    titlemodlist.append(str(input.compare_venn_peplength_pick())+"mers")
                else:
                    A=A
                    B=B
                a=set(A["PEP.StrippedSequence"])
                b=set(B["PEP.StrippedSequence"])
                titlemod="Stripped Peptides"
            if titlemodlist==[]:
                titlemodlist=""
            else:
                titlemodlist=" ("+", ".join(titlemodlist)+")"
            fig,ax=plt.subplots(figsize=(input.compare_venn_width()*(1/plt.rcParams['figure.dpi']),input.compare_venn_height()*(1/plt.rcParams['figure.dpi'])))
            Ab=len(a-b)
            aB=len(b-a)
            AB=len(a&b)
            venn2(subsets=(Ab,aB,AB),set_labels=(label1,label2),set_colors=(color1,color2),ax=ax)
            venn2_circles(subsets=(Ab,aB,AB),linestyle="dashed",linewidth=0.5)
            plt.title("Venn Diagram for "+input.compare_venn_run_list()+" "+titlemod+titlemodlist)
            fig.set_tight_layout(True)
            imagedownload_twosoftware("venndiagram_"+input.compare_venn_plotproperty())
    @render.download(filename=lambda: f"VennList_A-{input.software1()}_vs_B-{input.software2()}_{input.compare_venn_plotproperty()}.csv")
    def compare_venn_download():
        file1,file2=compare_metadata_update()
        if input.compare_venn_conditionorrun()=="condition":
            A=file1[file1["R.Condition"]==input.compare_venn_run_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            B=file2[file2["R.Condition"]==input.compare_venn_run_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
        if input.compare_venn_conditionorrun()=="individual":
            A=file1[file1["Cond_Rep"]==input.compare_venn_run_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)
            B=file2[file2["Cond_Rep"]==input.compare_venn_run_list()][["PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge","PEP.StrippedSequence","Peptide Length"]].drop_duplicates().reset_index(drop=True)

        #some software return EG.ModifiedPeptide with _seq_, need to strip the underscores in order for them to be able to be compared
        A_strippedmodpeplist=[]
        B_strippedmodpeplist=[]
        for ele1 in A["EG.ModifiedPeptide"]:
            if ele1.find("_")==-1:
                A_strippedmodpeplist.append(ele1)
            else:
                A_strippedmodpeplist.append(ele1.split("_")[1]) 
        for ele2 in B["EG.ModifiedPeptide"]:
            if ele2.find("_")==-1:
                B_strippedmodpeplist.append(ele2)
            else:
                B_strippedmodpeplist.append(ele2.split("_")[1])
        A["EG.ModifiedPeptide"]=A_strippedmodpeplist
        B["EG.ModifiedPeptide"]=B_strippedmodpeplist

        A["pep_charge"]=A["EG.ModifiedPeptide"]+A["FG.Charge"].astype(str)
        B["pep_charge"]=B["EG.ModifiedPeptide"]+B["FG.Charge"].astype(str)

        if input.compare_venn_plotproperty()=="proteingroups":
            a=set(A["PG.ProteinGroups"])
            b=set(B["PG.ProteinGroups"])
        if input.compare_venn_plotproperty()=="peptides":
            if input.compare_venn_ptm()==True:
                ptm=input.compare_venn_foundptms()
                A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
            if input.compare_venn_specific_length()==True:
                A=A[A["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                B=B[B["Peptide Length"]==int(input.compare_venn_peplength_pick())]
            else:
                A=A
                B=B
            a=set(A["EG.ModifiedPeptide"])
            b=set(B["EG.ModifiedPeptide"])
            titlemod="Peptides"
        if input.compare_venn_plotproperty()=="precursors":
            if input.compare_venn_ptm()==True:
                ptm=input.compare_venn_foundptms()
                A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
            if input.compare_venn_specific_length()==True:
                A=A[A["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                B=B[B["Peptide Length"]==int(input.compare_venn_peplength_pick())]
            else:
                A=A
                B=B
            a=set(A["pep_charge"])
            b=set(B["pep_charge"])
        if input.compare_venn_plotproperty()=="peptides_stripped":
            if input.compare_venn_ptm()==True:
                ptm=input.compare_venn_foundptms()
                A=A[A["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
                B=B[B["EG.ModifiedPeptide"].str.contains(ptm,regex=False)]
            if input.compare_venn_specific_length()==True:
                A=A[A["Peptide Length"]==int(input.compare_venn_peplength_pick())]
                B=B[B["Peptide Length"]==int(input.compare_venn_peplength_pick())]
            else:
                A=A
                B=B
            a=set(A["PEP.StrippedSequence"])
            b=set(B["PEP.StrippedSequence"])

        df=pd.DataFrame()
        Ab=list(a-b)
        aB=list(b-a)
        AB=list(a&b)
        df=pd.concat([df,pd.Series(Ab,name=input.software1())],axis=1)
        df=pd.concat([df,pd.Series(aB,name=input.software2())],axis=1)
        df=pd.concat([df,pd.Series(AB,name="Both")],axis=1)
        with io.BytesIO() as buf:
            df.to_csv(buf,index=False)
            yield buf.getvalue() 

#endregion

# ============================================================================= Export Tables 
#region 
    # ====================================== Export Tables
    #download table of peptide IDs
    @render.download(filename=lambda: f"Peptide List_{input.searchreport()[0]['name']}.csv")
    def peptidelist():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

        columnlist=["R.Condition","R.Replicate","PG.Genes","PG.ProteinAccessions","PG.ProteinNames","PG.ProteinGroups","EG.ModifiedPeptide","FG.Charge"]
        for column in columnlist:
            if column not in searchoutput.columns:
                columnlist.remove(column)
            else:
                pass
        peptidetable=searchoutput[columnlist].drop_duplicates().reset_index(drop=True)
        with io.BytesIO() as buf:
            peptidetable.to_csv(buf,index=False)
            yield buf.getvalue()

    #condition or run dropdown for stripped peptide list export
    @render.ui
    def peptidelist_dropdown():
        if input.peptidelist_condition_or_run()=="condition":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=sampleconditions
            return ui.input_selectize("peptidelist_dropdown_pick","Pick run to export peptide list from",choices=opts)
        if input.peptidelist_condition_or_run()=="individual":
            searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
            opts=runnames
            return ui.input_selectize("peptidelist_dropdown_pick","Pick run to export peptide list from",choices=opts)

    #download list of stripped peptide IDs of specific length
    @render.download(filename=lambda: f"Stripped Peptide List_{input.peptidelist_dropdown_pick()}_{input.strippedpeptidelength()}-mers.csv")
    def strippedpeptidelist():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        if input.peptidelist_condition_or_run()=="condition":
            placeholder=searchoutput[searchoutput["R.Condition"]==input.peptidelist_dropdown_pick()]["PEP.StrippedSequence"].drop_duplicates().reset_index(drop=True)

            df=pd.DataFrame()
            df["PEP.StrippedSequence"]=placeholder
            lengths=[]
            for pep in placeholder:
                lengths.append(len(pep))
            df["Peptide Length"]=lengths
            outputdf=df[df["Peptide Length"]==input.strippedpeptidelength()].drop(columns=["Peptide Length"])
            with io.BytesIO() as buf:
                outputdf.to_csv(buf,header=False,index=False)
                yield buf.getvalue()
        if input.peptidelist_condition_or_run()=="individual":
            placeholder=searchoutput[searchoutput["Cond_Rep"]==input.peptidelist_dropdown_pick()]["PEP.StrippedSequence"].drop_duplicates().reset_index(drop=True)

            df=pd.DataFrame()
            df["PEP.StrippedSequence"]=placeholder
            lengths=[]
            for pep in placeholder:
                lengths.append(len(pep))
            df["Peptide Length"]=lengths
            outputdf=df[df["Peptide Length"]==input.strippedpeptidelength()].drop(columns=["Peptide Length"])
            with io.BytesIO() as buf:
                outputdf.to_csv(buf,header=False,index=False)
                yield buf.getvalue()

    #download table of protein ID metrics/CVs
    @render.download(filename=lambda: f"Protein CV Table_{input.searchreport()[0]['name']}.csv")
    def proteinidmetrics_download():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        columnlist=["R.Condition","R.Replicate","PG.ProteinGroups","PG.MS2Quantity"]
        try:
            cvproteingroup=searchoutput[columnlist].drop_duplicates().reset_index(drop=True)
            cvproteinmean=cvproteingroup.drop(columns="R.Replicate").groupby(["R.Condition","PG.ProteinGroups"]).mean().rename(columns={"PG.MS2Quantity":"Mean"})
            cvproteinstdev=cvproteingroup.drop(columns="R.Replicate").groupby(["R.Condition","PG.ProteinGroups"]).std().rename(columns={"PG.MS2Quantity":"Stdev"})
            cvproteincount=cvproteingroup.drop(columns="R.Replicate").groupby(["R.Condition","PG.ProteinGroups"]).size().reset_index(drop=True)
            cvproteintable=pd.concat([cvproteinmean,cvproteinstdev],axis=1).reindex(cvproteinmean.index)
            cvproteintable["CV"]=cvproteintable["Stdev"]/cvproteintable["Mean"]*100
            cvproteintable["# replicates observed"]=cvproteincount.tolist()
            cvproteintable=cvproteintable.reset_index()
            with io.BytesIO() as buf:
                cvproteintable.to_csv(buf,index=False)
                yield buf.getvalue()
        except KeyError:
            raise KeyError("PG.MS2Quantity not present in search results") from None

    @render.download(filename=lambda: f"Protein ID Matrix_{input.searchreport()[0]['name']}.csv")
    def proteinidmatrix_download():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        try:
            pg_matrix=searchoutput[["R.Condition","R.Replicate","PG.ProteinGroups","PG.Genes","PG.MS2Quantity"]].drop_duplicates().pivot(columns=["R.Condition","R.Replicate"],index=["PG.ProteinGroups","PG.Genes"]).reset_index()
            with io.BytesIO() as buf:
                pg_matrix.to_csv(buf,header=True,index=False)
                yield buf.getvalue()
        except KeyError:
            raise KeyError("PG.MS2Quantity not present in search results") from None

    #download table of precursor ID metrics/CVs
    @render.download(filename=lambda: f"Precursor CV Table_{input.searchreport()[0]['name']}.csv")
    def precursoridmetrics_download():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        columnlist=["R.Condition","R.Replicate","EG.ModifiedPeptide","FG.Charge","FG.MS2Quantity"]
        try:
            cvprecursorgroup=searchoutput[columnlist].drop_duplicates().reset_index(drop=True)

            cvprecursormean=cvprecursorgroup.drop(columns="R.Replicate").groupby(["R.Condition","EG.ModifiedPeptide","FG.Charge"]).mean().rename(columns={"FG.MS2Quantity":"Mean"})
            cvprecursorstdev=cvprecursorgroup.drop(columns="R.Replicate").groupby(["R.Condition","EG.ModifiedPeptide","FG.Charge"]).std().rename(columns={"FG.MS2Quantity":"Stdev"})
            cvprecursorcount=cvprecursorgroup.drop(columns="R.Replicate").groupby(["R.Condition","EG.ModifiedPeptide","FG.Charge"]).size().reset_index(drop=True)
            cvprecursortable=pd.concat([cvprecursormean,cvprecursorstdev],axis=1).reindex(cvprecursormean.index)
            cvprecursortable["CV"]=cvprecursortable["Stdev"]/cvprecursortable["Mean"]*100
            cvprecursortable["# replicates observed"]=cvprecursorcount.tolist()
            cvprecursortable=cvprecursortable.reset_index()
            with io.BytesIO() as buf:
                cvprecursortable.to_csv(buf,index=False)
                yield buf.getvalue()
        except KeyError:
            raise KeyError("FG.MS2Quantity not present in search results") from None
    
    #download table of MOMA precursors for a specified run
    @render.download(filename=lambda: f"MOMA Table_{input.searchreport()[0]['name']}.csv")
    def moma_download():
        #RT tolerance in s
        rttolerance=input.rttolerance()
        #MZ tolerance in m/z
        mztolerance=input.mztolerance()
        #IM tolerance in 1/K0
        imtolerance=input.imtolerance()

        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()

        sample=input.cond_rep()

        columns=["EG.ModifiedPeptide","FG.Charge","EG.IonMobility","EG.ApexRT","FG.PrecMz","FG.MS2Quantity","PG.ProteinGroups","PG.ProteinNames"]
        df=searchoutput[searchoutput["Cond_Rep"]==sample][["EG.ModifiedPeptide","FG.Charge","EG.IonMobility","EG.ApexRT","FG.PrecMz","FG.MS2Quantity","PG.ProteinGroups","PG.ProteinNames"]].sort_values(["EG.ApexRT"]).reset_index(drop=True)
        coelutingpeptides=pd.DataFrame(columns=columns)
        for i in range(len(df)):
            if i+1 not in range(len(df)):
                break
            #rtpercentdiff=(abs(df["EG.ApexRT"][i]-df["EG.ApexRT"][i+1])/df["EG.ApexRT"][i])*100
            rtdiff=abs(df["EG.ApexRT"][i]-df["EG.ApexRT"][i+1])
            mzdiff=abs(df["FG.PrecMz"][i]-df["FG.PrecMz"][i+1])
            imdiff=abs(df["EG.IonMobility"][i]-df["EG.IonMobility"][i+1])
            if rtdiff <= rttolerance and mzdiff <= mztolerance and imdiff >= imtolerance:
                coelutingpeptides.loc[len(coelutingpeptides)]=df.iloc[i].tolist()
                coelutingpeptides.loc[len(coelutingpeptides)]=df.iloc[i+1].tolist()

        #adding a column for a rough group number for each group of peptides detected
        for i in range(len(coelutingpeptides)):
            if i+1 not in range(len(coelutingpeptides)):
                break
            #rtpercentdiff=(abs(coelutingpeptides["EG.ApexRT"][i]-coelutingpeptides["EG.ApexRT"][i+1])/coelutingpeptides["EG.ApexRT"][i])*100
            rtdiff=abs(coelutingpeptides["EG.ApexRT"][i]-coelutingpeptides["EG.ApexRT"][i+1])
            mzdiff=abs(coelutingpeptides["FG.PrecMz"][i]-coelutingpeptides["FG.PrecMz"][i+1])
            imdiff=abs(coelutingpeptides["EG.IonMobility"][i]-coelutingpeptides["EG.IonMobility"][i+1])
            if rtdiff <= rttolerance and mzdiff <= mztolerance and imdiff >= imtolerance:
                coelutingpeptides.loc[coelutingpeptides.index[i],"Group"]=i
        coelutingpeptides=coelutingpeptides[["Group","EG.ModifiedPeptide","FG.Charge","FG.PrecMz","EG.ApexRT","EG.IonMobility","FG.MS2Quantity","PG.ProteinGroups","PG.ProteinNames"]]
        with io.BytesIO() as buf:
            coelutingpeptides.to_csv(buf,index=False)
            yield buf.getvalue()

    #download table of PTMs per precursor
    @render.download(filename=lambda: f"PTM List_{input.searchreport()[0]['name']}.csv")
    def ptmlist_download():
        searchoutput,sampleconditions,maxreplicatelist,numconditions,numsamples,runnames=variables()
        ptmdf=pd.DataFrame()
        for condition in sampleconditions:
            df=searchoutput[searchoutput["R.Condition"]==condition][["EG.ModifiedPeptide","FG.Charge"]].drop_duplicates().reset_index(drop=True)
            dfptmlist=[]
            numptms=[]
            for i in df["EG.ModifiedPeptide"]:
                foundptms=re.findall(r"[^[]*\[([^]]*)\]",i)
                dfptmlist.append(foundptms)
                numptms.append(len(foundptms))
            dfptmlist=pd.Series(dfptmlist).value_counts().to_frame().reset_index().rename(columns={"index":condition,"count":condition+"_count"})
            ptmdf=pd.concat([ptmdf,dfptmlist],axis=1)
        with io.BytesIO() as buf:
            ptmdf.to_csv(buf,index=False)
            yield buf.getvalue()

    # ====================================== Figures to PPT
    #generation of keyword lists for output
    @render.ui
    def keys_ui():
        if input.download_path():
            image_dir=input.download_path()
        elif input.twosoftware_download_path():
            image_dir=input.twosoftware_download_path()
        else:
            return ui.p("Figure output directory not specified in either File Import or Two-Software Comparison tabs")
        
        os.chdir(image_dir)
        cwd=os.getcwd()

        #get list of files in the given directory that are .pngs
        pnglist=[file for file in os.listdir() if ".png" in file]

        software_split=pd.DataFrame(pnglist)[0].str.split("_",n=1,expand=True)
        searchfile_split=software_split[1].str.split(".",n=1,expand=True)
        searchfile_split[1]=searchfile_split[1].str.split(".png",expand=True)[0].str.split("_",n=1,expand=True)[1]
        searchfile_split[0]=searchfile_split[0]+"."+software_split[1].str.split(".",expand=True)[1].str.split("_",n=1,expand=True)[0]

        #keywords to download specific plots, especially if the direcotry where the figures are stored has other timsplot figures
        key_software=software_split[0].drop_duplicates().tolist()
        key_searchfile=searchfile_split[0].drop_duplicates().tolist()
        key_plots=searchfile_split[1].str.split(".png",expand=True)[0].drop_duplicates().tolist()
        keyword_df=pd.concat([software_split[[0]].rename(columns={0:"software"}),searchfile_split],axis=1).rename(columns={0:"searchfile",1:"plot"})

        #row and column setup to confine all keywords to a small card structure
        return ui.row(ui.column(4,ui.input_checkbox_group("software_keys","Software keywords:",choices=key_software)),ui.column(4,ui.input_checkbox_group("searchfile_keys","Search report keywords:",choices=key_searchfile)),ui.column(4,ui.input_checkbox_group("plots_keys","Plot type keywords:",choices=key_plots)))
    #show list of what files will be exported based on the selected keywords
    @render.text
    def selected_images():
        if input.download_path():
            image_dir=input.download_path()
        elif input.twosoftware_download_path():
            image_dir=input.twosoftware_download_path()
        else:
            return ""
        os.chdir(image_dir)
        cwd=os.getcwd()

        #get list of files in the given directory that are .pngs
        pnglist=[file for file in os.listdir() if ".png" in file]

        software_split=pd.DataFrame(pnglist)[0].str.split("_",n=1,expand=True)
        searchfile_split=software_split[1].str.split(".",n=1,expand=True)
        searchfile_split[1]=searchfile_split[1].str.split(".png",expand=True)[0].str.split("_",n=1,expand=True)[1]
        searchfile_split[0]=searchfile_split[0]+"."+software_split[1].str.split(".",expand=True)[1].str.split("_",n=1,expand=True)[0]

        #keywords to download specific plots, especially if the direcotry where the figures are stored has other timsplot figures
        key_software=software_split[0].drop_duplicates().tolist()
        key_searchfile=searchfile_split[0].drop_duplicates().tolist()
        key_plots=searchfile_split[1].str.split(".png",expand=True)[0].drop_duplicates().tolist()
        keyword_df=pd.concat([software_split[[0]].rename(columns={0:"software"}),searchfile_split],axis=1).rename(columns={0:"searchfile",1:"plot"})

        software_keys=input.software_keys()
        searchfile_keys=input.searchfile_keys()
        plots_keys=input.plots_keys()
        #hierarchy of selection: software>searchfile>plots
        #first check if software keys are selected
        if software_keys:
            if searchfile_keys:
                if plots_keys:
                    #all 3 filters
                    image_export_df=keyword_df[(keyword_df["software"].isin(software_keys))&(keyword_df["searchfile"].isin(searchfile_keys))&(keyword_df["plot"].isin(plots_keys))]
                else:
                    #software and searchfile filters if no plot keys
                    image_export_df=keyword_df[(keyword_df["software"].isin(software_keys))&(keyword_df["searchfile"].isin(searchfile_keys))]
            else:
                #software keys but no searchfile keys
                if plots_keys:
                    #software and plots keys
                    image_export_df=keyword_df[(keyword_df["software"].isin(software_keys))&(keyword_df["plot"].isin(plots_keys))]
                else:
                    #only software keys
                    image_export_df=keyword_df[(keyword_df["software"].isin(software_keys))]
        #if no software keys
        else:
            if searchfile_keys:
                if plots_keys:
                    #no software keys
                    image_export_df=keyword_df[(keyword_df["searchfile"].isin(searchfile_keys))&(keyword_df["plot"].isin(plots_keys))]
                else:
                    #no software or plots keys
                    image_export_df=keyword_df[(keyword_df["searchfile"].isin(searchfile_keys))]
            else:
                if plots_keys:
                    #just plots keys
                    image_export_df=keyword_df[(keyword_df["plot"].isin(plots_keys))]
                else:
                    #no selection
                    image_export_df=pd.DataFrame()
        if image_export_df.empty:
            return "No keywords selected"
        else:
            image_export_df=image_export_df.reset_index(drop=True)
            image_export_list=[]
            for i in range(len(image_export_df)):
                image_export_list.append(image_export_df["software"][i]+"_"+image_export_df["searchfile"][i]+"_"+image_export_df["plot"][i]+".png")
            return image_export_list
    @render.download(filename="timsplot_figures.pptx")
    def pptx_download():
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        import PIL.Image

        os_var=os_check()
        if input.download_path():
            image_dir=input.download_path()
        elif input.twosoftware_download_path():
            image_dir=input.twosoftware_download_path()

        prs=Presentation()
        slide_layout=prs.slide_layouts[6] #blank slide

        slideheight=7.5
        slidewidth=13.33
        prs.slide_height=Inches(slideheight)
        prs.slide_width=Inches(slidewidth)

        os.chdir(image_dir)
        cwd=os.getcwd()

        #get list of files in the given directory that are .pngs
        pnglist=[file for file in os.listdir() if ".png" in file]

        software_split=pd.DataFrame(pnglist)[0].str.split("_",n=1,expand=True)
        searchfile_split=software_split[1].str.split(".",n=1,expand=True)
        searchfile_split[1]=searchfile_split[1].str.split(".png",expand=True)[0].str.split("_",n=1,expand=True)[1]
        searchfile_split[0]=searchfile_split[0]+"."+software_split[1].str.split(".",expand=True)[1].str.split("_",n=1,expand=True)[0]

        #keywords to download specific plots, especially if the direcotry where the figures are stored has other timsplot figures
        key_software=software_split[0].drop_duplicates().tolist()
        key_searchfile=searchfile_split[0].drop_duplicates().tolist()
        key_plots=searchfile_split[1].str.split(".png",expand=True)[0].drop_duplicates().tolist()
        keyword_df=pd.concat([software_split[[0]].rename(columns={0:"software"}),searchfile_split],axis=1).rename(columns={0:"searchfile",1:"plot"})

        software_keys=input.software_keys()
        searchfile_keys=input.searchfile_keys()
        plots_keys=input.plots_keys()
        #hierarchy of selection: software>searchfile>plots
        #first check if software keys are selected
        if software_keys:
            if searchfile_keys:
                if plots_keys:
                    #all 3 filters
                    image_export_df=keyword_df[(keyword_df["software"].isin(software_keys))&(keyword_df["searchfile"].isin(searchfile_keys))&(keyword_df["plot"].isin(plots_keys))]
                else:
                    #software and searchfile filters if no plot keys
                    image_export_df=keyword_df[(keyword_df["software"].isin(software_keys))&(keyword_df["searchfile"].isin(searchfile_keys))]
            else:
                #software keys but no searchfile keys
                if plots_keys:
                    #software and plots keys
                    image_export_df=keyword_df[(keyword_df["software"].isin(software_keys))&(keyword_df["plot"].isin(plots_keys))]
                else:
                    #only software keys
                    image_export_df=keyword_df[(keyword_df["software"].isin(software_keys))]
        #if no software keys
        else:
            if searchfile_keys:
                if plots_keys:
                    #no software keys
                    image_export_df=keyword_df[(keyword_df["searchfile"].isin(searchfile_keys))&(keyword_df["plot"].isin(plots_keys))]
                else:
                    #no software or plots keys
                    image_export_df=keyword_df[(keyword_df["searchfile"].isin(searchfile_keys))]
            else:
                if plots_keys:
                    #just plots keys
                    image_export_df=keyword_df[(keyword_df["plot"].isin(plots_keys))]
                else:
                    #no selection
                    image_export_df=pd.DataFrame()
        image_export_df=image_export_df.reset_index(drop=True)
        image_export_list=[]
        for i in range(len(image_export_df)):
            image_export_list.append(image_export_df["software"][i]+"_"+image_export_df["searchfile"][i]+"_"+image_export_df["plot"][i]+".png")

        border=0.5
        slidewidth_border=Inches(slidewidth-border)
        slideheight_border=Inches(slideheight-border)

        for image in image_export_list:
            image_path=image_dir+os_var+image
            figureimage=PIL.Image.open(image_path)
            width=Inches(figureimage.size[0]*(1/plt.rcParams["figure.dpi"])) #width in inches
            height=Inches(figureimage.size[1]*(1/plt.rcParams["figure.dpi"])) #height in inches

            slide=prs.slides.add_slide(slide_layout)
            
            left=top=Inches(border)
            if width>slidewidth_border or height>slideheight_border:
                #if too wide but fits vertically
                if width>slidewidth_border and height<=slideheight_border:
                    add_picture=slide.shapes.add_picture(image_path,left=left,top=top,width=slidewidth_border)
                #if too tall but fits horizontally
                elif width<=slidewidth_border and height>slideheight_border:
                    add_picture=slide.shapes.add_picture(image_path,left=left,top=top,height=slideheight_border)
                #if too big in both dimensions, rescale based on height
                elif width>slidewidth_border and height>slideheight_border:
                    add_picture=slide.shapes.add_picture(image_path,left=left,top=top,height=slideheight_border)
            else:
                add_picture=slide.shapes.add_picture(image_path,left=left,top=top,width=width,height=height)
        with io.BytesIO() as buf:
            prs.save(buf)
            yield buf.getvalue()

#endregion

# ============================================================================= Raw Data
#region
    # ====================================== Multi-File Import
    #choose whether to import data from individual file names or from a directory
    @render.ui
    def rawfile_input_ui():
        if input.file_or_folder()=="individual":
            return ui.input_text_area("rawfile_input","Paste the path for each .d file you want to upload (note: do not leave whitespace at the end):",width="1500px",autoresize=True,placeholder="ex. - C:\\Users\\Data\\K562_500ng_1_Slot1-49_1_3838.d")
        if input.file_or_folder()=="directory":
            return ui.input_text_area("rawfile_input","Paste the path for the directory containing the raw files to upload (note: do not leave whitespace at the end):",width="1500px",autoresize=True,placeholder="ex. - C:\\Users\\Data")

    #separate function for downloading files from this section since it's independent from the File Import section
    #either exports to the directory that's input or the parent directory of the first raw file that's input
    def imagedownload_rawfile(plotname):
        if input.rawfile_download_parent_path()==True:
            os_var=os_check()
            if input.file_or_folder()=="individual":
                firstfile=list(input.rawfile_input().split("\n"))[0]
                directory=(os_var).join(firstfile.split(os_var)[:-1])+os_var
                plt.savefig(directory+plotname+".png")
            if input.file_or_folder()=="directory":
                directory=input.rawfile_input()
                plt.savefig(directory+os_var+plotname+".png")
        else:
            pass

    #take text input for data paths and make dictionaries of frame data
    @reactive.calc
    def rawfile_list():
        import alphatims.bruker as atb
        import alphatims.plotting as atp
        os_var=os_check()
        if input.file_or_folder()=="individual":
            filelist=list(input.rawfile_input().split("\n"))
            MSframedict=dict()
            precursordict=dict()
            samplenames=[]
            for run in filelist:
                frames=pd.DataFrame(atb.read_bruker_sql(run)[2])
                MSframedict[run]=frames[frames["MsMsType"]==0].reset_index(drop=True)
                precursordict[run]=pd.DataFrame(atb.read_bruker_sql(run)[3])
                samplenames.append(run.split(os_var)[-1])
            return MSframedict,precursordict,samplenames
        if input.file_or_folder()=="directory":
            os.chdir(input.rawfile_input())
            cwd=os.getcwd()
            filelist=[]
            for file in os.listdir():
                if ".d" in file:
                    filelist.append(cwd+os_var+file)
            MSframedict=dict()
            precursordict=dict()
            samplenames=[]
            for run in filelist:
                frames=pd.DataFrame(atb.read_bruker_sql(run)[2])
                MSframedict[run]=frames[frames["MsMsType"]==0].reset_index(drop=True)
                precursordict[run]=pd.DataFrame(atb.read_bruker_sql(run)[3])
                samplenames.append(run.split(os_var)[-1])
            return MSframedict,precursordict,samplenames

    #for directory input, show what files were found in the directory
    @render.text
    def uploadedfiles():
        if input.file_or_folder()=="directory":
            try:
                MSframedict,precursordict,samplenames=rawfile_list()
                printsamplenames="\n".join(str(el) for el in samplenames)
            except:
                return ""
            return printsamplenames

    # ====================================== TIC Plot
    #render ui for checkboxes to plot specific runs
    @render.ui
    def rawfile_checkboxes_tic():
        MSframedict,precursordict,samplenames=rawfile_list()
        opts=dict()
        keys=MSframedict.keys()
        labels=samplenames
        for x,y in zip(keys,labels):
            opts[x]=y
        return ui.input_checkbox_group("rawfile_pick_tic","Pick files to plot data for:",choices=opts,width="800px")
    #plot TIC from raw data
    @reactive.effect
    def _():
        @render.plot(width=input.tic_width(),height=input.tic_height())
        def TIC_plot():
            MSframedict,precursordict,samplenames=rawfile_list()
            checkgroup=input.rawfile_pick_tic()
            os_var=os_check()
            colors=list(mcolors.TABLEAU_COLORS)

            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.stacked_tic()==True:
                fig,ax=plt.subplots(nrows=len(checkgroup),sharex=True,figsize=(input.tic_width()*(1/plt.rcParams['figure.dpi']),input.tic_height()*(1/plt.rcParams['figure.dpi'])))
                for i,run in enumerate(checkgroup):
                    x=MSframedict[run]["Time"]/60
                    y=MSframedict[run]["SummedIntensities"]
                    ax[i].plot(x,y,label=run.split(os_var)[-1],linewidth=1.5,color=colors[i])
                    ax[i].set_ylabel("Intensity",fontsize=axisfont)
                    ax[0].set_title("Total Ion Chromatogram",fontsize=titlefont)
                    ax[i].set_axisbelow(True)
                    ax[i].grid(linestyle="--")
                    ax[i].xaxis.set_minor_locator(MultipleLocator(1))
                    ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                    legend=ax[i].legend(loc="upper left",prop={'size':legendfont})
                    for z in legend.legend_handles:
                        z.set_linewidth(5)
                ax[i].set_xlabel("Time (min)",fontsize=axisfont)
            else:
                fig,ax=plt.subplots(figsize=(input.tic_width()*(1/plt.rcParams['figure.dpi']),input.tic_height()*(1/plt.rcParams['figure.dpi'])))
                for run in checkgroup:
                    x=MSframedict[run]["Time"]/60
                    y=MSframedict[run]["SummedIntensities"]
                    ax.plot(x,y,label=run.split(os_var)[-1],linewidth=0.75)
                ax.set_xlabel("Time (min)",fontsize=axisfont)
                ax.set_ylabel("Intensity",fontsize=axisfont)
                ax.set_title("Total Ion Chromatogram",fontsize=titlefont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.xaxis.set_minor_locator(MultipleLocator(1))
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                legend=ax.legend(loc="upper left",prop={'size':legendfont})
                for z in legend.legend_handles:
                    z.set_linewidth(5)
            fig.set_tight_layout(True)
            if checkgroup:
                imagedownload_rawfile("TIC")

    # ====================================== BPC Plot
    #render ui for checkboxes to plot specific runs
    @render.ui
    def rawfile_checkboxes_bpc():
        MSframedict,precursordict,samplenames=rawfile_list()
        opts=dict()
        keys=MSframedict.keys()
        labels=samplenames
        for x,y in zip(keys,labels):
            opts[x]=y
        return ui.input_checkbox_group("rawfile_pick_bpc","Pick files to plot data for:",choices=opts,width="800px")
    #plot BPC from raw data
    @reactive.effect
    def _():
        @render.plot(width=input.bpc_width(),height=input.bpc_height())
        def BPC_plot():
            MSframedict,precursordict,samplenames=rawfile_list()
            checkgroup=input.rawfile_pick_bpc()
            os_var=os_check()
            colors=list(mcolors.TABLEAU_COLORS)
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.stacked_bpc()==True:
                fig,ax=plt.subplots(nrows=len(checkgroup),sharex=True,figsize=(input.bpc_width()*(1/plt.rcParams['figure.dpi']),input.bpc_height()*(1/plt.rcParams['figure.dpi'])))
                for i,run in enumerate(checkgroup):
                    x=MSframedict[run]["Time"]/60
                    y=MSframedict[run]["MaxIntensity"]
                    ax[i].plot(x,y,label=run.split(os_var)[-1],linewidth=1.5,color=colors[i])
                    ax[i].set_ylabel("Intensity",fontsize=axisfont)
                    ax[0].set_title("Base Peak Chromatogram",fontsize=titlefont)
                    ax[i].set_axisbelow(True)
                    ax[i].grid(linestyle="--")
                    ax[i].xaxis.set_minor_locator(MultipleLocator(1))
                    ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                    legend=ax[i].legend(loc="upper left",prop={'size':legendfont})
                    for z in legend.legend_handles:
                        z.set_linewidth(5)
                ax[i].set_xlabel("Time (min)",fontsize=axisfont)
            else:
                fig,ax=plt.subplots(figsize=(input.bpc_width()*(1/plt.rcParams['figure.dpi']),input.bpc_height()*(1/plt.rcParams['figure.dpi'])))
                for run in checkgroup:
                    x=MSframedict[run]["Time"]/60
                    y=MSframedict[run]["MaxIntensity"]
                    ax.plot(x,y,label=run.split(os_var)[-1],linewidth=0.75)
                ax.set_xlabel("Time (min)",fontsize=axisfont)
                ax.set_ylabel("Intensity",fontsize=axisfont)
                ax.set_title("Base Peak Chromatogram",fontsize=titlefont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.xaxis.set_minor_locator(MultipleLocator(1))
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                legend=ax.legend(loc='upper left',prop={'size':legendfont})
                for z in legend.legend_handles:
                    z.set_linewidth(5)
            fig.set_tight_layout(True)
            if checkgroup:
                imagedownload_rawfile("BPC")

    # ====================================== Accumulation Time
    #render ui for checkboxes to plot specific runs
    @render.ui
    def rawfile_checkboxes_accutime():
        MSframedict,precursordict,samplenames=rawfile_list()
        opts=dict()
        keys=MSframedict.keys()
        labels=samplenames
        for x,y in zip(keys,labels):
            opts[x]=y
        return ui.input_checkbox_group("rawfile_pick_accutime","Pick files to plot data for:",choices=opts,width="800px")
    #plot accumulation time from raw data
    @reactive.effect
    def _():
        @render.plot(width=input.accutime_width(),height=input.accutime_height())
        def accutime_plot():
            MSframedict,precursordict,samplenames=rawfile_list()
            checkgroup=input.rawfile_pick_accutime()
            os_var=os_check()
            colors=list(mcolors.TABLEAU_COLORS)
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            if input.stacked_accutime()==True:
                fig,ax=plt.subplots(nrows=len(checkgroup),sharex=True,figsize=(input.accutime_width()*(1/plt.rcParams['figure.dpi']),input.accutime_height()*(1/plt.rcParams['figure.dpi'])))
                for i,run in enumerate(checkgroup):
                    x=MSframedict[run]["Time"]/60
                    y=MSframedict[run]["AccumulationTime"]
                    ax[i].plot(x,y,label=run.split(os_var)[-1],linewidth=1.5,color=colors[i])
                    ax[i].set_ylabel("Accumulation Time (ms)",fontsize=axisfont)
                    ax[0].set_title("Accumulation Time Chromatogram",fontsize=titlefont)
                    ax[i].set_axisbelow(True)
                    ax[i].grid(linestyle="--")
                    ax[i].xaxis.set_minor_locator(MultipleLocator(1))
                    ax[i].tick_params(axis="both",labelsize=axisfont_labels)
                    legend=ax[i].legend(loc="upper left",prop={'size':legendfont})
                    for z in legend.legend_handles:
                        z.set_linewidth(5)
                ax[i].set_xlabel("Time (min)",fontsize=axisfont)
            else:
                fig,ax=plt.subplots(figsize=(input.accutime_width()*(1/plt.rcParams['figure.dpi']),input.accutime_height()*(1/plt.rcParams['figure.dpi'])))
                for run in checkgroup:
                    x=MSframedict[run]["Time"]/60
                    y=MSframedict[run]["AccumulationTime"]
                    ax.plot(x,y,label=run.split(os_var)[-1],linewidth=0.75)
                ax.set_xlabel("Time (min)",fontsize=axisfont)
                ax.set_ylabel("Accumulation Time (ms)",fontsize=axisfont)
                ax.set_title("Accumulation Time Chromatogram",fontsize=titlefont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.xaxis.set_minor_locator(MultipleLocator(1))
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                legend=ax.legend(loc='center left',bbox_to_anchor=(1,0.5),prop={'size':legendfont})
                for z in legend.legend_handles:
                    z.set_linewidth(5)
            fig.set_tight_layout(True)
            if checkgroup:
                imagedownload_rawfile("AccumulationTime")

    # ====================================== EIC Plot
    @render.ui
    def rawfile_buttons_eic():
        MSframedict,precursordict,samplenames=rawfile_list()
        opts=dict()
        keys=MSframedict.keys()
        labels=samplenames
        for x,y in zip(keys,labels):
            opts[x]=y
        return ui.input_radio_buttons("rawfile_pick_eic","Pick file to plot data for:",choices=opts,width="800px")
    #input for mobility add-on for EICs
    @render.ui
    def mobility_input():
        if input.include_mobility()==True:
            return ui.input_text("mobility_input_value","Input mobility for EIC:"),ui.input_text("mobility_input_window","Input mobility window (1/k0) for EIC:")
    @reactive.calc
    @reactive.event(input.eic_load_rawfile)
    def eic_rawfile_import():
        import alphatims.bruker as atb
        import alphatims.plotting as atp
        rawfile=atb.TimsTOF(input.rawfile_pick_eic())
        return rawfile
    #plot EIC for input mass
    @reactive.effect
    def _():
        @render.plot(width=input.eic_width(),height=input.eic_height())
        def eic():
            rawfile=eic_rawfile_import()
            os_var=os_check()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            try:
                mz=float(input.eic_mz_input())
                ppm_error=float(input.eic_ppm_input())
                
                low_mz=mz/(1+ppm_error/10**6)
                high_mz=mz*(1+ppm_error/10**6)

                if input.include_mobility()==True:
                    mobility=float(input.mobility_input_value())
                    window=float(input.mobility_input_window())
                    low_mobility=mobility-window
                    high_mobility=mobility+window
                    eic_df=rawfile[:,low_mobility: high_mobility,0,low_mz: high_mz]
                else:
                    eic_df=rawfile[:,:,0,low_mz: high_mz]

                fig,ax=plt.subplots(figsize=(input.eic_width()*(1/plt.rcParams['figure.dpi']),input.eic_height()*(1/plt.rcParams['figure.dpi'])))
                ax.plot(eic_df["rt_values_min"],eic_df["intensity_values"],linewidth=0.5)
                if input.include_mobility()==True:
                    ax.set_title(input.rawfile_pick_eic().split(os_var)[-1]+"\n"+"EIC: "+str(input.eic_mz_input())+", Mobility: "+str(input.mobility_input_value()),fontsize=titlefont)
                else:
                    ax.set_title(input.rawfile_pick_eic().split(os_var)[-1]+"\n"+"EIC: "+str(input.eic_mz_input()),fontsize=titlefont)
                ax.xaxis.set_minor_locator(MultipleLocator(1))
                ax.set_xlabel("Time (min)",fontsize=axisfont)
                ax.set_ylabel("Intensity",fontsize=axisfont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                fig.set_tight_layout(True)
                imagedownload_rawfile("EIC_"+str(input.eic_mz_input()))
            except:
                fig,ax=plt.subplots()
                ax.set_xlabel("Time (min)",fontsize=axisfont)
                ax.set_ylabel("Intensity",fontsize=axisfont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)

    # ====================================== EIM Plot
    @render.ui
    def rawfile_buttons_eim():
        MSframedict,precursordict,samplenames=rawfile_list()
        opts=dict()
        keys=MSframedict.keys()
        labels=samplenames
        for x,y in zip(keys,labels):
            opts[x]=y
        return ui.input_radio_buttons("rawfile_pick_eim","Pick file to plot data for:",choices=opts,width="800px")
    @reactive.calc
    @reactive.event(input.eim_load_rawfile)
    def eim_rawfile_import():
        import alphatims.bruker as atb
        import alphatims.plotting as atp
        rawfile=atb.TimsTOF(input.rawfile_pick_eim())
        return rawfile
    #plot EIM for input mass
    @reactive.effect
    def _():
        @render.plot(width=input.eim_width(),height=input.eim_height())
        def eim():
            rawfile=eim_rawfile_import()
            os_var=os_check()
            titlefont=input.titlefont()
            axisfont=input.axisfont()
            axisfont_labels=input.axisfont_labels()
            labelfont=input.labelfont()
            legendfont=input.legendfont()
            y_padding=input.ypadding()
            x_label_rotation=input.xaxis_label_rotation()

            try:
                mz=float(input.eim_mz_input())
                ppm_error=float(input.eim_ppm_input())
                low_mz=mz/(1+ppm_error/10**6)
                high_mz=mz*(1+ppm_error/10**6)

                eim_df=rawfile[:,:,0,low_mz: high_mz].sort_values("mobility_values")

                fig,ax=plt.subplots(figsize=(input.eim_width()*(1/plt.rcParams['figure.dpi']),input.eim_height()*(1/plt.rcParams['figure.dpi'])))
                ax.plot(eim_df["mobility_values"],eim_df["intensity_values"],linewidth=0.5)
                ax.set_title(input.rawfile_pick_eim().split(os_var)[-1]+"\n"+"EIM: "+str(input.eim_mz_input()),fontsize=titlefont)
                ax.xaxis.set_minor_locator(MultipleLocator(0.01))
                ax.set_xlabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
                ax.set_ylabel("Intensity",fontsize=axisfont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)
                fig.set_tight_layout(True)
                imagedownload_rawfile("EIM_"+str(input.eim_mz_input()))
            except:
                fig,ax=plt.subplots()
                ax.set_xlabel("Ion Mobility ($1/K_{0}$)",fontsize=axisfont)
                ax.set_ylabel("Intensity",fontsize=axisfont)
                ax.set_axisbelow(True)
                ax.grid(linestyle="--")
                ax.tick_params(axis="both",labelsize=axisfont_labels)

#endregion

#endregion

app=App(app_ui,server)
